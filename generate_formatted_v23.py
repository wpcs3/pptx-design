"""Generate complete presentation v23 with updated formatting.

Changes from v22:
1. Cropped section images to fit 11"x8.5" aspect ratio
2. Standardized bullet points (not HEADER-body format)
3. Updated table formatting:
   - Vertical center alignment for all rows
   - Header: 16pt, 0.5" row height
   - Data: 14pt, 0.4" row height
   - First column: 0.1" left margin
   - Other margins: 0.0"
   - Auto-fit column widths
"""

import json
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================================
# FORMATTING SPECIFICATIONS
# ============================================================================

COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),      # #051C2C - titles, table headers
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),    # #309CE7 - metric boxes
    'white': RGBColor(0xFF, 0xFF, 0xFF),          # #FFFFFF - text on dark
    'dark_text': RGBColor(0x06, 0x1F, 0x32),      # #061F32 - body text
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),     # Alternating row fill
}

# Layout indices in slide master
LAYOUTS = {
    'frontpage': 0,
    'default': 1,
    '1_default': 2,
    'top_left_title': 3,
    'section_breaker': 4,
    '1_4_grey': 5,
    '1_3_grey': 6,
    '1_2_grey': 7,
    '3_4_grey': 8,
    '1_1_2_grey': 9,
    'end': 10,
}

# Positioning (in inches)
POS = {
    'title': {'left': 0.40, 'top': 0.40, 'width': 10.20, 'height': 1.00},
    'subtitle': {'left': 0.40, 'top': 1.70, 'width': 10.20, 'height': 0.80},
    'content': {'left': 0.40, 'top': 2.80, 'width': 10.20, 'height': 4.70},
    'footer': {'left': 0.40, 'top': 7.50, 'width': 10.20, 'height': 0.40},
    'table': {'left': 0.40, 'top': 2.80, 'width': 10.20},
    'chart': {'left': 0.40, 'top': 2.80, 'width': 10.20, 'height': 4.00},
}

# Font specifications
FONTS = {
    'title_slide_title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': 'white'},
    'title_slide_subtitle': {'name': 'Arial', 'size': 20, 'bold': True, 'color': 'white'},
    'title_slide_date': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'white'},
    'content_title': {'name': 'Arial', 'size': 36, 'bold': True, 'color': 'dark_text'},
    'content_title_long': {'name': 'Arial', 'size': 32, 'bold': True, 'color': 'dark_text'},
    'subtitle': {'name': 'Arial', 'size': 20, 'bold': True, 'color': 'dark_text'},
    'body': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'dark_text'},
    'body_header': {'name': 'Arial', 'size': 14, 'bold': True, 'color': 'dark_text'},
    'metric_value': {'name': 'Arial', 'size': 28, 'bold': True, 'color': 'white'},
    'metric_label': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'white'},
    'table_header': {'name': 'Arial', 'size': 16, 'bold': True, 'color': 'white'},
    'table_data': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'dark_text'},
    'section_title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': 'white'},
    'footnote': {'name': 'Arial', 'size': 10, 'bold': False, 'color': 'dark_text'},
}

# Table formatting
TABLE_FORMAT = {
    'header_row_height': 0.50,
    'data_row_height': 0.40,
    'first_col_left_margin': Inches(0.1),
    'other_margins': Inches(0.0),
}

# Cropped image paths
CROPPED_IMAGE_DIR = Path("cache/images/cropped")

SECTION_IMAGE_MAP = {
    "Executive Summary": "title_slide.png",
    "Market Fundamentals": "market_fundamentals.png",
    "Target Markets": "target_markets.png",
    "Target Market Analysis": "target_markets.png",
    "Demand Drivers": "demand_drivers.png",
    "Structural Demand Drivers": "demand_drivers.png",
    "Investment Strategy": "investment_strategy.png",
    "Competitive Positioning": "competitive_positioning.png",
    "Risk Management": "risk_management.png",
    "Risk Factors": "risk_management.png",
    "Risk Factors & Mitigants": "risk_management.png",
    "ESG Strategy": "esg_strategy.png",
    "ESG & Sustainability": "esg_strategy.png",
    "JV Structure": "jv_structure.png",
    "JV Structure & Governance": "jv_structure.png",
    "Conclusion": "conclusion.png",
}


def apply_font(run, font_spec):
    """Apply font specification to a run."""
    font = run.font
    font.name = font_spec['name']
    font.size = Pt(font_spec['size'])
    font.bold = font_spec['bold']
    if font_spec['color'] in COLORS:
        font.color.rgb = COLORS[font_spec['color']]


def format_table(table):
    """Apply table formatting with proper margins and alignment."""
    num_cols = len(table.columns)

    for row_idx, row in enumerate(table.rows):
        # Set row height
        if row_idx == 0:
            row.height = Inches(TABLE_FORMAT['header_row_height'])
        else:
            row.height = Inches(TABLE_FORMAT['data_row_height'])

        for col_idx, cell in enumerate(row.cells):
            # Vertical alignment: center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins
            cell.margin_top = TABLE_FORMAT['other_margins']
            cell.margin_bottom = TABLE_FORMAT['other_margins']
            cell.margin_right = TABLE_FORMAT['other_margins']

            # First column gets 0.1" left margin, others get 0.0
            if col_idx == 0:
                cell.margin_left = TABLE_FORMAT['first_col_left_margin']
            else:
                cell.margin_left = TABLE_FORMAT['other_margins']

            # Apply fill colors
            if row_idx == 0:
                # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                # Even data rows (2, 4, 6...) - light gray
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                # Odd data rows (1, 3, 5...) - white
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            # Format text
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    if row_idx == 0:
                        apply_font(run, FONTS['table_header'])
                    else:
                        apply_font(run, FONTS['table_data'])


def calculate_column_widths(headers, rows, table_width_inches=10.2):
    """Calculate optimal column widths based on content to avoid wrapping."""
    num_cols = len(headers)

    # Measure max content length per column
    max_lengths = []
    for col_idx in range(num_cols):
        max_len = len(str(headers[col_idx]))
        for row in rows:
            if col_idx < len(row):
                max_len = max(max_len, len(str(row[col_idx])))
        max_lengths.append(max_len)

    total_chars = sum(max_lengths)
    if total_chars == 0:
        # Equal distribution if no content
        return [table_width_inches / num_cols] * num_cols

    # Distribute width proportionally, with minimum width of 1"
    min_width = 1.0
    available_width = table_width_inches - (min_width * num_cols)

    widths = []
    for max_len in max_lengths:
        proportion = max_len / total_chars
        width = min_width + (available_width * proportion)
        widths.append(width)

    # Normalize to exactly table_width_inches
    scale = table_width_inches / sum(widths)
    return [w * scale for w in widths]


def add_standardized_bullets(text_frame, bullets):
    """Add bullets using standard bullet point formatting."""
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        para.level = 0
        para.space_before = Pt(6)
        para.space_after = Pt(3)

        # Add the bullet text
        run = para.add_run()
        run.text = bullet
        apply_font(run, FONTS['body'])

        # Set bullet character via XML
        try:
            pPr = para._p.get_or_add_pPr()
            # Remove any existing bullet settings
            for child in list(pPr):
                if 'bu' in child.tag.lower():
                    pPr.remove(child)

            # Add bullet character
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '\u2022')  # Standard bullet

            # Set left margin and indent for proper bullet spacing
            para._p.get_or_add_pPr().set(qn('a:marL'), str(Emu(Inches(0.25))))
            para._p.get_or_add_pPr().set(qn('a:indent'), str(Emu(Inches(-0.15))))
        except:
            pass


def main():
    """Generate the formatted presentation v23."""
    print("Generating formatted presentation v23...")
    print("=" * 60)

    # Load the user's edited PPTX as template
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v21_CS_edits_v2.pptx")
    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    # Load the original outline for content
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    if not outline_path.exists():
        print(f"Error: Outline not found: {outline_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    print(f"Template loaded: {len(prs.slides)} slides")

    # Flatten outline to get all slides
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({
                'section': section.get('name', ''),
                **slide
            })

    print(f"Outline has {len(all_slides)} slides")

    existing_slide_count = len(prs.slides)
    print(f"\nAdding slides {existing_slide_count + 1} to {len(all_slides)}...")

    # Get layouts
    default_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['default']]
    section_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['section_breaker']]

    # Add remaining slides
    for idx in range(existing_slide_count, len(all_slides)):
        slide_data = all_slides[idx]
        slide_type = slide_data.get('slide_type', 'title_content')
        content = slide_data.get('content', {})
        section_name = slide_data.get('section', '')

        print(f"  Adding slide {idx + 1}: {slide_type} - {content.get('title', 'N/A')[:40]}...")

        # Choose layout
        if slide_type == 'section_divider':
            layout = section_layout
        else:
            layout = default_layout

        slide = prs.slides.add_slide(layout)

        # Find placeholders
        title_shape = None
        subtitle_shape = None
        content_shape = None

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # TITLE
                    title_shape = shape
                elif ph_type == 4:  # SUBTITLE
                    subtitle_shape = shape
                elif ph_type == 7:  # OBJECT/CONTENT
                    content_shape = shape

        # Set title
        if title_shape and content.get('title'):
            title_text = content['title']
            tf = title_shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = title_text
            if len(title_text) > 40:
                apply_font(run, FONTS['content_title_long'])
            else:
                apply_font(run, FONTS['content_title'])

        # Set subtitle/takeaway
        if subtitle_shape:
            subtitle_text = content.get('takeaway') or content.get('subtitle', '')
            if subtitle_text:
                tf = subtitle_shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = subtitle_text
                apply_font(run, FONTS['subtitle'])

        # Handle different slide types
        if slide_type in ['title_content', 'content']:
            bullets = content.get('bullets', [])
            if bullets and content_shape:
                add_standardized_bullets(content_shape.text_frame, bullets)

        elif slide_type == 'table_slide':
            # Get table data (both formats)
            table_data = content.get('table', {})
            if table_data:
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
            else:
                headers = content.get('headers', [])
                rows = content.get('data', [])

            if headers and rows:
                num_cols = len(headers)
                num_rows = len(rows) + 1

                # Calculate optimal column widths
                col_widths = calculate_column_widths(headers, rows)

                # Calculate table height
                table_height = TABLE_FORMAT['header_row_height'] + TABLE_FORMAT['data_row_height'] * len(rows)

                # Add table
                table_shape = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(POS['table']['left']),
                    Inches(POS['table']['top']),
                    Inches(POS['table']['width']),
                    Inches(table_height)
                )
                table = table_shape.table

                # Set column widths
                for col_idx, width in enumerate(col_widths):
                    table.columns[col_idx].width = Inches(width)

                # Populate header row
                for c_idx, header in enumerate(headers):
                    cell = table.cell(0, c_idx)
                    cell.text = header

                # Populate data rows
                for r_idx, row in enumerate(rows):
                    for c_idx, cell_val in enumerate(row):
                        if c_idx < num_cols:
                            cell = table.cell(r_idx + 1, c_idx)
                            cell.text = str(cell_val) if cell_val else ''

                # Apply formatting
                format_table(table)

        elif slide_type == 'data_chart':
            chart_data = content.get('chart_data', {})
            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = f"[Chart: {chart_data.get('type', 'N/A')}]"
                apply_font(run, FONTS['body'])

        elif slide_type == 'key_metrics':
            metrics = content.get('metrics', [])
            if metrics:
                # Add metrics as bullet points for now
                bullet_text = [f"{m.get('value', '')} - {m.get('label', '')}" for m in metrics]
                if content_shape:
                    add_standardized_bullets(content_shape.text_frame, bullet_text)

        elif slide_type == 'section_divider':
            # Add background image if available
            image_file = SECTION_IMAGE_MAP.get(section_name)
            if image_file:
                image_path = CROPPED_IMAGE_DIR / image_file
                if image_path.exists():
                    # Add as background (cover entire slide)
                    slide.shapes.add_picture(
                        str(image_path),
                        Inches(0), Inches(0),
                        Inches(11), Inches(8.5)
                    )
                    # Move to back
                    spTree = slide.shapes._spTree
                    pic = spTree[-1]
                    spTree.remove(pic)
                    spTree.insert(2, pic)

    # Save the presentation
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v23.pptx")
    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"Generated presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
