"""Add SEC-compliant source footnotes to presentation slides.

This script adds size 6, bottom-right-aligned footnotes to slides
that display specific data sourced from research publications.
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Source abbreviations for footnotes (SEC-friendly format)
SOURCES = {
    'rclco': 'RCLCO ODCE and NPI Results, Q2 2025',
    'clarion': 'Clarion Partners Industrial Real Estate Outlook, 2025',
    'commercialcafe': 'CommercialCafe National Industrial Report, December 2025',
    'cushman': 'Cushman & Wakefield Q2 2025 US Industrial MarketBeat',
    'cbre_cap': 'CBRE Cap Rate Survey, H1 2024',
    'cbre_rates': 'CBRE Impact of Interest Rate Cuts on Real Estate Cap Rates, 2024',
    'nashville': 'REBusinessOnline Nashville Industrial Market Analysis, 2024',
    'tampa': 'WareCRE Tampa Warehouse Market Report, 2025',
    'raleigh': 'Savills Raleigh-Durham Q4 2024 Industrial Market Report',
    'partners': 'Partners Real Estate Market Reports (DFW, Atlanta, San Antonio), 2024',
    'colliers': 'Colliers Phoenix Industrial Market Analysis, 2024',
    'kidder': 'Kidder Mathews Phoenix Industrial Market Report, 2024',
    'nasra': 'NASRA Public Pension Plan Investment Return Assumptions, FY 2024',
    'gresb': 'GRESB 2025 Real Estate Assessment Results',
    'naiop': 'NAIOP Nearshoring/Reshoring Analysis, 2024',
}

# Slide-to-source mapping (1-indexed slide numbers)
# Each slide lists the source keys that should be cited
SLIDE_SOURCES = {
    # Slide 3: Investment Thesis Summary (12.4% returns, 3.4% vacancy)
    3: ['rclco', 'commercialcafe'],

    # Slide 5: Repricing Context (Fed rates, 2.5B SF delivered)
    5: ['commercialcafe'],

    # Slide 6: US Industrial Market Overview table
    6: ['commercialcafe', 'cushman'],

    # Slide 7: Light Industrial vs Bulk Logistics comparison
    7: ['commercialcafe'],

    # Slide 8: Supply Pipeline Contraction
    8: ['commercialcafe'],

    # Slide 9: Cap Rate Evolution
    9: ['cbre_cap'],

    # Slide 11: Target Market Matrix table
    11: ['partners', 'nashville', 'tampa', 'raleigh', 'colliers'],

    # Slide 12: Tier 1 Deep Dive (Nashville, Tampa, Raleigh-Durham)
    12: ['nashville', 'tampa', 'raleigh'],

    # Slide 13: Tier 2 & 3 Markets
    13: ['partners', 'colliers', 'kidder'],

    # Slide 15: Structural Demand Drivers (e-commerce 23.2%, nearshoring)
    15: ['clarion', 'naiop'],

    # Slide 16: Manufacturing Construction Renaissance
    16: ['naiop'],

    # Slide 17: Property Type Interest Rate Sensitivity
    17: ['cbre_rates', 'rclco'],

    # Slide 23: Sensitivity Analysis (6.91% pension hurdle)
    23: ['nasra'],

    # Slide 24: Return Expectations by Strategy
    24: ['rclco'],

    # Slide 25: Historical Returns chart
    25: ['rclco'],

    # Slide 27: Institutional Industrial Landscape
    27: ['clarion'],

    # Slide 35: ESG Value Creation table
    35: ['gresb'],

    # Slide 36: ESG Integration & GRESB Roadmap
    36: ['gresb'],
}

# Footnote formatting
FOOTNOTE_FONT = {
    'name': 'Arial',
    'size': 6,
    'bold': False,
    'color': RGBColor(0x06, 0x1F, 0x32),  # Dark text
}

FOOTNOTE_POSITION = {
    'left': 0.40,
    'top': 7.90,  # Near bottom
    'width': 10.20,
    'height': 0.30,
}


def add_footnote_to_slide(slide, sources):
    """Add a footnote text box to a slide with source citations."""
    # Format source text
    source_keys = sources if isinstance(sources, list) else [sources]
    source_texts = [SOURCES[key] for key in source_keys if key in SOURCES]

    if not source_texts:
        return

    footnote_text = "Sources: " + "; ".join(source_texts) + "."

    # Check if there's already a footer placeholder at the bottom
    footer_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # Type 6 is FOOTER, Type 12 is DATE
            if ph_type in [6, 12]:
                # Check if near bottom of slide
                if shape.top.inches > 7.0:
                    footer_shape = shape
                    break
        elif hasattr(shape, 'top') and shape.top.inches > 7.5:
            # Found an existing text box near bottom
            if shape.has_text_frame:
                footer_shape = shape
                break

    # If no footer found, create a text box
    if footer_shape is None:
        footer_shape = slide.shapes.add_textbox(
            Inches(FOOTNOTE_POSITION['left']),
            Inches(FOOTNOTE_POSITION['top']),
            Inches(FOOTNOTE_POSITION['width']),
            Inches(FOOTNOTE_POSITION['height'])
        )

    # Set the footnote text
    tf = footer_shape.text_frame
    tf.clear()
    tf.word_wrap = True

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT

    run = para.add_run()
    run.text = footnote_text

    # Apply font formatting
    font = run.font
    font.name = FOOTNOTE_FONT['name']
    font.size = Pt(FOOTNOTE_FONT['size'])
    font.bold = FOOTNOTE_FONT['bold']
    font.color.rgb = FOOTNOTE_FONT['color']


def main():
    """Add footnotes to presentation."""
    print("Adding source footnotes to presentation...")
    print("=" * 60)

    # Load the v23 presentation
    input_path = Path("pptx_generator/output/Light_Industrial_Thesis_v23.pptx")
    if not input_path.exists():
        print(f"Error: Input not found: {input_path}")
        return

    prs = Presentation(str(input_path))
    print(f"Loaded presentation: {len(prs.slides)} slides")

    # Add footnotes to mapped slides
    footnotes_added = 0
    for slide_num, source_keys in SLIDE_SOURCES.items():
        if slide_num <= len(prs.slides):
            slide = prs.slides[slide_num - 1]  # Convert to 0-indexed
            add_footnote_to_slide(slide, source_keys)

            # Get slide title for logging
            title = "N/A"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            title = para.text[:40]
                            break
                    break

            print(f"  Slide {slide_num}: Added {len(source_keys)} source(s) - {title}...")
            footnotes_added += 1

    print(f"\nAdded footnotes to {footnotes_added} slides")

    # Save as v24
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v24.pptx")
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
