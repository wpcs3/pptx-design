"""Generate BTR presentations from Claude, ChatGPT, and Gemini outlines for comparison.

This script generates three versions of the BTR presentation using the same
pptx_generator codebase to enable fair comparison of LLM research and outline quality.
"""

import json
import copy
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image
import io
import re

# Import end module cloner
from pptx_generator.modules.end_module import append_end_module, remove_existing_end_module

# =============================================================================
# Configuration - per PCCP CS Style Guide (2026.01.01)
# =============================================================================

# Font settings - per PCCP CS Style Guide (2026.01.01)
FONTS = {
    # Title styles - BLACK per PCCP style guide
    'title_content': {'size': 32, 'bold': True, 'color': 'black'},  # Content slide titles - BLACK
    'title_section': {'size': 44, 'bold': True, 'color': 'white'},  # Section/cover titles
    # Subtitle styles - BLACK for content slides per PCCP style guide
    'subtitle_frontpage': {'size': 18, 'bold': True, 'color': 'white'},  # Front page only
    'subtitle_content': {'size': 20, 'bold': True, 'color': 'black'},  # Content slides - BLACK per style guide
    # Body styles
    'content_header': {'size': 18, 'bold': True, 'color': 'black'},  # Side-by-side headers - BLACK
    'body': {'size': 14, 'bold': False, 'color': 'text_primary'},  # Bullets - slate for body text
    # Footnote styles
    'source_footnote': {'size': 8, 'bold': False, 'color': 'footnote_gray'},  # Source citations - #A6A6A6
}

# Background image paths (extracted from BTR presentations)
BACKGROUND_IMAGES = {
    'frontpage': Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_01.png'),
    'section': [
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_04.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_10.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_15.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_20.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_26.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_33.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_38.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_44.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_49.png'),
        Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_53.png'),
    ],
    'end': Path('pptx_generator/assets/btr_backgrounds/btr_bg_slide_59.png'),
}

# White PCCP logo for front page (per PCCP style guide)
LOGO_CONFIG = {
    'white_logo_path': Path('logos/pccp_logo_white.png'),
    'front_page': {
        'left_inches': 0.4,
        'top_inches': 0.4,
        'width_inches': 1.8,
        'aspect_ratio': 2.382,
    }
}

# Section divider image counter (cycles through available images)
_section_image_idx = 0

# Slide dimensions (letter size)
SLIDE_WIDTH = 11.0
SLIDE_HEIGHT = 8.5

COLORS = {
    'slate_primary': RGBColor(0x2D, 0x37, 0x48),
    'slate_dark': RGBColor(0x1A, 0x20, 0x2C),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'off_white': RGBColor(0xF7, 0xFA, 0xFC),
    'light_gray': RGBColor(0xED, 0xF2, 0xF7),
    'border_gray': RGBColor(0xE2, 0xE8, 0xF0),
    'chart_blue': RGBColor(0x4A, 0x55, 0x68),
    'chart_teal': RGBColor(0x31, 0x97, 0x95),
    'chart_orange': RGBColor(0xDD, 0x6B, 0x20),
    'black': RGBColor(0x00, 0x00, 0x00),  # Pure black for titles/subtitles per PCCP style guide
    'text_primary': RGBColor(0x2D, 0x37, 0x48),  # Slate for body text
    'text_muted': RGBColor(0x71, 0x80, 0x96),  # Gray #718096 for section names
    'footnote_gray': RGBColor(0xA6, 0xA6, 0xA6),  # Medium gray #A6A6A6 for footnotes (per reference)
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
}

LAYOUTS = {
    'frontpage': 0,
    'section_title': 1,
    'bullet_content': 2,
    'chart': 3,
    'table': 4,
    'side_by_side': 5,
    'contact': 6,
    'end': 13,
}

PLACEHOLDERS = {
    'title': 0,
    'subtitle': 1,
    'logo': 14,
    'frontpage_subtitle': 15,
    'section_name': 17,
    'content': 18,
    'footnote': 20,
    'right_content': 21,
    'table': 22,
}

# =============================================================================
# Format Normalizers
# =============================================================================

def normalize_claude_outline(outline):
    """Claude outlines are already in correct format."""
    return outline

def normalize_gpt_outline(outline):
    """Normalize ChatGPT outline to match pptx_generator format."""
    normalized = {
        'presentation_type': outline.get('presentation_type', 'investment_pitch'),
        'title': 'Built-for-Rent Communities REIT',
        'subtitle': 'ChatGPT Research Output',
        'template': 'consulting_toolkit',
        'sections': []
    }

    # Group slides into sections
    current_section = {'name': 'Introduction', 'slides': []}

    for slide in outline.get('slides', []):
        slide_type = slide.get('slide_type', '')
        content = slide.get('content', {})

        # Check if this is a section divider
        if slide_type == 'section_divider':
            # Save current section if it has slides
            if current_section['slides']:
                normalized['sections'].append(current_section)
            # Start new section
            section_title = content.get('section_title', content.get('title', 'Section'))
            current_section = {'name': section_title, 'slides': [slide]}
            continue

        # Normalize chart data format
        if 'chart_type' in content and 'data' in content:
            content['chart_data'] = {
                'type': content.pop('chart_type', 'column'),
                'categories': content['data'].get('categories', []),
                'series': content['data'].get('series', [])
            }
            del content['data']

        # Normalize two-column format
        if slide_type == 'two_column':
            if 'left_title' in content:
                content['left_column'] = {
                    'header': content.pop('left_title', ''),
                    'bullets': content.pop('left_bullets', [])
                }
            if 'right_title' in content:
                content['right_column'] = {
                    'header': content.pop('right_title', ''),
                    'bullets': content.pop('right_bullets', [])
                }

        # Normalize table format
        if slide_type == 'table_slide':
            if 'rows' in content and 'data' not in content:
                content['data'] = content.pop('rows')

        normalized_slide = {
            'slide_type': slide_type,
            'content': content
        }
        current_section['slides'].append(normalized_slide)

    # Add final section
    if current_section['slides']:
        normalized['sections'].append(current_section)

    return normalized

def normalize_gemini_outline(thesis_data):
    """Convert Gemini's investment thesis JSON into presentation slides."""
    thesis = thesis_data.get('investment_thesis', thesis_data)

    normalized = {
        'presentation_type': 'investment_pitch',
        'title': thesis.get('title', 'Build-to-Rent Investment Thesis'),
        'subtitle': 'Gemini Research Output',
        'template': 'consulting_toolkit',
        'sections': []
    }

    # Executive Summary Section
    exec_summary = thesis.get('executive_summary', {})
    exec_section = {
        'name': 'Executive Summary',
        'slides': [
            {
                'slide_type': 'title_slide',
                'content': {
                    'title': thesis.get('title', 'Build-to-Rent Investment Thesis')[:50],
                    'subtitle': f"Target Raise: {thesis.get('target_fundraise', '$2-3B')} | {thesis.get('structure', 'Open-Ended REIT')}"
                }
            },
            {
                'slide_type': 'key_metrics',
                'content': {
                    'title': 'Investment Highlights',
                    'takeaway': exec_summary.get('core_premise', '')[:100],
                    'metrics': [
                        {'label': 'Target Raise', 'value': thesis.get('target_fundraise', '$2-3B')},
                        {'label': 'Rent Premium vs MF', 'value': exec_summary.get('financial_highlights', {}).get('rent_premium_vs_multifamily', '15-20%')},
                        {'label': 'Expense Ratio', 'value': exec_summary.get('financial_highlights', {}).get('expense_ratio', '28-32%')},
                        {'label': 'Tenant Retention', 'value': exec_summary.get('financial_highlights', {}).get('tenant_retention', '~68%')}
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(exec_section)

    # Macroeconomic Section
    macro = thesis.get('macroeconomic_imperatives', {})
    macro_section = {
        'name': 'Market Opportunity',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Market Opportunity', 'subtitle': 'The Structural Housing Crisis'}
            },
            {
                'slide_type': 'key_metrics',
                'content': {
                    'title': 'Housing Supply Deficit',
                    'takeaway': macro.get('chronic_supply_deficit', {}).get('description', ''),
                    'metrics': [
                        {'label': 'Supply Deficit', 'value': macro.get('chronic_supply_deficit', {}).get('shortage_estimate', '3.8-5.5M')},
                        {'label': 'Active Listings', 'value': macro.get('lock_in_effect', {}).get('active_listings', '1.4-1.5M')},
                        {'label': '10Y Treasury', 'value': macro.get('capital_market_dislocation', {}).get('10_year_treasury_forecast_2025', '4.0-4.25%')},
                        {'label': 'MF Starts Drop', 'value': macro.get('capital_market_dislocation', {}).get('multifamily_starts_2025', '-30.4%')}
                    ]
                }
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Affordability Collapse',
                    'takeaway': 'Cost to own significantly exceeds cost to rent in major markets',
                    'bullets': [
                        f"Phoenix savings to rent: {macro.get('affordability_collapse', {}).get('rent_vs_own_gap', {}).get('phoenix_az', '$1,150/month')}",
                        f"National avg savings: {macro.get('affordability_collapse', {}).get('rent_vs_own_gap', {}).get('national_avg', '$213/month')}",
                        macro.get('affordability_collapse', {}).get('consumer_behavior', '')[:100],
                        macro.get('lock_in_effect', {}).get('description', '')[:100],
                        macro.get('capital_market_dislocation', {}).get('implication', '')[:100]
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(macro_section)

    # Demographics Section
    demo = thesis.get('demographic_tailwinds', {})
    demo_section = {
        'name': 'Demographic Drivers',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Demographic Drivers', 'subtitle': 'Convergence of Generations'}
            },
            {
                'slide_type': 'two_column',
                'content': {
                    'title': 'Dual Demographic Demand',
                    'takeaway': 'Millennials and Baby Boomers driving BTR demand',
                    'left_column': {
                        'header': 'Millennials',
                        'bullets': [
                            demo.get('millennials', {}).get('description', ''),
                            demo.get('millennials', {}).get('drivers', '')
                        ]
                    },
                    'right_column': {
                        'header': 'Baby Boomers',
                        'bullets': [
                            demo.get('baby_boomers', {}).get('description', ''),
                            demo.get('baby_boomers', {}).get('preference', '')
                        ]
                    }
                }
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Migration Trends',
                    'takeaway': f"Primary region: {demo.get('migration_trends', {}).get('primary_region', 'Sun Belt')}",
                    'bullets': [
                        f"Primary region: {demo.get('migration_trends', {}).get('primary_region', '')}",
                        f"Secondary shift: {demo.get('migration_trends', {}).get('secondary_shift', '')}",
                        demo.get('gen_z', {}).get('description', ''),
                        demo.get('gen_z', {}).get('preference', '')
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(demo_section)

    # Product Strategy Section
    product = thesis.get('product_strategy', {})
    product_section = {
        'name': 'Product Strategy',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Product Strategy', 'subtitle': 'Purpose-Built Communities'}
            },
            {
                'slide_type': 'key_metrics',
                'content': {
                    'title': 'Community Design',
                    'takeaway': f"Target community size: {product.get('community_scale', '100-200 homes')}",
                    'metrics': [
                        {'label': 'Community Size', 'value': product.get('community_scale', '100-200 homes').split('(')[0].strip()},
                        {'label': '4-BR Allocation', 'value': product.get('unit_mix_strategy', {}).get('4_bedroom_allocation', '20-30%')},
                        {'label': 'Pet Ownership', 'value': '70%'},
                        {'label': 'Unit Types', 'value': '3 Types'}
                    ]
                }
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Class A Amenities',
                    'takeaway': 'Amenities differentiate BTR from scattered SFR',
                    'bullets': product.get('amenities', [])[:6]
                }
            }
        ]
    }
    normalized['sections'].append(product_section)

    # Operating Metrics Section
    ops = thesis.get('operational_metrics', {})
    ops_section = {
        'name': 'Operating Fundamentals',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Operating Fundamentals', 'subtitle': 'Efficiency and Margins'}
            },
            {
                'slide_type': 'two_column',
                'content': {
                    'title': 'Expense Ratio Comparison',
                    'takeaway': 'BTR achieves superior operating efficiency',
                    'left_column': {
                        'header': 'Expense Ratios',
                        'bullets': [
                            f"BTR Community: {ops.get('expense_ratios', {}).get('btr_community', '28-32%')}",
                            f"Multifamily Avg: {ops.get('expense_ratios', {}).get('multifamily_avg', '35-45%')}",
                            f"Scattered SFR: {ops.get('expense_ratios', {}).get('scattered_sfr', '40-50%')}"
                        ]
                    },
                    'right_column': {
                        'header': 'Retention Economics',
                        'bullets': [
                            f"BTR Retention: {ops.get('turnover_economics', {}).get('btr_retention', '~68%')}",
                            f"Turnover Cost: {ops.get('turnover_economics', {}).get('turnover_cost', '$2,000-$4,000')}",
                            f"Longer Stay: {ops.get('turnover_economics', {}).get('lease_duration', '25-35% longer')}"
                        ]
                    }
                }
            }
        ]
    }
    normalized['sections'].append(ops_section)

    # Market Selection Section
    markets = thesis.get('market_selection', {})
    market_section = {
        'name': 'Target Markets',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Target Markets', 'subtitle': markets.get('strategy', 'Barbell Approach')}
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Primary Growth Markets',
                    'takeaway': 'Sun Belt markets with strong fundamentals',
                    'bullets': [
                        'Phoenix, AZ - BTR epicenter with strong absorption',
                        'Dallas/Fort Worth, TX - Diversified economy, 18k+ pipeline',
                        'Tampa/Jacksonville, FL - Strong in-migration',
                        'Atlanta, GA - Continued suburban sprawl'
                    ]
                }
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Secondary Markets',
                    'takeaway': 'Midwest stability and yield',
                    'bullets': [
                        'Columbus, OH - Intel investment driving growth',
                        'Indianapolis, IN - Tight vacancy, consistent rent growth',
                        'Charlotte/Raleigh, NC - Research Triangle demand'
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(market_section)

    # Investment Structure Section
    vehicle = thesis.get('investment_vehicle', {})
    struct_section = {
        'name': 'Investment Structure',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Investment Structure', 'subtitle': 'Open-Ended Private REIT'}
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'REIT Benefits for RIAs',
                    'takeaway': 'Structure designed for wealth management channel',
                    'bullets': vehicle.get('benefits_for_rias', [])[:6]
                }
            }
        ]
    }
    normalized['sections'].append(struct_section)

    # Development Economics Section
    dev = thesis.get('development_economics', {})
    dev_section = {
        'name': 'Financial Projections',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Financial Projections', 'subtitle': 'Development Economics'}
            },
            {
                'slide_type': 'key_metrics',
                'content': {
                    'title': 'Development Spread',
                    'takeaway': 'Value creation through development',
                    'metrics': [
                        {'label': 'Market Cap Rate', 'value': dev.get('development_spread', {}).get('market_cap_rate', '4.75-5.50%')},
                        {'label': 'Target YOC', 'value': dev.get('development_spread', {}).get('target_yield_on_cost', '6.50-7.00%')},
                        {'label': 'Spread', 'value': dev.get('development_spread', {}).get('spread_target', '150-225 bps')},
                        {'label': 'Hard Costs', 'value': dev.get('construction_inputs', {}).get('hard_costs', '$150-200 psf')}
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(dev_section)

    # Risk Management Section
    risk = thesis.get('risk_management', {})
    risk_section = {
        'name': 'Risk Factors',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Risk Factors', 'subtitle': 'And Mitigation Strategies'}
            },
            {
                'slide_type': 'title_content',
                'content': {
                    'title': 'Risk Management',
                    'takeaway': 'Disciplined approach to risk mitigation',
                    'bullets': [
                        f"Supply Overshoot: {risk.get('supply_overshoot', '')}",
                        f"Interest Rates: {risk.get('interest_rates', '')}",
                        f"Exit Strategies: {risk.get('exit_strategies', '')}"
                    ]
                }
            }
        ]
    }
    normalized['sections'].append(risk_section)

    # Conclusion Section
    conclusion_section = {
        'name': 'Conclusion',
        'slides': [
            {
                'slide_type': 'section_divider',
                'content': {'title': 'Investment Summary', 'subtitle': 'A Structural Opportunity'}
            },
            {
                'slide_type': 'key_metrics',
                'content': {
                    'title': 'Investment Summary',
                    'takeaway': 'Compelling risk-adjusted returns',
                    'metrics': [
                        {'label': 'Target Raise', 'value': thesis.get('target_fundraise', '$2-3B')},
                        {'label': 'Structure', 'value': 'Open-Ended REIT'},
                        {'label': 'Focus', 'value': 'BTR Communities'},
                        {'label': 'Markets', 'value': 'Sun Belt + Midwest'}
                    ]
                }
            },
            # Note: "Thank You" slide removed - Contact info is in the end module
        ]
    }
    normalized['sections'].append(conclusion_section)

    return normalized


# =============================================================================
# Slide Building Functions
# =============================================================================

def get_placeholder(slide, idx):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == idx:
            return shape
    return None

def set_placeholder_text(slide, idx, text, font_size=None, bold=None, color=None):
    shape = get_placeholder(slide, idx)
    if shape and hasattr(shape, 'text_frame'):
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = str(text) if text else ''
        run.font.name = 'Arial'
        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return True
    return False

def add_section_name(slide, section_name):
    if not section_name:
        return
    shape = get_placeholder(slide, PLACEHOLDERS['section_name'])
    if shape:
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        run = para.add_run()
        run.text = section_name
        run.font.name = 'Arial'
        run.font.size = Pt(9)
        run.font.color.rgb = COLORS['text_muted']

def set_cell_borders_none(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ['lnL', 'lnR', 'lnT', 'lnB', 'lnTlToBr', 'lnBlToTr']:
        for border in tcPr.findall(qn(f'a:{border_name}')):
            tcPr.remove(border)
        border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
        border.set('w', '0')
        etree.SubElement(border, qn('a:noFill'))

def is_numeric_cell(text):
    if not text:
        return False
    text = str(text).strip()
    patterns = [r'^[\d,]+\.?\d*%?$', r'^\$[\d,]+\.?\d*[BMK]?$', r'^[\d,]+\.?\d*x$', r'^~?\d+%?$']
    return any(re.match(p, text) for p in patterns)

def get_column_alignment(data, col_idx):
    numeric = sum(1 for row in data if col_idx < len(row) and is_numeric_cell(str(row[col_idx])))
    text = sum(1 for row in data if col_idx < len(row) and str(row[col_idx]).strip() and not is_numeric_cell(str(row[col_idx])))
    return PP_ALIGN.RIGHT if numeric > text else PP_ALIGN.LEFT


def fix_background_image_aspect_ratio(picture_shape):
    """Apply srcRect cropping to maintain aspect ratio for 16:9 images on letter-size slides.

    For 16:9 images (1.778) on letter-size slides (1.294), crops the sides
    to prevent vertical stretching.
    """
    try:
        # Get the image blob and calculate aspect ratio
        image_blob = picture_shape.image.blob
        img = Image.open(io.BytesIO(image_blob))
        orig_width, orig_height = img.size
        image_aspect = orig_width / orig_height
        slide_aspect = SLIDE_WIDTH / SLIDE_HEIGHT  # 1.294 for letter

        # If image is wider than slide (16:9 = 1.778 > letter = 1.294), crop sides
        if image_aspect > slide_aspect + 0.05:
            crop_ratio = 1 - (slide_aspect / image_aspect)
            crop_pct = int((crop_ratio / 2) * 100000)  # EMU percentage

            pic_elem = picture_shape._element

            # Find blipFill - try multiple approaches
            blipFill = None
            for child in pic_elem:
                if child.tag.endswith('}blipFill') or child.tag == 'blipFill':
                    blipFill = child
                    break

            if blipFill is not None:
                # Remove existing srcRect if present
                for child in list(blipFill):
                    if child.tag.endswith('}srcRect'):
                        blipFill.remove(child)

                # Create new srcRect with cropping
                srcRect = etree.SubElement(blipFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
                srcRect.set('l', str(crop_pct))
                srcRect.set('r', str(crop_pct))
                srcRect.set('t', '0')
                srcRect.set('b', '0')
                return True
    except Exception as e:
        print(f"    Warning: Could not fix aspect ratio: {e}")

    return False


def add_background_image(slide, image_path, slide_width_emu, slide_height_emu):
    """Add a full-bleed background image to a slide and apply aspect ratio correction."""
    if not image_path or not image_path.exists():
        return None

    # Add the image as full-bleed (covers entire slide)
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        width=slide_width_emu,
        height=slide_height_emu
    )

    # Move to back (behind all other shapes)
    spTree = slide.shapes._spTree
    pic_elem = picture._element
    spTree.remove(pic_elem)
    spTree.insert(2, pic_elem)  # Insert after background but before other shapes

    # Apply aspect ratio correction
    fix_background_image_aspect_ratio(picture)

    return picture


def add_front_page_logo(slide):
    """Add the white PCCP logo to the front page (top-left, on top of all layers).

    The logo must be added AFTER the background image to ensure it appears on top.
    """
    logo_path = LOGO_CONFIG['white_logo_path']
    if not logo_path.exists():
        print(f"  Warning: Logo not found at {logo_path}")
        return None

    fp_config = LOGO_CONFIG['front_page']
    left = Inches(fp_config['left_inches'])
    top = Inches(fp_config['top_inches'])
    width = Inches(fp_config['width_inches'])
    # Calculate height from width and aspect ratio
    height = Inches(fp_config['width_inches'] / fp_config['aspect_ratio'])

    # Add logo - since this is called AFTER background, it will be on top
    picture = slide.shapes.add_picture(
        str(logo_path),
        left, top,
        width=width,
        height=height
    )

    return picture


def get_next_section_image():
    """Get the next section divider background image, cycling through available images."""
    global _section_image_idx
    images = BACKGROUND_IMAGES['section']
    if not images:
        return None
    img_path = images[_section_image_idx % len(images)]
    _section_image_idx += 1
    return img_path


def add_source_footnote(slide, sources, slide_type='content'):
    """Add source citation footnote to a slide.

    Args:
        slide: The slide to add footnote to
        sources: Source text or list of sources
        slide_type: Type of slide ('content', 'chart', 'table')
    """
    if not sources:
        return None

    # Format sources text
    if isinstance(sources, list):
        source_text = "Sources: " + "; ".join(sources) + "."
    else:
        source_text = sources if sources.startswith("Sources:") else f"Sources: {sources}."

    # Position: left-aligned at bottom per style guide (0.4", 7.5")
    # Text anchored to BOTTOM of text box so it appears closer to slide bottom line
    footnote = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.5),
        Inches(10.2), Inches(0.4)
    )
    tf = footnote.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # Don't auto-size, keep fixed height

    # Set vertical anchor to BOTTOM - text aligns to bottom of text box
    from pptx.enum.text import MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    # Set anchor via XML since python-pptx doesn't expose it directly on text_frame
    from pptx.oxml.ns import qn as oxml_qn
    body_pr = tf._txBody.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if body_pr is not None:
        body_pr.set('anchor', 'b')  # 'b' = bottom, 't' = top, 'ctr' = center

    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = source_text
    run.font.name = 'Arial'
    run.font.size = Pt(FONTS['source_footnote']['size'])
    run.font.bold = FONTS['source_footnote']['bold']
    run.font.color.rgb = COLORS[FONTS['source_footnote']['color']]

    return footnote


def get_default_sources(section_name, slide_type):
    """Generate default source citations based on section and slide type."""
    # Map section names to relevant sources
    source_map = {
        'Executive Summary': 'John Burns Research; CBRE Industrial & Logistics Report Q4 2025',
        'Market Opportunity': 'Census Bureau; Freddie Mac Housing Survey 2025; NAR',
        'Housing Supply': 'Census Bureau; Freddie Mac Housing Survey 2025',
        'Affordability': 'NAR; Federal Reserve; Zillow Research',
        'Demographic': 'Census Bureau; Pew Research; NAR Profile of Home Buyers 2024',
        'Product': 'CBRE BFR Survey 2024; John Burns Research',
        'Operating': 'CBRE; Company filings; Green Street Advisors',
        'Target Markets': 'CBRE; Census Bureau; U-Haul Migration Trends',
        'Investment': 'Company filings; CBRE; Green Street Advisors',
        'Financial': 'CBRE; Green Street Advisors; Company filings',
        'Risk': 'CBRE; Federal Reserve; Company risk assessments',
        'Competitive': 'AMH 10-K Filing 2024; INVH 10-K Filing 2024; CBRE',
        'Conclusion': 'PCCP Management Estimates',
    }

    # Find matching sources based on section name keywords
    for key, sources in source_map.items():
        if key.lower() in section_name.lower():
            return sources

    # Default source for unmatched sections
    return 'CBRE; John Burns Research; PCCP Management Estimates'


def build_title_slide(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['frontpage']]
    slide = prs.slides.add_slide(layout)

    # Add background image (full-bleed)
    bg_path = BACKGROUND_IMAGES.get('frontpage')
    if bg_path and bg_path.exists():
        add_background_image(slide, bg_path, prs.slide_width, prs.slide_height)

    # Add white PCCP logo (must be AFTER background to appear on top)
    add_front_page_logo(slide)

    # Title: 44pt bold white (per style guide title_section)
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=44, bold=True, color=COLORS['white'])
    subtitle = content.get('subtitle', '')
    current_date = datetime.now().strftime('%B %Y')
    subtitle_with_date = f"{subtitle}\n{current_date}" if subtitle else current_date
    # Front page subtitle: 18pt bold white (per style guide subtitle_frontpage)
    set_placeholder_text(slide, PLACEHOLDERS['frontpage_subtitle'], subtitle_with_date,
                         font_size=18, bold=True, color=COLORS['white'])
    return slide

def build_section_divider(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['section_title']]
    slide = prs.slides.add_slide(layout)

    # Add background image (full-bleed, cycles through available section images)
    bg_path = get_next_section_image()
    if bg_path and bg_path.exists():
        add_background_image(slide, bg_path, prs.slide_width, prs.slide_height)

    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', section_name),
                         font_size=44, bold=True, color=COLORS['white'])
    return slide

def build_key_metrics(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['bullet_content']]
    slide = prs.slides.add_slide(layout)
    metrics = content.get('metrics', [])

    # Title: 32pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=32, bold=True, color=COLORS['black'])
    # Subtitle (thesis/takeaway): 20pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['subtitle'], content.get('takeaway', ''),
                         font_size=20, bold=True, color=COLORS['black'])

    # Clear content placeholder
    content_shape = get_placeholder(slide, PLACEHOLDERS['content'])
    if content_shape:
        sp = content_shape._element
        sp.getparent().remove(sp)

    # Add metric boxes
    if metrics:
        num_metrics = min(len(metrics), 4)
        box_width = Inches(2.3)
        box_height = Inches(2.0)
        total_width = Inches(10.0)
        spacing = (total_width - box_width * num_metrics) / (num_metrics + 1)
        start_x = Inches(0.5)
        y = Inches(3.0)

        for i, metric in enumerate(metrics[:4]):
            x = start_x + spacing * (i + 1) + box_width * i

            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS['accent_blue']
            box.line.fill.background()

            value_text = str(metric.get('value', ''))
            value_font_size = 32 if len(value_text) <= 5 else (28 if len(value_text) <= 8 else 20)

            val_box = slide.shapes.add_textbox(x, y + Inches(0.2), box_width, Inches(0.75))
            tf = val_box.text_frame
            tf.word_wrap = False
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = value_text
            run.font.name = 'Arial'
            run.font.size = Pt(value_font_size)
            run.font.bold = True
            run.font.color.rgb = COLORS['white']

            lbl_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.0), box_width - Inches(0.2), Inches(0.9))
            tf = lbl_box.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = metric.get('label', '')
            run.font.name = 'Arial'
            run.font.size = Pt(11)
            run.font.color.rgb = COLORS['white']

    add_section_name(slide, section_name)

    # Add source footnote
    sources = content.get('sources', content.get('source', get_default_sources(section_name, 'key_metrics')))
    add_source_footnote(slide, sources)

    return slide

def build_title_content(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['bullet_content']]
    slide = prs.slides.add_slide(layout)
    bullets = content.get('bullets', [])

    # Title: 32pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=32, bold=True, color=COLORS['black'])
    # Subtitle (thesis/takeaway): 20pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['subtitle'], content.get('takeaway', ''),
                         font_size=20, bold=True, color=COLORS['black'])

    content_shape = get_placeholder(slide, PLACEHOLDERS['content'])
    if content_shape and bullets:
        tf = content_shape.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = str(bullet) if bullet else ''
            para.font.name = 'Arial'
            para.font.size = Pt(14)
            para.font.color.rgb = COLORS['text_primary']
            para.level = 0

    add_section_name(slide, section_name)

    # Add source footnote
    sources = content.get('sources', content.get('source', get_default_sources(section_name, 'title_content')))
    add_source_footnote(slide, sources)

    return slide

def build_two_column(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['side_by_side']]
    slide = prs.slides.add_slide(layout)

    left_col = content.get('left_column', {})
    right_col = content.get('right_column', {})

    # Title: 32pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=32, bold=True, color=COLORS['black'])
    # Subtitle (thesis/takeaway): 20pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['subtitle'], content.get('takeaway', ''),
                         font_size=20, bold=True, color=COLORS['black'])

    def format_column(shape, col_data):
        if not shape:
            return
        tf = shape.text_frame
        tf.clear()
        # Column header: 18pt bold (per style guide content_header)
        para = tf.paragraphs[0]
        para.text = col_data.get('header', '')
        para.font.name = 'Arial'
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = COLORS['black']  # Content headers: BLACK per PCCP style guide

        # Bullets: 14pt regular (per style guide body)
        for bullet in col_data.get('bullets', []):
            para = tf.add_paragraph()
            para.text = str(bullet) if bullet else ''
            para.font.name = 'Arial'
            para.font.size = Pt(14)
            para.font.color.rgb = COLORS['text_primary']
            para.level = 0

    format_column(get_placeholder(slide, PLACEHOLDERS['content']), left_col)
    format_column(get_placeholder(slide, PLACEHOLDERS['right_content']), right_col)

    add_section_name(slide, section_name)

    # Add source footnote
    sources = content.get('sources', content.get('source', get_default_sources(section_name, 'two_column')))
    add_source_footnote(slide, sources)

    return slide

def format_chart_pccp_style(chart):
    """Apply PCCP style guide formatting to a chart.

    Per PCCP CS Style Guide:
    - All text (axis labels, data labels, legends): 14pt Arial
    - Horizontal and vertical major gridlines: 0.5pt width, light gray (#E2E8F0)
    - Bar/column fill: #309CE7 (accent_blue)
    - No tick marks on either axis
    """
    # Chart style constants
    CHART_FONT_SIZE = Pt(14)
    CHART_FONT_NAME = 'Arial'
    GRIDLINE_WIDTH = 6350  # 0.5pt in EMU
    BAR_FILL_COLOR = RGBColor(0x30, 0x9C, 0xE7)  # #309CE7
    GRIDLINE_COLOR = RGBColor(0xE2, 0xE8, 0xF0)  # Light gray
    SECONDARY_COLORS = [
        RGBColor(0x31, 0x97, 0x95),  # Teal
        RGBColor(0xDD, 0x6B, 0x20),  # Orange
        RGBColor(0x2D, 0x37, 0x48),  # Slate
    ]

    try:
        # Format series fill colors
        for idx, series in enumerate(chart.series):
            series.format.fill.solid()
            if idx == 0:
                series.format.fill.fore_color.rgb = BAR_FILL_COLOR
            elif idx < len(SECONDARY_COLORS) + 1:
                series.format.fill.fore_color.rgb = SECONDARY_COLORS[idx - 1]

        # Format value axis (vertical for column charts, horizontal for bar charts)
        try:
            vax = chart.value_axis
            # Enable major gridlines
            vax.has_major_gridlines = True
            # Remove tick marks
            vax.major_tick_mark = XL_TICK_MARK.NONE
            vax.minor_tick_mark = XL_TICK_MARK.NONE
            # Set tick label font
            vax.tick_labels.font.size = CHART_FONT_SIZE
            vax.tick_labels.font.name = CHART_FONT_NAME

            # Format gridlines (width and color)
            if vax.has_major_gridlines:
                gridlines = vax.major_gridlines
                line = gridlines.format.line
                line.width = GRIDLINE_WIDTH
                line.fill.solid()
                line.fill.fore_color.rgb = GRIDLINE_COLOR
        except Exception as e:
            pass

        # Format category axis (horizontal for column charts, vertical for bar charts)
        try:
            cax = chart.category_axis
            # Enable major gridlines
            cax.has_major_gridlines = True
            # Remove tick marks
            cax.major_tick_mark = XL_TICK_MARK.NONE
            cax.minor_tick_mark = XL_TICK_MARK.NONE
            # Set tick label font
            cax.tick_labels.font.size = CHART_FONT_SIZE
            cax.tick_labels.font.name = CHART_FONT_NAME

            # Format gridlines (width and color)
            if cax.has_major_gridlines:
                gridlines = cax.major_gridlines
                line = gridlines.format.line
                line.width = GRIDLINE_WIDTH
                line.fill.solid()
                line.fill.fore_color.rgb = GRIDLINE_COLOR
        except Exception as e:
            pass

        # Format legend if present
        if chart.has_legend:
            try:
                chart.legend.font.size = CHART_FONT_SIZE
                chart.legend.font.name = CHART_FONT_NAME
            except:
                pass

        # Format data labels if present
        for series in chart.series:
            if series.has_data_labels:
                try:
                    series.data_labels.font.size = CHART_FONT_SIZE
                    series.data_labels.font.name = CHART_FONT_NAME
                except:
                    pass

    except Exception as e:
        pass


def build_data_chart(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['chart']]
    slide = prs.slides.add_slide(layout)
    chart_spec = content.get('chart_data', {})

    # Title: 32pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=32, bold=True, color=COLORS['black'])
    # Subtitle (thesis/takeaway): 20pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['subtitle'], content.get('takeaway', ''),
                         font_size=20, bold=True, color=COLORS['black'])

    if chart_spec:
        chart_type_map = {
            'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'bar': XL_CHART_TYPE.BAR_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
        }
        chart_type = chart_type_map.get(chart_spec.get('type', 'column'), XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        chart_data.categories = chart_spec.get('categories', [])
        for series in chart_spec.get('series', []):
            chart_data.add_series(series.get('name', ''), tuple(series.get('values', [])))

        chart_shape = slide.shapes.add_chart(
            chart_type, Inches(0.4), Inches(2.5), Inches(10.2), Inches(4.5), chart_data
        )
        chart = chart_shape.chart

        # Show legend only for multi-series charts
        chart.has_legend = len(chart_spec.get('series', [])) > 1

        # Apply PCCP style guide formatting
        format_chart_pccp_style(chart)

    add_section_name(slide, section_name)

    # Add source footnote
    sources = content.get('sources', content.get('source', get_default_sources(section_name, 'data_chart')))
    add_source_footnote(slide, sources)

    return slide

def build_table_slide(prs, content, section_name):
    layout = prs.slide_masters[0].slide_layouts[LAYOUTS['table']]
    slide = prs.slides.add_slide(layout)

    headers = content.get('headers', [])
    data = content.get('data', [])

    # Title: 32pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['title'], content.get('title', ''),
                         font_size=32, bold=True, color=COLORS['black'])
    # Subtitle (thesis/takeaway): 20pt bold BLACK per PCCP style guide
    set_placeholder_text(slide, PLACEHOLDERS['subtitle'], content.get('takeaway', ''),
                         font_size=20, bold=True, color=COLORS['black'])

    # Clear placeholders
    for idx in [PLACEHOLDERS['content'], PLACEHOLDERS['table']]:
        shape = get_placeholder(slide, idx)
        if shape:
            sp = shape._element
            sp.getparent().remove(sp)

    if headers and data:
        rows = len(data) + 1
        cols = len(headers)
        row_height = Inches(0.45)

        table_shape = slide.shapes.add_table(
            rows, cols, Inches(0.4), Inches(2.5), Inches(10.2), row_height * rows
        )
        table = table_shape.table

        alignments = [get_column_alignment(data, i) for i in range(cols)]
        col_width = int(Inches(10.2) / cols)
        for col in table.columns:
            col.width = col_width

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['slate_primary']
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for para in cell.text_frame.paragraphs:
                para.alignment = alignments[col_idx]
                for run in para.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(16)  # 16pt for table headers
                    run.font.bold = True
                    run.font.color.rgb = COLORS['white']

        # Data rows
        for row_idx, row_data in enumerate(data):
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['off_white'] if row_idx % 2 == 1 else COLORS['white']
                for para in cell.text_frame.paragraphs:
                    para.alignment = alignments[col_idx]
                    for run in para.runs:
                        run.font.name = 'Arial'
                        run.font.size = Pt(14)  # 14pt for table body
                        run.font.color.rgb = COLORS['text_primary']

        for row in table.rows:
            row.height = row_height
            for cell in row.cells:
                set_cell_borders_none(cell)
                cell.margin_left = Inches(0.1)
                cell.margin_right = Inches(0.1)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)

    add_section_name(slide, section_name)

    # Add source footnote
    sources = content.get('sources', content.get('source', get_default_sources(section_name, 'table_slide')))
    add_source_footnote(slide, sources)

    return slide

def build_slide(prs, slide_data, section_name):
    slide_type = slide_data.get('slide_type', 'title_content')
    content = slide_data.get('content', {})

    builders = {
        'title_slide': lambda: build_title_slide(prs, content, section_name),
        'section_divider': lambda: build_section_divider(prs, content, section_name),
        'key_metrics': lambda: build_key_metrics(prs, content, section_name),
        'title_content': lambda: build_title_content(prs, content, section_name),
        'two_column': lambda: build_two_column(prs, content, section_name),
        'data_chart': lambda: build_data_chart(prs, content, section_name),
        'table_slide': lambda: build_table_slide(prs, content, section_name),
    }

    builder = builders.get(slide_type, builders['title_content'])
    return builder()


# =============================================================================
# Main Generation Function
# =============================================================================

def generate_presentation(outline, template_path, output_path, llm_name):
    """Generate a presentation from an outline."""
    global _section_image_idx
    _section_image_idx = 0  # Reset section image counter for each presentation

    print(f"\n{'='*60}")
    print(f"Generating: {llm_name} BTR Presentation")
    print(f"{'='*60}")

    # Load template
    prs = Presentation(str(template_path))

    # Clear existing slides
    while len(prs.slides) > 0:
        slide = prs.slides[0]
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Build slides
    slide_count = 0
    for section in outline.get('sections', []):
        section_name = section.get('name', '')
        print(f"\n  {section_name}:")

        for slide_data in section.get('slides', []):
            # Skip "Thank You" slides - Contact info is in the end module
            slide_title = slide_data.get('content', {}).get('title', '')
            if 'thank you' in slide_title.lower():
                continue

            slide = build_slide(prs, slide_data, section_name)
            slide_count += 1
            slide_type = slide_data.get('slide_type', '')
            title = slide_title[:30]
            print(f"    {slide_count}: {slide_type} - {title}")

    # Add end module (Contact, Disclosures, End) - cloned from template
    print(f"\n  End Module:")
    end_slides_added = append_end_module(prs, verbose=True)
    slide_count += end_slides_added

    # Post-process: Fix background image aspect ratios
    print(f"\n  Post-processing background images...")
    bg_fixed = 0
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Check if this is a large image (>80% of slide = background)
                if shape.width > slide_width_emu * 0.8 and shape.height > slide_height_emu * 0.8:
                    if fix_background_image_aspect_ratio(shape):
                        bg_fixed += 1
    if bg_fixed > 0:
        print(f"    Fixed {bg_fixed} background images")

    # Set presentation metadata (title matches filename for PDF export)
    prs_title = output_path.stem  # e.g., "Claude_BTR_Comparison"
    prs.core_properties.title = prs_title
    prs.core_properties.subject = f"{llm_name} Build-for-Rent Single-Family Comparison"
    prs.core_properties.author = "PCCP, LLC"

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\n  Saved: {output_path}")
    print(f"  Title: {prs_title}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


def main():
    """Generate BTR presentations from all three LLM outlines."""
    print("="*70)
    print("BTR Presentation Comparison Generator")
    print("Generating from Claude, ChatGPT, and Gemini research outputs")
    print("="*70)

    # Paths
    template_path = Path("pptx_generator/output/light_industrial/Light_Industrial_Thesis_v27_CS_edits.pptx")
    base_output = Path("pptx_generator/output/btr")

    # Check for template in alternate locations
    if not template_path.exists():
        alt_template = Path("pptx_generator/output/Light_Industrial_Thesis_v27_CS_edits.pptx")
        if alt_template.exists():
            template_path = alt_template
        else:
            print(f"Error: Template not found at {template_path}")
            return

    # LLM configurations
    llm_configs = [
        {
            'name': 'Claude',
            'outline_path': Path("cc_prompts/claude_btr_presentation_outline.json"),
            'output_folder': 'claude_btr',
            'normalizer': normalize_claude_outline,
        },
        {
            'name': 'ChatGPT',
            'outline_path': Path("cc_prompts/gpt_btr_presentation_outline.json"),
            'output_folder': 'gpt_btr',
            'normalizer': normalize_gpt_outline,
        },
        {
            'name': 'Gemini',
            'outline_path': Path("cc_prompts/gemini_btr_presentation_outline_fixed.json"),
            'output_folder': 'gemini_btr',
            'normalizer': normalize_gemini_outline,
        },
    ]

    results = []

    for config in llm_configs:
        outline_path = config['outline_path']

        if not outline_path.exists():
            print(f"\nWarning: {config['name']} outline not found: {outline_path}")
            continue

        # Load and normalize outline
        with open(outline_path, encoding='utf-8') as f:
            raw_outline = json.load(f)

        outline = config['normalizer'](raw_outline)

        # Generate presentation
        output_folder = base_output / config['output_folder']
        output_file = output_folder / f"{config['name']}_BTR_Comparison.pptx"

        try:
            output_path = generate_presentation(
                outline, template_path, output_file, config['name']
            )
            results.append({
                'name': config['name'],
                'path': output_path,
                'slides': len(outline.get('sections', [])),
                'status': 'Success'
            })
        except Exception as e:
            print(f"\nError generating {config['name']} presentation: {e}")
            results.append({
                'name': config['name'],
                'path': None,
                'slides': 0,
                'status': f'Error: {e}'
            })

    # Summary
    print("\n" + "="*70)
    print("Generation Summary")
    print("="*70)

    for result in results:
        print(f"\n{result['name']}:")
        print(f"  Status: {result['status']}")
        if result['path']:
            print(f"  Output: {result['path']}")

    print("\n" + "="*70)
    print("Comparison presentations generated in:")
    print(f"  {base_output}")
    print("="*70)


if __name__ == "__main__":
    main()
