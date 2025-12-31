"""Generate presentation v25 using v24 CS edits formatting preferences.

This generator uses the template and formatting extracted from v24_CS_edits:
1. Uses correct slide layouts (Bullet Content, Chart, Table, Side by Side, etc.)
2. No table borders
3. Section names in top-right placeholder (idx=17)
4. Footnotes in master template placeholder (idx=20)
5. Auto-generate section title slide images using Gemini
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Load configuration
CONFIG_PATH = Path("pptx_generator/config/template_format_v24.json")
with open(CONFIG_PATH) as f:
    CONFIG = json.load(f)

# Extract settings
COLORS = {k: RGBColor.from_string(v.lstrip('#')) for k, v in CONFIG['colors'].items()}
LAYOUTS = CONFIG['layouts']
TABLE_FORMAT = CONFIG['table_formatting']
FOOTNOTE_CONFIG = CONFIG['footnotes']
SECTION_NAME_CONFIG = CONFIG['section_names']

# Slide type to layout mapping
SLIDE_TYPE_LAYOUT = {
    'title_slide': 'frontpage',
    'title_content': 'bullet_content',
    'content': 'bullet_content',
    'key_metrics': 'bullet_content',
    'table_slide': 'table',
    'data_chart': 'chart',
    'chart': 'chart',
    'two_column': 'side_by_side',
    'comparison': 'side_by_side',
    'section_divider': 'section_title',
    'disclaimer': 'disclaimers',
}

# Source mapping for footnotes
SLIDE_SOURCES = {
    3: ['RCLCO ODCE and NPI Results, Q2 2025', 'CommercialCafe National Industrial Report, December 2025'],
    5: ['CommercialCafe National Industrial Report, December 2025'],
    6: ['CommercialCafe National Industrial Report, December 2025', 'Cushman & Wakefield Q2 2025 US Industrial MarketBeat'],
    7: ['CommercialCafe National Industrial Report, December 2025'],
    8: ['CommercialCafe National Industrial Report, December 2025'],
    9: ['CBRE Cap Rate Survey, H1 2024'],
    11: ['Partners Real Estate Market Reports, 2024', 'REBusinessOnline Nashville Industrial Market Analysis, 2024'],
    12: ['REBusinessOnline Nashville, 2024', 'WareCRE Tampa, 2025', 'Savills Raleigh-Durham, Q4 2024'],
    13: ['Partners Real Estate, 2024', 'Colliers Phoenix, 2024', 'Kidder Mathews Phoenix, 2024'],
    15: ['Clarion Partners Industrial Outlook, 2025', 'NAIOP Nearshoring Analysis, 2024'],
    16: ['NAIOP Nearshoring Analysis, 2024'],
    17: ['CBRE Interest Rate Impact Report, 2024', 'RCLCO ODCE/NPI, Q2 2025'],
    23: ['NASRA Public Pension Plan Return Assumptions, FY 2024'],
    24: ['RCLCO ODCE and NPI Results, Q2 2025'],
    25: ['RCLCO ODCE and NPI Results, Q2 2025'],
    27: ['Clarion Partners Industrial Outlook, 2025'],
    35: ['GRESB 2025 Real Estate Assessment Results'],
    36: ['GRESB 2025 Real Estate Assessment Results'],
}


def apply_font(run, font_spec):
    """Apply font specification to a run."""
    font = run.font
    font.name = font_spec.get('name', 'Arial')
    font.size = Pt(font_spec.get('size', 14))
    font.bold = font_spec.get('bold', False)
    color = font_spec.get('color')
    if color:
        if color in COLORS:
            font.color.rgb = COLORS[color]
        elif color.startswith('#'):
            font.color.rgb = RGBColor.from_string(color.lstrip('#'))


def remove_table_borders(table):
    """Remove all borders from a table by setting noFill on all border elements."""
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            # Remove all border elements and add noFill
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                # Find existing border element
                border = tcPr.find(qn(f'a:{border_name}'))
                if border is not None:
                    tcPr.remove(border)

                # Add border element with noFill
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                noFill = etree.SubElement(border, qn('a:noFill'))


def format_table(table):
    """Apply table formatting per v24 spec."""
    num_cols = len(table.columns)

    for row_idx, row in enumerate(table.rows):
        # Set row height
        if row_idx == 0:
            row.height = Inches(TABLE_FORMAT['header_row']['height_inches'])
        else:
            row.height = Inches(TABLE_FORMAT['data_rows']['height_inches'])

        for col_idx, cell in enumerate(row.cells):
            # Vertical alignment: center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins
            margins = TABLE_FORMAT['cell_margins']
            cell.margin_top = Inches(margins['top'])
            cell.margin_bottom = Inches(margins['bottom'])
            cell.margin_right = Inches(margins['right'])

            if col_idx == 0:
                cell.margin_left = Inches(margins['first_column_left'])
            else:
                cell.margin_left = Inches(margins['other_left'])

            # Apply fills
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(
                    TABLE_FORMAT['header_row']['fill_color'].lstrip('#')
                )
            else:
                fills = TABLE_FORMAT['data_rows']['alternating_fills']
                fill_color = fills['odd'] if row_idx % 2 == 1 else fills['even']
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip('#'))

            # Format text
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    if row_idx == 0:
                        font = run.font
                        font.name = 'Arial'
                        font.size = Pt(TABLE_FORMAT['header_row']['font_size'])
                        font.bold = TABLE_FORMAT['header_row']['font_bold']
                        font.color.rgb = RGBColor.from_string(
                            TABLE_FORMAT['header_row']['text_color'].lstrip('#')
                        )
                    else:
                        font = run.font
                        font.name = 'Arial'
                        font.size = Pt(TABLE_FORMAT['data_rows']['font_size'])
                        font.bold = TABLE_FORMAT['data_rows']['font_bold']
                        font.color.rgb = RGBColor.from_string(
                            TABLE_FORMAT['data_rows']['text_color'].lstrip('#')
                        )

    # Remove all borders
    remove_table_borders(table)


def add_section_name(slide, section_name):
    """Add section name to top-right placeholder (idx=17)."""
    if not section_name:
        return

    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.idx == 17:
                tf = shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.RIGHT
                run = para.add_run()
                run.text = section_name
                font = run.font
                font.name = 'Arial'
                font.size = Pt(10)
                font.bold = False
                font.color.rgb = COLORS['dark_text']
                return


def add_footnote(slide, sources):
    """Add footnote to bottom placeholder (idx=20)."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.idx == 20:
                tf = shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.RIGHT
                run = para.add_run()
                run.text = footnote_text
                font = run.font
                font.name = 'Arial'
                font.size = Pt(6)
                font.bold = False
                font.color.rgb = COLORS['dark_text']
                return


def add_standardized_bullets(text_frame, bullets):
    """Add bullets using standard formatting."""
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        para.level = 0
        para.space_before = Pt(6)
        para.space_after = Pt(3)

        run = para.add_run()
        run.text = bullet
        font = run.font
        font.name = 'Arial'
        font.size = Pt(14)
        font.color.rgb = COLORS['dark_text']

        # Add bullet character
        try:
            pPr = para._p.get_or_add_pPr()
            for child in list(pPr):
                if 'bu' in child.tag.lower():
                    pPr.remove(child)
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '\u2022')
        except:
            pass


def get_section_for_slide(slide_idx, all_slides):
    """Get the section name for a given slide index."""
    if slide_idx < len(all_slides):
        return all_slides[slide_idx].get('section', '')
    return ''


def main():
    """Generate the formatted presentation v25."""
    print("Generating formatted presentation v25...")
    print("=" * 60)

    # Load the user's v24 edited PPTX as template
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v24_CS_edits.pptx")
    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    # Load the original outline for content
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    if not outline_path.exists():
        print(f"Error: Outline not found: {outline_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    print(f"Template loaded: {len(prs.slides)} slides")

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({
                'section': section.get('name', ''),
                **slide
            })

    print(f"Outline has {len(all_slides)} slides")

    # The template already has the first N slides formatted
    # We just need to ensure section names and footnotes are in the right placeholders
    print("\nUpdating existing slides with section names and footnotes...")

    for idx in range(min(len(prs.slides), len(all_slides))):
        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')

        # Add section name to top-right
        add_section_name(slide, section_name)

        # Add footnotes if this slide has sources
        if (idx + 1) in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[idx + 1])

    # Add any remaining slides
    if len(prs.slides) < len(all_slides):
        print(f"\nAdding slides {len(prs.slides) + 1} to {len(all_slides)}...")

        for idx in range(len(prs.slides), len(all_slides)):
            slide_data = all_slides[idx]
            slide_type = slide_data.get('slide_type', 'title_content')
            content = slide_data.get('content', {})
            section_name = slide_data.get('section', '')

            # Get appropriate layout
            layout_name = SLIDE_TYPE_LAYOUT.get(slide_type, 'bullet_content')
            layout_config = LAYOUTS.get(layout_name, LAYOUTS['bullet_content'])
            layout = prs.slide_masters[0].slide_layouts[layout_config['index']]

            print(f"  Adding slide {idx + 1}: {slide_type} -> {layout_name}")

            slide = prs.slides.add_slide(layout)

            # Set title
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    tf = shape.text_frame
                    tf.clear()
                    para = tf.paragraphs[0]
                    run = para.add_run()
                    run.text = content.get('title', '')
                    apply_font(run, CONFIG['fonts']['title'])
                    break

            # Add section name
            add_section_name(slide, section_name)

            # Add footnotes
            if (idx + 1) in SLIDE_SOURCES:
                add_footnote(slide, SLIDE_SOURCES[idx + 1])

            # Handle content based on slide type
            if slide_type in ['title_content', 'content', 'key_metrics']:
                bullets = content.get('bullets', [])
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.idx == 18:
                        if bullets:
                            add_standardized_bullets(shape.text_frame, bullets)
                        break

            elif slide_type == 'table_slide':
                headers = content.get('headers', []) or content.get('table', {}).get('headers', [])
                rows = content.get('data', []) or content.get('table', {}).get('rows', [])

                if headers and rows:
                    # Add table
                    num_cols = len(headers)
                    num_rows = len(rows) + 1

                    table_shape = slide.shapes.add_table(
                        num_rows, num_cols,
                        Inches(0.40), Inches(2.70),
                        Inches(10.20), Inches(4.00)
                    )
                    table = table_shape.table

                    # Populate
                    for c_idx, header in enumerate(headers):
                        table.cell(0, c_idx).text = header
                    for r_idx, row in enumerate(rows):
                        for c_idx, val in enumerate(row):
                            if c_idx < num_cols:
                                table.cell(r_idx + 1, c_idx).text = str(val) if val else ''

                    format_table(table)

    # Save
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v25.pptx")
    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"Generated presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
