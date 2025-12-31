"""Generate presentation v28 with all formatting fixes from v27_CS_edits.

Fixes applied:
1. Table borders - explicit noFill on all border elements
2. Cell margins - 0.1" left/right, 0.05" top/bottom
3. Column alignment - majority-based (numbers right, text left), headers match column
4. White logo on title and end slides (from logos/pccp_logo_white.png)
5. Table sizing - 10.2" wide, 0.4" from left edge, equal column distribution
6. Section names - 9pt, #A6A6A6, right-aligned
7. Footnotes - 6pt, #A6A6A6, left-aligned
8. CONFIDENTIAL watermark inherited from master
9. Updated layout indices per v27_CS_edits
10. Historical returns chart data for slide 25
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from lxml import etree

# Colors
COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text': RGBColor(0x06, 0x1F, 0x32),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'medium_gray': RGBColor(0xA6, 0xA6, 0xA6),
}

# Updated layout indices for v27_CS_edits
LAYOUTS = {
    'frontpage': 0,
    'section_title': 1,
    'bullet_content': 2,
    'chart': 3,
    'table': 4,
    'side_by_side': 5,
    'contact': 6,
    'disclaimers': 7,
    'end': 13,
}

# Section to image mapping
SECTION_IMAGES = {
    'Executive Summary': 'title_slide.png',
    'Market Fundamentals': 'market_fundamentals.png',
    'Target Markets': 'target_markets.png',
    'Demand Drivers': 'demand_drivers.png',
    'Investment Strategy': 'investment_strategy.png',
    'Competitive Positioning': 'competitive_positioning.png',
    'Risk Management': 'risk_management.png',
    'ESG Strategy': 'esg_strategy.png',
    'JV Structure': 'jv_structure.png',
    'Conclusion': 'conclusion.png',
}

# Source mapping
SLIDE_SOURCES = {
    2: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    3: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    5: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    6: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    7: ['CommercialCafe Dec 2025'],
    8: ['CommercialCafe Dec 2025'],
    9: ['CBRE Cap Rate Survey H1 2024'],
    11: ['Partners RE 2024', 'REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024', 'Colliers Phoenix 2024'],
    12: ['REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024'],
    13: ['Partners RE DFW/Atlanta/San Antonio 2024', 'Colliers Phoenix 2024', 'Kidder Mathews Phoenix 2024'],
    15: ['Clarion Partners Industrial Outlook 2025', 'NAIOP Nearshoring Analysis 2024'],
    16: ['NAIOP Nearshoring Analysis 2024'],
    17: ['CBRE Interest Rate Impact 2024', 'RCLCO ODCE/NPI Q2 2025'],
    19: ['Clarion Partners Industrial Outlook 2025'],
    20: ['REBusinessOnline Nashville 2024'],
    21: ['REBusinessOnline Nashville 2024'],
    22: ['RCLCO ODCE/NPI Q2 2025'],
    23: ['NASRA FY 2024', 'RCLCO ODCE/NPI Q2 2025'],
    24: ['RCLCO ODCE/NPI Q2 2025'],
    25: ['RCLCO ODCE/NPI Q2 2025'],
    27: ['Clarion Partners Industrial Outlook 2025'],
    28: ['Clarion Partners Industrial Outlook 2025'],
    30: ['Clarion Partners Industrial Outlook 2025'],
    31: ['Clarion Partners Industrial Outlook 2025'],
    32: ['Cushman & Wakefield Q2 2025'],
    33: ['Cushman & Wakefield Q2 2025'],
    35: ['GRESB 2025 Results'],
    36: ['GRESB 2025 Results'],
    38: ['Clarion Partners Industrial Outlook 2025'],
    39: ['Clarion Partners Industrial Outlook 2025'],
    40: ['Clarion Partners Industrial Outlook 2025'],
    42: ['RCLCO ODCE/NPI Q2 2025', 'Clarion Partners Industrial Outlook 2025'],
    43: ['RCLCO ODCE/NPI Q2 2025'],
}

# Historical returns data for slide 25 chart (10-year annualized returns)
# Based on NCREIF NPI data - Industrial outperformed by 260-590 bps
HISTORICAL_RETURNS_DATA = {
    'categories': ['Industrial', 'Apartment', 'Retail', 'Office'],
    'returns': [12.4, 9.8, 8.4, 6.5],  # 10-year annualized returns
}

# PCCP Contact Info
PCCP_CONTACT = {
    'name': 'PCCP, LLC',
    'headquarters': {
        'address': '10100 Santa Monica Blvd., Suite 1000',
        'city_state_zip': 'Los Angeles, CA 90067',
        'phone': '310.414.7870',
    },
    'email': 'Inquiry@pccpllc.com',
    'website': 'pccpllc.com',
    'other_offices': [
        {'city': 'New York', 'phone': '646.308.2100'},
        {'city': 'San Francisco', 'phone': '415.732.7645'},
        {'city': 'Atlanta', 'phone': '404.947.6080'},
    ]
}


def is_numeric_cell(text):
    """Check if cell content is numeric."""
    if not text:
        return False
    text = text.strip()
    patterns = [
        r'^[\d,]+\.?\d*%?$',
        r'^\$[\d,]+\.?\d*[BMK]?$',
        r'^[\d,]+\.?\d*x$',
        r'^[\d,]+-[\d,]+%?$',
        r'^\d+\.\d+%$',
        r'^-?\d+\.?\d*%?$',
        r'^\(\d+\.?\d*%?\)$',
    ]
    return any(re.match(p, text) for p in patterns)


def get_column_alignment(table, col_idx):
    """Determine alignment for a column based on majority of cell values."""
    numeric_count = 0
    text_count = 0

    # Skip header row (row 0), analyze data rows
    for row_idx in range(1, len(table.rows)):
        cell = table.cell(row_idx, col_idx)
        cell_text = ''
        for para in cell.text_frame.paragraphs:
            cell_text += para.text

        if is_numeric_cell(cell_text.strip()):
            numeric_count += 1
        elif cell_text.strip():
            text_count += 1

    # Return RIGHT for numeric columns, LEFT for text columns
    return PP_ALIGN.RIGHT if numeric_count > text_count else PP_ALIGN.LEFT


def remove_table_borders_completely(table):
    """Remove ALL borders from table cells with explicit noFill."""
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            # Remove existing borders and add noFill
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                # Remove any existing border element
                existing = tcPr.find(qn(f'a:{border_name}'))
                if existing is not None:
                    tcPr.remove(existing)

                # Create new border element with noFill
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                border.set('w', '0')  # Width = 0
                etree.SubElement(border, qn('a:noFill'))


def format_table_v28(table):
    """Apply v28 table formatting with all fixes."""
    num_cols = len(table.columns)
    num_rows = len(table.rows)

    # Calculate column alignments based on majority
    column_alignments = []
    for col_idx in range(num_cols):
        column_alignments.append(get_column_alignment(table, col_idx))

    # Set equal column widths
    total_width = Inches(10.2)
    col_width = total_width // num_cols
    for col in table.columns:
        col.width = col_width

    for row_idx, row in enumerate(table.rows):
        row.height = Inches(0.50 if row_idx == 0 else 0.40)

        for col_idx, cell in enumerate(row.cells):
            # Vertical alignment - center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins - 0.1" left/right, 0.05" top/bottom
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            # Fill colors
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            # Apply consistent column alignment (header matches data)
            for para in cell.text_frame.paragraphs:
                para.alignment = column_alignments[col_idx]

                for run in para.runs:
                    run.font.name = 'Arial'
                    if row_idx == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = COLORS['white']
                    else:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rgb = COLORS['dark_text']

    # Remove all borders
    remove_table_borders_completely(table)


def add_section_name(slide, section_name):
    """Add section name (9pt, #A6A6A6, right-aligned)."""
    if not section_name:
        return
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 17:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = section_name
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = COLORS['medium_gray']
            return


def add_footnote(slide, sources):
    """Add footnote (6pt, #A6A6A6, left-aligned)."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx in [20, 16]:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = footnote_text
            run.font.name = 'Arial'
            run.font.size = Pt(6)
            run.font.color.rgb = COLORS['medium_gray']
            return


def remove_all_pictures_from_slide(slide):
    """Remove all picture shapes from a slide."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_section_image(slide, section_name, image_dir):
    """Add background image to section divider."""
    remove_all_pictures_from_slide(slide)

    image_file = SECTION_IMAGES.get(section_name)
    if not image_file:
        return False

    image_path = Path(image_dir) / image_file
    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return False

    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)
    return True


def add_logo_to_slide(slide, logo_path, position, size):
    """Add logo image to slide at specified position."""
    if not logo_path.exists():
        print(f"    Warning: Logo not found: {logo_path}")
        return None

    pic = slide.shapes.add_picture(
        str(logo_path),
        Inches(position['left']),
        Inches(position['top']),
        Inches(size['width']),
        Inches(size['height'])
    )
    return pic


def update_historical_returns_chart(slide):
    """Add/update the chart on slide 25 with historical returns data."""
    # First check if chart already exists
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = HISTORICAL_RETURNS_DATA['categories']
            chart_data.add_series('10-Year Annualized Returns', HISTORICAL_RETURNS_DATA['returns'])
            chart.replace_data(chart_data)
            print("    Updated existing chart data")
            return True

    # If no chart exists, add one to the chart placeholder or as a new shape
    chart_data = CategoryChartData()
    chart_data.categories = HISTORICAL_RETURNS_DATA['categories']
    chart_data.add_series('10-Year Annualized Returns (%)',
                          tuple(HISTORICAL_RETURNS_DATA['returns']))

    # Add chart to slide
    x, y = Inches(0.4), Inches(2.7)
    cx, cy = Inches(10.2), Inches(3.8)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Style the chart
    chart.has_legend = False

    # Color the bars - Industrial in accent blue, others in gray
    plot = chart.plots[0]
    series = plot.series[0]

    # Set data labels
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.number_format = '0.0"%"'
    data_labels.font.size = Pt(12)
    data_labels.font.bold = True

    print("    Added historical returns column chart")
    return True


def update_historical_returns_textbox(slide):
    """Update the text box on slide 25 with returns narrative."""
    narrative = (
        "Industrial delivered 12.4% 10-year annualized returns, outperforming "
        "Apartment (+260 bps), Retail (+400 bps), and Office (+590 bps)."
    )

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 18:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = narrative
            run.font.name = 'Arial'
            run.font.size = Pt(11)
            run.font.color.rgb = COLORS['dark_text']
            print("    Updated historical returns text box")
            return True

    return False


def update_contact_slide(slide):
    """Update contact slide with PCCP information."""
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx in [18, 7, 21]:
            tf = shape.text_frame
            tf.clear()

            contact_lines = [
                f"{PCCP_CONTACT['name']}",
                f"{PCCP_CONTACT['headquarters']['address']}",
                f"{PCCP_CONTACT['headquarters']['city_state_zip']}",
                f"Phone: {PCCP_CONTACT['headquarters']['phone']}",
                f"Email: {PCCP_CONTACT['email']}",
                f"Web: {PCCP_CONTACT['website']}",
                "",
                "Additional Offices:",
            ]

            for office in PCCP_CONTACT['other_offices']:
                contact_lines.append(f"{office['city']}: {office['phone']}")

            for i, line in enumerate(contact_lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = line
                run.font.name = 'Arial'
                run.font.size = Pt(14)
                run.font.color.rgb = COLORS['dark_text']
            return


def add_end_slide(prs, image_dir, logo_path):
    """Add End slide with background image and centered logo."""
    try:
        end_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['end']]
    except:
        end_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['section_title']]

    slide = prs.slides.add_slide(end_layout)

    # Add background image
    image_path = Path(image_dir) / 'end_slide.png'
    if image_path.exists():
        pic = slide.shapes.add_picture(
            str(image_path),
            Inches(0), Inches(0),
            Inches(11), Inches(8.5)
        )
        # Move to back
        spTree = slide.shapes._spTree
        sp = spTree[-1]
        spTree.remove(sp)
        spTree.insert(2, sp)

    # Add centered white logo
    if logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            Inches(4.6), Inches(3.6),
            Inches(1.8), Inches(1.3)
        )

    return slide


def main():
    """Generate presentation v28."""
    print("Generating presentation v28 with all formatting fixes...")
    print("=" * 70)

    # Paths
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v27_CS_edits.pptx")
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    image_dir = Path("cache/images/cropped")
    logo_path = Path("logos/pccp_logo_white.png")
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v28.pptx")

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({'section': section.get('name', ''), **slide})

    print(f"Template: {len(prs.slides)} slides")
    print(f"Image directory: {image_dir}")
    print(f"Logo: {logo_path}")

    # Process each slide
    print("\nProcessing slides...")

    for idx in range(len(prs.slides)):
        if idx >= len(all_slides):
            break

        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')
        slide_type = slide_data.get('slide_type', '')
        slide_num = idx + 1

        # Add section name (9pt, #A6A6A6, right-aligned)
        add_section_name(slide, section_name)

        # Add footnote (6pt, #A6A6A6, left-aligned)
        if slide_num in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[slide_num])

        # Handle section dividers
        if slide_type == 'section_divider':
            if add_section_image(slide, section_name, image_dir):
                print(f"  Slide {slide_num}: Updated section image for '{section_name}'")

        # Add white logo to title slide
        if slide_num == 1:
            # Remove existing logo placeholder content and add white logo
            add_logo_to_slide(
                slide, logo_path,
                {'left': 0.4, 'top': 0.4},
                {'width': 2.5, 'height': 1.8}
            )
            print(f"  Slide {slide_num}: Added white logo to title slide")

        # Update slide 25 historical returns
        if slide_num == 25:
            updated_chart = update_historical_returns_chart(slide)
            updated_text = update_historical_returns_textbox(slide)
            if updated_chart or updated_text:
                print(f"  Slide {slide_num}: Updated historical returns data")

        # Update contact slide (44)
        if slide_num == 44:
            update_contact_slide(slide)
            print(f"  Slide {slide_num}: Updated contact information")

    # Format all tables with v28 specifications
    print("\nFormatting tables (v28 spec)...")
    tables_updated = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                # Ensure table position
                shape.left = Inches(0.4)
                shape.width = Inches(10.2)

                format_table_v28(shape.table)
                tables_updated += 1
                print(f"  Slide {slide_idx + 1}: Formatted table")

    print(f"  Total tables formatted: {tables_updated}")

    # Check if End slide exists, add if not
    # Look at last slide to determine if it's an End slide
    last_slide = prs.slides[-1]
    has_end_slide = False
    for shape in last_slide.shapes:
        if hasattr(shape, 'text') and 'PCCP' in shape.text:
            has_end_slide = True
            break

    if not has_end_slide:
        print("\nAdding End slide...")
        add_end_slide(prs, image_dir, logo_path)
        print("  Added End slide with white logo")

    # Save
    prs.save(str(output_path))

    print(f"\n{'=' * 70}")
    print(f"Generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless", "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
