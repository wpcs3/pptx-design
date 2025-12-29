"""
PPTEval - Presentation Evaluation Framework

Evaluates generated presentations across three dimensions:
1. Content Quality - Topic coverage, accuracy, completeness
2. Visual Design - Element count, spacing, consistency
3. Logical Coherence - Flow, section structure, narrative

Inspired by PPTAgent's PPTEval framework.

Phase 2 Enhancement (2025-12-29):
- Content evaluation (bullet count, word limits, completeness)
- Design evaluation (element density, spacing, visual consistency)
- Coherence evaluation (slide flow, section structure)
- Overall scoring with detailed feedback
"""

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


# =============================================================================
# Evaluation Constants (from PPTAgent best practices)
# =============================================================================

class EvaluationLimits:
    """Best practice limits for slide content."""
    MAX_ELEMENTS_PER_SLIDE = 6
    MAX_BULLETS_PER_SLIDE = 6
    MAX_WORDS_PER_BULLET = 15
    MIN_WORDS_PER_BULLET = 3
    MAX_TITLE_WORDS = 8
    MIN_SLIDE_COUNT = 5
    MAX_SLIDE_COUNT = 30
    IDEAL_TEXT_DENSITY = 0.6  # 60% of element space
    MIN_FONT_SIZE_PT = 14
    MAX_FONT_SIZE_PT = 44


# =============================================================================
# Data Classes
# =============================================================================

@dataclass
class ContentScore:
    """Content quality evaluation results."""
    score: float  # 0.0 - 1.0
    issues: List[str] = field(default_factory=list)
    details: Dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return self.score >= 0.7


@dataclass
class DesignScore:
    """Visual design evaluation results."""
    score: float  # 0.0 - 1.0
    issues: List[str] = field(default_factory=list)
    details: Dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return self.score >= 0.7


@dataclass
class CoherenceScore:
    """Logical coherence evaluation results."""
    score: float  # 0.0 - 1.0
    issues: List[str] = field(default_factory=list)
    details: Dict[str, Any] = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return self.score >= 0.7


@dataclass
class EvaluationResult:
    """Complete evaluation result for a presentation."""
    presentation_path: str
    content: ContentScore
    design: DesignScore
    coherence: CoherenceScore
    overall_score: float = 0.0
    grade: str = "F"
    recommendations: List[str] = field(default_factory=list)

    def __post_init__(self):
        """Calculate overall score and grade."""
        # Weighted average: Content 40%, Design 30%, Coherence 30%
        self.overall_score = (
            self.content.score * 0.4 +
            self.design.score * 0.3 +
            self.coherence.score * 0.3
        )
        self.grade = self._calculate_grade()
        self.recommendations = self._generate_recommendations()

    def _calculate_grade(self) -> str:
        """Convert score to letter grade."""
        if self.overall_score >= 0.9:
            return "A"
        elif self.overall_score >= 0.8:
            return "B"
        elif self.overall_score >= 0.7:
            return "C"
        elif self.overall_score >= 0.6:
            return "D"
        else:
            return "F"

    def _generate_recommendations(self) -> List[str]:
        """Generate improvement recommendations."""
        recs = []

        # Content recommendations
        if self.content.score < 0.7:
            recs.extend(self.content.issues[:3])

        # Design recommendations
        if self.design.score < 0.7:
            recs.extend(self.design.issues[:3])

        # Coherence recommendations
        if self.coherence.score < 0.7:
            recs.extend(self.coherence.issues[:3])

        return recs[:5]  # Top 5 recommendations

    def summary(self) -> str:
        """Generate human-readable summary."""
        lines = [
            f"PPTEval Results: {Path(self.presentation_path).name}",
            "=" * 50,
            f"Overall Score: {self.overall_score:.2f} (Grade: {self.grade})",
            "",
            f"Content Quality:   {self.content.score:.2f} {'PASS' if self.content.passed else 'FAIL'}",
            f"Visual Design:     {self.design.score:.2f} {'PASS' if self.design.passed else 'FAIL'}",
            f"Logical Coherence: {self.coherence.score:.2f} {'PASS' if self.coherence.passed else 'FAIL'}",
            "",
        ]

        if self.recommendations:
            lines.append("Top Recommendations:")
            for i, rec in enumerate(self.recommendations, 1):
                lines.append(f"  {i}. {rec}")

        return "\n".join(lines)

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            "presentation_path": self.presentation_path,
            "overall_score": self.overall_score,
            "grade": self.grade,
            "content": {
                "score": self.content.score,
                "passed": self.content.passed,
                "issues": self.content.issues,
                "details": self.content.details
            },
            "design": {
                "score": self.design.score,
                "passed": self.design.passed,
                "issues": self.design.issues,
                "details": self.design.details
            },
            "coherence": {
                "score": self.coherence.score,
                "passed": self.coherence.passed,
                "issues": self.coherence.issues,
                "details": self.coherence.details
            },
            "recommendations": self.recommendations
        }


# =============================================================================
# Main Evaluator Class
# =============================================================================

class PresentationEvaluator:
    """
    Evaluates presentations using the PPTEval framework.

    Usage:
        evaluator = PresentationEvaluator()
        result = evaluator.evaluate("presentation.pptx")
        print(result.summary())

        # Or evaluate with context
        result = evaluator.evaluate("presentation.pptx", context={
            "topic": "Investment Pitch",
            "expected_sections": ["Overview", "Market", "Financials"]
        })
    """

    def __init__(self, limits: EvaluationLimits = None):
        """
        Initialize evaluator.

        Args:
            limits: Custom evaluation limits (uses defaults if not specified)
        """
        self.limits = limits or EvaluationLimits()

    def evaluate(
        self,
        pptx_path: str,
        context: Optional[Dict[str, Any]] = None
    ) -> EvaluationResult:
        """
        Evaluate a presentation file.

        Args:
            pptx_path: Path to the PPTX file
            context: Optional context for evaluation (topic, expected sections, etc.)

        Returns:
            EvaluationResult with scores and recommendations
        """
        prs = Presentation(pptx_path)

        # Extract slide data
        slide_data = self._extract_slide_data(prs)

        # Evaluate each dimension
        content = self._evaluate_content(slide_data, context or {})
        design = self._evaluate_design(slide_data, prs)
        coherence = self._evaluate_coherence(slide_data, context or {})

        return EvaluationResult(
            presentation_path=pptx_path,
            content=content,
            design=design,
            coherence=coherence
        )

    def _extract_slide_data(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Extract structured data from all slides."""
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "index": i,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "shapes": [],
                "text_content": [],
                "title": None,
                "bullet_count": 0,
                "word_count": 0,
                "has_chart": False,
                "has_table": False,
                "has_image": False,
            }

            for shape in slide.shapes:
                shape_info = {
                    "type": shape.shape_type,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }

                # Extract text content
                if shape.has_text_frame:
                    text_parts = []
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)
                            words = text.split()
                            slide_info["word_count"] += len(words)

                            # Count bullets (lines starting with bullet character)
                            if text.startswith("•") or text.startswith("-"):
                                slide_info["bullet_count"] += 1

                    if text_parts:
                        shape_info["text"] = "\n".join(text_parts)
                        slide_info["text_content"].extend(text_parts)

                # Check for title
                if shape == slide.shapes.title and shape.has_text_frame:
                    slide_info["title"] = shape.text_frame.text.strip()

                # Check for charts
                if shape.has_chart:
                    slide_info["has_chart"] = True
                    shape_info["chart_type"] = str(shape.chart.chart_type)

                # Check for tables
                if shape.has_table:
                    slide_info["has_table"] = True
                    shape_info["table_rows"] = len(shape.table.rows)
                    shape_info["table_cols"] = len(shape.table.columns)

                # Check for images
                if hasattr(shape, "image"):
                    slide_info["has_image"] = True

                slide_info["shapes"].append(shape_info)

            slides_data.append(slide_info)

        return slides_data

    def _evaluate_content(
        self,
        slide_data: List[Dict[str, Any]],
        context: Dict[str, Any]
    ) -> ContentScore:
        """
        Evaluate content quality.

        Checks:
        - Slide count within reasonable range
        - Bullet counts per slide
        - Word counts per element
        - Title presence
        - Data visualization presence
        """
        issues = []
        details = {
            "slide_count": len(slide_data),
            "total_bullets": 0,
            "total_words": 0,
            "slides_over_bullet_limit": 0,
            "slides_missing_title": 0,
            "slides_with_visuals": 0,
        }

        penalties = 0.0
        max_penalties = 10.0

        # Check slide count
        if len(slide_data) < self.limits.MIN_SLIDE_COUNT:
            issues.append(f"Too few slides ({len(slide_data)}). Consider adding more content.")
            penalties += 1.0
        elif len(slide_data) > self.limits.MAX_SLIDE_COUNT:
            issues.append(f"Too many slides ({len(slide_data)}). Consider condensing content.")
            penalties += 1.0

        for slide in slide_data:
            details["total_bullets"] += slide["bullet_count"]
            details["total_words"] += slide["word_count"]

            # Check bullet count
            if slide["bullet_count"] > self.limits.MAX_BULLETS_PER_SLIDE:
                details["slides_over_bullet_limit"] += 1

            # Check for title
            if not slide["title"]:
                details["slides_missing_title"] += 1

            # Check for visuals
            if slide["has_chart"] or slide["has_table"] or slide["has_image"]:
                details["slides_with_visuals"] += 1

        # Penalty for too many bullets
        if details["slides_over_bullet_limit"] > 0:
            pct = details["slides_over_bullet_limit"] / len(slide_data)
            issues.append(
                f"{details['slides_over_bullet_limit']} slides exceed {self.limits.MAX_BULLETS_PER_SLIDE} bullets. "
                "Consider splitting content."
            )
            penalties += pct * 2.0

        # Penalty for missing titles
        if details["slides_missing_title"] > 1:  # Allow 1 (title slide may not have "title")
            pct = (details["slides_missing_title"] - 1) / len(slide_data)
            issues.append(
                f"{details['slides_missing_title']} slides missing titles."
            )
            penalties += pct * 1.5

        # Bonus for visuals
        visual_ratio = details["slides_with_visuals"] / len(slide_data) if slide_data else 0
        if visual_ratio < 0.2:
            issues.append("Consider adding more visual elements (charts, tables, images).")
            penalties += 0.5

        # Calculate score
        score = max(0.0, 1.0 - (penalties / max_penalties))

        return ContentScore(score=score, issues=issues, details=details)

    def _evaluate_design(
        self,
        slide_data: List[Dict[str, Any]],
        prs: Presentation
    ) -> DesignScore:
        """
        Evaluate visual design quality.

        Checks:
        - Element count per slide
        - Consistent positioning
        - Font size appropriateness
        - Visual density
        """
        issues = []
        details = {
            "slides_overcrowded": 0,
            "avg_elements_per_slide": 0,
            "max_elements": 0,
            "position_consistency": 0.0,
        }

        penalties = 0.0
        max_penalties = 10.0

        total_elements = 0
        max_elements = 0
        element_counts = []

        for slide in slide_data:
            # Count meaningful shapes (text, charts, tables, images)
            meaningful_shapes = len([
                s for s in slide["shapes"]
                if s.get("text") or s.get("chart_type") or s.get("table_rows")
            ])

            element_counts.append(meaningful_shapes)
            total_elements += meaningful_shapes
            max_elements = max(max_elements, meaningful_shapes)

            if meaningful_shapes > self.limits.MAX_ELEMENTS_PER_SLIDE:
                details["slides_overcrowded"] += 1

        details["avg_elements_per_slide"] = (
            total_elements / len(slide_data) if slide_data else 0
        )
        details["max_elements"] = max_elements

        # Penalty for overcrowded slides
        if details["slides_overcrowded"] > 0:
            pct = details["slides_overcrowded"] / len(slide_data)
            issues.append(
                f"{details['slides_overcrowded']} slides have too many elements "
                f"(>{self.limits.MAX_ELEMENTS_PER_SLIDE}). Consider simplifying."
            )
            penalties += pct * 3.0

        # Check element count variance (consistency)
        if element_counts:
            import statistics
            if len(element_counts) > 1:
                variance = statistics.variance(element_counts)
                details["element_variance"] = variance
                if variance > 10:
                    issues.append("Inconsistent element density across slides.")
                    penalties += 1.0

        # Calculate position consistency (check title alignment)
        title_positions = []
        for slide in slide_data:
            for shape in slide["shapes"]:
                if shape.get("text") and slide["title"] and shape["text"].startswith(slide["title"][:20]):
                    title_positions.append((shape["left"], shape["top"]))

        if title_positions and len(title_positions) > 1:
            left_values = [p[0] for p in title_positions]
            top_values = [p[1] for p in title_positions]
            left_var = statistics.variance(left_values) if len(left_values) > 1 else 0
            top_var = statistics.variance(top_values) if len(top_values) > 1 else 0

            # High variance = inconsistent positioning
            if left_var > 100000 or top_var > 100000:  # EMU units
                issues.append("Title positions are inconsistent across slides.")
                penalties += 1.0
            else:
                details["position_consistency"] = 1.0

        # Calculate score
        score = max(0.0, 1.0 - (penalties / max_penalties))

        return DesignScore(score=score, issues=issues, details=details)

    def _evaluate_coherence(
        self,
        slide_data: List[Dict[str, Any]],
        context: Dict[str, Any]
    ) -> CoherenceScore:
        """
        Evaluate logical coherence and flow.

        Checks:
        - Has title slide
        - Has section structure
        - Logical progression
        - Has conclusion/summary
        """
        issues = []
        details = {
            "has_title_slide": False,
            "has_section_breaks": False,
            "has_conclusion": False,
            "section_count": 0,
        }

        penalties = 0.0
        max_penalties = 10.0

        if not slide_data:
            return CoherenceScore(score=0.0, issues=["No slides found"], details=details)

        # Check for title slide (first slide)
        first_slide = slide_data[0]
        if first_slide["layout_name"] in ["Frontpage", "Title Slide", "title_slide"]:
            details["has_title_slide"] = True
        elif first_slide["title"] and first_slide["word_count"] < 20:
            details["has_title_slide"] = True  # Likely a title slide
        else:
            issues.append("Missing clear title slide.")
            penalties += 1.0

        # Check for section structure
        section_keywords = ["Section", "Chapter", "Part", "Overview", "Summary", "Conclusion"]
        section_count = 0
        for slide in slide_data:
            layout = slide["layout_name"].lower()
            title = (slide["title"] or "").lower()

            if "section" in layout or "breaker" in layout:
                section_count += 1
            elif any(kw.lower() in title for kw in section_keywords):
                section_count += 1

        details["section_count"] = section_count
        if section_count >= 2:
            details["has_section_breaks"] = True
        else:
            if len(slide_data) > 10:
                issues.append("Consider adding section dividers for better organization.")
                penalties += 0.5

        # Check for conclusion/summary
        last_slides = slide_data[-3:] if len(slide_data) >= 3 else slide_data
        conclusion_keywords = ["summary", "conclusion", "takeaway", "next step", "contact", "thank", "question"]
        for slide in last_slides:
            title = (slide["title"] or "").lower()
            layout = slide["layout_name"].lower()
            if any(kw in title or kw in layout for kw in conclusion_keywords):
                details["has_conclusion"] = True
                break

        if not details["has_conclusion"]:
            issues.append("Consider adding a conclusion or summary slide.")
            penalties += 0.5

        # Check for logical flow (titles should not repeat exactly)
        titles = [s["title"] for s in slide_data if s["title"]]
        unique_titles = set(titles)
        if len(titles) > len(unique_titles) + 1:  # Allow 1 duplicate
            issues.append("Multiple slides have duplicate titles. Consider making titles more specific.")
            penalties += 0.5

        # Check expected sections if provided in context
        expected_sections = context.get("expected_sections", [])
        if expected_sections:
            found_sections = 0
            all_text = " ".join(
                " ".join(s["text_content"]) for s in slide_data
            ).lower()

            for section in expected_sections:
                if section.lower() in all_text:
                    found_sections += 1

            coverage = found_sections / len(expected_sections) if expected_sections else 1
            details["section_coverage"] = coverage
            if coverage < 0.7:
                issues.append(f"Missing expected sections. Found {found_sections}/{len(expected_sections)}.")
                penalties += (1 - coverage) * 2.0

        # Calculate score
        score = max(0.0, 1.0 - (penalties / max_penalties))

        return CoherenceScore(score=score, issues=issues, details=details)


# =============================================================================
# Convenience Functions
# =============================================================================

def evaluate_presentation(
    pptx_path: str,
    context: Optional[Dict[str, Any]] = None
) -> EvaluationResult:
    """
    Quick evaluation of a presentation.

    Args:
        pptx_path: Path to PPTX file
        context: Optional evaluation context

    Returns:
        EvaluationResult
    """
    evaluator = PresentationEvaluator()
    return evaluator.evaluate(pptx_path, context)


def quick_score(pptx_path: str) -> float:
    """
    Get quick overall score for a presentation.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        Overall score (0.0-1.0)
    """
    result = evaluate_presentation(pptx_path)
    return result.overall_score


# =============================================================================
# CLI
# =============================================================================

def main():
    """CLI for presentation evaluation."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="PPTEval - Presentation Evaluation")
    parser.add_argument("pptx_path", help="Path to PPTX file")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    parser.add_argument("--context", type=str, help="JSON context file")

    args = parser.parse_args()

    # Load context if provided
    context = {}
    if args.context:
        with open(args.context, "r") as f:
            context = json.load(f)

    # Evaluate
    result = evaluate_presentation(args.pptx_path, context)

    # Output
    if args.json:
        print(json.dumps(result.to_dict(), indent=2))
    else:
        print(result.summary())


if __name__ == "__main__":
    main()
