"""
Round-Trip Testing Module

Verifies extraction → generation fidelity by:
1. Extracting content from source PPTX
2. Regenerating a new PPTX from extracted content
3. Comparing original vs regenerated using multiple metrics

Phase 2 Enhancement (2025-12-29):
- Text content comparison
- Visual similarity (SSIM) comparison
- Structural comparison (slide count, element count)
- Comprehensive fidelity reporting
"""

import json
import logging
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


# =============================================================================
# Data Classes
# =============================================================================

@dataclass
class SlideComparison:
    """Comparison results for a single slide."""
    slide_index: int
    text_similarity: float  # 0.0 - 1.0
    element_count_match: bool
    title_match: bool
    layout_match: bool
    issues: List[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return (
            self.text_similarity >= 0.8 and
            self.element_count_match and
            self.title_match
        )


@dataclass
class RoundtripResult:
    """Complete round-trip test result."""
    source_path: str
    generated_path: str
    slide_comparisons: List[SlideComparison]
    overall_text_similarity: float
    overall_visual_similarity: float
    slide_count_match: bool
    structural_fidelity: float  # 0.0 - 1.0
    issues: List[str] = field(default_factory=list)

    @property
    def passed(self) -> bool:
        return (
            self.overall_text_similarity >= 0.85 and
            self.slide_count_match and
            self.structural_fidelity >= 0.8
        )

    @property
    def fidelity_score(self) -> float:
        """Overall fidelity score (0.0 - 1.0)."""
        return (
            self.overall_text_similarity * 0.4 +
            self.overall_visual_similarity * 0.3 +
            self.structural_fidelity * 0.3
        )

    def summary(self) -> str:
        """Generate human-readable summary."""
        status = "PASS" if self.passed else "FAIL"
        lines = [
            f"Round-Trip Test: {status}",
            "=" * 50,
            f"Source: {Path(self.source_path).name}",
            f"Generated: {Path(self.generated_path).name}",
            "",
            f"Fidelity Score: {self.fidelity_score:.2%}",
            f"  Text Similarity: {self.overall_text_similarity:.2%}",
            f"  Visual Similarity: {self.overall_visual_similarity:.2%}",
            f"  Structural Fidelity: {self.structural_fidelity:.2%}",
            "",
            f"Slide Count Match: {'Yes' if self.slide_count_match else 'No'}",
            "",
            "Slide Details:",
        ]

        for comp in self.slide_comparisons:
            status = "OK" if comp.passed else "ISSUE"
            lines.append(
                f"  Slide {comp.slide_index + 1}: [{status}] "
                f"Text={comp.text_similarity:.2%}, "
                f"Title={'OK' if comp.title_match else 'DIFF'}, "
                f"Elements={'OK' if comp.element_count_match else 'DIFF'}"
            )

        if self.issues:
            lines.append("")
            lines.append("Issues:")
            for issue in self.issues[:5]:
                lines.append(f"  - {issue}")

        return "\n".join(lines)

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for JSON serialization."""
        return {
            "source_path": self.source_path,
            "generated_path": self.generated_path,
            "passed": self.passed,
            "fidelity_score": self.fidelity_score,
            "overall_text_similarity": self.overall_text_similarity,
            "overall_visual_similarity": self.overall_visual_similarity,
            "structural_fidelity": self.structural_fidelity,
            "slide_count_match": self.slide_count_match,
            "slide_comparisons": [
                {
                    "slide_index": c.slide_index,
                    "text_similarity": c.text_similarity,
                    "element_count_match": c.element_count_match,
                    "title_match": c.title_match,
                    "layout_match": c.layout_match,
                    "passed": c.passed,
                    "issues": c.issues
                }
                for c in self.slide_comparisons
            ],
            "issues": self.issues
        }


# =============================================================================
# Extraction Functions
# =============================================================================

def extract_slide_content(slide) -> Dict[str, Any]:
    """
    Extract content from a single slide.

    Args:
        slide: python-pptx Slide object

    Returns:
        Dictionary with slide content
    """
    content = {
        "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
        "title": None,
        "text_content": [],
        "shapes_count": 0,
        "has_chart": False,
        "has_table": False,
        "has_image": False,
    }

    for shape in slide.shapes:
        content["shapes_count"] += 1

        # Extract title
        if shape == slide.shapes.title and shape.has_text_frame:
            content["title"] = shape.text_frame.text.strip()

        # Extract text content
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    content["text_content"].append(text)

        # Check for charts
        if shape.has_chart:
            content["has_chart"] = True

        # Check for tables
        if shape.has_table:
            content["has_table"] = True
            content["table_rows"] = len(shape.table.rows)
            content["table_cols"] = len(shape.table.columns)

        # Check for images
        if hasattr(shape, "image"):
            content["has_image"] = True

    return content


def extract_presentation_content(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Extract content from all slides in a presentation.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        List of slide content dictionaries
    """
    prs = Presentation(pptx_path)
    slides_content = []

    for slide in prs.slides:
        content = extract_slide_content(slide)
        slides_content.append(content)

    return slides_content


# =============================================================================
# Comparison Functions
# =============================================================================

def calculate_text_similarity(text1: List[str], text2: List[str]) -> float:
    """
    Calculate text similarity between two lists of text.

    Uses Jaccard similarity on word sets.

    Args:
        text1: First text list
        text2: Second text list

    Returns:
        Similarity score (0.0 - 1.0)
    """
    # Combine and normalize text
    words1 = set()
    words2 = set()

    for t in text1:
        words1.update(t.lower().split())
    for t in text2:
        words2.update(t.lower().split())

    # Remove common stopwords
    stopwords = {"the", "a", "an", "is", "are", "was", "were", "be", "been", "to", "of", "and", "in", "for", "on", "with"}
    words1 -= stopwords
    words2 -= stopwords

    if not words1 and not words2:
        return 1.0
    if not words1 or not words2:
        return 0.0

    # Jaccard similarity
    intersection = words1 & words2
    union = words1 | words2

    return len(intersection) / len(union) if union else 0.0


def compare_slides(
    slide1_content: Dict[str, Any],
    slide2_content: Dict[str, Any],
    slide_index: int
) -> SlideComparison:
    """
    Compare two slides.

    Args:
        slide1_content: First slide content
        slide2_content: Second slide content
        slide_index: Index of the slides

    Returns:
        SlideComparison result
    """
    issues = []

    # Text similarity
    text_sim = calculate_text_similarity(
        slide1_content.get("text_content", []),
        slide2_content.get("text_content", [])
    )

    # Title match
    title1 = (slide1_content.get("title") or "").strip().lower()
    title2 = (slide2_content.get("title") or "").strip().lower()
    title_match = title1 == title2 or (title1 in title2 or title2 in title1)

    if not title_match and title1 and title2:
        issues.append(f"Title mismatch: '{title1[:30]}' vs '{title2[:30]}'")

    # Element count match (within tolerance)
    count1 = slide1_content.get("shapes_count", 0)
    count2 = slide2_content.get("shapes_count", 0)
    element_count_match = abs(count1 - count2) <= 2  # Allow 2 elements difference

    if not element_count_match:
        issues.append(f"Element count: {count1} vs {count2}")

    # Layout match
    layout1 = slide1_content.get("layout_name", "").lower()
    layout2 = slide2_content.get("layout_name", "").lower()
    layout_match = layout1 == layout2 or layout1 in layout2 or layout2 in layout1

    # Check for chart/table presence
    if slide1_content.get("has_chart") != slide2_content.get("has_chart"):
        issues.append("Chart presence mismatch")
    if slide1_content.get("has_table") != slide2_content.get("has_table"):
        issues.append("Table presence mismatch")

    return SlideComparison(
        slide_index=slide_index,
        text_similarity=text_sim,
        element_count_match=element_count_match,
        title_match=title_match,
        layout_match=layout_match,
        issues=issues
    )


# =============================================================================
# Round-Trip Tester
# =============================================================================

class RoundtripTester:
    """
    Tests round-trip fidelity of extraction → generation.

    Usage:
        tester = RoundtripTester()
        result = tester.test("source.pptx", "regenerated.pptx")
        print(result.summary())

        # Or with automatic regeneration
        result = tester.test_with_regeneration("source.pptx", template="consulting_toolkit")
    """

    def __init__(self, visual_comparison: bool = False):
        """
        Initialize tester.

        Args:
            visual_comparison: Whether to perform visual (SSIM) comparison
        """
        self.visual_comparison = visual_comparison

    def test(
        self,
        source_path: str,
        generated_path: str
    ) -> RoundtripResult:
        """
        Test round-trip fidelity between source and generated presentations.

        Args:
            source_path: Path to original PPTX
            generated_path: Path to regenerated PPTX

        Returns:
            RoundtripResult
        """
        issues = []

        # Extract content from both presentations
        source_content = extract_presentation_content(source_path)
        generated_content = extract_presentation_content(generated_path)

        # Check slide count
        slide_count_match = len(source_content) == len(generated_content)
        if not slide_count_match:
            issues.append(
                f"Slide count mismatch: {len(source_content)} vs {len(generated_content)}"
            )

        # Compare slides
        slide_comparisons = []
        min_slides = min(len(source_content), len(generated_content))

        for i in range(min_slides):
            comparison = compare_slides(
                source_content[i],
                generated_content[i],
                i
            )
            slide_comparisons.append(comparison)

        # Calculate overall text similarity
        all_text_sims = [c.text_similarity for c in slide_comparisons]
        overall_text_sim = sum(all_text_sims) / len(all_text_sims) if all_text_sims else 0.0

        # Calculate structural fidelity
        structural_points = 0
        structural_max = 0

        for comp in slide_comparisons:
            structural_max += 3
            if comp.title_match:
                structural_points += 1
            if comp.element_count_match:
                structural_points += 1
            if comp.layout_match:
                structural_points += 1

        # Penalize for missing slides
        if len(source_content) != len(generated_content):
            structural_max += abs(len(source_content) - len(generated_content))

        structural_fidelity = structural_points / structural_max if structural_max > 0 else 0.0

        # Visual comparison (if enabled)
        overall_visual_sim = 0.0
        if self.visual_comparison:
            try:
                overall_visual_sim = self._calculate_visual_similarity(
                    source_path, generated_path
                )
            except Exception as e:
                issues.append(f"Visual comparison failed: {e}")
                overall_visual_sim = 0.5  # Default to neutral

        return RoundtripResult(
            source_path=source_path,
            generated_path=generated_path,
            slide_comparisons=slide_comparisons,
            overall_text_similarity=overall_text_sim,
            overall_visual_similarity=overall_visual_sim,
            slide_count_match=slide_count_match,
            structural_fidelity=structural_fidelity,
            issues=issues
        )

    def _calculate_visual_similarity(
        self,
        source_path: str,
        generated_path: str
    ) -> float:
        """
        Calculate visual similarity using SSIM.

        Requires VisualTester from testing module.
        """
        try:
            from .testing import VisualTester

            tester = VisualTester()

            with tempfile.TemporaryDirectory() as tmp_dir:
                tmp_path = Path(tmp_dir)

                # Convert both to PNG
                source_png = tester.pptx_to_png(
                    Path(source_path),
                    tmp_path,
                    slide_index=0
                )
                generated_png = tester.pptx_to_png(
                    Path(generated_path),
                    tmp_path,
                    slide_index=0
                )

                if source_png and generated_png:
                    score, _ = tester.compare_images(source_png, generated_png)
                    return score

        except ImportError:
            logger.warning("VisualTester not available for visual comparison")
        except Exception as e:
            logger.warning(f"Visual comparison failed: {e}")

        return 0.0

    def test_with_regeneration(
        self,
        source_path: str,
        template: str = "consulting_toolkit",
        output_path: Optional[str] = None
    ) -> RoundtripResult:
        """
        Extract from source, regenerate, and test fidelity.

        Args:
            source_path: Path to source PPTX
            template: Template to use for regeneration
            output_path: Optional path for regenerated file

        Returns:
            RoundtripResult
        """
        try:
            from .builder import PresentationBuilder
        except ImportError:
            raise ImportError("PresentationBuilder required for regeneration")

        # Extract content
        source_content = extract_presentation_content(source_path)

        # Create output path
        if output_path is None:
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                output_path = f.name

        # Regenerate presentation
        builder = PresentationBuilder(template)

        for i, slide_content in enumerate(source_content):
            title = slide_content.get("title", f"Slide {i + 1}")
            text = slide_content.get("text_content", [])

            # Determine slide type based on content
            layout = slide_content.get("layout_name", "").lower()

            if "title" in layout or "frontpage" in layout:
                # Title slide
                subtitle = text[0] if text else ""
                builder.add_title_slide(title, subtitle)
            elif "section" in layout or "breaker" in layout:
                # Section divider
                builder.add_section(title)
            else:
                # Content slide
                bullets = text if text else []
                builder.add_content_slide(title, bullets=bullets)

        builder.save(output_path)

        # Test fidelity
        return self.test(source_path, output_path)


# =============================================================================
# Convenience Functions
# =============================================================================

def test_roundtrip(source_path: str, generated_path: str) -> RoundtripResult:
    """
    Quick round-trip test.

    Args:
        source_path: Path to original PPTX
        generated_path: Path to regenerated PPTX

    Returns:
        RoundtripResult
    """
    tester = RoundtripTester()
    return tester.test(source_path, generated_path)


def quick_fidelity_score(source_path: str, generated_path: str) -> float:
    """
    Get quick fidelity score for round-trip.

    Args:
        source_path: Path to original PPTX
        generated_path: Path to regenerated PPTX

    Returns:
        Fidelity score (0.0 - 1.0)
    """
    result = test_roundtrip(source_path, generated_path)
    return result.fidelity_score


# =============================================================================
# CLI
# =============================================================================

def main():
    """CLI for round-trip testing."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Round-Trip Fidelity Test")
    parser.add_argument("source", help="Source PPTX path")
    parser.add_argument("generated", nargs="?", help="Generated PPTX path (optional)")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    parser.add_argument("--template", default="consulting_toolkit", help="Template for regeneration")
    parser.add_argument("--visual", action="store_true", help="Enable visual comparison")

    args = parser.parse_args()

    tester = RoundtripTester(visual_comparison=args.visual)

    if args.generated:
        # Test existing pair
        result = tester.test(args.source, args.generated)
    else:
        # Regenerate and test
        result = tester.test_with_regeneration(args.source, template=args.template)

    if args.json:
        print(json.dumps(result.to_dict(), indent=2))
    else:
        print(result.summary())


if __name__ == "__main__":
    main()
