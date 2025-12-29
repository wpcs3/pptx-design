"""
Template Registry - Catalog of available templates with metadata.

Provides template discovery, layout information, and styling details.
"""

import json
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation

logger = logging.getLogger(__name__)


class TemplateRegistry:
    """
    Registry of available PowerPoint templates.

    Maintains a catalog of templates with their layouts, color palettes,
    fonts, and recommended use cases.
    """

    def __init__(self, registry_path: Optional[Path] = None):
        """
        Initialize the registry.

        Args:
            registry_path: Path to registry JSON file (auto-discovered if not provided)
        """
        self._templates_dir = Path(__file__).parent.parent / "pptx_templates"
        self._registry_path = registry_path or (
            Path(__file__).parent.parent / "config" / "template_registry.json"
        )
        self._registry: Dict[str, Any] = {}
        self._load_registry()

    def _load_registry(self) -> None:
        """Load registry from JSON file or build from templates."""
        if self._registry_path.exists():
            try:
                with open(self._registry_path, "r", encoding="utf-8") as f:
                    self._registry = json.load(f)
                logger.info(f"Loaded registry: {len(self._registry.get('templates', {}))} templates")
                return
            except Exception as e:
                logger.warning(f"Could not load registry: {e}")

        # Build registry from available templates
        self._build_registry()

    def _build_registry(self) -> None:
        """Build registry by scanning template directories."""
        self._registry = {"templates": {}, "version": "1.0"}

        if not self._templates_dir.exists():
            logger.warning(f"Templates directory not found: {self._templates_dir}")
            return

        for template_dir in self._templates_dir.iterdir():
            if not template_dir.is_dir():
                continue

            # Find .pptx files in the directory
            pptx_files = list(template_dir.glob("*.pptx"))
            if not pptx_files:
                continue

            template_path = pptx_files[0]
            template_name = self._normalize_name(template_dir.name)

            try:
                template_info = self._extract_template_info(template_path)
                template_info["path"] = str(template_path)
                template_info["directory"] = str(template_dir)
                self._registry["templates"][template_name] = template_info
                logger.debug(f"Indexed template: {template_name}")
            except Exception as e:
                logger.warning(f"Could not index {template_path}: {e}")

        logger.info(f"Built registry: {len(self._registry['templates'])} templates")

    def _normalize_name(self, name: str) -> str:
        """Normalize template name for lookup."""
        # Remove common prefixes
        name = name.replace("pptx_template_", "").replace("template_", "")
        # Convert to snake_case
        name = name.lower().replace("-", "_").replace(" ", "_")
        return name

    def _extract_template_info(self, template_path: Path) -> Dict[str, Any]:
        """Extract layout and style information from a template."""
        prs = Presentation(str(template_path))

        # Extract layouts
        layouts = []
        for layout in prs.slide_layouts:
            layout_info = {
                "name": layout.name,
                "placeholders": []
            }
            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type).split(".")[-1].strip("()")
                layout_info["placeholders"].append({
                    "idx": ph.placeholder_format.idx,
                    "type": ph_type,
                })
            layouts.append(layout_info)

        # Extract dimensions
        width_inches = round(prs.slide_width / 914400, 2)
        height_inches = round(prs.slide_height / 914400, 2)

        return {
            "layouts": layouts,
            "layout_names": [l["name"] for l in layouts],
            "dimensions": {
                "width_inches": width_inches,
                "height_inches": height_inches,
            },
            "slide_count": len(prs.slides),
        }

    def get_template(self, name: str) -> Optional[Dict[str, Any]]:
        """
        Get template information by name.

        Args:
            name: Template name (case-insensitive, partial match supported)

        Returns:
            Template info dict or None if not found
        """
        normalized = self._normalize_name(name)

        # Exact match
        if normalized in self._registry.get("templates", {}):
            return self._registry["templates"][normalized]

        # Partial match
        for template_name, info in self._registry.get("templates", {}).items():
            if normalized in template_name or template_name in normalized:
                return info

        return None

    def list_templates(self) -> List[str]:
        """List all available template names."""
        return list(self._registry.get("templates", {}).keys())

    def get_layouts(self, template_name: str) -> List[str]:
        """Get available layout names for a template."""
        info = self.get_template(template_name)
        if info:
            return info.get("layout_names", [])
        return []

    def find_template(
        self,
        use_case: Optional[str] = None,
        min_layouts: int = 0
    ) -> Optional[str]:
        """
        Find a suitable template based on criteria.

        Args:
            use_case: Desired use case (e.g., "consulting", "pitch", "report")
            min_layouts: Minimum number of layouts required

        Returns:
            Template name or None
        """
        for name, info in self._registry.get("templates", {}).items():
            # Check layout count
            if len(info.get("layouts", [])) < min_layouts:
                continue

            # Check use case
            if use_case:
                template_use_cases = info.get("use_cases", [])
                if use_case.lower() not in [u.lower() for u in template_use_cases]:
                    # Also check name
                    if use_case.lower() not in name.lower():
                        continue

            return name

        return None

    def save_registry(self, output_path: Optional[Path] = None) -> Path:
        """
        Save the registry to a JSON file.

        Args:
            output_path: Path to save to (default: config/template_registry.json)

        Returns:
            Path to saved file
        """
        output_path = output_path or self._registry_path
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(self._registry, f, indent=2)

        logger.info(f"Saved registry: {output_path}")
        return output_path

    def refresh(self) -> None:
        """Rebuild the registry from templates."""
        self._build_registry()

    def add_metadata(
        self,
        template_name: str,
        palette: List[str] = None,
        fonts: Dict[str, str] = None,
        use_cases: List[str] = None
    ) -> None:
        """
        Add metadata to a template entry.

        Args:
            template_name: Template to update
            palette: List of hex colors
            fonts: Dict of font styles (e.g., {"title": "Arial Bold 28pt"})
            use_cases: List of use case tags
        """
        normalized = self._normalize_name(template_name)
        if normalized not in self._registry.get("templates", {}):
            logger.warning(f"Template not found: {template_name}")
            return

        info = self._registry["templates"][normalized]
        if palette:
            info["palette"] = palette
        if fonts:
            info["fonts"] = fonts
        if use_cases:
            info["use_cases"] = use_cases

    def __len__(self) -> int:
        return len(self._registry.get("templates", {}))

    def __contains__(self, name: str) -> bool:
        return self.get_template(name) is not None

    def __repr__(self) -> str:
        return f"TemplateRegistry({len(self)} templates)"


def build_registry() -> TemplateRegistry:
    """
    Build and save a fresh template registry.

    Returns:
        TemplateRegistry instance
    """
    registry = TemplateRegistry()
    registry.refresh()
    registry.save_registry()
    return registry
