"""Generate presentation v26 with new Gemini images and complete source footnotes.

Features:
1. New Gemini-generated realistic section images (cropped to 11x8.5)
2. Complete source footnotes on all slides with sourced data
3. Section names in top-right placeholder
4. Table formatting with no borders
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Colors
COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text': RGBColor(0x06, 0x1F, 0x32),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
}

# Layout indices
LAYOUTS = {
    'frontpage': 0,
    'bullet_content': 1,
    'disclaimers': 2,
    'chart': 3,
    'table': 4,
    'side_by_side': 5,
    'section_title': 6,
}

# Slide type to layout mapping
SLIDE_TYPE_LAYOUT = {
    'title_slide': 'frontpage',
    'title_content': 'bullet_content',
    'content': 'bullet_content',
    'key_metrics': 'bullet_content',
    'table_slide': 'table',
    'data_chart': 'chart',
    'two_column': 'side_by_side',
    'section_divider': 'section_title',
    'disclaimer': 'disclaimers',
}

# Section to image mapping
SECTION_IMAGES = {
    'Executive Summary': 'title_slide.png',
    'Market Fundamentals': 'market_fundamentals.png',
    'Target Markets': 'target_markets.png',
    'Demand Drivers': 'demand_drivers.png',
    'Investment Strategy': 'investment_strategy.png',
    'Competitive Positioning': 'competitive_positioning.png',
    'Risk Management': 'risk_management.png',
    'ESG Strategy': 'esg_strategy.png',
    'JV Structure': 'jv_structure.png',
    'Conclusion': 'conclusion.png',
}

# COMPREHENSIVE SOURCE MAPPING - Based on research data in each slide
# Format: slide_number (1-indexed): [list of source citations]
SLIDE_SOURCES = {
    # Executive Summary
    2: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],  # Key metrics
    3: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],  # Thesis summary (12.4% returns, 3.4% vacancy)

    # Market Fundamentals
    5: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],  # Repricing context
    6: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],  # Market overview table
    7: ['CommercialCafe Dec 2025'],  # Light vs bulk comparison
    8: ['CommercialCafe Dec 2025'],  # Supply pipeline chart
    9: ['CBRE Cap Rate Survey H1 2024'],  # Cap rate evolution

    # Target Markets
    11: ['Partners RE 2024', 'REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025',
         'Savills Raleigh-Durham Q4 2024', 'Colliers Phoenix 2024'],  # Market matrix
    12: ['REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024'],  # Tier 1
    13: ['Partners RE DFW/Atlanta/San Antonio 2024', 'Colliers Phoenix 2024', 'Kidder Mathews Phoenix 2024'],  # Tier 2&3

    # Demand Drivers
    15: ['Clarion Partners Industrial Outlook 2025', 'NAIOP Nearshoring Analysis 2024'],  # Three drivers
    16: ['NAIOP Nearshoring Analysis 2024'],  # Manufacturing construction
    17: ['CBRE Interest Rate Impact 2024', 'RCLCO ODCE/NPI Q2 2025'],  # Rate sensitivity

    # Investment Strategy
    19: ['Clarion Partners Industrial Outlook 2025'],  # Portfolio construction
    20: ['REBusinessOnline Nashville 2024'],  # Nashville acquisition
    21: ['REBusinessOnline Nashville 2024'],  # Nashville case study
    22: ['RCLCO ODCE/NPI Q2 2025'],  # Return sensitivity
    23: ['NASRA FY 2024', 'RCLCO ODCE/NPI Q2 2025'],  # Sensitivity observations
    24: ['RCLCO ODCE/NPI Q2 2025'],  # Return expectations
    25: ['RCLCO ODCE/NPI Q2 2025'],  # Historical returns

    # Competitive Positioning
    27: ['Clarion Partners Industrial Outlook 2025'],  # Institutional landscape
    28: ['Clarion Partners Industrial Outlook 2025'],  # GP differentiation

    # Risk Management
    30: ['Clarion Partners Industrial Outlook 2025'],  # Risk matrix
    31: ['Clarion Partners Industrial Outlook 2025'],  # Defensive characteristics
    32: ['Cushman & Wakefield Q2 2025'],  # Mark-to-market
    33: ['Cushman & Wakefield Q2 2025'],  # NOI analysis

    # ESG Strategy
    35: ['GRESB 2025 Results'],  # ESG value creation
    36: ['GRESB 2025 Results'],  # ESG roadmap

    # JV Structure
    38: ['Clarion Partners Industrial Outlook 2025'],  # JV structure
    39: ['Clarion Partners Industrial Outlook 2025'],  # GP capabilities
    40: ['Clarion Partners Industrial Outlook 2025'],  # GP track record

    # Conclusion
    42: ['RCLCO ODCE/NPI Q2 2025', 'Clarion Partners Industrial Outlook 2025'],  # Investment summary
    43: ['RCLCO ODCE/NPI Q2 2025'],  # Partnership invitation
}


def remove_table_borders(table):
    """Remove all borders from a table."""
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                border = tcPr.find(qn(f'a:{border_name}'))
                if border is not None:
                    tcPr.remove(border)
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                etree.SubElement(border, qn('a:noFill'))


def format_table(table):
    """Apply table formatting."""
    for row_idx, row in enumerate(table.rows):
        row.height = Inches(0.50 if row_idx == 0 else 0.40)

        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Inches(0)
            cell.margin_bottom = Inches(0)
            cell.margin_right = Inches(0)
            cell.margin_left = Inches(0.1) if col_idx == 0 else Inches(0)

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = 'Arial'
                    if row_idx == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = COLORS['white']
                    else:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rgb = COLORS['dark_text']

    remove_table_borders(table)


def add_section_name(slide, section_name):
    """Add section name to idx=17 placeholder."""
    if not section_name:
        return
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 17:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = section_name
            run.font.name = 'Arial'
            run.font.size = Pt(10)
            run.font.color.rgb = COLORS['dark_text']
            return


def add_footnote(slide, sources):
    """Add footnote to idx=20 placeholder."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 20:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = footnote_text
            run.font.name = 'Arial'
            run.font.size = Pt(6)
            run.font.color.rgb = COLORS['dark_text']
            return


def add_section_image(slide, section_name, image_dir):
    """Add background image to section divider slide."""
    image_file = SECTION_IMAGES.get(section_name)
    if not image_file:
        return

    image_path = Path(image_dir) / image_file
    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return

    # Add as background covering full slide
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_bullets(text_frame, bullets):
    """Add standardized bullets."""
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.level = 0
        para.space_before = Pt(6)
        para.space_after = Pt(3)
        run = para.add_run()
        run.text = bullet
        run.font.name = 'Arial'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['dark_text']

        try:
            pPr = para._p.get_or_add_pPr()
            for child in list(pPr):
                if 'bu' in child.tag.lower():
                    pPr.remove(child)
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '\u2022')
        except:
            pass


def main():
    """Generate presentation v26."""
    print("Generating presentation v26 with new images and complete footnotes...")
    print("=" * 70)

    # Paths
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v24_CS_edits.pptx")
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    image_dir = Path("cache/images/cropped")
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v26.pptx")

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({'section': section.get('name', ''), **slide})

    print(f"Template: {len(prs.slides)} slides")
    print(f"Outline: {len(all_slides)} slides")
    print(f"Image directory: {image_dir}")

    # Process each slide
    print("\nProcessing slides...")
    footnotes_added = 0
    images_added = 0

    for idx in range(len(prs.slides)):
        if idx >= len(all_slides):
            break

        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')
        slide_type = slide_data.get('slide_type', '')
        slide_num = idx + 1

        # Add section name
        add_section_name(slide, section_name)

        # Add footnote if this slide has sources
        if slide_num in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[slide_num])
            footnotes_added += 1

        # Add section image for section dividers
        if slide_type == 'section_divider':
            add_section_image(slide, section_name, image_dir)
            images_added += 1
            print(f"  Slide {slide_num}: Added section image for '{section_name}'")

    # Update any tables to remove borders
    print("\nRemoving table borders...")
    tables_updated = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                format_table(shape.table)
                tables_updated += 1

    print(f"  Updated {tables_updated} tables")

    # Save
    prs.save(str(output_path))

    print(f"\n{'=' * 70}")
    print(f"Generated: {output_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Footnotes added: {footnotes_added}")
    print(f"Section images added: {images_added}")
    print(f"Tables formatted: {tables_updated}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless", "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
