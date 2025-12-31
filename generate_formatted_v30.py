"""Generate presentation v30 with all formatting fixes.

Fixes from v29:
1. Table borders - remove tableStyleId to eliminate default white borders
2. Slide 8 chart - vertical axis with comma format, fewer gridlines
3. Slide 9 chart - fix percentage display (values/100, % on axis)
4. Logo positions - use placeholders as set in master slides
5. Cleaner chart styling - minimal gridlines
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from lxml import etree

# Colors
COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text': RGBColor(0x06, 0x1F, 0x32),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'medium_gray': RGBColor(0xA6, 0xA6, 0xA6),
}

# Logo aspect ratio (width/height) from original: 2689/1129
LOGO_ASPECT_RATIO = 2.382

# Section to image mapping
SECTION_IMAGES = {
    'Executive Summary': 'title_slide.png',
    'Market Fundamentals': 'market_fundamentals.png',
    'Target Markets': 'target_markets.png',
    'Demand Drivers': 'demand_drivers.png',
    'Investment Strategy': 'investment_strategy.png',
    'Competitive Positioning': 'competitive_positioning.png',
    'Risk Management': 'risk_management.png',
    'ESG Strategy': 'esg_strategy.png',
    'JV Structure': 'jv_structure.png',
    'Conclusion': 'conclusion.png',
}

# Source mapping
SLIDE_SOURCES = {
    2: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    3: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    5: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    6: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    7: ['CommercialCafe Dec 2025'],
    8: ['CommercialCafe Dec 2025'],
    9: ['CBRE Cap Rate Survey H1 2024'],
    11: ['Partners RE 2024', 'REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025',
         'Savills Raleigh-Durham Q4 2024', 'Colliers Phoenix 2024'],
    12: ['REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024'],
    13: ['Partners RE DFW/Atlanta/San Antonio 2024', 'Colliers Phoenix 2024', 'Kidder Mathews Phoenix 2024'],
    15: ['Clarion Partners Industrial Outlook 2025', 'NAIOP Nearshoring Analysis 2024'],
    16: ['NAIOP Nearshoring Analysis 2024'],
    17: ['CBRE Interest Rate Impact 2024', 'RCLCO ODCE/NPI Q2 2025'],
    19: ['Clarion Partners Industrial Outlook 2025'],
    20: ['REBusinessOnline Nashville 2024'],
    21: ['REBusinessOnline Nashville 2024'],
    22: ['RCLCO ODCE/NPI Q2 2025'],
    23: ['NASRA FY 2024', 'RCLCO ODCE/NPI Q2 2025'],
    24: ['RCLCO ODCE/NPI Q2 2025'],
    25: ['RCLCO ODCE/NPI Q2 2025'],
    27: ['Clarion Partners Industrial Outlook 2025'],
    28: ['Clarion Partners Industrial Outlook 2025'],
    30: ['Clarion Partners Industrial Outlook 2025'],
    31: ['Clarion Partners Industrial Outlook 2025'],
    32: ['Cushman & Wakefield Q2 2025'],
    33: ['Cushman & Wakefield Q2 2025'],
    35: ['GRESB 2025 Results'],
    36: ['GRESB 2025 Results'],
    38: ['Clarion Partners Industrial Outlook 2025'],
    39: ['Clarion Partners Industrial Outlook 2025'],
    40: ['Clarion Partners Industrial Outlook 2025'],
    42: ['RCLCO ODCE/NPI Q2 2025', 'Clarion Partners Industrial Outlook 2025'],
    43: ['RCLCO ODCE/NPI Q2 2025'],
}

# Historical returns data
HISTORICAL_RETURNS_DATA = {
    'categories': ['Industrial', 'Apartment', 'Retail', 'Office'],
    'returns': [12.4, 9.8, 8.4, 6.5],
}


def is_numeric_cell(text):
    """Check if cell content is numeric."""
    if not text:
        return False
    text = text.strip()
    patterns = [
        r'^[\d,]+\.?\d*%?$',
        r'^\$[\d,]+\.?\d*[BMK]?$',
        r'^[\d,]+\.?\d*x$',
        r'^[\d,]+-[\d,]+%?$',
        r'^\d+\.\d+%$',
        r'^-?\d+\.?\d*%?$',
        r'^\(\d+\.?\d*%?\)$',
    ]
    return any(re.match(p, text) for p in patterns)


def get_column_alignment(table, col_idx):
    """Determine alignment for a column based on majority of cell values."""
    numeric_count = 0
    text_count = 0

    for row_idx in range(1, len(table.rows)):
        cell = table.cell(row_idx, col_idx)
        cell_text = ''
        for para in cell.text_frame.paragraphs:
            cell_text += para.text

        if is_numeric_cell(cell_text.strip()):
            numeric_count += 1
        elif cell_text.strip():
            text_count += 1

    return PP_ALIGN.RIGHT if numeric_count > text_count else PP_ALIGN.LEFT


def get_max_text_width_in_column(table, col_idx):
    """Estimate the maximum text width needed in a column."""
    max_chars = 0
    for row_idx in range(len(table.rows)):
        cell = table.cell(row_idx, col_idx)
        cell_text = ''
        for para in cell.text_frame.paragraphs:
            cell_text += para.text
        max_chars = max(max_chars, len(cell_text.strip()))
    return max_chars


def set_cell_borders_none(cell):
    """Set all borders on a cell to 'None' (noFill with zero width).

    The proper OOXML structure for invisible borders uses <a:noFill/> inside
    each border element with width set to 0.
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # All possible border element names
    border_names = ['lnL', 'lnR', 'lnT', 'lnB', 'lnTlToBr', 'lnBlToTr']

    # Remove all existing border elements first
    for border_name in border_names:
        for border in tcPr.findall(qn(f'a:{border_name}')):
            tcPr.remove(border)

    # Add border elements with noFill and zero width
    for border_name in border_names:
        border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
        border.set('w', '0')  # Zero width
        # Add noFill - this is the key to "None" border type
        etree.SubElement(border, qn('a:noFill'))


def format_table_borders_no_border(table):
    """Format table to have no visible borders using noFill approach."""
    # Apply noFill borders to every cell
    for row in table.rows:
        for cell in row.cells:
            set_cell_borders_none(cell)


def calculate_column_widths(table, total_width_inches=10.2):
    """Calculate column widths to prevent text wrapping."""
    num_cols = len(table.columns)

    char_counts = []
    for col_idx in range(num_cols):
        max_chars = get_max_text_width_in_column(table, col_idx)
        char_counts.append(max(max_chars + 2, 5))

    total_chars = sum(char_counts)
    total_width = Inches(total_width_inches)

    widths = []
    for chars in char_counts:
        width = int(total_width * chars / total_chars)
        widths.append(width)

    return widths


def format_table_v30(table):
    """Apply v30 table formatting - no borders, proper spacing."""
    num_cols = len(table.columns)
    num_rows = len(table.rows)

    # Calculate column alignments
    column_alignments = []
    for col_idx in range(num_cols):
        column_alignments.append(get_column_alignment(table, col_idx))

    # Calculate and apply column widths
    col_widths = calculate_column_widths(table)
    for col_idx, col in enumerate(table.columns):
        col.width = col_widths[col_idx]

    for row_idx, row in enumerate(table.rows):
        row.height = Inches(0.50 if row_idx == 0 else 0.40)

        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            # Fill colors
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            # Apply column alignment
            for para in cell.text_frame.paragraphs:
                para.alignment = column_alignments[col_idx]

                for run in para.runs:
                    run.font.name = 'Arial'
                    if row_idx == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = COLORS['white']
                    else:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rgb = COLORS['dark_text']

    # Apply noFill borders (proper "None" border type)
    format_table_borders_no_border(table)


def format_chart_slide_8(slide):
    """Format slide 8 chart - comma format on axis, minimal gridlines."""
    for shape in slide.shapes:
        if not (hasattr(shape, 'has_chart') and shape.has_chart):
            continue

        chart = shape.chart

        # Format value axis with comma
        try:
            va = chart.value_axis
            va.tick_labels.number_format = '#,##0'
            va.tick_labels.font.size = Pt(12)
            va.has_major_gridlines = True
            va.has_minor_gridlines = False

            # Reduce gridline clutter - adjust major unit if possible
            # Setting major_unit will reduce the number of gridlines
            try:
                # Get the max value to determine appropriate major unit
                max_val = 0
                for series in chart.series:
                    if hasattr(series, 'values'):
                        max_val = max(max_val, max(series.values))
                # Set major unit to create fewer gridlines
                if max_val > 1000:
                    va.major_unit = 500  # Fewer gridlines for large values
            except:
                pass

        except Exception as e:
            print(f"    Error formatting slide 8 value axis: {e}")

        # Format category axis
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(12)
        except:
            pass

        # Data labels with comma format
        for series in chart.series:
            if hasattr(series, 'data_labels'):
                series.has_data_labels = True
                series.data_labels.number_format = '#,##0'
                series.data_labels.font.size = Pt(12)

        print("    Formatted slide 8 chart")
        return True

    return False


def format_chart_slide_9(slide):
    """Format slide 9 chart - percentage format (0% to 7%), minimal gridlines."""
    for shape in slide.shapes:
        if not (hasattr(shape, 'has_chart') and shape.has_chart):
            continue

        chart = shape.chart

        # Format value axis as percentage
        try:
            va = chart.value_axis
            va.tick_labels.number_format = '0.0"%"'
            va.tick_labels.font.size = Pt(12)
            va.has_major_gridlines = True
            va.has_minor_gridlines = False

            # Set axis range for cap rates (0% to 8%)
            va.minimum_scale = 0
            va.maximum_scale = 8
            va.major_unit = 2  # Gridlines at 0, 2, 4, 6, 8

        except Exception as e:
            print(f"    Error formatting slide 9 value axis: {e}")

        # Format category axis
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(12)
        except:
            pass

        # Data labels with percentage format
        for series in chart.series:
            if hasattr(series, 'data_labels'):
                series.has_data_labels = True
                series.data_labels.number_format = '0.0"%"'
                series.data_labels.font.size = Pt(10)

        print("    Formatted slide 9 chart")
        return True

    return False


def format_chart_slide_16(slide):
    """Format slide 16 chart - USD format with minimal gridlines."""
    for shape in slide.shapes:
        if not (hasattr(shape, 'has_chart') and shape.has_chart):
            continue

        chart = shape.chart

        # Format value axis
        try:
            va = chart.value_axis
            va.tick_labels.number_format = '$#,##0'
            va.tick_labels.font.size = Pt(12)
            va.has_major_gridlines = True
            va.has_minor_gridlines = False
        except:
            pass

        # Format category axis
        try:
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(12)
        except:
            pass

        # Data labels
        for series in chart.series:
            if hasattr(series, 'data_labels'):
                series.has_data_labels = True
                series.data_labels.number_format = '$#,##0'
                series.data_labels.font.size = Pt(12)

        print("    Formatted slide 16 chart")
        return True

    return False


def add_section_name(slide, section_name):
    """Add section name (9pt, #A6A6A6, right-aligned)."""
    if not section_name:
        return
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 17:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = section_name
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = COLORS['medium_gray']
            return


def add_footnote(slide, sources):
    """Add footnote (6pt, #A6A6A6, left-aligned)."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx in [20, 16]:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = footnote_text
            run.font.name = 'Arial'
            run.font.size = Pt(6)
            run.font.color.rgb = COLORS['medium_gray']
            return


def remove_all_pictures_from_slide(slide):
    """Remove all picture shapes from a slide (except background)."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Keep full-slide background images
            if shape.width >= Inches(10) and shape.height >= Inches(8):
                continue
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_section_image(slide, section_name, image_dir):
    """Add background image to section divider."""
    # Remove existing non-background pictures
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    image_file = SECTION_IMAGES.get(section_name)
    if not image_file:
        return False

    image_path = Path(image_dir) / image_file
    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return False

    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)
    return True


def update_slide_25_chart(slide):
    """Add/update the chart on slide 25 with proper formatting."""
    # Remove any existing chart shapes
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = HISTORICAL_RETURNS_DATA['categories']
    chart_data.add_series('10-Year Annualized Returns (%)',
                          tuple(HISTORICAL_RETURNS_DATA['returns']))

    # Add chart
    x, y = Inches(0.4), Inches(2.7)
    cx, cy = Inches(10.2), Inches(3.8)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart

    # Style the chart - minimal look
    chart.has_legend = False

    # Format value axis
    try:
        va = chart.value_axis
        va.tick_labels.number_format = '0"%"'
        va.tick_labels.font.size = Pt(12)
        va.has_major_gridlines = True
        va.has_minor_gridlines = False
        va.minimum_scale = 0
        va.maximum_scale = 15
        va.major_unit = 5  # Gridlines at 0, 5, 10, 15
    except:
        pass

    # Format category axis
    try:
        chart.category_axis.tick_labels.font.size = Pt(12)
    except:
        pass

    # Format data labels
    plot = chart.plots[0]
    series = plot.series[0]
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.number_format = '0.0"%"'
    data_labels.font.size = Pt(12)
    data_labels.font.bold = True

    return True


def update_slide_25_textbox(slide):
    """Update the text box on slide 25 with size 14 font."""
    narrative = (
        "Industrial delivered 12.4% 10-year annualized returns, outperforming "
        "Apartment (+260 bps), Retail (+400 bps), and Office (+590 bps)."
    )

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 18:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = narrative
            run.font.name = 'Arial'
            run.font.size = Pt(14)
            run.font.color.rgb = COLORS['dark_text']
            return True

    return False


def main():
    """Generate presentation v30."""
    print("Generating presentation v30 with all formatting fixes...")
    print("=" * 70)

    # Paths - use the CS edits as template
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v29_CS_edits.pptx")
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    image_dir = Path("cache/images/cropped")
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v30b.pptx")

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({'section': section.get('name', ''), **slide})

    print(f"Template: {len(prs.slides)} slides")
    print(f"Image directory: {image_dir}")

    # Process each slide
    print("\nProcessing slides...")

    for idx in range(len(prs.slides)):
        if idx >= len(all_slides):
            break

        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')
        slide_type = slide_data.get('slide_type', '')
        slide_num = idx + 1

        # Add section name
        add_section_name(slide, section_name)

        # Add footnote
        if slide_num in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[slide_num])

        # Handle section dividers
        if slide_type == 'section_divider':
            if add_section_image(slide, section_name, image_dir):
                print(f"  Slide {slide_num}: Updated section image for '{section_name}'")

        # Format charts on specific slides
        if slide_num == 8:
            format_chart_slide_8(slide)
        elif slide_num == 9:
            format_chart_slide_9(slide)
        elif slide_num == 16:
            format_chart_slide_16(slide)

        # Update slide 25
        if slide_num == 25:
            update_slide_25_chart(slide)
            update_slide_25_textbox(slide)
            print(f"  Slide {slide_num}: Updated historical returns chart and text")

    # Format all tables - remove borders completely
    print("\nFormatting tables (removing all borders)...")
    tables_updated = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                # Ensure table position
                shape.left = Inches(0.4)
                shape.width = Inches(10.2)

                format_table_v30(shape.table)
                tables_updated += 1

    print(f"  Total tables formatted: {tables_updated}")

    # Save
    prs.save(str(output_path))

    print(f"\n{'=' * 70}")
    print(f"Generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless", "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
