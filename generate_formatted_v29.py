"""Generate presentation v29 with all formatting fixes.

Fixes from v28:
1. Logo aspect ratio - maintain 2.382 ratio (width/height) from original
2. Table column spacing - auto-size to prevent text wrapping
3. Table borders - use different approach to ensure no borders
4. Chart number formatting - commas for thousands, % for percentages, $ for USD
5. Slide 25 chart formatting - match other bar charts
6. Contact slide 44 - no modifications (keep original)
7. End slide - no duplicates, white logo with correct aspect ratio
8. Delete blue logos, keep only white
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from lxml import etree
import copy

# Colors
COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text': RGBColor(0x06, 0x1F, 0x32),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'medium_gray': RGBColor(0xA6, 0xA6, 0xA6),
}

# Logo aspect ratio (width/height) from original blue logo: 2689/1129
LOGO_ASPECT_RATIO = 2.382

# Section to image mapping
SECTION_IMAGES = {
    'Executive Summary': 'title_slide.png',
    'Market Fundamentals': 'market_fundamentals.png',
    'Target Markets': 'target_markets.png',
    'Demand Drivers': 'demand_drivers.png',
    'Investment Strategy': 'investment_strategy.png',
    'Competitive Positioning': 'competitive_positioning.png',
    'Risk Management': 'risk_management.png',
    'ESG Strategy': 'esg_strategy.png',
    'JV Structure': 'jv_structure.png',
    'Conclusion': 'conclusion.png',
}

# Source mapping
SLIDE_SOURCES = {
    2: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    3: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    5: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    6: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    7: ['CommercialCafe Dec 2025'],
    8: ['CommercialCafe Dec 2025'],
    9: ['CBRE Cap Rate Survey H1 2024'],
    11: ['Partners RE 2024', 'REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025',
         'Savills Raleigh-Durham Q4 2024', 'Colliers Phoenix 2024'],
    12: ['REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024'],
    13: ['Partners RE DFW/Atlanta/San Antonio 2024', 'Colliers Phoenix 2024', 'Kidder Mathews Phoenix 2024'],
    15: ['Clarion Partners Industrial Outlook 2025', 'NAIOP Nearshoring Analysis 2024'],
    16: ['NAIOP Nearshoring Analysis 2024'],
    17: ['CBRE Interest Rate Impact 2024', 'RCLCO ODCE/NPI Q2 2025'],
    19: ['Clarion Partners Industrial Outlook 2025'],
    20: ['REBusinessOnline Nashville 2024'],
    21: ['REBusinessOnline Nashville 2024'],
    22: ['RCLCO ODCE/NPI Q2 2025'],
    23: ['NASRA FY 2024', 'RCLCO ODCE/NPI Q2 2025'],
    24: ['RCLCO ODCE/NPI Q2 2025'],
    25: ['RCLCO ODCE/NPI Q2 2025'],
    27: ['Clarion Partners Industrial Outlook 2025'],
    28: ['Clarion Partners Industrial Outlook 2025'],
    30: ['Clarion Partners Industrial Outlook 2025'],
    31: ['Clarion Partners Industrial Outlook 2025'],
    32: ['Cushman & Wakefield Q2 2025'],
    33: ['Cushman & Wakefield Q2 2025'],
    35: ['GRESB 2025 Results'],
    36: ['GRESB 2025 Results'],
    38: ['Clarion Partners Industrial Outlook 2025'],
    39: ['Clarion Partners Industrial Outlook 2025'],
    40: ['Clarion Partners Industrial Outlook 2025'],
    42: ['RCLCO ODCE/NPI Q2 2025', 'Clarion Partners Industrial Outlook 2025'],
    43: ['RCLCO ODCE/NPI Q2 2025'],
}

# Historical returns data for slide 25 chart (10-year annualized returns)
HISTORICAL_RETURNS_DATA = {
    'categories': ['Industrial', 'Apartment', 'Retail', 'Office'],
    'returns': [12.4, 9.8, 8.4, 6.5],
}


def is_numeric_cell(text):
    """Check if cell content is numeric."""
    if not text:
        return False
    text = text.strip()
    patterns = [
        r'^[\d,]+\.?\d*%?$',
        r'^\$[\d,]+\.?\d*[BMK]?$',
        r'^[\d,]+\.?\d*x$',
        r'^[\d,]+-[\d,]+%?$',
        r'^\d+\.\d+%$',
        r'^-?\d+\.?\d*%?$',
        r'^\(\d+\.?\d*%?\)$',
    ]
    return any(re.match(p, text) for p in patterns)


def get_column_alignment(table, col_idx):
    """Determine alignment for a column based on majority of cell values."""
    numeric_count = 0
    text_count = 0

    for row_idx in range(1, len(table.rows)):
        cell = table.cell(row_idx, col_idx)
        cell_text = ''
        for para in cell.text_frame.paragraphs:
            cell_text += para.text

        if is_numeric_cell(cell_text.strip()):
            numeric_count += 1
        elif cell_text.strip():
            text_count += 1

    return PP_ALIGN.RIGHT if numeric_count > text_count else PP_ALIGN.LEFT


def get_max_text_width_in_column(table, col_idx):
    """Estimate the maximum text width needed in a column."""
    max_chars = 0
    for row_idx in range(len(table.rows)):
        cell = table.cell(row_idx, col_idx)
        cell_text = ''
        for para in cell.text_frame.paragraphs:
            cell_text += para.text
        max_chars = max(max_chars, len(cell_text.strip()))
    return max_chars


def remove_all_table_borders(table):
    """Remove ALL borders from table using XML manipulation."""
    tbl = table._tbl

    # Remove table-level borders
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        # Remove any existing border elements at table level
        for elem in tblPr.findall('.//' + qn('a:ln')):
            elem.getparent().remove(elem)

    # Process each cell
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            # Remove all existing border elements
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                for border in tcPr.findall(qn(f'a:{border_name}')):
                    tcPr.remove(border)

            # Add explicit no-border elements
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                border.set('w', '0')
                border.set('cap', 'flat')
                border.set('cmpd', 'sng')
                border.set('algn', 'ctr')
                noFill = etree.SubElement(border, qn('a:noFill'))


def calculate_column_widths(table, total_width_inches=10.2):
    """Calculate column widths to prevent text wrapping."""
    num_cols = len(table.columns)

    # Get character counts for each column
    char_counts = []
    for col_idx in range(num_cols):
        max_chars = get_max_text_width_in_column(table, col_idx)
        # Minimum of 5 chars, add padding
        char_counts.append(max(max_chars + 2, 5))

    total_chars = sum(char_counts)
    total_width = Inches(total_width_inches)

    # Distribute width proportionally
    widths = []
    for chars in char_counts:
        width = int(total_width * chars / total_chars)
        widths.append(width)

    return widths


def format_table_v29(table):
    """Apply v29 table formatting with improved column spacing."""
    num_cols = len(table.columns)
    num_rows = len(table.rows)

    # Calculate column alignments
    column_alignments = []
    for col_idx in range(num_cols):
        column_alignments.append(get_column_alignment(table, col_idx))

    # Calculate and apply column widths to prevent wrapping
    col_widths = calculate_column_widths(table)
    for col_idx, col in enumerate(table.columns):
        col.width = col_widths[col_idx]

    for row_idx, row in enumerate(table.rows):
        row.height = Inches(0.50 if row_idx == 0 else 0.40)

        for col_idx, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            # Fill colors
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            # Apply column alignment
            for para in cell.text_frame.paragraphs:
                para.alignment = column_alignments[col_idx]

                for run in para.runs:
                    run.font.name = 'Arial'
                    if row_idx == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = COLORS['white']
                    else:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rgb = COLORS['dark_text']

    # Remove all borders
    remove_all_table_borders(table)


def format_chart_numbers(slide, slide_num):
    """Format chart numbers based on slide content."""
    for shape in slide.shapes:
        if not (hasattr(shape, 'has_chart') and shape.has_chart):
            continue

        chart = shape.chart

        # Slide 8: Large numbers need commas
        if slide_num == 8:
            for series in chart.series:
                if hasattr(series, 'data_labels'):
                    series.has_data_labels = True
                    series.data_labels.number_format = '#,##0'
                    series.data_labels.font.size = Pt(12)

        # Slide 9: Percentages
        elif slide_num == 9:
            for series in chart.series:
                if hasattr(series, 'data_labels'):
                    series.has_data_labels = True
                    series.data_labels.number_format = '0.0%'
                    series.data_labels.font.size = Pt(12)

        # Slide 16: USD values
        elif slide_num == 16:
            for series in chart.series:
                if hasattr(series, 'data_labels'):
                    series.has_data_labels = True
                    series.data_labels.number_format = '$#,##0'
                    series.data_labels.font.size = Pt(12)

            # Try to format value axis
            try:
                if chart.value_axis:
                    chart.value_axis.tick_labels.number_format = '$#,##0'
                    chart.value_axis.tick_labels.font.size = Pt(12)
            except:
                pass


def add_section_name(slide, section_name):
    """Add section name (9pt, #A6A6A6, right-aligned)."""
    if not section_name:
        return
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 17:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = section_name
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = COLORS['medium_gray']
            return


def add_footnote(slide, sources):
    """Add footnote (6pt, #A6A6A6, left-aligned)."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx in [20, 16]:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = footnote_text
            run.font.name = 'Arial'
            run.font.size = Pt(6)
            run.font.color.rgb = COLORS['medium_gray']
            return


def remove_all_pictures_from_slide(slide):
    """Remove all picture shapes from a slide."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_section_image(slide, section_name, image_dir):
    """Add background image to section divider."""
    remove_all_pictures_from_slide(slide)

    image_file = SECTION_IMAGES.get(section_name)
    if not image_file:
        return False

    image_path = Path(image_dir) / image_file
    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return False

    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)
    return True


def add_logo_with_aspect_ratio(slide, logo_path, left, top, width):
    """Add logo maintaining correct aspect ratio."""
    if not logo_path.exists():
        print(f"    Warning: Logo not found: {logo_path}")
        return None

    # Calculate height based on aspect ratio
    height = width / LOGO_ASPECT_RATIO

    pic = slide.shapes.add_picture(
        str(logo_path),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    return pic


def update_slide_25_chart(slide):
    """Add/update the chart on slide 25 with proper formatting."""
    # Remove any existing chart shapes first
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = HISTORICAL_RETURNS_DATA['categories']
    chart_data.add_series('10-Year Annualized Returns (%)',
                          tuple(HISTORICAL_RETURNS_DATA['returns']))

    # Add chart
    x, y = Inches(0.4), Inches(2.7)
    cx, cy = Inches(10.2), Inches(3.8)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = chart_shape.chart

    # Style the chart
    chart.has_legend = False

    # Format data labels
    plot = chart.plots[0]
    series = plot.series[0]
    series.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.number_format = '0.0"%"'
    data_labels.font.size = Pt(12)
    data_labels.font.bold = True
    data_labels.font.color.rgb = COLORS['dark_text']

    # Format category axis labels
    try:
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.color.rgb = COLORS['dark_text']
    except:
        pass

    # Format value axis
    try:
        chart.value_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.number_format = '0"%"'
        chart.value_axis.has_major_gridlines = True
    except:
        pass

    return True


def update_slide_25_textbox(slide):
    """Update the text box on slide 25 with size 14 font."""
    narrative = (
        "Industrial delivered 12.4% 10-year annualized returns, outperforming "
        "Apartment (+260 bps), Retail (+400 bps), and Office (+590 bps)."
    )

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 18:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = narrative
            run.font.name = 'Arial'
            run.font.size = Pt(14)  # Size 14 per Chart master slide
            run.font.color.rgb = COLORS['dark_text']
            return True

    return False


def main():
    """Generate presentation v29."""
    print("Generating presentation v29 with all formatting fixes...")
    print("=" * 70)

    # Paths
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v27_CS_edits.pptx")
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    image_dir = Path("cache/images/cropped")
    logo_path = Path("logos/pccp_logo_white.png")
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v29.pptx")

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({'section': section.get('name', ''), **slide})

    print(f"Template: {len(prs.slides)} slides")
    print(f"Image directory: {image_dir}")
    print(f"Logo: {logo_path}")

    # Process each slide
    print("\nProcessing slides...")

    for idx in range(len(prs.slides)):
        if idx >= len(all_slides):
            break

        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')
        slide_type = slide_data.get('slide_type', '')
        slide_num = idx + 1

        # Add section name
        add_section_name(slide, section_name)

        # Add footnote
        if slide_num in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[slide_num])

        # Handle section dividers
        if slide_type == 'section_divider':
            if add_section_image(slide, section_name, image_dir):
                print(f"  Slide {slide_num}: Updated section image for '{section_name}'")

        # Add white logo to title slide (slide 1)
        if slide_num == 1:
            # Remove any existing pictures that might be logos (small ones at top-left)
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Check if it's in the logo area (top-left, small)
                    if shape.left < Inches(3) and shape.top < Inches(2) and shape.width < Inches(4):
                        shapes_to_remove.append(shape)

            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)

            # Add white logo with correct aspect ratio
            add_logo_with_aspect_ratio(slide, logo_path, 0.4, 0.4, 2.5)
            print(f"  Slide {slide_num}: Added white logo to title slide")

        # Format charts on specific slides
        if slide_num in [8, 9, 16]:
            format_chart_numbers(slide, slide_num)
            print(f"  Slide {slide_num}: Formatted chart numbers")

        # Update slide 25 with chart and textbox
        if slide_num == 25:
            update_slide_25_chart(slide)
            update_slide_25_textbox(slide)
            print(f"  Slide {slide_num}: Updated historical returns chart and text")

        # Contact slide 44 - DO NOT MODIFY (keep original)
        # Skipping any modifications

    # Format all tables
    print("\nFormatting tables...")
    tables_updated = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                # Ensure table position
                shape.left = Inches(0.4)
                shape.width = Inches(10.2)

                format_table_v29(shape.table)
                tables_updated += 1

    print(f"  Total tables formatted: {tables_updated}")

    # Handle End slide (slide 46)
    # First, check if we need to update the logo
    if len(prs.slides) >= 46:
        end_slide = prs.slides[45]  # 0-indexed, slide 46

        # Remove any existing logo images (small pictures)
        shapes_to_remove = []
        for shape in end_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Check if it's a small image (logo-sized), not the background
                if shape.width < Inches(5) and shape.height < Inches(3):
                    shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # Add white logo centered with correct aspect ratio
        # Center position: (11 - logo_width) / 2 = (11 - 2.5) / 2 = 4.25
        add_logo_with_aspect_ratio(end_slide, logo_path, 4.25, 3.6, 2.5)
        print("  Slide 46: Updated End slide with white logo")

    # Delete slide 47 if it exists (duplicate End slide)
    if len(prs.slides) > 46:
        print("\nDeleting duplicate End slide (47)...")
        # Get the slide to delete
        slide_to_delete = prs.slides[46]  # 0-indexed, slide 47
        rId = prs.part.relate_to(slide_to_delete.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[46]
        print("  Deleted slide 47")

    # Save
    prs.save(str(output_path))

    print(f"\n{'=' * 70}")
    print(f"Generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless", "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
