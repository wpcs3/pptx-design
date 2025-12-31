"""Generate presentation v27 with all formatting fixes.

Fixes:
1. Section dividers: Remove old images, add new small bay industrial images
2. Section names: top-right aligned, 9pt, #A6A6A6
3. Footnotes: bottom-left aligned (0.4" from left), 6pt, #A6A6A6
4. Table columns: left-align text, right-align numbers, center vertically
5. Add End slide with Gemini image and centered white PCCP logo
6. Add PCCP logo to title slide (top-left, white)
7. Update slide 44 with PCCP contact info
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# Colors
COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'dark_text': RGBColor(0x06, 0x1F, 0x32),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'medium_gray': RGBColor(0xA6, 0xA6, 0xA6),  # For section names and footnotes
}

# Section to image mapping
SECTION_IMAGES = {
    'Executive Summary': 'title_slide.png',
    'Market Fundamentals': 'market_fundamentals.png',
    'Target Markets': 'target_markets.png',
    'Demand Drivers': 'demand_drivers.png',
    'Investment Strategy': 'investment_strategy.png',
    'Competitive Positioning': 'competitive_positioning.png',
    'Risk Management': 'risk_management.png',
    'ESG Strategy': 'esg_strategy.png',
    'JV Structure': 'jv_structure.png',
    'Conclusion': 'conclusion.png',
}

# Source mapping
SLIDE_SOURCES = {
    2: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    3: ['RCLCO ODCE/NPI Q2 2025', 'CommercialCafe Dec 2025'],
    5: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    6: ['CommercialCafe Dec 2025', 'Cushman & Wakefield Q2 2025'],
    7: ['CommercialCafe Dec 2025'],
    8: ['CommercialCafe Dec 2025'],
    9: ['CBRE Cap Rate Survey H1 2024'],
    11: ['Partners RE 2024', 'REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024', 'Colliers Phoenix 2024'],
    12: ['REBusinessOnline Nashville 2024', 'WareCRE Tampa 2025', 'Savills Raleigh-Durham Q4 2024'],
    13: ['Partners RE DFW/Atlanta/San Antonio 2024', 'Colliers Phoenix 2024', 'Kidder Mathews Phoenix 2024'],
    15: ['Clarion Partners Industrial Outlook 2025', 'NAIOP Nearshoring Analysis 2024'],
    16: ['NAIOP Nearshoring Analysis 2024'],
    17: ['CBRE Interest Rate Impact 2024', 'RCLCO ODCE/NPI Q2 2025'],
    19: ['Clarion Partners Industrial Outlook 2025'],
    20: ['REBusinessOnline Nashville 2024'],
    21: ['REBusinessOnline Nashville 2024'],
    22: ['RCLCO ODCE/NPI Q2 2025'],
    23: ['NASRA FY 2024', 'RCLCO ODCE/NPI Q2 2025'],
    24: ['RCLCO ODCE/NPI Q2 2025'],
    25: ['RCLCO ODCE/NPI Q2 2025'],
    27: ['Clarion Partners Industrial Outlook 2025'],
    28: ['Clarion Partners Industrial Outlook 2025'],
    30: ['Clarion Partners Industrial Outlook 2025'],
    31: ['Clarion Partners Industrial Outlook 2025'],
    32: ['Cushman & Wakefield Q2 2025'],
    33: ['Cushman & Wakefield Q2 2025'],
    35: ['GRESB 2025 Results'],
    36: ['GRESB 2025 Results'],
    38: ['Clarion Partners Industrial Outlook 2025'],
    39: ['Clarion Partners Industrial Outlook 2025'],
    40: ['Clarion Partners Industrial Outlook 2025'],
    42: ['RCLCO ODCE/NPI Q2 2025', 'Clarion Partners Industrial Outlook 2025'],
    43: ['RCLCO ODCE/NPI Q2 2025'],
}

# PCCP Contact Info
PCCP_CONTACT = {
    'name': 'PCCP, LLC',
    'headquarters': {
        'address': '10100 Santa Monica Blvd., Suite 1000',
        'city_state_zip': 'Los Angeles, CA 90067',
        'phone': '310.414.7870',
    },
    'email': 'Inquiry@pccpllc.com',
    'website': 'pccpllc.com',
    'other_offices': [
        {'city': 'New York', 'phone': '646.308.2100'},
        {'city': 'San Francisco', 'phone': '415.732.7645'},
        {'city': 'Atlanta', 'phone': '404.947.6080'},
    ]
}


def is_numeric_cell(text):
    """Check if cell content is numeric (for right-alignment)."""
    if not text:
        return False
    text = text.strip()
    # Match numbers, percentages, currency, multiples
    patterns = [
        r'^[\d,]+\.?\d*%?$',  # Numbers with optional %
        r'^\$[\d,]+\.?\d*[BMK]?$',  # Currency
        r'^[\d,]+\.?\d*x$',  # Multiples (1.5x)
        r'^[\d,]+-[\d,]+%?$',  # Ranges (10-15%)
        r'^\d+\.\d+%$',  # Decimal percentages
    ]
    return any(re.match(p, text) for p in patterns)


def remove_table_borders(table):
    """Remove all borders from a table."""
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                border = tcPr.find(qn(f'a:{border_name}'))
                if border is not None:
                    tcPr.remove(border)
                border = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                etree.SubElement(border, qn('a:noFill'))


def format_table(table):
    """Apply table formatting with proper alignment."""
    num_cols = len(table.columns)

    for row_idx, row in enumerate(table.rows):
        row.height = Inches(0.50 if row_idx == 0 else 0.40)

        for col_idx, cell in enumerate(row.cells):
            # Always center vertically
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Cell margins
            cell.margin_top = Inches(0)
            cell.margin_bottom = Inches(0)
            cell.margin_right = Inches(0)
            cell.margin_left = Inches(0.1) if col_idx == 0 else Inches(0)

            # Fill colors
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

            # Format text with proper alignment
            for para in cell.text_frame.paragraphs:
                cell_text = para.text.strip() if para.text else ''

                # Determine alignment: numbers right, text left
                if row_idx == 0:
                    # Header row - left align
                    para.alignment = PP_ALIGN.LEFT
                elif is_numeric_cell(cell_text):
                    para.alignment = PP_ALIGN.RIGHT
                else:
                    para.alignment = PP_ALIGN.LEFT

                for run in para.runs:
                    run.font.name = 'Arial'
                    if row_idx == 0:
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = COLORS['white']
                    else:
                        run.font.size = Pt(14)
                        run.font.bold = False
                        run.font.color.rgb = COLORS['dark_text']

    remove_table_borders(table)


def add_section_name(slide, section_name):
    """Add section name to idx=17 placeholder (top-right, 9pt, #A6A6A6)."""
    if not section_name:
        return
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 17:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = section_name
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = COLORS['medium_gray']
            return


def add_footnote(slide, sources):
    """Add footnote to idx=20 placeholder (bottom-left, 6pt, #A6A6A6)."""
    if not sources:
        return

    footnote_text = "Sources: " + "; ".join(sources) + "."

    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 20:
            tf = shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT  # Changed to left align
            run = para.add_run()
            run.text = footnote_text
            run.font.name = 'Arial'
            run.font.size = Pt(6)
            run.font.color.rgb = COLORS['medium_gray']
            return


def remove_all_pictures_from_slide(slide):
    """Remove all picture shapes from a slide."""
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def add_section_image(slide, section_name, image_dir):
    """Add background image to section divider (after removing old images)."""
    # First remove any existing pictures
    remove_all_pictures_from_slide(slide)

    image_file = SECTION_IMAGES.get(section_name)
    if not image_file:
        return False

    image_path = Path(image_dir) / image_file
    if not image_path.exists():
        print(f"    Warning: Image not found: {image_path}")
        return False

    # Add as background covering full slide
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)
    return True


def add_pccp_logo_text(slide, position='top_left', color='white'):
    """Add PCCP, LLC text logo to slide."""
    if position == 'top_left':
        left, top = Inches(0.4), Inches(0.4)
    elif position == 'center':
        left, top = Inches(4.0), Inches(4.0)

    # Find if there's already a text box we can use, otherwise add one
    textbox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
    tf = textbox.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT if position == 'top_left' else PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "PCCP, LLC"
    run.font.name = 'Arial'
    run.font.size = Pt(24) if position == 'top_left' else Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['white'] if color == 'white' else COLORS['dark_text']


def update_contact_slide(slide):
    """Update slide 44 with PCCP contact information."""
    # Find the content placeholder
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx in [18, 7]:
            tf = shape.text_frame
            tf.clear()

            # Add contact info as bullets
            contact_lines = [
                f"{PCCP_CONTACT['name']}",
                f"{PCCP_CONTACT['headquarters']['address']}",
                f"{PCCP_CONTACT['headquarters']['city_state_zip']}",
                f"Phone: {PCCP_CONTACT['headquarters']['phone']}",
                f"Email: {PCCP_CONTACT['email']}",
                f"Web: {PCCP_CONTACT['website']}",
                "",
                "Additional Offices:",
            ]

            for office in PCCP_CONTACT['other_offices']:
                contact_lines.append(f"{office['city']}: {office['phone']}")

            for i, line in enumerate(contact_lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.alignment = PP_ALIGN.LEFT
                run = para.add_run()
                run.text = line
                run.font.name = 'Arial'
                run.font.size = Pt(14)
                run.font.color.rgb = COLORS['dark_text']
            return


def add_end_slide(prs, image_dir):
    """Add End slide with background image and centered PCCP logo."""
    # Get the End layout (index 12 based on earlier analysis)
    try:
        end_layout = prs.slide_masters[0].slide_layouts[12]
    except:
        # Fall back to section title layout
        end_layout = prs.slide_masters[0].slide_layouts[6]

    slide = prs.slides.add_slide(end_layout)

    # Add background image
    image_path = Path(image_dir) / 'end_slide.png'
    if image_path.exists():
        pic = slide.shapes.add_picture(
            str(image_path),
            Inches(0), Inches(0),
            Inches(11), Inches(8.5)
        )
        # Move to back
        spTree = slide.shapes._spTree
        sp = spTree[-1]
        spTree.remove(sp)
        spTree.insert(2, sp)

    # Add centered PCCP logo
    textbox = slide.shapes.add_textbox(Inches(0), Inches(3.5), Inches(11), Inches(1.5))
    tf = textbox.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "PCCP, LLC"
    run.font.name = 'Arial'
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    return slide


def main():
    """Generate presentation v27."""
    print("Generating presentation v27 with all formatting fixes...")
    print("=" * 70)

    # Paths
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v24_CS_edits.pptx")
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    image_dir = Path("cache/images/cropped")
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v27.pptx")

    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    prs = Presentation(str(template_path))
    with open(outline_path) as f:
        outline = json.load(f)

    # Flatten outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({'section': section.get('name', ''), **slide})

    print(f"Template: {len(prs.slides)} slides")
    print(f"Image directory: {image_dir}")

    # Process each slide
    print("\nProcessing slides...")

    for idx in range(len(prs.slides)):
        if idx >= len(all_slides):
            break

        slide = prs.slides[idx]
        slide_data = all_slides[idx]
        section_name = slide_data.get('section', '')
        slide_type = slide_data.get('slide_type', '')
        slide_num = idx + 1

        # Add section name (top-right, 9pt, #A6A6A6)
        add_section_name(slide, section_name)

        # Add footnote (bottom-left, 6pt, #A6A6A6)
        if slide_num in SLIDE_SOURCES:
            add_footnote(slide, SLIDE_SOURCES[slide_num])

        # Handle section dividers - remove old images, add new
        if slide_type == 'section_divider':
            if add_section_image(slide, section_name, image_dir):
                print(f"  Slide {slide_num}: Updated section image for '{section_name}'")

        # Add PCCP logo to title slide
        if slide_num == 1:
            add_pccp_logo_text(slide, 'top_left', 'white')
            print(f"  Slide {slide_num}: Added PCCP logo to title slide")

        # Update contact slide (44)
        if slide_num == 44:
            update_contact_slide(slide)
            print(f"  Slide {slide_num}: Updated contact information")

    # Format all tables
    print("\nFormatting tables...")
    tables_updated = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'has_table') and shape.has_table:
                format_table(shape.table)
                tables_updated += 1
    print(f"  Updated {tables_updated} tables")

    # Add End slide
    print("\nAdding End slide...")
    add_end_slide(prs, image_dir)
    print("  Added End slide with PCCP logo")

    # Save
    prs.save(str(output_path))

    print(f"\n{'=' * 70}")
    print(f"Generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless", "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
