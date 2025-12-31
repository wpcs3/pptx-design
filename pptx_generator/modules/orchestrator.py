"""
Presentation Orchestrator Module

Main workflow controller that coordinates all components to generate presentations.

Phase 2 Enhancements (2025-12-29):
- Integrated LayoutCascade for automatic layout selection
- Added PresentationEvaluator for quality feedback
- Added functional slide auto-insertion (TOC, section headers)

Phase 4 Enhancements (2025-12-29):
- Added SlidePool for reference slide matching and cloning (PPTAgent-inspired)
- Added IterativeRefiner for quality-based outline refinement
- Clone-and-edit workflow for better visual fidelity
"""

import asyncio
import json
import logging
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation

from .outline_generator import OutlineGenerator
from .research_agent import ResearchAgent
from .slide_library import SlideLibrary
from .template_renderer import TemplateRenderer
from .layout_cascade import LayoutCascade, LayoutType, analyze_content

# Phase 4: SlidePool and IterativeRefiner
try:
    from .slide_pool import (
        SlidePool, IndexedSlide, EditAction, EditActionType,
        FunctionalType, create_edit_actions, build_slide_pool
    )
    from .iterative_refiner import (
        IterativeRefiner, OutlineRefiner, RefinementHistory,
        RefinementAction, RefinementType
    )
    HAS_SLIDE_POOL = True
except ImportError as e:
    logger.warning(f"SlidePool/IterativeRefiner not available: {e}")
    HAS_SLIDE_POOL = False
    SlidePool = None
    IterativeRefiner = None

# Import evaluation if available
try:
    from pptx_design.evaluation import PresentationEvaluator, EvaluationResult
    HAS_EVALUATION = True
except ImportError:
    HAS_EVALUATION = False
    PresentationEvaluator = None
    EvaluationResult = None

logger = logging.getLogger(__name__)


@dataclass
class GenerationOptions:
    """Options for presentation generation."""
    auto_layout: bool = True  # Use LayoutCascade for layout selection
    evaluate_after: bool = True  # Run evaluation after generation
    auto_toc: bool = False  # Auto-insert table of contents
    auto_section_headers: bool = True  # Auto-insert section dividers
    auto_ending: bool = False  # Auto-insert ending slide
    min_quality_score: float = 0.0  # Minimum quality score (0 = no threshold)
    # Phase 4: SlidePool options
    use_slide_pool: bool = False  # Disabled: clone-and-edit copies template content. Use clean layouts instead.
    # Phase 4: Iterative refinement options
    iterative_refinement: bool = False  # Enable iterative refinement
    target_grade: str = "B"  # Target quality grade (A, B, C, D)
    max_iterations: int = 3  # Maximum refinement iterations


@dataclass
class GenerationResult:
    """Result of presentation generation with quality metrics."""
    presentation: Presentation
    output_path: Optional[Path] = None
    evaluation: Optional[Any] = None  # EvaluationResult if available
    layout_decisions: List[Dict[str, Any]] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    # Phase 4: Refinement history
    refinement_history: Optional[Any] = None  # RefinementHistory if iterative refinement used
    slides_cloned: int = 0  # Number of slides created via clone-and-edit


class PresentationOrchestrator:
    """Main workflow controller for presentation generation."""

    def __init__(
        self,
        config_dir: str,
        templates_dir: str,
        output_dir: str = None,
        options: GenerationOptions = None
    ):
        """
        Initialize the orchestrator with all components.

        Args:
            config_dir: Path to configuration files directory
            templates_dir: Path to PPTX templates directory
            output_dir: Path to output directory
            options: Generation options (Phase 2)
        """
        self.config_dir = Path(config_dir)
        self.templates_dir = Path(templates_dir)
        self.output_dir = Path(output_dir) if output_dir else self.config_dir.parent / "output"
        self.options = options or GenerationOptions()

        # Load configuration files
        self.style_guide = self._load_config("style_guide.json")
        self.slide_catalog = self._load_config("slide_catalog.json")
        self.content_patterns = self._load_config("content_patterns.json")

        # Initialize components
        self.outline_gen = OutlineGenerator(self.content_patterns, self.slide_catalog)
        self.slide_lib = SlideLibrary(
            str(self.templates_dir),
            str(self.config_dir / "slide_catalog.json")
        )
        self.research = ResearchAgent(
            self.content_patterns,
            cache_dir=str(self.config_dir.parent / "cache" / "research")
        )
        # Initialize template renderer with formatting and library integration
        base_template = self._find_base_template()
        if base_template:
            self.renderer = TemplateRenderer(base_template, use_library=True)
        else:
            raise FileNotFoundError("No base template found in templates directory")

        # Phase 2: Initialize layout cascade
        self.layout_cascade = LayoutCascade()

        # Phase 2: Initialize evaluator if available
        self.evaluator = PresentationEvaluator() if HAS_EVALUATION else None

        # Phase 4: Initialize slide pool if available
        self.slide_pool = None
        if HAS_SLIDE_POOL and self.options.use_slide_pool:
            try:
                index_path = self.config_dir.parent / "cache" / "slide_pool_index.json"
                self.slide_pool = SlidePool(index_path=str(index_path))
                # Index templates if pool is empty
                if not self.slide_pool.slides:
                    logger.info("Building slide pool from templates...")
                    self.slide_pool.index_templates_directory(str(self.templates_dir))
                    self.slide_pool.save_index()
                logger.info(f"SlidePool initialized with {len(self.slide_pool.slides)} slides")
            except Exception as e:
                logger.warning(f"Failed to initialize SlidePool: {e}")
                self.slide_pool = None

        # Phase 4: Initialize iterative refiner if available
        self.iterative_refiner = None
        if HAS_SLIDE_POOL and self.options.iterative_refinement:
            self.iterative_refiner = IterativeRefiner(
                evaluator=self.evaluator,
                llm_manager=None  # Can be set later via set_llm_manager()
            )

        # Ensure output directory exists
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _load_config(self, filename: str) -> dict:
        """Load a configuration file."""
        path = self.config_dir / filename
        if path.exists():
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        logger.warning(f"Config file not found: {path}")
        return {}

    def _find_base_template(self) -> Optional[str]:
        """Find a base template to use for styling."""
        for template_path in self.templates_dir.rglob("*.pptx"):
            return str(template_path)
        return None

    async def create_presentation(self, user_request: str) -> "Workflow":
        """
        Start the full presentation generation workflow.

        Args:
            user_request: Natural language description of the presentation

        Returns:
            Workflow object for tracking progress and iterations
        """
        workflow = Workflow(self, user_request)
        await workflow.start()
        return workflow

    def generate_outline(self, user_request: str) -> dict:
        """Generate a presentation outline from a user request."""
        return self.outline_gen.generate_outline(user_request)

    def refine_outline(self, outline: dict, feedback: str) -> dict:
        """Refine an outline based on user feedback."""
        return self.outline_gen.refine_outline(outline, feedback)

    async def assemble_content(self, outline: dict) -> dict:
        """
        Assemble content for all sections in the outline.

        Args:
            outline: Approved presentation outline

        Returns:
            Outline with content populated for each slide
        """
        enriched_outline = outline.copy()

        for section in enriched_outline.get("sections", []):
            content_source = section.get("content_source", "user_input")

            if content_source == "reusable":
                # Mark for slide library copy
                section["use_library"] = True
                logger.info(f"Section '{section['name']}' will use reusable slides")

            elif content_source == "research":
                # Perform research
                topics = section.get("research_topics", [section["name"]])
                context = outline.get("context", {})

                research_result = await self.research.research_section(
                    section["name"],
                    topics,
                    context
                )

                # Merge research content into slides
                research_slides = research_result.get("slides", [])
                for i, slide in enumerate(section.get("slides", [])):
                    if i < len(research_slides):
                        slide["content"] = research_slides[i].get("content", {})
                        slide["sources"] = research_slides[i].get("sources", [])

            else:
                # User input - keep placeholder content
                logger.info(f"Section '{section['name']}' requires user input")

        return enriched_outline

    def generate_pptx(self, outline: dict) -> Presentation:
        """
        Generate a PPTX file from an enriched outline.

        Args:
            outline: Outline with content populated

        Returns:
            Generated Presentation object
        """
        # Get presentation title for metadata
        presentation_title = outline.get("title", "Presentation")
        prs = self.renderer.create_presentation(title=presentation_title)
        layout_decisions = []

        for section_idx, section in enumerate(outline.get("sections", [])):
            # Phase 2: Auto-insert section headers
            if self.options.auto_section_headers and section_idx > 0:
                section_name = section.get("name", f"Section {section_idx + 1}")
                self.renderer.create_slide(
                    prs,
                    "section_divider",
                    {"title": section_name}
                )

            # Render all slides through TemplateRenderer with proper formatting
            # The renderer handles ComponentLibrary integration for charts/tables
            for slide_spec in section.get("slides", []):
                content = slide_spec.get("content", {})

                # Phase 2: Use LayoutCascade for automatic layout selection
                if self.options.auto_layout:
                    explicit_type = slide_spec.get("slide_type")
                    if not explicit_type or explicit_type == "auto":
                        # Let cascade decide the layout
                        selected_layout = self.layout_cascade.select_layout(content)
                        slide_type = selected_layout.value
                        layout_decisions.append({
                            "content_title": content.get("title", "Untitled"),
                            "selected_layout": slide_type,
                            "explicit_type": explicit_type
                        })
                    else:
                        slide_type = explicit_type
                else:
                    slide_type = slide_spec.get("slide_type", "title_content")

                # Phase 4: Try clone-and-edit workflow if SlidePool available
                slide_created = False
                if self.slide_pool and self.options.use_slide_pool:
                    # Add slide_type to content for matching
                    match_content = {**content, "slide_type": slide_type}
                    match = self.slide_pool.find_best_match(match_content, min_score=0.5)

                    if match:
                        try:
                            # Generate edit actions from content
                            edits = create_edit_actions(content)
                            # Clone and edit
                            self.slide_pool.clone_and_edit(match, prs, edits)
                            slide_created = True
                            layout_decisions.append({
                                "content_title": content.get("title", "Untitled"),
                                "selected_layout": slide_type,
                                "clone_source": match.slide_id,
                                "clone_template": Path(match.template_path).name,
                                "match_score": match.matches_content(match_content)
                            })
                            logger.debug(f"Cloned slide from {match.slide_id}")
                        except Exception as e:
                            logger.debug(f"Clone-and-edit failed, falling back: {e}")

                # Fall back to regular rendering
                if not slide_created:
                    self.renderer.create_slide(prs, slide_type, content)

        # Phase 2: Auto-insert ending slide
        if self.options.auto_ending:
            self.renderer.create_slide(
                prs,
                "title_slide",
                {"title": "Thank You", "subtitle": "Questions?"}
            )

        # Store layout decisions for later inspection
        self._last_layout_decisions = layout_decisions

        return prs

    def generate_pptx_with_evaluation(
        self,
        outline: dict,
        context: Optional[Dict[str, Any]] = None
    ) -> GenerationResult:
        """
        Generate a PPTX and run evaluation (Phase 2).

        Args:
            outline: Outline with content populated
            context: Optional context for evaluation

        Returns:
            GenerationResult with presentation and quality metrics
        """
        # Generate the presentation
        prs = self.generate_pptx(outline)
        result = GenerationResult(
            presentation=prs,
            layout_decisions=getattr(self, '_last_layout_decisions', [])
        )

        # Run evaluation if enabled and available
        if self.options.evaluate_after and self.evaluator:
            # Save temporarily for evaluation
            temp_path = self.output_dir / "_temp_eval.pptx"
            prs.save(str(temp_path))

            try:
                eval_context = context or {}
                eval_context["outline"] = outline
                result.evaluation = self.evaluator.evaluate(str(temp_path), eval_context)

                # Check quality threshold
                if self.options.min_quality_score > 0:
                    if result.evaluation.overall_score < self.options.min_quality_score:
                        result.warnings.append(
                            f"Quality score {result.evaluation.overall_score:.1f} "
                            f"below threshold {self.options.min_quality_score:.1f}"
                        )

                logger.info(
                    f"Evaluation: {result.evaluation.grade} "
                    f"({result.evaluation.overall_score:.1f}/100)"
                )
            except Exception as e:
                logger.warning(f"Evaluation failed: {e}")
                result.warnings.append(f"Evaluation failed: {e}")
            finally:
                # Clean up temp file
                if temp_path.exists():
                    temp_path.unlink()

        return result

    def generate_pptx_with_refinement(
        self,
        outline: dict,
        context: Optional[Dict[str, Any]] = None
    ) -> GenerationResult:
        """
        Generate a PPTX with iterative refinement (Phase 4).

        Runs a refinement loop that evaluates the presentation and
        improves the outline until quality thresholds are met.

        Args:
            outline: Outline with content populated
            context: Optional context for evaluation

        Returns:
            GenerationResult with presentation, evaluation, and refinement history
        """
        if not self.iterative_refiner or not self.options.iterative_refinement:
            # Fall back to standard generation with evaluation
            return self.generate_pptx_with_evaluation(outline, context)

        def generator_fn(o: dict) -> str:
            """Generate PPTX and return path."""
            prs = self.generate_pptx(o)
            temp_path = self.output_dir / "_temp_refine.pptx"
            prs.save(str(temp_path))
            return str(temp_path)

        # Run iterative refinement
        refined_outline, history = self.iterative_refiner.refine(
            outline,
            generator_fn,
            target_grade=self.options.target_grade,
            max_iterations=self.options.max_iterations
        )

        # Generate final presentation with refined outline
        prs = self.generate_pptx(refined_outline)

        result = GenerationResult(
            presentation=prs,
            layout_decisions=getattr(self, '_last_layout_decisions', []),
            refinement_history=history
        )

        # Run final evaluation
        if self.options.evaluate_after and self.evaluator:
            temp_path = self.output_dir / "_temp_eval.pptx"
            prs.save(str(temp_path))

            try:
                eval_context = context or {}
                eval_context["outline"] = refined_outline
                result.evaluation = self.evaluator.evaluate(str(temp_path), eval_context)

                logger.info(
                    f"Final Evaluation: {result.evaluation.grade} "
                    f"({result.evaluation.overall_score:.2f})"
                )
                logger.info(history.summary())
            except Exception as e:
                logger.warning(f"Final evaluation failed: {e}")
            finally:
                if temp_path.exists():
                    temp_path.unlink()

        # Clean up temp refinement files
        temp_refine = self.output_dir / "_temp_refine.pptx"
        if temp_refine.exists():
            temp_refine.unlink()

        return result

    def set_llm_manager(self, llm_manager) -> None:
        """
        Set the LLM manager for AI-powered refinement.

        Args:
            llm_manager: LLMManager instance from llm_provider module
        """
        if self.iterative_refiner:
            self.iterative_refiner.llm_manager = llm_manager
            self.iterative_refiner.outline_refiner.llm_manager = llm_manager

    def get_slide_pool_stats(self) -> Optional[Dict[str, Any]]:
        """Get statistics about the slide pool."""
        if self.slide_pool:
            return self.slide_pool.get_stats()
        return None

    def get_layout_explanation(self, content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Explain why a particular layout would be selected for content.

        Args:
            content: Slide content dictionary

        Returns:
            Explanation dictionary with analysis
        """
        return self.layout_cascade.explain_selection(content)

    def _copy_reusable_section(self, presentation: Presentation, section: dict) -> None:
        """Copy slides from the reusable slide library."""
        section_id = section.get("reusable_section_id")
        if not section_id:
            return

        reusable_sections = self.content_patterns.get("reusable_sections", {})
        if section_id not in reusable_sections:
            logger.warning(f"Reusable section not found: {section_id}")
            return

        section_def = reusable_sections[section_id]
        source_template = section_def.get("source_template")
        source_slides = section_def.get("source_slides", [])

        for slide_idx in source_slides:
            self.slide_lib.copy_slide(source_template, slide_idx, presentation)

    def export_pptx(self, presentation: Presentation, filename: str = None) -> Path:
        """
        Save a presentation to a PPTX file.

        Args:
            presentation: Presentation object to save
            filename: Optional filename (will be auto-generated if not provided)

        Returns:
            Path to the saved file
        """
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"presentation_{timestamp}.pptx"

        output_path = self.output_dir / filename
        presentation.save(str(output_path))

        logger.info(f"Saved presentation to: {output_path}")
        return output_path

    def refine_presentation(
        self,
        presentation: Presentation,
        feedback: str,
        slide_number: int = None
    ) -> Presentation:
        """
        Refine a presentation based on user feedback.

        Args:
            presentation: Existing Presentation object
            feedback: User's modification request
            slide_number: Optional specific slide to modify

        Returns:
            Modified Presentation object
        """
        feedback_lower = feedback.lower()

        # Parse the type of edit
        if slide_number is not None:
            # Edit specific slide
            self._edit_slide(presentation, slide_number, feedback)
        elif "add" in feedback_lower and "slide" in feedback_lower:
            # Add new slide
            self._add_slide_from_feedback(presentation, feedback)
        elif "remove" in feedback_lower or "delete" in feedback_lower:
            # Remove slide(s)
            self._remove_slides_from_feedback(presentation, feedback)
        else:
            # General style or content edit
            self._apply_general_edit(presentation, feedback)

        return presentation

    def _edit_slide(
        self,
        presentation: Presentation,
        slide_number: int,
        feedback: str
    ) -> None:
        """Edit a specific slide based on feedback."""
        if slide_number < 1 or slide_number > len(presentation.slides):
            logger.warning(f"Invalid slide number: {slide_number}")
            return

        slide = presentation.slides[slide_number - 1]
        feedback_lower = feedback.lower()

        # Handle common edit types
        if "simplify" in feedback_lower:
            self._simplify_slide(slide)
        elif "concise" in feedback_lower:
            self._make_concise(slide)
        # Add more edit handlers as needed

    def _simplify_slide(self, slide) -> None:
        """Simplify a slide by reducing content."""
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Reduce text
                for para in shape.text_frame.paragraphs:
                    if len(para.runs) > 0 and len(para.text) > 100:
                        # Truncate long text
                        para.runs[0].text = para.text[:100] + "..."

    def _make_concise(self, slide) -> None:
        """Make slide content more concise."""
        self._simplify_slide(slide)

    def _add_slide_from_feedback(self, presentation: Presentation, feedback: str) -> None:
        """Add a new slide based on feedback."""
        # Extract the topic from feedback
        import re
        match = re.search(r'add\s+(?:a\s+)?(?:slide\s+(?:about|on|for)\s+)?(.+)', feedback, re.I)
        if match:
            topic = match.group(1).strip()
            self.renderer.create_slide(
                presentation,
                "title_content",
                {"title": topic, "body": f"Content about {topic}"}
            )

    def _remove_slides_from_feedback(self, presentation: Presentation, feedback: str) -> None:
        """Remove slides based on feedback."""
        import re
        # Extract slide number(s)
        numbers = re.findall(r'slide\s*(\d+)', feedback, re.I)
        for num_str in numbers:
            try:
                num = int(num_str) - 1
                if 0 <= num < len(presentation.slides):
                    # Note: python-pptx doesn't support direct slide removal
                    # This would require XML manipulation
                    logger.info(f"Would remove slide {num + 1}")
            except ValueError:
                pass

    def _apply_general_edit(self, presentation: Presentation, feedback: str) -> None:
        """Apply a general edit to the presentation."""
        logger.info(f"General edit requested: {feedback}")
        # Implement general editing logic


class Workflow:
    """Tracks the state of a presentation generation workflow."""

    def __init__(self, orchestrator: PresentationOrchestrator, request: str):
        """
        Initialize a workflow.

        Args:
            orchestrator: Parent orchestrator
            request: Original user request
        """
        self.orchestrator = orchestrator
        self.request = request
        self.outline: Optional[dict] = None
        self.enriched_outline: Optional[dict] = None
        self.presentation: Optional[Presentation] = None
        self.output_path: Optional[Path] = None
        self.status = "initialized"
        self.history: list[dict] = []
        # Phase 2: Track generation result with evaluation
        self.generation_result: Optional[GenerationResult] = None

    async def start(self) -> None:
        """Start the workflow by generating an outline."""
        self.status = "generating_outline"
        self.outline = self.orchestrator.generate_outline(self.request)
        self.status = "awaiting_outline_approval"
        self._log_action("outline_generated")

    def get_outline_preview(self) -> str:
        """Get a human-readable preview of the current outline."""
        if not self.outline:
            return "No outline generated yet."
        return self.orchestrator.outline_gen.outline_to_text(self.outline)

    def approve_outline(self) -> None:
        """Mark the current outline as approved."""
        if self.status != "awaiting_outline_approval":
            logger.warning("Cannot approve outline in current state")
            return
        self.status = "outline_approved"
        self._log_action("outline_approved")

    def modify_outline(self, feedback: str) -> None:
        """Modify the outline based on feedback."""
        if not self.outline:
            return
        self.outline = self.orchestrator.refine_outline(self.outline, feedback)
        self._log_action("outline_modified", {"feedback": feedback})

    async def assemble_content(self) -> None:
        """Assemble content for all sections."""
        if self.status != "outline_approved":
            logger.warning("Outline must be approved before assembling content")
            return

        self.status = "assembling_content"
        self.enriched_outline = await self.orchestrator.assemble_content(self.outline)
        self.status = "content_ready"
        self._log_action("content_assembled")

    def generate_presentation(self, with_evaluation: bool = True) -> None:
        """
        Generate the PPTX from the enriched outline.

        Args:
            with_evaluation: If True, run evaluation after generation (Phase 2)
        """
        if not self.enriched_outline:
            logger.warning("Content must be assembled before generating")
            return

        self.status = "generating_pptx"

        if with_evaluation and self.orchestrator.evaluator:
            # Phase 2: Generate with evaluation
            self.generation_result = self.orchestrator.generate_pptx_with_evaluation(
                self.enriched_outline,
                context={"request": self.request}
            )
            self.presentation = self.generation_result.presentation

            # Log evaluation results
            if self.generation_result.evaluation:
                self._log_action("presentation_evaluated", {
                    "grade": self.generation_result.evaluation.grade,
                    "score": self.generation_result.evaluation.overall_score,
                    "recommendations": self.generation_result.evaluation.recommendations[:3]
                })
        else:
            self.presentation = self.orchestrator.generate_pptx(self.enriched_outline)

        self.status = "draft_ready"
        self._log_action("presentation_generated")

    def get_evaluation_summary(self) -> Optional[str]:
        """Get a summary of the evaluation results (Phase 2)."""
        if not self.generation_result or not self.generation_result.evaluation:
            return None

        eval_result = self.generation_result.evaluation
        lines = [
            f"Quality Grade: {eval_result.grade} ({eval_result.overall_score * 100:.1f}/100)",
            "",
            "Scores:",
            f"  Content:   {eval_result.content.score * 100:.1f}/100",
            f"  Design:    {eval_result.design.score * 100:.1f}/100",
            f"  Coherence: {eval_result.coherence.score * 100:.1f}/100",
            "",
            "Recommendations:"
        ]
        for rec in eval_result.recommendations[:5]:
            lines.append(f"  • {rec}")

        return "\n".join(lines)

    def get_layout_decisions(self) -> List[Dict[str, Any]]:
        """Get the layout decisions made during generation (Phase 2)."""
        if self.generation_result:
            return self.generation_result.layout_decisions
        return getattr(self.orchestrator, '_last_layout_decisions', [])

    def export(self, filename: str = None) -> Path:
        """Export the presentation to a file."""
        if not self.presentation:
            logger.warning("No presentation to export")
            return None

        self.output_path = self.orchestrator.export_pptx(self.presentation, filename)
        self._log_action("presentation_exported", {"path": str(self.output_path)})
        return self.output_path

    def refine(self, feedback: str, slide_number: int = None) -> None:
        """Refine the presentation based on feedback."""
        if not self.presentation:
            return

        self.presentation = self.orchestrator.refine_presentation(
            self.presentation,
            feedback,
            slide_number
        )
        self._log_action("presentation_refined", {
            "feedback": feedback,
            "slide": slide_number
        })

    def finalize(self, filename: str = None) -> Path:
        """Finalize and export the presentation."""
        self.status = "finalized"
        return self.export(filename)

    def _log_action(self, action: str, details: dict = None) -> None:
        """Log an action to the workflow history."""
        self.history.append({
            "action": action,
            "timestamp": datetime.now().isoformat(),
            "status": self.status,
            "details": details or {}
        })


def main():
    """Test the orchestrator."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Presentation Orchestrator")
    parser.add_argument(
        "--config-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/config",
        help="Path to config directory"
    )
    parser.add_argument(
        "--templates-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_templates",
        help="Path to templates directory"
    )
    parser.add_argument(
        "--output-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/output",
        help="Path to output directory"
    )
    parser.add_argument(
        "--request",
        default="Create a pitch deck for a $200M industrial logistics fund",
        help="User request"
    )
    parser.add_argument(
        "--no-eval",
        action="store_true",
        help="Skip evaluation after generation"
    )
    parser.add_argument(
        "--auto-layout",
        action="store_true",
        default=True,
        help="Use automatic layout selection"
    )

    args = parser.parse_args()

    async def run():
        # Phase 2: Configure generation options
        options = GenerationOptions(
            auto_layout=args.auto_layout,
            evaluate_after=not args.no_eval,
            auto_section_headers=True
        )

        orchestrator = PresentationOrchestrator(
            args.config_dir,
            args.templates_dir,
            args.output_dir,
            options=options
        )

        # Generate outline
        print("\n=== Generating Outline ===")
        workflow = await orchestrator.create_presentation(args.request)
        print(workflow.get_outline_preview())

        # Approve and build
        print("\n=== Approving Outline ===")
        workflow.approve_outline()

        print("\n=== Assembling Content ===")
        await workflow.assemble_content()

        print("\n=== Generating Presentation ===")
        workflow.generate_presentation(with_evaluation=not args.no_eval)

        # Phase 2: Show evaluation results
        eval_summary = workflow.get_evaluation_summary()
        if eval_summary:
            print("\n=== Quality Evaluation ===")
            print(eval_summary)

        # Phase 2: Show layout decisions
        layout_decisions = workflow.get_layout_decisions()
        if layout_decisions:
            print("\n=== Layout Decisions ===")
            for decision in layout_decisions[:5]:
                print(f"  {decision['content_title']}: {decision['selected_layout']}")
            if len(layout_decisions) > 5:
                print(f"  ... and {len(layout_decisions) - 5} more")

        print("\n=== Exporting ===")
        output_path = workflow.finalize()
        print(f"Saved to: {output_path}")

    asyncio.run(run())


if __name__ == "__main__":
    main()
