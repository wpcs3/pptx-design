"""
Template-Based Slide Renderer

Generates slides using actual template master layouts and styling.
Integrates with ComponentLibrary for reusable charts, tables, and diagrams.
Uses StyleGuideConfig for centralized formatting specifications.

Fixed version: Uses actual placeholder positions and proper auto-sizing.
"""

import json
import logging
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

# Import ComponentLibrary and LibraryEnhancer for reusable components
try:
    from .component_library import ComponentLibrary
    from .library_enhancer import LibraryEnhancer
    LIBRARY_AVAILABLE = True
except ImportError:
    LIBRARY_AVAILABLE = False
    ComponentLibrary = None
    LibraryEnhancer = None

# Import StyleGuideConfig for centralized formatting
try:
    from .style_guide_config import StyleGuideConfig, get_style_config
    STYLE_CONFIG_AVAILABLE = True
except ImportError:
    STYLE_CONFIG_AVAILABLE = False
    StyleGuideConfig = None
    get_style_config = None

logger = logging.getLogger(__name__)


class TemplateRenderer:
    """Renders slides using template master layouts and styling."""

    # Layout mappings for different slide types
    # Using template master layouts for professional styling
    # "Default" layout has title at top=0.39" (ABOVE the horizontal line)
    # "Top left title" has title at top=1.87" (BELOW the line) - avoid for standard content
    LAYOUT_MAP = {
        "title_slide": "Frontpage",
        "frontpage": "Frontpage",
        "section_divider": "Section breaker",
        "section_breaker": "Section breaker",
        "title_content": "Default",  # Title above horizontal line (top=0.39")
        "default": "Default",
        "top_left": "Default",  # Use Default for proper title placement
        "content": "Default",
        "two_column": "1/2 grey",
        "comparison": "1/2 grey",
        "sidebar_left": "1/3 grey",
        "sidebar_right": "2/3 grey",
        "agenda": "Agenda",
        "end_slide": "End",
        "blank": "Blank",
        "data_chart": "Default",  # Title above line
        "table_slide": "Default",  # Title above line
        "key_metrics": "Default",  # Title above line
    }

    # Slide dimensions (standard 16:9)
    SLIDE_WIDTH = 13.333  # inches
    SLIDE_HEIGHT = 7.5    # inches

    # Margins
    MARGIN_LEFT = 0.61
    MARGIN_RIGHT = 0.6
    MARGIN_TOP = 0.39
    MARGIN_BOTTOM = 0.5

    # Bullet character for lists
    BULLET_CHAR = "•"

    # Colors extracted from templates
    COLORS = {
        "primary": RGBColor(0x3C, 0x96, 0xB4),      # Teal
        "secondary": RGBColor(0xE5, 0x54, 0x6C),    # Coral/Red
        "accent1": RGBColor(0x05, 0x1C, 0x2C),      # Dark blue
        "accent2": RGBColor(0x00, 0xB0, 0x50),      # Green
        "black": RGBColor(0x00, 0x00, 0x00),        # Pure black for titles/subtitles (per PCCP style guide)
        "text_dark": RGBColor(0x06, 0x1F, 0x32),    # Dark navy (legacy)
        "text_body": RGBColor(0x2D, 0x37, 0x48),    # Slate #2D3748 for body text
        "text_light": RGBColor(0x71, 0x80, 0x96),   # Gray #718096 for section names
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_gray": RGBColor(0xE6, 0xE6, 0xE6),
        "background": RGBColor(0xF2, 0xF2, 0xF2),
        "overlay_navy": RGBColor(0x05, 0x1C, 0x2C), # Dark navy for overlays
        "footnote": RGBColor(0xA6, 0xA6, 0xA6),     # Medium gray #A6A6A6 for footnotes (per PCCP style guide)
        "chart_bar": RGBColor(0x30, 0x9C, 0xE7),    # Blue #309CE7 for bar charts
        "chart_line_primary": RGBColor(0x05, 0x1C, 0x2C),  # Dark navy #051C2C for line charts
        "chart_line_secondary": RGBColor(0x22, 0x22, 0xF6),  # Blue #2222F6 for secondary lines
        "chart_gridline": RGBColor(0xE2, 0xE8, 0xF0),  # Light gray #E2E8F0 for gridlines (per PCCP style guide)
    }

    # Line widths
    LINE_WIDTH_CHART = Pt(3)      # 3pt for chart lines
    LINE_WIDTH_GRIDLINE = Pt(0.5)  # 0.5pt for gridlines

    # Footer configuration
    FOOTER_COMPANY = "PCCP, LLC"

    # Overlay transparency (0-100000, where 100000 = fully transparent)
    OVERLAY_TRANSPARENCY = 30000  # 30% transparent (70% opaque)

    # Font settings - per PCCP CS Style Guide (2026.01.01)
    # Title (content slides): 32pt bold black
    # Title (section/cover): 44pt bold white
    # Subtitle (content): 20pt bold black
    # Subtitle (frontpage): 18pt bold white
    # Body/Bullets: 14pt regular black
    # Content header (side-by-side): 18pt bold black
    FONTS = {
        # Title styles
        "cover_title": {"name": "Arial", "size": Pt(44), "bold": True, "color": "white"},
        "section_title": {"name": "Arial", "size": Pt(44), "bold": True, "color": "white"},
        "title": {"name": "Arial", "size": Pt(32), "bold": True, "color": "black"},  # Black per PCCP style guide
        "section": {"name": "Arial", "size": Pt(44), "bold": True, "color": "white"},
        # Subtitle styles
        "subtitle": {"name": "Arial", "size": Pt(18), "bold": True, "color": "white"},  # Front page
        "subtitle_content": {"name": "Arial", "size": Pt(20), "bold": True, "color": "black"},  # Black per PCCP style guide
        "takeaway": {"name": "Arial", "size": Pt(18), "bold": True, "italic": False, "color": "black"},
        # Body/content styles
        "heading": {"name": "Arial", "size": Pt(16), "bold": True, "color": "black"},
        "content_header": {"name": "Arial", "size": Pt(18), "bold": True, "color": "black"},  # Side-by-side headers
        "body": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},
        "body_large": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},
        "bullet": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},
        "caption": {"name": "Arial", "size": Pt(9), "bold": False, "color": "text_light"},
        "footnote": {"name": "Arial", "size": Pt(8), "bold": False, "color": "footnote"},
        # Component styles
        "metric_value": {"name": "Arial", "size": Pt(28), "bold": True, "color": "white"},
        "metric_label": {"name": "Arial", "size": Pt(12), "bold": False, "color": "white"},
        "table_header": {"name": "Arial", "size": Pt(16), "bold": True, "color": "white"},  # 16pt per BTR standard
        "table_cell": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},  # 14pt per BTR standard
        "chart_label": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},  # 14pt per PCCP style guide
        "chart_axis": {"name": "Arial", "size": Pt(14), "bold": False, "color": "text_body"},  # 14pt per PCCP style guide
    }

    def __init__(self, template_path: str, use_library: bool = True, style_guide_version: str = None):
        """
        Initialize renderer with a template.

        Args:
            template_path: Path to the PPTX template file
            use_library: Whether to use ComponentLibrary for charts/tables
            style_guide_version: Specific style guide version to use (e.g., "2026.01.01")
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load style guide configuration
        self.style_config = None
        if STYLE_CONFIG_AVAILABLE:
            try:
                self.style_config = get_style_config(style_guide_version)
                logger.info(f"Loaded style guide: {self.style_config.name} v{self.style_config.version}")
            except Exception as e:
                logger.warning(f"Could not load style guide config: {e}")

        # Load template to get layouts
        self._template_prs = Presentation(str(self.template_path))
        self._build_layout_index()
        self._extract_layout_dimensions()

        # Initialize component library and enhancer
        self.library = None
        self.enhancer = None
        if use_library and LIBRARY_AVAILABLE:
            try:
                self.library = ComponentLibrary()
                if self.library.is_available:
                    stats = self.library.get_stats()
                    total = sum(stats['components'].values())
                    logger.info(f"ComponentLibrary loaded with {total} components")

                    # Initialize enhancer for smart matching
                    if LibraryEnhancer:
                        self.enhancer = LibraryEnhancer(self.library)
                else:
                    self.library = None
            except Exception as e:
                logger.warning(f"Could not initialize ComponentLibrary: {e}")
                self.library = None

    def _build_layout_index(self) -> None:
        """Build index of available layouts."""
        self.layouts = {}
        for layout in self._template_prs.slide_layouts:
            self.layouts[layout.name] = layout
        logger.info(f"Indexed {len(self.layouts)} layouts from template")

    def _extract_layout_dimensions(self) -> None:
        """Extract actual placeholder dimensions from each layout."""
        self.layout_dimensions = {}

        for layout in self._template_prs.slide_layouts:
            dims = {
                "title": None,
                "subtitle": None,
                "body": None,
                "content_top": 1.8,  # Default content start
            }

            max_placeholder_bottom = self.MARGIN_TOP

            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type)
                left = ph.left / 914400  # EMUs to inches
                top = ph.top / 914400
                width = ph.width / 914400
                height = ph.height / 914400
                bottom = top + height

                if "TITLE" in ph_type:
                    dims["title"] = {"left": left, "top": top, "width": width, "height": height}
                    max_placeholder_bottom = max(max_placeholder_bottom, bottom)
                elif "SUBTITLE" in ph_type:
                    dims["subtitle"] = {"left": left, "top": top, "width": width, "height": height}
                    max_placeholder_bottom = max(max_placeholder_bottom, bottom)
                elif "BODY" in ph_type and ph.placeholder_format.idx != 17:
                    # idx 17 is typically a small header body, skip it
                    dims["body"] = {"left": left, "top": top, "width": width, "height": height}

            # Content area starts after title/subtitle placeholders
            dims["content_top"] = max_placeholder_bottom + 0.2
            self.layout_dimensions[layout.name] = dims

    def _get_content_area(self, layout_name: str) -> Dict[str, float]:
        """Get the content area dimensions for a layout."""
        dims = self.layout_dimensions.get(layout_name, {})
        content_top = dims.get("content_top", 1.8)

        return {
            "left": self.MARGIN_LEFT,
            "top": content_top,
            "width": self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT,
            "height": self.SLIDE_HEIGHT - content_top - self.MARGIN_BOTTOM,
        }

    def create_presentation(self, title: str = None) -> Presentation:
        """Create a new empty presentation with template layouts.

        Args:
            title: Optional title for the presentation (used in metadata)

        Returns:
            Empty Presentation object with layouts from template
        """
        prs = Presentation(str(self.template_path))

        # Remove all existing slides from the template
        while len(prs.slides) > 0:
            slide = prs.slides[0]
            rId = prs.part.relate_to(slide.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            sldIdLst = prs._element.get_or_add_sldIdLst()
            sldId = sldIdLst.sldId_lst[0]
            sldIdLst.remove(sldId)
            prs.part.drop_rel(sldId.rId)

        # Set presentation metadata
        if title:
            self.set_presentation_metadata(prs, title=title)

        logger.info(f"Created empty presentation with {len(prs.slide_layouts)} layouts")
        return prs

    def set_presentation_metadata(self, prs: Presentation, title: str = None, author: str = None, subject: str = None) -> None:
        """Set presentation metadata (title, author, subject).

        This metadata appears in PDF properties when converted.

        Args:
            prs: Presentation object
            title: Document title
            author: Document author
            subject: Document subject
        """
        core_props = prs.core_properties
        if title:
            core_props.title = title
        if author:
            core_props.author = author
        if subject:
            core_props.subject = subject

    def get_layout(self, prs: Presentation, slide_type: str):
        """Get the appropriate layout for a slide type."""
        layout_name = self.LAYOUT_MAP.get(slide_type, "Default")

        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout

        # Fallback to Default
        for layout in prs.slide_layouts:
            if layout.name == "Default":
                return layout

        return prs.slide_layouts[1]

    def create_slide(
        self,
        prs: Presentation,
        slide_type: str,
        content: dict,
        use_placeholders: bool = True
    ) -> Slide:
        """
        Create a slide using template layout.

        Args:
            prs: Presentation object
            slide_type: Type of slide to create
            content: Content dictionary
            use_placeholders: If True, fill master layout placeholders directly

        Returns:
            Created Slide object
        """
        layout = self.get_layout(prs, slide_type)
        slide = prs.slides.add_slide(layout)
        layout_name = layout.name

        # Clear ALL placeholder content from the template
        # This prevents template default text from appearing
        self._clear_placeholder_content(slide)

        # Handle title - add manually for Blank layout
        title = content.get("title", "")
        if title:
            if slide.shapes.title:
                self._set_title(slide.shapes.title, title, slide_type)
            elif layout_name == "Blank":
                # Add title manually for Blank layout
                self._add_title_to_blank_slide(slide, title, slide_type)

        # Route to appropriate renderer
        renderer_method = getattr(self, f"_render_{slide_type}", None)
        if renderer_method:
            renderer_method(slide, content, layout_name)
        else:
            self._render_default(slide, content, layout_name)

        # Set footer text on all slides
        self._set_footer(slide)

        return slide

    def _set_footer(self, slide: Slide) -> None:
        """Set footer text to company name on slide."""
        for shape in slide.placeholders:
            ph_type = str(shape.placeholder_format.type)
            if "FOOTER" in ph_type:
                shape.text = self.FOOTER_COMPANY
                break

    def _add_title_to_blank_slide(self, slide: Slide, title: str, slide_type: str) -> None:
        """Add a title text box to a blank slide layout."""
        # Title position matching template style
        left = Inches(self.MARGIN_LEFT)
        top = Inches(0.39)
        width = Inches(self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT)
        height = Inches(0.6)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = tf.paragraphs[0]
        p.text = title
        self._apply_font(p, "title")

        # Add a horizontal line below the title
        line_top = Inches(1.05)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, line_top,
            width, Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["text_dark"]
        line.line.fill.background()

    def _clear_placeholder_content(self, slide: Slide) -> None:
        """Clear all content shapes from a slide to remove template defaults.

        This AGGRESSIVELY removes ALL shapes except:
        - Title placeholder (we'll fill it)
        - Footer/slide number/date placeholders
        - Horizontal lines in footer area (decorative dividers)

        Also removes any text that looks like template section names
        (e.g., "Template overview", "Guide: How to write...", etc.)
        """
        shapes_to_remove = []

        # Get title shape reference before iterating
        title_shape = slide.shapes.title

        # Template section names to remove (partial matches)
        template_section_names = [
            "Template overview",
            "Guide:",
            "How to write",
            "Section Name",
            "Click to edit",
            "Add your",
        ]

        # Identify ALL shapes to remove
        for shape in slide.shapes:
            # Keep slide title placeholder only
            if shape == title_shape:
                continue

            # Check if it's a footer/slide number/date placeholder - keep these
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if any(keep in ph_type for keep in ["FOOTER", "SLIDE_NUMBER", "DATE"]):
                    continue

            # Check for footer divider lines (thin horizontal lines at bottom)
            try:
                shape_top = shape.top / 914400  # EMUs to inches
                shape_height = shape.height / 914400
                shape_width = shape.width / 914400

                # Keep only thin horizontal lines (height < 0.1") in footer area (bottom 1")
                footer_area = self.SLIDE_HEIGHT - 1.0
                is_horizontal_line = shape_height < 0.1 and shape_width > 5.0

                if shape_top > footer_area and is_horizontal_line:
                    continue
            except Exception:
                pass

            # Check if shape contains template section name text - remove it
            try:
                if hasattr(shape, 'text') and shape.text:
                    text = shape.text.strip()
                    if any(template_name in text for template_name in template_section_names):
                        shapes_to_remove.append(shape)
                        continue
            except Exception:
                pass

            # REMOVE EVERYTHING ELSE - no exceptions
            shapes_to_remove.append(shape)

        # Remove all identified shapes
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception as e:
                logger.warning(f"Failed to remove shape: {e}")

    def _set_title(self, title_shape, text: str, slide_type: str) -> None:
        """Set title text with proper formatting."""
        title_shape.text = text
        tf = title_shape.text_frame

        # Enable word wrap and auto-size
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Apply font
        style = "section" if slide_type == "section_divider" else "title"
        for para in tf.paragraphs:
            self._apply_font(para, style)

    def _create_text_box(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        auto_size: bool = True
    ) -> Any:
        """Create a text box with proper settings."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        if auto_size:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        return textbox

    def _add_bullets(
        self,
        slide: Slide,
        bullets: list,
        area: Dict[str, float],
        style: str = "bullet"
    ) -> None:
        """Add bullet points to a slide with proper PowerPoint bullet formatting.

        Paragraph formatting (from PowerPoint dialog):
        - Alignment: Left
        - Before text: 0.2"
        - Special: Hanging, By: 0.2"
        - Spacing Before: 0 pt
        - Spacing After: 6 pt
        - Line Spacing: Single
        """
        import re
        from lxml import etree

        if not bullets:
            return

        textbox = self._create_text_box(
            slide,
            area["left"],
            area["top"],
            area["width"],
            area["height"]
        )
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        font_style = self.FONTS.get(style, self.FONTS["bullet"])

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Set bullet level and alignment
            p.level = 0
            p.alignment = PP_ALIGN.LEFT

            # Apply proper PowerPoint bullet formatting via XML
            try:
                pPr = p._p.get_or_add_pPr()

                # Set indentation attributes on pPr element:
                # marL = left margin (where text wraps to on subsequent lines)
                # indent = first line indent relative to marL (negative = hanging)
                #
                # PowerPoint "Before text: 0.2" + Hanging: 0.2"" means:
                # - marL = 0.4" (text starts here)
                # - indent = -0.2" (bullet hangs 0.2" left of marL, so at 0.2")
                # This creates 0.2" spacing between bullet (at 0.2") and text (at 0.4")
                marL_emu = int(Inches(0.4))  # Text starts at 0.4"
                indent_emu = int(Inches(-0.2))  # Bullet at marL + indent = 0.2"
                pPr.set(qn('a:marL'), str(marL_emu))
                pPr.set(qn('a:indent'), str(indent_emu))

                # Remove ALL existing child elements to rebuild from scratch
                for child in list(pPr):
                    pPr.remove(child)

                # OOXML schema requires elements in specific order:
                # 1. lnSpc, spcBef, spcAft (spacing)
                # 2. buClr (bullet color)
                # 3. buSzPct (bullet size)
                # 4. buFont (bullet font)
                # 5. buChar (bullet character)

                # Line Spacing: Single (100%)
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
                spcPct.set('val', '100000')

                # Spacing Before: 0 pt
                spcBef = etree.SubElement(pPr, qn('a:spcBef'))
                spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
                spcPts.set('val', '0')

                # Spacing After: 6 pt (600 = 6pt * 100)
                spcAft = etree.SubElement(pPr, qn('a:spcAft'))
                spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
                spcPts.set('val', '600')

                # Bullet color (black)
                buClr = etree.SubElement(pPr, qn('a:buClr'))
                srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
                srgbClr.set('val', '000000')

                # Bullet size (100% of text size)
                buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
                buSzPct.set('val', '100000')

                # Bullet font (Arial)
                buFont = etree.SubElement(pPr, qn('a:buFont'))
                buFont.set('typeface', 'Arial')
                buFont.set('pitchFamily', '34')
                buFont.set('charset', '0')

                # Bullet character (standard round bullet)
                buChar = etree.SubElement(pPr, qn('a:buChar'))
                buChar.set('char', '•')

            except Exception as e:
                logger.debug(f"Could not set bullet formatting via XML: {e}")

            # Check if bullet has a header pattern (CAPS TEXT: rest of text)
            header_match = re.match(r'^([A-Z][A-Z\s\-&/]+:)\s*(.*)$', bullet)

            if header_match:
                # Split into header (bold) and body (normal)
                header_text = header_match.group(1)
                body_text = header_match.group(2)

                # Add bold header run
                run1 = p.add_run()
                run1.text = header_text + " "
                run1.font.name = font_style["name"]
                run1.font.size = font_style["size"]
                run1.font.bold = True
                run1.font.color.rgb = self.COLORS.get(font_style.get("color", "text_body"), self.COLORS["text_body"])

                # Add normal body run
                if body_text:
                    run2 = p.add_run()
                    run2.text = body_text
                    run2.font.name = font_style["name"]
                    run2.font.size = font_style["size"]
                    run2.font.bold = False
                    run2.font.color.rgb = self.COLORS.get(font_style.get("color", "text_body"), self.COLORS["text_body"])
            else:
                # No header pattern - apply normal styling
                run = p.add_run()
                run.text = bullet
                run.font.name = font_style["name"]
                run.font.size = font_style["size"]
                run.font.bold = False
                run.font.color.rgb = self.COLORS.get(font_style.get("color", "text_body"), self.COLORS["text_body"])

    def _add_takeaway(self, slide: Slide, takeaway: str) -> float:
        """
        Add a takeaway/subheader text below the horizontal line.

        The takeaway provides a quick summary of the slide's main point,
        making it easy for readers to get the gist when scanning the deck.

        Args:
            slide: Slide to add takeaway to
            takeaway: The takeaway text (1-2 sentences)

        Returns:
            The height consumed by the takeaway (for adjusting content area)
        """
        if not takeaway:
            return 0.0

        # Position and size from PowerPoint specifications:
        # Horizontal: 0.6", Vertical: 1.5", Width: 12.2", Height: 0.7"
        takeaway_left = 0.6
        takeaway_top = 1.5
        takeaway_width = 12.2
        takeaway_height = 0.7

        textbox = self._create_text_box(
            slide,
            left=takeaway_left,
            top=takeaway_top,
            width=takeaway_width,
            height=takeaway_height,
            auto_size=False
        )

        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        p = tf.paragraphs[0]
        p.text = takeaway
        p.alignment = PP_ALIGN.LEFT
        self._apply_font(p, "takeaway")

        # Return the space consumed (for adjusting content area)
        # Takeaway ends at 1.5" + 0.7" = 2.2", add small gap
        return (takeaway_top + takeaway_height) - 1.8 + 0.1  # Adjust relative to default content_top

    # ==================== Slide Renderers ====================

    def _render_title_slide(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a title/frontpage slide with left-aligned overlay boxes for readability."""
        import datetime

        title = content.get("title", "")
        subtitle = content.get("subtitle", "")
        background_image = content.get("background_image")

        # Add background image if provided
        if background_image:
            self._add_background_image(slide, background_image)

            # Clear existing title placeholder (we'll use overlay box instead)
            if slide.shapes.title:
                try:
                    sp = slide.shapes.title._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass

            # Clear subtitle placeholder too
            for shape in list(slide.placeholders):
                ph_type = str(shape.placeholder_format.type)
                if "SUBTITLE" in ph_type:
                    try:
                        sp = shape._element
                        sp.getparent().remove(sp)
                    except Exception:
                        pass

            # Add title overlay box (LEFT-ALIGNED)
            # Using 44pt font for cover title
            title_box_width = 8.0
            title_box_height = 1.5
            title_left = self.MARGIN_LEFT
            title_top = 2.5

            self._add_text_overlay_box(
                slide,
                left=title_left,
                top=title_top,
                width=title_box_width,
                height=title_box_height,
                text=title,
                font_style="cover_title",
                padding=0.3,
                align_left=True
            )

            # Add subtitle overlay box below title (LEFT-ALIGNED)
            if subtitle:
                subtitle_box_width = 6.0
                subtitle_box_height = 0.8
                subtitle_left = self.MARGIN_LEFT
                subtitle_top = title_top + title_box_height + 0.2

                self._add_text_overlay_box(
                    slide,
                    left=subtitle_left,
                    top=subtitle_top,
                    width=subtitle_box_width,
                    height=subtitle_box_height,
                    text=subtitle,
                    font_style="subtitle",
                    padding=0.2,
                    align_left=True
                )

                # Add date below subtitle
                date_top = subtitle_top + subtitle_box_height + 0.2
            else:
                # Add date below title if no subtitle
                date_top = title_top + title_box_height + 0.2

            # Add today's date
            today_str = datetime.date.today().strftime("%B %d, %Y")
            date_box_width = 3.0
            date_box_height = 0.5

            self._add_text_overlay_box(
                slide,
                left=self.MARGIN_LEFT,
                top=date_top,
                width=date_box_width,
                height=date_box_height,
                text=today_str,
                font_style="subtitle",
                padding=0.15,
                align_left=True
            )
        else:
            # No background image - use standard placeholder styling
            for shape in slide.placeholders:
                ph_type = str(shape.placeholder_format.type)
                if "SUBTITLE" in ph_type:
                    shape.text = subtitle
                    tf = shape.text_frame
                    tf.word_wrap = True
                    for para in tf.paragraphs:
                        self._apply_font(para, "subtitle")
                    break

    def _render_section_divider(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a section divider slide with left-aligned overlay box for readability."""
        title = content.get("title", "")
        background_image = content.get("background_image")

        # Add background image if provided
        if background_image:
            self._add_background_image(slide, background_image)

            # Clear existing title placeholder (we'll use overlay box instead)
            if slide.shapes.title:
                try:
                    sp = slide.shapes.title._element
                    sp.getparent().remove(sp)
                except Exception:
                    pass

            # Add section title overlay box (LEFT-ALIGNED)
            # Using 36pt font for section titles
            title_box_width = 7.0
            title_box_height = 1.3
            title_left = self.MARGIN_LEFT
            title_top = (self.SLIDE_HEIGHT - title_box_height) / 2

            self._add_text_overlay_box(
                slide,
                left=title_left,
                top=title_top,
                width=title_box_width,
                height=title_box_height,
                text=title,
                font_style="section_title",
                padding=0.25,
                align_left=True
            )

    def _add_background_image(self, slide: Slide, image_path: str) -> None:
        """Add a background image to a slide with proper aspect ratio handling.

        For 16:9 images on letter-size (11x8.5) slides, applies cropping to
        maintain aspect ratio and fill the slide without stretching.
        """
        from pptx.util import Inches
        from lxml import etree
        from PIL import Image as PILImage
        from io import BytesIO

        image_path = Path(image_path)
        if not image_path.exists():
            logger.warning(f"Background image not found: {image_path}")
            return

        # Add image to fill the entire slide
        left = Inches(0)
        top = Inches(0)
        width = Inches(self.SLIDE_WIDTH)
        height = Inches(self.SLIDE_HEIGHT)

        # Add the image
        picture = slide.shapes.add_picture(
            str(image_path), left, top, width, height
        )

        # Apply aspect ratio correction for 16:9 images on letter-size slides
        # This prevents vertical stretching
        try:
            slide_aspect = self.SLIDE_WIDTH / self.SLIDE_HEIGHT
            with open(image_path, 'rb') as f:
                img = PILImage.open(f)
                orig_width, orig_height = img.size
                orig_aspect = orig_width / orig_height

            # If image is wider than slide (16:9 = 1.778 > letter = 1.294), crop sides
            if orig_aspect > slide_aspect + 0.05:
                # Calculate crop percentage for each side
                crop_ratio = 1 - (slide_aspect / orig_aspect)
                crop_pct = int((crop_ratio / 2) * 100000)  # EMU percentage

                pic_elem = picture._element

                # Find blipFill - try multiple approaches for different PPTX structures
                blipFill = pic_elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
                if blipFill is None:
                    blipFill = pic_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')
                if blipFill is None:
                    # Try direct child with local name
                    for child in pic_elem:
                        if child.tag.endswith('}blipFill') or child.tag == 'blipFill':
                            blipFill = child
                            break

                if blipFill is not None:
                    # Remove existing srcRect if present
                    existing_srcRect = blipFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
                    if existing_srcRect is not None:
                        blipFill.remove(existing_srcRect)

                    # Create new srcRect with cropping
                    srcRect = etree.SubElement(blipFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
                    srcRect.set('l', str(crop_pct))
                    srcRect.set('r', str(crop_pct))
                    srcRect.set('t', '0')
                    srcRect.set('b', '0')
                    logger.debug(f"Applied {crop_pct/1000:.1f}% side crop to background image")
        except Exception as e:
            logger.debug(f"Could not apply aspect ratio correction: {e}")

        # Send to back (behind all other shapes)
        # Move to position 0 in the shape tree
        spTree = slide.shapes._spTree
        pic_element = picture._element
        spTree.remove(pic_element)
        spTree.insert(2, pic_element)  # Insert after nvGrpSpPr and grpSpPr

    def _add_text_overlay_box(
        self,
        slide: Slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_style: str = "cover_title",
        padding: float = 0.3,
        align_left: bool = False,
        background_color: str = "overlay_navy"
    ) -> Any:
        """
        Add a semi-transparent box with text.

        Args:
            slide: Slide to add overlay to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text to display
            font_style: Font style key from FONTS dict
            padding: Padding around text in inches
            align_left: If True, align text left; otherwise center
            background_color: Color key from COLORS dict

        Returns:
            The created shape
        """
        # Create the rectangle shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # Set fill with transparency
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS.get(background_color, self.COLORS["overlay_navy"])

        # Set transparency using alpha channel in XML
        # python-pptx doesn't directly support transparency, so we use XML
        fill_elem = shape.fill._xPr
        srgbClr = fill_elem.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgbClr is not None:
            from lxml import etree
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', str(100000 - self.OVERLAY_TRANSPARENCY))  # 70000 = 70% opaque

        # Remove border
        shape.line.fill.background()

        # Add text to the shape
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # Set margins/padding
        tf.margin_left = Inches(padding)
        tf.margin_right = Inches(padding)
        tf.margin_top = Inches(padding)
        tf.margin_bottom = Inches(padding)

        # Center text vertically
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT if align_left else PP_ALIGN.CENTER
        self._apply_font(p, font_style)

        return shape

    def _render_default(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a default content slide."""
        self._render_title_content(slide, content, layout_name)

    def _render_title_content(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a title + content slide with optional takeaway subheader."""
        bullets = content.get("bullets", [])
        body = content.get("body", "")
        takeaway = content.get("takeaway", "")

        # Add takeaway subheader if provided
        takeaway_offset = self._add_takeaway(slide, takeaway)

        if not bullets and not body:
            return

        area = self._get_content_area(layout_name)

        # Adjust content area if takeaway was added
        if takeaway_offset > 0:
            area["top"] += takeaway_offset
            area["height"] -= takeaway_offset

        if body:
            textbox = self._create_text_box(
                slide, area["left"], area["top"],
                area["width"], min(area["height"], 1.5)
            )
            tf = textbox.text_frame
            p = tf.paragraphs[0]
            p.text = body
            self._apply_font(p, "body_large")

            # Adjust area for bullets
            area["top"] += 1.6
            area["height"] -= 1.6

        if bullets:
            self._add_bullets(slide, bullets, area)

    def _render_two_column(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a two-column comparison slide with optional takeaway subheader."""
        left_col = content.get("left_column", content.get("left", {}))
        right_col = content.get("right_column", content.get("right", {}))
        takeaway = content.get("takeaway", "")

        # Add takeaway subheader if provided
        takeaway_offset = self._add_takeaway(slide, takeaway)

        area = self._get_content_area(layout_name)

        # Adjust content area if takeaway was added
        if takeaway_offset > 0:
            area["top"] += takeaway_offset
            area["height"] -= takeaway_offset

        col_width = (area["width"] - 0.4) / 2  # Gap between columns

        if left_col:
            left_area = {
                "left": area["left"],
                "top": area["top"],
                "width": col_width,
                "height": area["height"]
            }
            self._render_column(slide, left_col, left_area)

        if right_col:
            right_area = {
                "left": area["left"] + col_width + 0.4,
                "top": area["top"],
                "width": col_width,
                "height": area["height"]
            }
            self._render_column(slide, right_col, right_area)

    def _render_column(self, slide: Slide, col_content: dict, area: Dict[str, float]) -> None:
        """Render a single column with proper PowerPoint bullet formatting.

        Paragraph formatting (from PowerPoint dialog):
        - Alignment: Left
        - Before text: 0.2"
        - Special: Hanging, By: 0.2"
        - Spacing Before: 0 pt
        - Spacing After: 6 pt
        - Line Spacing: Single
        """
        import re
        from lxml import etree

        header = col_content.get("header", col_content.get("heading", ""))
        bullets = col_content.get("bullets", [])

        textbox = self._create_text_box(
            slide, area["left"], area["top"],
            area["width"], area["height"]
        )
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        font_style = self.FONTS["bullet"]

        # Header
        if header:
            p = tf.paragraphs[0]
            p.text = header
            self._apply_font(p, "heading")
            p.space_after = Pt(10)

        # Bullets with proper PowerPoint formatting
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if (i > 0 or header) else tf.paragraphs[0]

            # Set bullet level and alignment
            p.level = 0
            p.alignment = PP_ALIGN.LEFT

            # Apply proper PowerPoint bullet formatting via XML
            try:
                pPr = p._p.get_or_add_pPr()

                # Set indentation attributes on pPr element:
                # marL = left margin (where text wraps to on subsequent lines)
                # indent = first line indent relative to marL (negative = hanging)
                #
                # PowerPoint "Before text: 0.2" + Hanging: 0.2"" means:
                # - marL = 0.4" (text starts here)
                # - indent = -0.2" (bullet hangs 0.2" left of marL, so at 0.2")
                # This creates 0.2" spacing between bullet (at 0.2") and text (at 0.4")
                marL_emu = int(Inches(0.4))  # Text starts at 0.4"
                indent_emu = int(Inches(-0.2))  # Bullet at marL + indent = 0.2"
                pPr.set(qn('a:marL'), str(marL_emu))
                pPr.set(qn('a:indent'), str(indent_emu))

                # Remove ALL existing child elements to rebuild from scratch
                for child in list(pPr):
                    pPr.remove(child)

                # OOXML schema requires elements in specific order:
                # 1. lnSpc, spcBef, spcAft (spacing)
                # 2. buClr (bullet color)
                # 3. buSzPct (bullet size)
                # 4. buFont (bullet font)
                # 5. buChar (bullet character)

                # Line Spacing: Single (100%)
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
                spcPct.set('val', '100000')

                # Spacing Before: 0 pt
                spcBef = etree.SubElement(pPr, qn('a:spcBef'))
                spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
                spcPts.set('val', '0')

                # Spacing After: 6 pt (600 = 6pt * 100)
                spcAft = etree.SubElement(pPr, qn('a:spcAft'))
                spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
                spcPts.set('val', '600')

                # Bullet color (black)
                buClr = etree.SubElement(pPr, qn('a:buClr'))
                srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
                srgbClr.set('val', '000000')

                # Bullet size (100% of text size)
                buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
                buSzPct.set('val', '100000')

                # Bullet font (Arial)
                buFont = etree.SubElement(pPr, qn('a:buFont'))
                buFont.set('typeface', 'Arial')
                buFont.set('pitchFamily', '34')
                buFont.set('charset', '0')

                # Bullet character (standard round bullet)
                buChar = etree.SubElement(pPr, qn('a:buChar'))
                buChar.set('char', '•')

            except Exception:
                pass

            # Check if bullet has a header pattern (CAPS TEXT: rest of text)
            header_match = re.match(r'^([A-Z][A-Z\s\-&/]+:)\s*(.*)$', bullet)

            if header_match:
                header_text = header_match.group(1)
                body_text = header_match.group(2)

                run1 = p.add_run()
                run1.text = header_text + " "
                run1.font.name = font_style["name"]
                run1.font.size = font_style["size"]
                run1.font.bold = True
                run1.font.color.rgb = self.COLORS["text_body"]

                if body_text:
                    run2 = p.add_run()
                    run2.text = body_text
                    run2.font.name = font_style["name"]
                    run2.font.size = font_style["size"]
                    run2.font.bold = False
                    run2.font.color.rgb = self.COLORS["text_body"]
            else:
                run = p.add_run()
                run.text = bullet
                run.font.name = font_style["name"]
                run.font.size = font_style["size"]
                run.font.bold = False
                run.font.color.rgb = self.COLORS["text_body"]

    def _render_key_metrics(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a key metrics slide with KPI boxes and optional takeaway subheader."""
        metrics = content.get("metrics", [])
        takeaway = content.get("takeaway", "")

        # Add takeaway subheader if provided
        takeaway_offset = self._add_takeaway(slide, takeaway)

        if not metrics:
            return

        area = self._get_content_area(layout_name)

        # Adjust content area if takeaway was added
        if takeaway_offset > 0:
            area["top"] += takeaway_offset
            area["height"] -= takeaway_offset

        num_metrics = min(len(metrics), 5)
        box_margin = 0.15
        box_width = (area["width"] - (num_metrics - 1) * box_margin) / num_metrics
        box_height = 1.4

        # Center vertically in content area
        top = area["top"] + (area["height"] - box_height) / 2

        for i, metric in enumerate(metrics[:5]):
            left = area["left"] + i * (box_width + box_margin)

            # Create box shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(top),
                Inches(box_width), Inches(box_height)
            )

            # Style the box with #309CE7 blue color
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.COLORS["chart_bar"]  # #309CE7
            shape.line.fill.background()

            # Add text
            tf = shape.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # Value
            p = tf.paragraphs[0]
            p.text = str(metric.get("value", ""))
            p.alignment = PP_ALIGN.CENTER
            self._apply_font(p, "metric_value")

            # Label
            p2 = tf.add_paragraph()
            p2.text = metric.get("label", "")
            p2.alignment = PP_ALIGN.CENTER
            self._apply_font(p2, "metric_label")

    def _render_table_slide(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a slide with a data table using style guide specifications."""
        headers = content.get("headers", [])
        data = content.get("data", [])
        takeaway = content.get("takeaway", "")

        # Add takeaway subheader if provided
        takeaway_offset = self._add_takeaway(slide, takeaway)

        if not data and not headers:
            return

        area = self._get_content_area(layout_name)

        # Adjust content area if takeaway was added
        if takeaway_offset > 0:
            area["top"] += takeaway_offset
            area["height"] -= takeaway_offset

        rows = len(data) + (1 if headers else 0)
        cols = len(headers) if headers else (len(data[0]) if data else 0)

        if rows == 0 or cols == 0:
            return

        # Calculate table dimensions - fit to content area
        row_height = min(0.35, (area["height"] - 0.2) / rows)
        table_height = min(rows * row_height, area["height"] - 0.2)

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(area["left"]), Inches(area["top"]),
            Inches(area["width"]), Inches(table_height)
        )
        table = table_shape.table

        # Get style config values (with fallbacks)
        if self.style_config:
            header_color = self.style_config.table.header_rgb
            row_odd_color = self.style_config.table.row_odd_rgb
            row_even_color = self.style_config.table.row_even_rgb
        else:
            header_color = self.COLORS["overlay_navy"]
            row_odd_color = self.COLORS["white"]
            row_even_color = self.COLORS["background"]

        # Style header row
        if headers:
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                self._style_table_cell(cell, is_header=True)

        # Add data rows with alternating colors (per style guide)
        start_row = 1 if headers else 0
        for i, row in enumerate(data):
            for j, value in enumerate(row):
                if j < cols:  # Ensure we don't exceed columns
                    cell = table.cell(start_row + i, j)
                    cell.text = str(value)

                    # Alternate row colors (per style guide: #FFFFFF / #F5F5F5)
                    cell.fill.solid()
                    if i % 2 == 0:
                        cell.fill.fore_color.rgb = row_odd_color
                    else:
                        cell.fill.fore_color.rgb = row_even_color

                    self._style_table_cell(cell, is_header=False)

    def _style_table_cell(self, cell, is_header: bool = False) -> None:
        """Style a table cell using style guide specifications."""
        # Get style config values (with fallbacks)
        if self.style_config:
            margin_lr = self.style_config.table.cell_margin_lr
            margin_tb = self.style_config.table.cell_margin_tb
            header_text_color = self.style_config.table.header_text_rgb
            row_text_color = self.style_config.table.row_text_rgb
            header_font_size = Pt(self.style_config.table.header_font_size_pt)
            row_font_size = Pt(self.style_config.table.row_font_size_pt)
            font_name = self.style_config.table.header_font_name
        else:
            margin_lr = Inches(0.1)
            margin_tb = Inches(0.05)
            header_text_color = self.COLORS["white"]
            row_text_color = self.COLORS["text_body"]
            header_font_size = self.FONTS["table_header"]["size"]
            row_font_size = self.FONTS["table_cell"]["size"]
            font_name = "Arial"

        # Apply cell margins (per style guide: 0.1" L/R, 0.05" T/B)
        cell.margin_left = margin_lr
        cell.margin_right = margin_lr
        cell.margin_top = margin_tb
        cell.margin_bottom = margin_tb

        # Vertically center text in cell
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = cell.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Apply font styling based on header/data row
        for para in tf.paragraphs:
            para.font.name = font_name
            if is_header:
                para.font.size = header_font_size
                para.font.bold = True
                para.font.color.rgb = header_text_color
            else:
                para.font.size = row_font_size
                para.font.bold = False
                para.font.color.rgb = row_text_color

    def _render_data_chart(self, slide: Slide, content: dict, layout_name: str) -> None:
        """Render a slide with a chart and optional takeaway subheader."""
        chart_data = content.get("chart_data", {})
        narrative = content.get("narrative", "")
        takeaway = content.get("takeaway", "")

        # Add takeaway subheader if provided
        takeaway_offset = self._add_takeaway(slide, takeaway)

        area = self._get_content_area(layout_name)

        # Adjust content area if takeaway was added
        if takeaway_offset > 0:
            area["top"] += takeaway_offset
            area["height"] -= takeaway_offset

        if chart_data:
            # Reserve space for narrative if present
            chart_height = area["height"] - 0.6 if narrative else area["height"]
            self._add_chart(slide, chart_data, {
                "left": area["left"],
                "top": area["top"],
                "width": area["width"],
                "height": chart_height
            })

        # Add narrative/source text
        if narrative:
            textbox = self._create_text_box(
                slide,
                area["left"],
                area["top"] + area["height"] - 0.4,
                area["width"],
                0.4,
                auto_size=False
            )
            p = textbox.text_frame.paragraphs[0]
            p.text = narrative
            self._apply_font(p, "caption")

    def _add_chart(self, slide: Slide, chart_data: dict, area: Dict[str, float]) -> None:
        """Add a chart to the slide."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        chart_type = chart_data.get("type", "column")
        categories = chart_data.get("categories", [])
        series_list = chart_data.get("series", [])

        if not categories or not series_list:
            return

        # Map chart types
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Create chart data
        data = CategoryChartData()
        data.categories = categories

        for series in series_list:
            data.add_series(series.get("name", "Series"), series.get("values", []))

        # Adjust size for pie charts
        chart_width = area["width"] * 0.75 if chart_type == "pie" else area["width"]

        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            Inches(area["left"]), Inches(area["top"]),
            Inches(chart_width), Inches(area["height"]),
            data
        )
        chart = chart_shape.chart

        # Style the chart based on type
        if chart_type == "pie":
            self._style_pie_chart(chart)
        elif chart_type == "line":
            self._style_line_chart(chart)
        else:
            self._style_bar_chart(chart)

    def _style_bar_chart(self, chart) -> None:
        """Style a bar/column chart using style guide specifications."""
        from pptx.enum.chart import XL_TICK_MARK

        # Get style config values (with fallbacks)
        if self.style_config:
            gridline_color = self.style_config.chart.gridline_rgb
            gridline_width = self.style_config.chart.gridline_width
            gridline_enabled = self.style_config.chart.gridline_enabled
            primary_color = self.style_config.chart.primary_series_rgb
            secondary_color = self.style_config.chart.secondary_series_rgb
            axis_color = self.style_config.chart.axis_label_rgb
            axis_font_size = Pt(self.style_config.chart.axis_label_font_size_pt)
            data_label_size = Pt(self.style_config.chart.data_label_font_size_pt)
        else:
            gridline_color = self.COLORS["chart_gridline"]
            gridline_width = self.LINE_WIDTH_GRIDLINE
            gridline_enabled = True
            primary_color = self.COLORS["chart_bar"]
            secondary_color = self.COLORS["secondary"]
            axis_color = self.COLORS["text_body"]
            axis_font_size = self.FONTS["chart_axis"]["size"]
            data_label_size = self.FONTS["chart_label"]["size"]

        try:
            plot = chart.plots[0]
            if hasattr(plot, 'series'):
                for i, series in enumerate(plot.series):
                    series.format.fill.solid()
                    if i == 0:
                        series.format.fill.fore_color.rgb = primary_color
                    else:
                        series.format.fill.fore_color.rgb = secondary_color

            # Add data labels
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = data_label_size
            data_labels.font.color.rgb = axis_color

            # Style category axis with no tick marks (per style guide)
            if hasattr(chart, 'category_axis'):
                chart.category_axis.tick_labels.font.size = axis_font_size
                chart.category_axis.tick_labels.font.color.rgb = axis_color
                chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE

                # Add gridlines to category axis (per PCCP style guide)
                if gridline_enabled and hasattr(chart.category_axis, 'major_gridlines'):
                    chart.category_axis.has_major_gridlines = True
                    gridlines = chart.category_axis.major_gridlines
                    gridlines.format.line.color.rgb = gridline_color
                    gridlines.format.line.width = gridline_width

            # Style value axis and gridlines with no tick marks (per style guide)
            if hasattr(chart, 'value_axis'):
                chart.value_axis.tick_labels.font.size = axis_font_size
                chart.value_axis.tick_labels.font.color.rgb = axis_color
                chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE

                # Style major gridlines (per PCCP style guide: 0.5pt, #E2E8F0)
                if gridline_enabled and hasattr(chart.value_axis, 'major_gridlines'):
                    chart.value_axis.has_major_gridlines = True
                    gridlines = chart.value_axis.major_gridlines
                    gridlines.format.line.color.rgb = gridline_color
                    gridlines.format.line.width = gridline_width

        except Exception as e:
            logger.debug(f"Error styling bar chart: {e}")

    def _style_line_chart(self, chart) -> None:
        """Style a line chart using style guide specifications."""
        from pptx.enum.chart import XL_DATA_LABEL_POSITION, XL_TICK_MARK

        # Get style config values (with fallbacks)
        if self.style_config:
            gridline_color = self.style_config.chart.gridline_rgb
            gridline_width = self.style_config.chart.gridline_width
            gridline_enabled = self.style_config.chart.gridline_enabled
            primary_color = self.style_config.chart.primary_series_rgb
            secondary_color = self.style_config.chart.secondary_series_rgb
            line_width = self.style_config.chart.line_width
            axis_color = self.style_config.chart.axis_label_rgb
            axis_font_size = Pt(self.style_config.chart.axis_label_font_size_pt)
            data_label_size = Pt(self.style_config.chart.data_label_font_size_pt)
        else:
            gridline_color = self.COLORS["chart_gridline"]
            gridline_width = self.LINE_WIDTH_GRIDLINE
            gridline_enabled = True
            primary_color = self.COLORS["chart_line_primary"]
            secondary_color = self.COLORS["chart_line_secondary"]
            line_width = self.LINE_WIDTH_CHART
            axis_color = self.COLORS["text_body"]
            axis_font_size = self.FONTS["chart_axis"]["size"]
            data_label_size = self.FONTS["chart_label"]["size"]

        try:
            plot = chart.plots[0]
            if hasattr(plot, 'series'):
                for i, series in enumerate(plot.series):
                    # Set line color
                    if i == 0:
                        series.format.line.color.rgb = primary_color
                    else:
                        series.format.line.color.rgb = secondary_color

                    # Set line width
                    series.format.line.width = line_width

            # Add data labels positioned ABOVE
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = data_label_size
            data_labels.font.color.rgb = axis_color
            data_labels.position = XL_DATA_LABEL_POSITION.ABOVE

            # Style category axis with no tick marks (per style guide)
            if hasattr(chart, 'category_axis'):
                chart.category_axis.tick_labels.font.size = axis_font_size
                chart.category_axis.tick_labels.font.color.rgb = axis_color
                chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE

                # Add gridlines to category axis (per PCCP style guide)
                if gridline_enabled and hasattr(chart.category_axis, 'major_gridlines'):
                    chart.category_axis.has_major_gridlines = True
                    gridlines = chart.category_axis.major_gridlines
                    gridlines.format.line.color.rgb = gridline_color
                    gridlines.format.line.width = gridline_width

            # Style value axis and gridlines with no tick marks (per style guide)
            if hasattr(chart, 'value_axis'):
                chart.value_axis.tick_labels.font.size = axis_font_size
                chart.value_axis.tick_labels.font.color.rgb = axis_color
                chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
                chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE

                # Style major gridlines (per PCCP style guide: 0.5pt, #E2E8F0)
                if gridline_enabled and hasattr(chart.value_axis, 'major_gridlines'):
                    chart.value_axis.has_major_gridlines = True
                    gridlines = chart.value_axis.major_gridlines
                    gridlines.format.line.color.rgb = gridline_color
                    gridlines.format.line.width = gridline_width

        except Exception as e:
            logger.debug(f"Error styling line chart: {e}")

    def _style_pie_chart(self, chart) -> None:
        """Style a pie chart with minimum 12pt labels."""
        from pptx.enum.chart import XL_LEGEND_POSITION

        pie_colors = [
            self.COLORS["primary"],
            self.COLORS["secondary"],
            self.COLORS["accent2"],
            RGBColor(0xFF, 0xB8, 0x00),  # Orange
            RGBColor(0x7C, 0x4D, 0xC4),  # Purple
            self.COLORS["accent1"],
        ]

        try:
            plot = chart.plots[0]
            if hasattr(plot, 'series') and len(plot.series) > 0:
                series = plot.series[0]
                for i, point in enumerate(series.points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = pie_colors[i % len(pie_colors)]

            # Add legend (minimum 12pt)
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
            chart.legend.font.size = self.FONTS["chart_label"]["size"]  # 12pt
        except Exception:
            pass

    def _apply_font(self, paragraph, style_name: str) -> None:
        """Apply font styling to a paragraph."""
        style = self.FONTS.get(style_name, self.FONTS["body"])

        if paragraph.runs:
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()

        run.font.name = style["name"]
        run.font.size = style["size"]
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)

        color_key = style.get("color", "text_body")
        run.font.color.rgb = self.COLORS.get(color_key, self.COLORS["text_body"])

    # ==================== Component Library Integration ====================

    def find_library_chart(
        self,
        chart_type: str = "column",
        category: Optional[str] = None,
        min_series: int = 1
    ) -> Optional[dict]:
        """Find a matching chart component from the library."""
        if not self.library:
            return None

        type_map = {
            "column": "COLUMN",
            "bar": "BAR",
            "line": "LINE",
            "pie": "PIE",
            "area": "AREA",
        }
        lib_type = type_map.get(chart_type.lower(), "COLUMN")

        results = self.library.search_charts(
            chart_type=lib_type,
            min_series=min_series,
            limit=5
        )

        return results[0] if results else None

    def find_library_table(
        self,
        rows: int = 3,
        cols: int = 3,
        category: Optional[str] = None
    ) -> Optional[dict]:
        """Find a matching table component from the library."""
        if not self.library:
            return None

        results = self.library.search_tables(
            category=category,
            min_rows=max(1, rows - 2),
            max_rows=rows + 5,
            min_cols=max(1, cols - 1),
            max_cols=cols + 2,
            limit=5
        )

        if results:
            results.sort(key=lambda t: abs(t.get('rows', 0) - rows) + abs(t.get('cols', 0) - cols))
            return results[0]

        return None


def main():
    """Test the template renderer."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Template Renderer Test")
    parser.add_argument(
        "--template",
        default="pptx_templates/pptx_template_business_consulting_toolkit/template_business_consulting_toolkit.pptx",
        help="Path to template"
    )
    parser.add_argument(
        "--output",
        default="pptx_generator/output/template_render_test.pptx",
        help="Output file"
    )

    args = parser.parse_args()

    renderer = TemplateRenderer(args.template)
    prs = renderer.create_presentation()

    # Create test slides
    renderer.create_slide(prs, "title_slide", {
        "title": "Test Presentation",
        "subtitle": "Template Renderer Demo"
    })

    renderer.create_slide(prs, "section_divider", {
        "title": "Section One"
    })

    renderer.create_slide(prs, "title_content", {
        "title": "Key Points",
        "bullets": ["First important point", "Second important point", "Third important point"]
    })

    renderer.create_slide(prs, "key_metrics", {
        "title": "Key Metrics",
        "metrics": [
            {"label": "Revenue", "value": "$1.2M"},
            {"label": "Growth", "value": "25%"},
            {"label": "Users", "value": "10K"},
            {"label": "NPS", "value": "72"}
        ]
    })

    renderer.create_slide(prs, "table_slide", {
        "title": "Data Table",
        "headers": ["Name", "Value", "Status"],
        "data": [
            ["Item 1", "100", "Active"],
            ["Item 2", "200", "Pending"],
            ["Item 3", "300", "Complete"]
        ]
    })

    prs.save(args.output)
    print(f"Saved to: {args.output}")


if __name__ == "__main__":
    main()
