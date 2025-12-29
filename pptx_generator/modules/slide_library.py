"""
Slide Library Module

Manages reusable slide retrieval and copying between presentations.
"""

import json
import logging
from copy import deepcopy
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.slide import Slide

logger = logging.getLogger(__name__)


class SlideLibrary:
    """Manages reusable slide retrieval and copying."""

    def __init__(self, templates_dir: str, catalog_path: str):
        """
        Initialize the slide library.

        Args:
            templates_dir: Path to directory containing PPTX templates
            catalog_path: Path to slide_catalog.json
        """
        self.templates_dir = Path(templates_dir)
        self.catalog_path = Path(catalog_path)
        self.catalog: dict = {}
        self.template_cache: dict[str, Presentation] = {}
        self.slide_index: list[dict] = []

        self._load_catalog()
        self._build_index()

    def _load_catalog(self) -> None:
        """Load the slide catalog."""
        if self.catalog_path.exists():
            with open(self.catalog_path, "r", encoding="utf-8") as f:
                self.catalog = json.load(f)
            logger.info(f"Loaded catalog with {len(self.catalog.get('slide_types', []))} slide types")
        else:
            logger.warning(f"Catalog not found at {self.catalog_path}")
            self.catalog = {"slide_types": []}

    def _build_index(self) -> None:
        """Build searchable index from catalog."""
        self.slide_index = []

        for slide_type in self.catalog.get("slide_types", []):
            for example in slide_type.get("examples", []):
                entry = {
                    "template": example["template"],
                    "slide_index": example["slide_index"],
                    "text_preview": example.get("text_preview", ""),
                    "slide_type_id": slide_type["id"],
                    "slide_type_name": slide_type["name"],
                    "master_layout": slide_type["master_layout"],
                    "keywords": self._extract_keywords(example, slide_type)
                }
                self.slide_index.append(entry)

        logger.info(f"Built index with {len(self.slide_index)} slides")

    def _extract_keywords(self, example: dict, slide_type: dict) -> list[str]:
        """Extract searchable keywords from slide information."""
        keywords = []

        # From text preview
        text = example.get("text_preview", "").lower()
        keywords.extend(text.split())

        # From slide type
        keywords.append(slide_type["id"].lower())
        keywords.append(slide_type["name"].lower())
        keywords.extend(slide_type.get("description", "").lower().split())

        # Clean and deduplicate
        keywords = [k.strip(",.;:()[]") for k in keywords if len(k) > 2]
        return list(set(keywords))

    def _get_template(self, template_name: str) -> Optional[Presentation]:
        """Get or load a template presentation."""
        if template_name in self.template_cache:
            return self.template_cache[template_name]

        # Search for template file
        template_path = None
        for path in self.templates_dir.rglob(template_name):
            template_path = path
            break

        if not template_path or not template_path.exists():
            logger.error(f"Template not found: {template_name}")
            return None

        try:
            prs = Presentation(str(template_path))
            self.template_cache[template_name] = prs
            return prs
        except Exception as e:
            logger.error(f"Error loading template {template_name}: {e}")
            return None

    def search(
        self,
        query: str,
        slide_types: Optional[list[str]] = None,
        limit: int = 10
    ) -> list[dict]:
        """
        Search for slides matching a query.

        Args:
            query: Search query string
            slide_types: Optional list of slide type IDs to filter by
            limit: Maximum number of results

        Returns:
            List of matching slides with relevance scores
        """
        query_terms = query.lower().split()
        results = []

        for entry in self.slide_index:
            # Filter by slide type if specified
            if slide_types and entry["slide_type_id"] not in slide_types:
                continue

            # Calculate relevance score
            score = 0
            for term in query_terms:
                if term in entry["keywords"]:
                    score += 2
                if term in entry["text_preview"].lower():
                    score += 3
                if term in entry["slide_type_name"].lower():
                    score += 1

            if score > 0:
                results.append({
                    **entry,
                    "relevance_score": score
                })

        # Sort by relevance
        results.sort(key=lambda x: x["relevance_score"], reverse=True)
        return results[:limit]

    def get_by_type(self, slide_type_id: str, limit: int = 5) -> list[dict]:
        """Get slides of a specific type."""
        results = [
            entry for entry in self.slide_index
            if entry["slide_type_id"] == slide_type_id
        ]
        return results[:limit]

    def get_slide_type_info(self, slide_type_id: str) -> Optional[dict]:
        """Get detailed information about a slide type."""
        for slide_type in self.catalog.get("slide_types", []):
            if slide_type["id"] == slide_type_id:
                return slide_type
        return None

    def list_slide_types(self) -> list[dict]:
        """List all available slide types."""
        return [
            {
                "id": st["id"],
                "name": st["name"],
                "description": st["description"],
                "occurrence_count": st.get("occurrence_count", 0)
            }
            for st in self.catalog.get("slide_types", [])
        ]

    def copy_slide(
        self,
        source_template: str,
        slide_index: int,
        target_presentation: Presentation
    ) -> Optional[Slide]:
        """
        Copy a slide from a source template to a target presentation.

        This creates a duplicate of the source slide including all shapes,
        text, and formatting.

        Args:
            source_template: Name of the source template file
            slide_index: Index of the slide to copy
            target_presentation: Target Presentation object

        Returns:
            The newly created slide, or None if copy failed
        """
        source_prs = self._get_template(source_template)
        if not source_prs:
            return None

        if slide_index >= len(source_prs.slides):
            logger.error(f"Slide index {slide_index} out of range for {source_template}")
            return None

        source_slide = source_prs.slides[slide_index]

        try:
            # Find matching layout in target presentation
            layout_name = source_slide.slide_layout.name
            target_layout = self._find_matching_layout(target_presentation, layout_name)

            if not target_layout:
                # Use blank layout as fallback
                target_layout = target_presentation.slide_layouts[6]  # Usually blank

            # Add new slide
            new_slide = target_presentation.slides.add_slide(target_layout)

            # Copy shapes from source to target
            self._copy_shapes(source_slide, new_slide)

            logger.info(f"Copied slide {slide_index} from {source_template}")
            return new_slide

        except Exception as e:
            logger.error(f"Error copying slide: {e}")
            return None

    def _find_matching_layout(
        self,
        presentation: Presentation,
        layout_name: str
    ) -> Optional[object]:
        """Find a matching slide layout by name."""
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None

    def _copy_shapes(self, source_slide: Slide, target_slide: Slide) -> None:
        """Copy all shapes from source slide to target slide."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Emu

        for shape in source_slide.shapes:
            try:
                # Skip placeholders that already exist in target
                if shape.is_placeholder:
                    self._copy_placeholder_content(shape, target_slide)
                    continue

                # Copy based on shape type
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    self._copy_text_box(shape, target_slide)
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    self._copy_auto_shape(shape, target_slide)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._copy_picture(shape, target_slide)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self._copy_table(shape, target_slide)
                elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    self._copy_line(shape, target_slide)
                # Add more shape types as needed

            except Exception as e:
                logger.warning(f"Could not copy shape: {e}")

    def _copy_placeholder_content(self, source_shape, target_slide: Slide) -> None:
        """Copy content to matching placeholder in target slide."""
        try:
            placeholder_idx = source_shape.placeholder_format.idx

            for shape in target_slide.placeholders:
                if shape.placeholder_format.idx == placeholder_idx:
                    if source_shape.has_text_frame:
                        self._copy_text_frame(source_shape.text_frame, shape.text_frame)
                    break
        except Exception:
            pass

    def _copy_text_box(self, source_shape, target_slide: Slide) -> None:
        """Copy a text box shape."""
        from pptx.util import Emu

        new_shape = target_slide.shapes.add_textbox(
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )

        if source_shape.has_text_frame:
            self._copy_text_frame(source_shape.text_frame, new_shape.text_frame)

    def _copy_auto_shape(self, source_shape, target_slide: Slide) -> None:
        """Copy an auto shape."""
        try:
            new_shape = target_slide.shapes.add_shape(
                source_shape.auto_shape_type,
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )

            # Copy fill
            if source_shape.fill.type is not None:
                self._copy_fill(source_shape.fill, new_shape.fill)

            # Copy line
            if source_shape.line:
                self._copy_line_format(source_shape.line, new_shape.line)

            # Copy text
            if source_shape.has_text_frame:
                self._copy_text_frame(source_shape.text_frame, new_shape.text_frame)

        except Exception as e:
            logger.warning(f"Could not copy auto shape: {e}")

    def _copy_picture(self, source_shape, target_slide: Slide) -> None:
        """Copy a picture shape."""
        try:
            # Get the image blob
            image = source_shape.image
            new_shape = target_slide.shapes.add_picture(
                image.blob,
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
        except Exception as e:
            logger.warning(f"Could not copy picture: {e}")

    def _copy_table(self, source_shape, target_slide: Slide) -> None:
        """Copy a table shape."""
        try:
            source_table = source_shape.table
            rows = len(source_table.rows)
            cols = len(source_table.columns)

            new_table = target_slide.shapes.add_table(
                rows, cols,
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            ).table

            # Copy cell contents
            for i, row in enumerate(source_table.rows):
                for j, cell in enumerate(row.cells):
                    target_cell = new_table.cell(i, j)
                    if cell.text_frame:
                        self._copy_text_frame(cell.text_frame, target_cell.text_frame)

        except Exception as e:
            logger.warning(f"Could not copy table: {e}")

    def _copy_line(self, source_shape, target_slide: Slide) -> None:
        """Copy a line shape."""
        try:
            from pptx.enum.shapes import MSO_CONNECTOR

            connector = target_slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                source_shape.begin_x,
                source_shape.begin_y,
                source_shape.end_x,
                source_shape.end_y
            )

            # Copy line formatting
            if source_shape.line:
                self._copy_line_format(source_shape.line, connector.line)

        except Exception as e:
            logger.warning(f"Could not copy line: {e}")

    def _copy_text_frame(self, source_tf, target_tf) -> None:
        """Copy text frame content and formatting."""
        # Clear existing paragraphs except first
        while len(target_tf.paragraphs) > 1:
            p = target_tf.paragraphs[-1]
            p._p.getparent().remove(p._p)

        for i, source_para in enumerate(source_tf.paragraphs):
            if i == 0:
                target_para = target_tf.paragraphs[0]
            else:
                target_para = target_tf.paragraphs[-1]._add_paragraph()

            # Copy paragraph properties
            try:
                target_para.level = source_para.level
                target_para.alignment = source_para.alignment
            except:
                pass

            # Copy runs
            for j, source_run in enumerate(source_para.runs):
                if j == 0 and len(target_para.runs) > 0:
                    target_run = target_para.runs[0]
                else:
                    target_run = target_para.add_run()

                target_run.text = source_run.text

                # Copy font properties
                try:
                    target_run.font.name = source_run.font.name
                    target_run.font.size = source_run.font.size
                    target_run.font.bold = source_run.font.bold
                    target_run.font.italic = source_run.font.italic
                    if source_run.font.color.type == 1:  # RGB
                        target_run.font.color.rgb = source_run.font.color.rgb
                except:
                    pass

    def _copy_fill(self, source_fill, target_fill) -> None:
        """Copy fill properties."""
        try:
            if source_fill.type == 1:  # Solid
                target_fill.solid()
                if source_fill.fore_color.type == 1:  # RGB
                    target_fill.fore_color.rgb = source_fill.fore_color.rgb
        except:
            pass

    def _copy_line_format(self, source_line, target_line) -> None:
        """Copy line formatting."""
        try:
            target_line.width = source_line.width
            if source_line.fill.type == 1:  # Solid
                target_line.fill.solid()
                if source_line.fill.fore_color.type == 1:  # RGB
                    target_line.fill.fore_color.rgb = source_line.fill.fore_color.rgb
        except:
            pass

    def copy_section(
        self,
        section_id: str,
        target_presentation: Presentation,
        reusable_sections: dict
    ) -> list[Slide]:
        """
        Copy a reusable section to a target presentation.

        Args:
            section_id: ID of the reusable section
            target_presentation: Target Presentation object
            reusable_sections: Dictionary of reusable section definitions

        Returns:
            List of copied slides
        """
        if section_id not in reusable_sections:
            logger.error(f"Section not found: {section_id}")
            return []

        section = reusable_sections[section_id]
        source_template = section.get("source_template")
        source_slides = section.get("source_slides", [])

        copied = []
        for slide_idx in source_slides:
            new_slide = self.copy_slide(source_template, slide_idx, target_presentation)
            if new_slide:
                copied.append(new_slide)

        return copied


def main():
    """Test the slide library."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Slide Library")
    parser.add_argument(
        "--templates-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_templates",
        help="Directory containing PPTX templates"
    )
    parser.add_argument(
        "--catalog",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/config/slide_catalog.json",
        help="Path to slide catalog"
    )
    parser.add_argument(
        "--search",
        type=str,
        help="Search query"
    )
    parser.add_argument(
        "--list-types",
        action="store_true",
        help="List all slide types"
    )

    args = parser.parse_args()

    library = SlideLibrary(args.templates_dir, args.catalog)

    if args.list_types:
        print("\nAvailable Slide Types:")
        print("-" * 60)
        for st in library.list_slide_types():
            print(f"  {st['id']}: {st['name']} ({st['occurrence_count']} slides)")

    if args.search:
        print(f"\nSearch results for '{args.search}':")
        print("-" * 60)
        results = library.search(args.search)
        for r in results:
            print(f"  [{r['slide_type_name']}] {r['template']}:slide {r['slide_index']}")
            print(f"    Preview: {r['text_preview'][:80]}...")


if __name__ == "__main__":
    main()
