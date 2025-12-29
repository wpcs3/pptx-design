"""
Slide Pool - Reference Slide Matching and Cloning

PPTAgent-inspired slide pool that indexes existing slides from templates,
enabling slide-level matching and cloning for better visual fidelity.

Key Features:
- Index slides by functional type (title, content, data, comparison, etc.)
- Match new content to best-fitting reference slides
- Clone slides and apply targeted edits while preserving formatting
- Support for edit actions (replace_title, update_bullets, swap_chart, etc.)

Phase 4 Enhancement (2025-12-29):
Based on PPTAgent's two-phase architecture (github.com/icip-cas/PPTAgent)
"""

import copy
import hashlib
import json
import logging
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Callable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide

logger = logging.getLogger(__name__)


# =============================================================================
# Functional Slide Types (from PPTAgent)
# =============================================================================

class FunctionalType(Enum):
    """Functional types for slides based on their purpose."""
    TITLE = "title"                    # Opening/title slide
    SECTION = "section"                # Section divider
    AGENDA = "agenda"                  # Table of contents/agenda
    CONTENT = "content"                # Text-heavy content
    CONTENT_WITH_VISUAL = "content_visual"  # Content + image/diagram
    DATA_CHART = "data_chart"          # Chart-focused slide
    DATA_TABLE = "data_table"          # Table-focused slide
    COMPARISON = "comparison"          # Two-column comparison
    METRICS = "metrics"                # KPI/metrics boxes
    TIMELINE = "timeline"              # Timeline/roadmap
    TEAM = "team"                      # Team/profile slide
    QUOTE = "quote"                    # Quote/testimonial
    CLOSING = "closing"                # Thank you/closing slide
    UNKNOWN = "unknown"


# =============================================================================
# Edit Actions
# =============================================================================

class EditActionType(Enum):
    """Types of edit actions that can be applied to cloned slides."""
    REPLACE_TITLE = "replace_title"
    REPLACE_SUBTITLE = "replace_subtitle"
    UPDATE_BULLETS = "update_bullets"
    REPLACE_BODY = "replace_body"
    SWAP_CHART_DATA = "swap_chart_data"
    UPDATE_TABLE = "update_table"
    REPLACE_IMAGE = "replace_image"
    UPDATE_METRICS = "update_metrics"
    CLEAR_PLACEHOLDER = "clear_placeholder"


@dataclass
class EditAction:
    """An edit action to apply to a cloned slide."""
    action_type: EditActionType
    target: str  # Placeholder name or shape index
    value: Any   # New value to apply
    options: Dict[str, Any] = field(default_factory=dict)

    def __repr__(self):
        return f"EditAction({self.action_type.value}, target={self.target})"


# =============================================================================
# Indexed Slide
# =============================================================================

@dataclass
class IndexedSlide:
    """A slide indexed from a template with metadata for matching."""
    slide_id: str
    template_path: str
    slide_index: int
    functional_type: FunctionalType

    # Content analysis
    title: Optional[str] = None
    subtitle: Optional[str] = None
    bullet_count: int = 0
    word_count: int = 0

    # Visual elements
    has_chart: bool = False
    chart_type: Optional[str] = None
    has_table: bool = False
    table_dimensions: Optional[Tuple[int, int]] = None  # (rows, cols)
    has_image: bool = False
    image_count: int = 0
    shape_count: int = 0

    # Layout info
    layout_name: str = ""
    placeholder_types: List[str] = field(default_factory=list)

    # Matching metadata
    content_signature: str = ""  # Hash for quick matching
    element_pattern: str = ""    # Pattern like "T-B6-C" (Title, 6 Bullets, Chart)

    # Quality score from original template (higher = better designed)
    quality_score: float = 1.0

    def matches_content(self, content: Dict[str, Any]) -> float:
        """
        Calculate match score (0-1) between this slide and new content.

        Higher score = better match.
        """
        score = 0.0
        max_score = 0.0

        # Match functional type (most important)
        max_score += 3.0
        content_type = content.get("slide_type", "").lower()
        type_mapping = {
            "title_slide": FunctionalType.TITLE,
            "section_divider": FunctionalType.SECTION,
            "agenda": FunctionalType.AGENDA,
            "title_content": FunctionalType.CONTENT,
            "content_slide": FunctionalType.CONTENT,
            "data_chart": FunctionalType.DATA_CHART,
            "chart_slide": FunctionalType.DATA_CHART,
            "table_slide": FunctionalType.DATA_TABLE,
            "two_column": FunctionalType.COMPARISON,
            "comparison": FunctionalType.COMPARISON,
            "key_metrics": FunctionalType.METRICS,
            "metrics": FunctionalType.METRICS,
            "timeline": FunctionalType.TIMELINE,
            "thank_you": FunctionalType.CLOSING,
            "closing": FunctionalType.CLOSING,
        }
        expected_type = type_mapping.get(content_type, FunctionalType.UNKNOWN)
        if self.functional_type == expected_type:
            score += 3.0
        elif self.functional_type == FunctionalType.CONTENT and expected_type == FunctionalType.UNKNOWN:
            score += 1.5  # Content is a reasonable fallback

        # Match visual elements
        max_score += 2.0
        if "chart_data" in content and self.has_chart:
            score += 2.0
        elif "chart_data" not in content and not self.has_chart:
            score += 1.0

        max_score += 2.0
        has_table = "headers" in content and "data" in content
        if has_table and self.has_table:
            score += 2.0
        elif not has_table and not self.has_table:
            score += 1.0

        # Match bullet count (approximate)
        max_score += 1.0
        content_bullets = len(content.get("bullets", []))
        if content_bullets > 0:
            bullet_diff = abs(self.bullet_count - content_bullets)
            if bullet_diff == 0:
                score += 1.0
            elif bullet_diff <= 2:
                score += 0.5

        # Match metrics
        max_score += 1.0
        has_metrics = "metrics" in content
        if has_metrics and self.functional_type == FunctionalType.METRICS:
            score += 1.0

        # Quality bonus
        score *= self.quality_score

        return score / max_score if max_score > 0 else 0.0


# =============================================================================
# Slide Pool
# =============================================================================

class SlidePool:
    """
    Pool of indexed reference slides for matching and cloning.

    Usage:
        pool = SlidePool()
        pool.index_template("template.pptx")

        # Find best matching slide for content
        match = pool.find_best_match({"slide_type": "data_chart", "chart_data": {...}})

        # Clone and edit
        new_slide = pool.clone_and_edit(match, prs, [
            EditAction(EditActionType.REPLACE_TITLE, "title", "New Title"),
            EditAction(EditActionType.SWAP_CHART_DATA, "chart", new_chart_data),
        ])
    """

    def __init__(self, index_path: Optional[str] = None):
        """
        Initialize slide pool.

        Args:
            index_path: Optional path to load/save slide index
        """
        self.slides: List[IndexedSlide] = []
        self.index_path = Path(index_path) if index_path else None
        self._templates: Dict[str, Presentation] = {}  # Cached template presentations

        if self.index_path and self.index_path.exists():
            self._load_index()

    def index_template(self, template_path: str, quality_score: float = 1.0) -> int:
        """
        Index all slides from a template.

        Args:
            template_path: Path to PPTX template
            quality_score: Quality score for slides from this template (0-1)

        Returns:
            Number of slides indexed
        """
        template_path = str(template_path)
        prs = Presentation(template_path)
        count = 0

        for idx, slide in enumerate(prs.slides):
            indexed = self._analyze_slide(slide, template_path, idx)
            indexed.quality_score = quality_score
            self.slides.append(indexed)
            count += 1

        logger.info(f"Indexed {count} slides from {template_path}")
        return count

    def index_templates_directory(self, templates_dir: str) -> int:
        """Index all templates in a directory."""
        templates_dir = Path(templates_dir)
        total = 0

        for pptx_file in templates_dir.glob("**/*.pptx"):
            if pptx_file.name.startswith("~"):  # Skip temp files
                continue
            try:
                count = self.index_template(str(pptx_file))
                total += count
            except Exception as e:
                logger.warning(f"Failed to index {pptx_file}: {e}")

        return total

    def find_best_match(
        self,
        content: Dict[str, Any],
        min_score: float = 0.3
    ) -> Optional[IndexedSlide]:
        """
        Find the best matching slide for given content.

        Args:
            content: Slide content dictionary
            min_score: Minimum match score to accept

        Returns:
            Best matching IndexedSlide or None if no good match
        """
        if not self.slides:
            return None

        best_match = None
        best_score = 0.0

        for slide in self.slides:
            score = slide.matches_content(content)
            if score > best_score and score >= min_score:
                best_score = score
                best_match = slide

        if best_match:
            logger.debug(f"Best match: {best_match.slide_id} (score: {best_score:.2f})")

        return best_match

    def find_matches(
        self,
        content: Dict[str, Any],
        limit: int = 5,
        min_score: float = 0.2
    ) -> List[Tuple[IndexedSlide, float]]:
        """
        Find top N matching slides with scores.

        Returns:
            List of (IndexedSlide, score) tuples sorted by score descending
        """
        matches = []

        for slide in self.slides:
            score = slide.matches_content(content)
            if score >= min_score:
                matches.append((slide, score))

        matches.sort(key=lambda x: x[1], reverse=True)
        return matches[:limit]

    def clone_and_edit(
        self,
        source: IndexedSlide,
        target_prs: Presentation,
        edits: List[EditAction]
    ) -> Slide:
        """
        Clone a source slide into target presentation and apply edits.

        Args:
            source: IndexedSlide to clone
            target_prs: Target presentation to add slide to
            edits: List of EditActions to apply

        Returns:
            The new slide (already added to target_prs)
        """
        # Load source presentation if not cached
        if source.template_path not in self._templates:
            self._templates[source.template_path] = Presentation(source.template_path)

        source_prs = self._templates[source.template_path]
        source_slide = source_prs.slides[source.slide_index]

        # Clone the slide
        new_slide = self._clone_slide(source_slide, target_prs)

        # Apply edits
        for edit in edits:
            self._apply_edit(new_slide, edit)

        return new_slide

    def _analyze_slide(
        self,
        slide: Slide,
        template_path: str,
        slide_index: int
    ) -> IndexedSlide:
        """Analyze a slide and create IndexedSlide metadata."""
        # Generate unique ID
        slide_id = hashlib.md5(
            f"{template_path}:{slide_index}".encode()
        ).hexdigest()[:12]

        # Extract content
        title = None
        subtitle = None
        bullet_count = 0
        word_count = 0
        has_chart = False
        chart_type = None
        has_table = False
        table_dims = None
        has_image = False
        image_count = 0
        placeholder_types = []

        for shape in slide.shapes:
            # Check for title
            if shape == slide.shapes.title and shape.has_text_frame:
                title = shape.text_frame.text.strip()

            # Check placeholders
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                placeholder_types.append(str(ph_type))

                # Subtitle detection
                if "SUBTITLE" in str(ph_type) and shape.has_text_frame:
                    subtitle = shape.text_frame.text.strip()

            # Count text content
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        word_count += len(text.split())
                        # Count bullets (indented or marked paragraphs)
                        if para.level > 0 or text.startswith(("•", "-", "*")):
                            bullet_count += 1

            # Check for charts
            if shape.has_chart:
                has_chart = True
                chart_type = str(shape.chart.chart_type)

            # Check for tables
            if shape.has_table:
                has_table = True
                table_dims = (len(shape.table.rows), len(shape.table.columns))

            # Check for images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                image_count += 1

        # Determine functional type
        functional_type = self._infer_functional_type(
            title=title,
            bullet_count=bullet_count,
            has_chart=has_chart,
            has_table=has_table,
            has_image=has_image,
            shape_count=len(slide.shapes),
            layout_name=slide.slide_layout.name if slide.slide_layout else ""
        )

        # Build element pattern
        pattern_parts = []
        if title:
            pattern_parts.append("T")
        if bullet_count > 0:
            pattern_parts.append(f"B{bullet_count}")
        if has_chart:
            pattern_parts.append("C")
        if has_table:
            pattern_parts.append("Tb")
        if has_image:
            pattern_parts.append(f"I{image_count}")
        element_pattern = "-".join(pattern_parts) or "Empty"

        # Content signature for quick matching
        sig_data = f"{functional_type.value}:{element_pattern}:{bullet_count}"
        content_signature = hashlib.md5(sig_data.encode()).hexdigest()[:8]

        return IndexedSlide(
            slide_id=slide_id,
            template_path=template_path,
            slide_index=slide_index,
            functional_type=functional_type,
            title=title,
            subtitle=subtitle,
            bullet_count=bullet_count,
            word_count=word_count,
            has_chart=has_chart,
            chart_type=chart_type,
            has_table=has_table,
            table_dimensions=table_dims,
            has_image=has_image,
            image_count=image_count,
            shape_count=len(slide.shapes),
            layout_name=slide.slide_layout.name if slide.slide_layout else "",
            placeholder_types=placeholder_types,
            content_signature=content_signature,
            element_pattern=element_pattern,
        )

    def _infer_functional_type(
        self,
        title: Optional[str],
        bullet_count: int,
        has_chart: bool,
        has_table: bool,
        has_image: bool,
        shape_count: int,
        layout_name: str
    ) -> FunctionalType:
        """Infer functional type from slide characteristics."""
        layout_lower = layout_name.lower()

        # Check layout name hints
        if "title" in layout_lower and "content" not in layout_lower:
            return FunctionalType.TITLE
        if "section" in layout_lower:
            return FunctionalType.SECTION
        if "comparison" in layout_lower or "two" in layout_lower:
            return FunctionalType.COMPARISON
        if "blank" in layout_lower and shape_count <= 2:
            return FunctionalType.TITLE

        # Check content characteristics
        if has_chart:
            return FunctionalType.DATA_CHART
        if has_table:
            return FunctionalType.DATA_TABLE
        if has_image and bullet_count <= 3:
            return FunctionalType.CONTENT_WITH_VISUAL

        # Check for special slides by title patterns
        if title:
            title_lower = title.lower()
            if any(w in title_lower for w in ["agenda", "contents", "overview", "outline"]):
                return FunctionalType.AGENDA
            if any(w in title_lower for w in ["thank", "question", "contact", "end"]):
                return FunctionalType.CLOSING
            if any(w in title_lower for w in ["team", "leadership", "about us"]):
                return FunctionalType.TEAM
            if any(w in title_lower for w in ["timeline", "roadmap", "milestone"]):
                return FunctionalType.TIMELINE
            if any(w in title_lower for w in ["metric", "kpi", "performance"]):
                return FunctionalType.METRICS

        # Default based on content density
        if bullet_count > 0:
            return FunctionalType.CONTENT
        if shape_count <= 3:
            return FunctionalType.SECTION

        return FunctionalType.UNKNOWN

    def _clone_slide(self, source_slide: Slide, target_prs: Presentation) -> Slide:
        """
        Clone a slide from source to target presentation.

        This creates a new slide with the same layout and copies all shapes.
        """
        # Try to find matching layout in target presentation
        source_layout_name = source_slide.slide_layout.name if source_slide.slide_layout else ""
        target_layout = None

        for layout in target_prs.slide_layouts:
            if layout.name == source_layout_name:
                target_layout = layout
                break

        # Fall back to first layout if no match
        if target_layout is None:
            target_layout = target_prs.slide_layouts[0]

        # Create new slide
        new_slide = target_prs.slides.add_slide(target_layout)

        # Copy shapes from source
        for shape in source_slide.shapes:
            self._copy_shape(shape, new_slide)

        return new_slide

    def _copy_shape(self, source_shape, target_slide: Slide):
        """Copy a shape to target slide."""
        try:
            # Handle text frames
            if source_shape.has_text_frame:
                # Find matching placeholder or create textbox
                if source_shape.is_placeholder:
                    ph_idx = source_shape.placeholder_format.idx
                    for shape in target_slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
                            self._copy_text_frame(source_shape.text_frame, shape.text_frame)
                            return
                else:
                    # Create textbox
                    textbox = target_slide.shapes.add_textbox(
                        source_shape.left,
                        source_shape.top,
                        source_shape.width,
                        source_shape.height
                    )
                    self._copy_text_frame(source_shape.text_frame, textbox.text_frame)

        except Exception as e:
            logger.debug(f"Could not copy shape: {e}")

    def _copy_text_frame(self, source_tf, target_tf):
        """Copy text frame content."""
        # Clear target
        for para in list(target_tf.paragraphs)[1:]:
            p = para._p
            p.getparent().remove(p)

        # Copy paragraphs
        for i, source_para in enumerate(source_tf.paragraphs):
            if i == 0:
                target_para = target_tf.paragraphs[0]
            else:
                target_para = target_tf.add_paragraph()

            target_para.text = source_para.text
            target_para.level = source_para.level

            # Copy font properties
            if source_para.font.size:
                target_para.font.size = source_para.font.size
            if source_para.font.bold is not None:
                target_para.font.bold = source_para.font.bold

    def _apply_edit(self, slide: Slide, edit: EditAction):
        """Apply an edit action to a slide."""
        try:
            if edit.action_type == EditActionType.REPLACE_TITLE:
                self._edit_title(slide, edit.value)

            elif edit.action_type == EditActionType.REPLACE_SUBTITLE:
                self._edit_subtitle(slide, edit.value)

            elif edit.action_type == EditActionType.UPDATE_BULLETS:
                self._edit_bullets(slide, edit.value, edit.options)

            elif edit.action_type == EditActionType.REPLACE_BODY:
                self._edit_body(slide, edit.value)

            elif edit.action_type == EditActionType.SWAP_CHART_DATA:
                self._edit_chart(slide, edit.value)

            elif edit.action_type == EditActionType.UPDATE_TABLE:
                self._edit_table(slide, edit.value)

            elif edit.action_type == EditActionType.CLEAR_PLACEHOLDER:
                self._clear_placeholder(slide, edit.target)

        except Exception as e:
            logger.warning(f"Failed to apply edit {edit}: {e}")

    def _edit_title(self, slide: Slide, new_title: str):
        """Replace slide title."""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide.shapes.title.text_frame.paragraphs[0].text = new_title

    def _edit_subtitle(self, slide: Slide, new_subtitle: str):
        """Replace slide subtitle."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if "SUBTITLE" in ph_type and shape.has_text_frame:
                    shape.text_frame.paragraphs[0].text = new_subtitle
                    return

    def _edit_bullets(self, slide: Slide, bullets: List[str], options: Dict = None):
        """Update bullet points in the slide body."""
        options = options or {}
        target_placeholder = options.get("placeholder", "BODY")

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Check if this is the body/content placeholder
                is_body = False
                if shape.is_placeholder:
                    ph_type = str(shape.placeholder_format.type)
                    is_body = target_placeholder in ph_type or "BODY" in ph_type or "OBJECT" in ph_type
                elif shape != slide.shapes.title:
                    # Non-placeholder text box that's not the title
                    is_body = True

                if is_body:
                    tf = shape.text_frame
                    # Clear existing paragraphs
                    for para in list(tf.paragraphs)[1:]:
                        p = para._p
                        p.getparent().remove(p)

                    # Add new bullets
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            para = tf.paragraphs[0]
                        else:
                            para = tf.add_paragraph()
                        para.text = bullet
                        para.level = 0
                    return

    def _edit_body(self, slide: Slide, new_text: str):
        """Replace body text."""
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                if shape.is_placeholder:
                    ph_type = str(shape.placeholder_format.type)
                    if "BODY" in ph_type or "OBJECT" in ph_type:
                        shape.text_frame.paragraphs[0].text = new_text
                        return

    def _edit_chart(self, slide: Slide, chart_data: Dict[str, Any]):
        """Update chart data."""
        from pptx.chart.data import CategoryChartData

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart

                # Build new chart data
                new_data = CategoryChartData()
                categories = chart_data.get("categories", [])
                new_data.categories = categories

                for series in chart_data.get("series", []):
                    new_data.add_series(
                        series.get("name", "Series"),
                        series.get("values", [])
                    )

                # Replace chart data
                chart.replace_data(new_data)
                return

    def _edit_table(self, slide: Slide, table_data: Dict[str, Any]):
        """Update table content."""
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                headers = table_data.get("headers", [])
                rows = table_data.get("data", [])

                # Update headers (first row)
                if len(table.rows) > 0:
                    for i, header in enumerate(headers):
                        if i < len(table.columns):
                            cell = table.cell(0, i)
                            cell.text = str(header)

                # Update data rows
                for row_idx, row_data in enumerate(rows):
                    table_row = row_idx + 1  # Skip header
                    if table_row < len(table.rows):
                        for col_idx, value in enumerate(row_data):
                            if col_idx < len(table.columns):
                                cell = table.cell(table_row, col_idx)
                                cell.text = str(value)
                return

    def _clear_placeholder(self, slide: Slide, placeholder_type: str):
        """Clear content from a placeholder."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if placeholder_type in ph_type and shape.has_text_frame:
                    shape.text_frame.clear()

    def get_stats(self) -> Dict[str, Any]:
        """Get statistics about the slide pool."""
        type_counts = {}
        for slide in self.slides:
            ft = slide.functional_type.value
            type_counts[ft] = type_counts.get(ft, 0) + 1

        templates = set(s.template_path for s in self.slides)

        return {
            "total_slides": len(self.slides),
            "templates": len(templates),
            "by_type": type_counts,
            "with_charts": sum(1 for s in self.slides if s.has_chart),
            "with_tables": sum(1 for s in self.slides if s.has_table),
            "with_images": sum(1 for s in self.slides if s.has_image),
        }

    def save_index(self, path: str = None):
        """Save slide index to JSON."""
        path = Path(path) if path else self.index_path
        if not path:
            raise ValueError("No index path specified")

        data = {
            "version": "1.0",
            "slides": [
                {
                    "slide_id": s.slide_id,
                    "template_path": s.template_path,
                    "slide_index": s.slide_index,
                    "functional_type": s.functional_type.value,
                    "title": s.title,
                    "subtitle": s.subtitle,
                    "bullet_count": s.bullet_count,
                    "word_count": s.word_count,
                    "has_chart": s.has_chart,
                    "chart_type": s.chart_type,
                    "has_table": s.has_table,
                    "table_dimensions": s.table_dimensions,
                    "has_image": s.has_image,
                    "image_count": s.image_count,
                    "shape_count": s.shape_count,
                    "layout_name": s.layout_name,
                    "placeholder_types": s.placeholder_types,
                    "content_signature": s.content_signature,
                    "element_pattern": s.element_pattern,
                    "quality_score": s.quality_score,
                }
                for s in self.slides
            ]
        }

        path.parent.mkdir(parents=True, exist_ok=True)
        with open(path, "w") as f:
            json.dump(data, f, indent=2)

        logger.info(f"Saved slide index to {path}")

    def _load_index(self):
        """Load slide index from JSON."""
        if not self.index_path or not self.index_path.exists():
            return

        with open(self.index_path) as f:
            data = json.load(f)

        for s in data.get("slides", []):
            self.slides.append(IndexedSlide(
                slide_id=s["slide_id"],
                template_path=s["template_path"],
                slide_index=s["slide_index"],
                functional_type=FunctionalType(s["functional_type"]),
                title=s.get("title"),
                subtitle=s.get("subtitle"),
                bullet_count=s.get("bullet_count", 0),
                word_count=s.get("word_count", 0),
                has_chart=s.get("has_chart", False),
                chart_type=s.get("chart_type"),
                has_table=s.get("has_table", False),
                table_dimensions=tuple(s["table_dimensions"]) if s.get("table_dimensions") else None,
                has_image=s.get("has_image", False),
                image_count=s.get("image_count", 0),
                shape_count=s.get("shape_count", 0),
                layout_name=s.get("layout_name", ""),
                placeholder_types=s.get("placeholder_types", []),
                content_signature=s.get("content_signature", ""),
                element_pattern=s.get("element_pattern", ""),
                quality_score=s.get("quality_score", 1.0),
            ))

        logger.info(f"Loaded {len(self.slides)} slides from index")


# =============================================================================
# Convenience Functions
# =============================================================================

def build_slide_pool(templates_dir: str, index_path: str = None) -> SlidePool:
    """
    Build a slide pool from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        index_path: Optional path to save/load index

    Returns:
        Populated SlidePool
    """
    pool = SlidePool(index_path=index_path)
    pool.index_templates_directory(templates_dir)

    if index_path:
        pool.save_index(index_path)

    return pool


def create_edit_actions(content: Dict[str, Any]) -> List[EditAction]:
    """
    Create edit actions from content dictionary.

    Args:
        content: Slide content dictionary

    Returns:
        List of EditActions to apply to a cloned slide
    """
    actions = []

    # Title
    if "title" in content:
        actions.append(EditAction(
            EditActionType.REPLACE_TITLE,
            "title",
            content["title"]
        ))

    # Subtitle
    if "subtitle" in content:
        actions.append(EditAction(
            EditActionType.REPLACE_SUBTITLE,
            "subtitle",
            content["subtitle"]
        ))

    # Bullets
    if "bullets" in content and content["bullets"]:
        actions.append(EditAction(
            EditActionType.UPDATE_BULLETS,
            "body",
            content["bullets"]
        ))

    # Chart data
    if "chart_data" in content:
        actions.append(EditAction(
            EditActionType.SWAP_CHART_DATA,
            "chart",
            content["chart_data"]
        ))

    # Table data
    if "headers" in content and "data" in content:
        actions.append(EditAction(
            EditActionType.UPDATE_TABLE,
            "table",
            {"headers": content["headers"], "data": content["data"]}
        ))

    return actions
