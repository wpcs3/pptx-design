"""
Enhanced Slide Pool Re-indexing Script

Re-indexes all templates with improved functional type classification:
1. Better heuristics for slide type inference
2. Optional ML-based classification enhancement
3. Template-specific quality scores
4. Persisted index for fast startup

Usage:
    python -m pptx_generator.modules.reindex_slides
    python -m pptx_generator.modules.reindex_slides --use-ml
"""

import argparse
import json
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .slide_pool import (
    SlidePool, IndexedSlide, FunctionalType
)

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)


# Enhanced functional type patterns
LAYOUT_NAME_PATTERNS = {
    # Title slides
    "title slide": FunctionalType.TITLE,
    "title only": FunctionalType.TITLE,
    "opening": FunctionalType.TITLE,
    "cover": FunctionalType.TITLE,

    # Section dividers
    "section": FunctionalType.SECTION,
    "divider": FunctionalType.SECTION,
    "chapter": FunctionalType.SECTION,
    "break": FunctionalType.SECTION,

    # Content layouts
    "title and content": FunctionalType.CONTENT,
    "content": FunctionalType.CONTENT,
    "bullet": FunctionalType.CONTENT,
    "text": FunctionalType.CONTENT,

    # Two column / comparison
    "two content": FunctionalType.COMPARISON,
    "comparison": FunctionalType.COMPARISON,
    "two column": FunctionalType.COMPARISON,
    "side by side": FunctionalType.COMPARISON,

    # Visual layouts
    "picture": FunctionalType.CONTENT_WITH_VISUAL,
    "image": FunctionalType.CONTENT_WITH_VISUAL,
    "photo": FunctionalType.CONTENT_WITH_VISUAL,
    "content with caption": FunctionalType.CONTENT_WITH_VISUAL,

    # Blank
    "blank": FunctionalType.UNKNOWN,
}

TITLE_PATTERNS = {
    # Agenda/TOC
    ("agenda", "contents", "outline", "overview", "topics", "today"): FunctionalType.AGENDA,

    # Closing
    ("thank", "question", "q&a", "contact", "next step", "conclusion", "summary"): FunctionalType.CLOSING,

    # Team
    ("team", "leadership", "management", "about us", "our people", "who we are"): FunctionalType.TEAM,

    # Timeline
    ("timeline", "roadmap", "milestone", "schedule", "plan", "phases"): FunctionalType.TIMELINE,

    # Metrics/KPIs
    ("metric", "kpi", "performance", "dashboard", "scorecard", "results"): FunctionalType.METRICS,

    # Financial
    ("financial", "revenue", "profit", "budget", "forecast", "p&l"): FunctionalType.DATA_CHART,

    # Market
    ("market", "competitor", "industry", "landscape", "swot"): FunctionalType.DATA_CHART,
}


def enhanced_infer_type(
    slide,
    layout_name: str,
    title: Optional[str],
    has_chart: bool,
    has_table: bool,
    has_image: bool,
    image_count: int,
    bullet_count: int,
    shape_count: int,
    word_count: int
) -> FunctionalType:
    """Enhanced functional type inference with better heuristics."""

    layout_lower = layout_name.lower()

    # 1. Check layout name patterns first
    for pattern, ftype in LAYOUT_NAME_PATTERNS.items():
        if pattern in layout_lower:
            # But override if we have specific content
            if ftype == FunctionalType.CONTENT:
                if has_chart:
                    return FunctionalType.DATA_CHART
                if has_table:
                    return FunctionalType.DATA_TABLE
            return ftype

    # 2. Check title patterns
    if title:
        title_lower = title.lower()
        for patterns, ftype in TITLE_PATTERNS.items():
            if any(p in title_lower for p in patterns):
                return ftype

    # 3. Check content characteristics
    if has_chart:
        return FunctionalType.DATA_CHART

    if has_table:
        return FunctionalType.DATA_TABLE

    # 4. Check for metrics boxes (multiple small text boxes with numbers)
    if shape_count >= 3 and word_count < 50 and bullet_count == 0:
        # Likely a metrics/KPI slide
        return FunctionalType.METRICS

    # 5. Check for comparison (two columns of content)
    if "two" in layout_lower or "comparison" in layout_lower:
        return FunctionalType.COMPARISON

    # 6. Check for content with visual
    if has_image and bullet_count > 0:
        return FunctionalType.CONTENT_WITH_VISUAL

    # 7. Check for pure content
    if bullet_count >= 2:
        return FunctionalType.CONTENT

    # 8. Check for title-only slides
    if title and bullet_count == 0 and shape_count <= 3:
        if word_count < 20:
            return FunctionalType.SECTION
        return FunctionalType.TITLE

    # 9. Check for closing slides
    if shape_count <= 4 and word_count < 30:
        if title and any(w in title.lower() for w in ["thank", "question", "contact"]):
            return FunctionalType.CLOSING

    # 10. Visual/infographic slides (title + many shapes, few bullets)
    if title and shape_count >= 10 and bullet_count <= 1:
        if has_image:
            return FunctionalType.CONTENT_WITH_VISUAL
        # High shape count suggests visual elements (icons, diagrams)
        return FunctionalType.CONTENT_WITH_VISUAL

    # 11. Image gallery slides
    if has_image and image_count >= 2:
        return FunctionalType.CONTENT_WITH_VISUAL

    # 12. Section/transition slides (title-only with moderate shapes)
    if title and bullet_count == 0 and not has_chart and not has_table:
        if shape_count >= 4 and word_count < 100:
            return FunctionalType.SECTION

    return FunctionalType.UNKNOWN


def analyze_slide_enhanced(slide, template_path: str, slide_index: int) -> Dict[str, Any]:
    """Enhanced slide analysis with better type inference."""
    import hashlib

    # Generate unique ID
    slide_id = hashlib.md5(f"{template_path}:{slide_index}".encode()).hexdigest()[:12]

    # Extract content
    title = None
    subtitle = None
    bullet_count = 0
    word_count = 0
    has_chart = False
    chart_type = None
    has_table = False
    table_dims = None
    has_image = False
    image_count = 0
    placeholder_types = []

    for shape in slide.shapes:
        # Title
        if shape == slide.shapes.title and shape.has_text_frame:
            title = shape.text_frame.text.strip()

        # Placeholders
        if shape.is_placeholder:
            ph_type = str(shape.placeholder_format.type)
            placeholder_types.append(ph_type)

            if "SUBTITLE" in ph_type and shape.has_text_frame:
                subtitle = shape.text_frame.text.strip()

        # Text content
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    word_count += len(text.split())
                    if para.level > 0 or text.startswith(("•", "-", "*", "●", "○")):
                        bullet_count += 1

        # Charts
        if shape.has_chart:
            has_chart = True
            chart_type = str(shape.chart.chart_type)

        # Tables
        if shape.has_table:
            has_table = True
            table_dims = (len(shape.table.rows), len(shape.table.columns))

        # Images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_image = True
            image_count += 1

    layout_name = slide.slide_layout.name if slide.slide_layout else ""

    # Use enhanced inference
    functional_type = enhanced_infer_type(
        slide=slide,
        layout_name=layout_name,
        title=title,
        has_chart=has_chart,
        has_table=has_table,
        has_image=has_image,
        image_count=image_count,
        bullet_count=bullet_count,
        shape_count=len(slide.shapes),
        word_count=word_count
    )

    # Build element pattern
    pattern_parts = []
    if title:
        pattern_parts.append("T")
    if bullet_count > 0:
        pattern_parts.append(f"B{bullet_count}")
    if has_chart:
        pattern_parts.append("C")
    if has_table:
        pattern_parts.append("Tb")
    if has_image:
        pattern_parts.append(f"I{image_count}")
    element_pattern = "-".join(pattern_parts) or "Empty"

    # Content signature
    sig_data = f"{functional_type.value}:{element_pattern}:{bullet_count}"
    content_signature = hashlib.md5(sig_data.encode()).hexdigest()[:8]

    return {
        "slide_id": slide_id,
        "template_path": template_path,
        "slide_index": slide_index,
        "functional_type": functional_type.value,
        "title": title,
        "subtitle": subtitle,
        "bullet_count": bullet_count,
        "word_count": word_count,
        "has_chart": has_chart,
        "chart_type": chart_type,
        "has_table": has_table,
        "table_dimensions": list(table_dims) if table_dims else None,
        "has_image": has_image,
        "image_count": image_count,
        "shape_count": len(slide.shapes),
        "layout_name": layout_name,
        "placeholder_types": placeholder_types,
        "content_signature": content_signature,
        "element_pattern": element_pattern,
        "quality_score": 1.0,
    }


def reindex_templates(
    templates_dir: str,
    output_path: str,
    use_ml: bool = False,
    template_quality_scores: Dict[str, float] = None
) -> Dict[str, Any]:
    """
    Re-index all templates with enhanced classification.

    Args:
        templates_dir: Directory containing template PPTX files
        output_path: Path to save the index JSON
        use_ml: Whether to use ML classifiers for enhancement
        template_quality_scores: Optional quality scores per template name

    Returns:
        Statistics about the indexing
    """
    templates_dir = Path(templates_dir)
    output_path = Path(output_path)

    # Default quality scores (higher = better designed template)
    default_scores = template_quality_scores or {
        "consulting": 1.0,
        "business_case": 0.95,
        "market_analysis": 0.9,
        "due_diligence": 0.85,
    }

    all_slides = []
    stats = {
        "templates": 0,
        "total_slides": 0,
        "by_type": {},
        "by_template": {},
    }

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob("*.pptx"))
    pptx_files = [f for f in pptx_files if not f.name.startswith("~")]

    logger.info(f"Found {len(pptx_files)} template files")

    for pptx_path in pptx_files:
        logger.info(f"Indexing: {pptx_path.name}")

        # Determine quality score for this template
        quality = 1.0
        path_lower = str(pptx_path).lower()
        for key, score in default_scores.items():
            if key in path_lower:
                quality = score
                break

        try:
            prs = Presentation(str(pptx_path))
            template_slides = 0

            for idx, slide in enumerate(prs.slides):
                slide_data = analyze_slide_enhanced(slide, str(pptx_path), idx)
                slide_data["quality_score"] = quality
                all_slides.append(slide_data)

                # Update stats
                ftype = slide_data["functional_type"]
                stats["by_type"][ftype] = stats["by_type"].get(ftype, 0) + 1
                template_slides += 1

            stats["by_template"][pptx_path.name] = template_slides
            stats["templates"] += 1

        except Exception as e:
            logger.error(f"Failed to index {pptx_path}: {e}")

    stats["total_slides"] = len(all_slides)

    # Optionally enhance with ML classifiers
    if use_ml:
        logger.info("Enhancing classification with ML (not implemented yet)...")
        # TODO: Use LayoutClassifier and SemanticAnalyzer to refine types

    # Save index
    output_path.parent.mkdir(parents=True, exist_ok=True)

    index_data = {
        "version": "2.0",
        "enhanced": True,
        "stats": stats,
        "slides": all_slides,
    }

    with open(output_path, "w") as f:
        json.dump(index_data, f, indent=2)

    logger.info(f"Saved index to: {output_path}")
    logger.info(f"Total slides: {stats['total_slides']}")
    logger.info(f"By type: {stats['by_type']}")

    # Show improvement over baseline
    unknown_count = stats["by_type"].get("unknown", 0)
    unknown_pct = (unknown_count / stats["total_slides"] * 100) if stats["total_slides"] > 0 else 0
    logger.info(f"Unknown slides: {unknown_count} ({unknown_pct:.1f}%)")

    return stats


def main():
    parser = argparse.ArgumentParser(description="Re-index slide pool with enhanced classification")
    parser.add_argument(
        "--templates-dir",
        default="pptx_templates",
        help="Directory containing template PPTX files"
    )
    parser.add_argument(
        "--output",
        default="cache/slide_pool_index.json",
        help="Output path for index JSON"
    )
    parser.add_argument(
        "--use-ml",
        action="store_true",
        help="Use ML classifiers for enhanced classification"
    )

    args = parser.parse_args()

    stats = reindex_templates(
        templates_dir=args.templates_dir,
        output_path=args.output,
        use_ml=args.use_ml
    )

    print("\n" + "=" * 50)
    print("Re-indexing Complete!")
    print("=" * 50)
    print(f"Templates: {stats['templates']}")
    print(f"Total slides: {stats['total_slides']}")
    print(f"\nBy functional type:")
    for ftype, count in sorted(stats['by_type'].items(), key=lambda x: -x[1]):
        pct = count / stats['total_slides'] * 100
        print(f"  {ftype:20s}: {count:4d} ({pct:5.1f}%)")


if __name__ == "__main__":
    main()
