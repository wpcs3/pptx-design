"""
Slide Renderer Module

Generates individual slides from specifications using the style guide.
"""

import json
import logging
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class SlideRenderer:
    """Generates slides from specifications using style guide formatting."""

    # Standard slide dimensions (widescreen 16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(
        self,
        style_guide: dict,
        slide_catalog: dict,
        template_path: Optional[str] = None
    ):
        """
        Initialize the slide renderer.

        Args:
            style_guide: Style guide dictionary
            slide_catalog: Slide catalog dictionary
            template_path: Optional path to base template for styling
        """
        self.style_guide = style_guide
        self.slide_catalog = slide_catalog
        self.template_path = template_path

        # Extract style shortcuts
        self.colors = style_guide.get("colors", {})
        self.fonts = style_guide.get("fonts", {})
        self.spacing = style_guide.get("spacing", {})

        # Build slide type lookup
        self.slide_types = {
            st["id"]: st
            for st in slide_catalog.get("slide_types", [])
        }

    def create_presentation(self) -> Presentation:
        """Create a new presentation, optionally based on template."""
        if self.template_path and Path(self.template_path).exists():
            return Presentation(self.template_path)
        return Presentation()

    def create_slide(
        self,
        presentation: Presentation,
        slide_type: str,
        content: dict
    ) -> Slide:
        """
        Create a slide of the given type with provided content.

        Args:
            presentation: Target Presentation object
            slide_type: ID from slide_catalog.json
            content: Content dictionary with keys like:
                - title: str
                - subtitle: str
                - body: str or list of bullets
                - bullets: list of str
                - chart_data: dict
                - images: list of image paths
                - left_column: dict
                - right_column: dict

        Returns:
            The newly created Slide
        """
        # Get slide type definition
        type_def = self.slide_types.get(slide_type)
        if not type_def:
            logger.warning(f"Unknown slide type: {slide_type}, using default")
            type_def = {"master_layout": "Default", "elements": []}

        # Find matching layout
        layout = self._find_layout(presentation, type_def.get("master_layout", "Default"))
        slide = presentation.slides.add_slide(layout)

        # Render based on slide type
        renderer_method = self._get_renderer_method(slide_type)
        renderer_method(slide, content, type_def)

        return slide

    def _find_layout(self, presentation: Presentation, layout_name: str):
        """Find a slide layout by name."""
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout

        # Fallback to index-based selection
        layout_map = {
            "title_slide": 0,
            "Title Slide": 0,
            "Frontpage": 0,
            "section_divider": 2,
            "Section breaker": 2,
            "default": 1,
            "Default": 1,
            "blank": 6,
            "Blank": 6,
        }

        idx = layout_map.get(layout_name, 1)
        if idx < len(presentation.slide_layouts):
            return presentation.slide_layouts[idx]
        return presentation.slide_layouts[0]

    def _get_renderer_method(self, slide_type: str):
        """Get the appropriate renderer method for a slide type."""
        renderers = {
            "title_slide": self._render_title_slide,
            "section_divider": self._render_section_divider,
            "title_content": self._render_title_content,
            "two_column": self._render_two_column,
            "data_chart": self._render_data_chart,
            "table_slide": self._render_table_slide,
            "key_metrics": self._render_key_metrics,
            "image_slide": self._render_image_slide,
        }

        # Check for exact match or pattern match
        for pattern, renderer in renderers.items():
            if pattern in slide_type.lower():
                return renderer

        return self._render_generic

    def _render_title_slide(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a title slide."""
        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        # Add title
        if title:
            self._add_title(slide, title, is_main_title=True)

        # Add subtitle
        if subtitle:
            self._add_subtitle(slide, subtitle)

    def _render_section_divider(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a section divider slide."""
        title = content.get("title", content.get("section_name", ""))

        if title:
            # Large centered title for section dividers
            left = Inches(0.5)
            top = Inches(3)
            width = self.SLIDE_WIDTH - Inches(1)
            height = Inches(1.5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER

            self._apply_font_style(p.runs[0].font, "title", size_override=Pt(44))

    def _render_title_content(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a title and content slide."""
        title = content.get("title", "")
        body = content.get("body", "")
        bullets = content.get("bullets", [])

        if title:
            self._add_title(slide, title)

        # Add body content
        if body or bullets:
            self._add_body_content(slide, body, bullets)

    def _render_two_column(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a two-column comparison slide."""
        title = content.get("title", "")
        left_col = content.get("left_column", content.get("left", {}))
        right_col = content.get("right_column", content.get("right", {}))

        if title:
            self._add_title(slide, title)

        # Left column
        self._add_column(slide, left_col, is_left=True)

        # Right column
        self._add_column(slide, right_col, is_left=False)

    def _render_data_chart(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a data visualization slide."""
        title = content.get("title", "")
        chart_data = content.get("chart_data", {})
        narrative = content.get("narrative", content.get("body", ""))

        if title:
            self._add_title(slide, title)

        # Add chart placeholder or data
        if chart_data:
            self._add_chart(slide, chart_data)
        else:
            # Add placeholder text
            self._add_chart_placeholder(slide)

        # Add narrative text
        if narrative:
            self._add_narrative(slide, narrative)

    def _render_table_slide(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a table slide."""
        title = content.get("title", "")
        table_data = content.get("table_data", content.get("data", []))
        headers = content.get("headers", [])

        if title:
            self._add_title(slide, title)

        if table_data:
            self._add_table(slide, headers, table_data)

    def _render_key_metrics(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a key metrics dashboard slide."""
        title = content.get("title", "Key Metrics")
        metrics = content.get("metrics", [])

        if title:
            self._add_title(slide, title)

        # Add metric boxes
        self._add_metric_boxes(slide, metrics)

    def _render_image_slide(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render an image slide."""
        title = content.get("title", "")
        image_path = content.get("image", content.get("image_path", ""))
        caption = content.get("caption", "")

        if title:
            self._add_title(slide, title)

        if image_path and Path(image_path).exists():
            self._add_image(slide, image_path)

        if caption:
            self._add_caption(slide, caption)

    def _render_generic(self, slide: Slide, content: dict, type_def: dict) -> None:
        """Render a generic slide with available content."""
        title = content.get("title", "")
        body = content.get("body", "")
        bullets = content.get("bullets", [])

        if title:
            self._add_title(slide, title)

        if body or bullets:
            self._add_body_content(slide, body, bullets)

    # Helper methods for adding elements

    def _add_title(
        self,
        slide: Slide,
        text: str,
        is_main_title: bool = False
    ) -> None:
        """Add a title to the slide."""
        # Try to use placeholder first
        if slide.shapes.title:
            slide.shapes.title.text = text
            if slide.shapes.title.text_frame.paragraphs:
                run = slide.shapes.title.text_frame.paragraphs[0].runs[0] if slide.shapes.title.text_frame.paragraphs[0].runs else None
                if run:
                    self._apply_font_style(run.font, "title")
            return

        # Otherwise add text box
        left = Inches(0.5)
        top = Inches(0.5) if not is_main_title else Inches(2.5)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER if is_main_title else PP_ALIGN.LEFT

        self._apply_font_style(p.runs[0].font, "title")

    def _add_subtitle(self, slide: Slide, text: str) -> None:
        """Add a subtitle to the slide."""
        left = Inches(0.5)
        top = Inches(3.8)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(0.75)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER

        self._apply_font_style(p.runs[0].font, "subtitle")

    def _add_body_content(
        self,
        slide: Slide,
        body: str,
        bullets: list
    ) -> None:
        """Add body content to the slide."""
        left = Inches(0.5)
        top = Inches(1.5)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        if body:
            p = tf.paragraphs[0]
            p.text = body
            self._apply_font_style(p.runs[0].font, "body")

        for i, bullet in enumerate(bullets):
            if i == 0 and not body:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = bullet
            p.level = 0
            if p.runs:
                self._apply_font_style(p.runs[0].font, "body")

    def _add_column(
        self,
        slide: Slide,
        column_content: dict,
        is_left: bool
    ) -> None:
        """Add a column to a two-column slide."""
        width = (self.SLIDE_WIDTH - Inches(1.5)) / 2
        left = Inches(0.5) if is_left else Inches(0.5) + width + Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Add header
        header = column_content.get("header", "")
        if header:
            p = tf.paragraphs[0]
            p.text = header
            if p.runs:
                self._apply_font_style(p.runs[0].font, "subtitle")

        # Add bullets
        bullets = column_content.get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            if p.runs:
                self._apply_font_style(p.runs[0].font, "body")

    def _add_chart_placeholder(self, slide: Slide) -> None:
        """Add a placeholder for chart data."""
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(4.5)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE6, 0xE6, 0xE6)
        shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Add text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[Chart Placeholder]"
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            self._apply_font_style(p.runs[0].font, "body")

    def _add_chart(self, slide: Slide, chart_data: dict) -> None:
        """Add a chart to the slide."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_type = chart_data.get("type", "column")
        categories = chart_data.get("categories", [])
        series_data = chart_data.get("series", [])

        if not categories or not series_data:
            self._add_chart_placeholder(slide)
            return

        # Map chart types
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }

        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Create chart data
        data = CategoryChartData()
        data.categories = categories

        for series in series_data:
            data.add_series(series.get("name", "Series"), series.get("values", []))

        # Add chart
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11)
        height = Inches(5)

        slide.shapes.add_chart(xl_chart_type, left, top, width, height, data)

    def _add_table(
        self,
        slide: Slide,
        headers: list,
        data: list
    ) -> None:
        """Add a table to the slide."""
        if not data:
            return

        rows = len(data) + (1 if headers else 0)
        cols = len(headers) if headers else (len(data[0]) if data else 0)

        if rows == 0 or cols == 0:
            return

        left = Inches(0.5)
        top = Inches(1.5)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set headers
        if headers:
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                if cell.text_frame.paragraphs[0].runs:
                    self._apply_font_style(
                        cell.text_frame.paragraphs[0].runs[0].font,
                        "body"
                    )
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True

        # Set data
        start_row = 1 if headers else 0
        for i, row in enumerate(data):
            for j, value in enumerate(row):
                cell = table.cell(start_row + i, j)
                cell.text = str(value)
                if cell.text_frame.paragraphs[0].runs:
                    self._apply_font_style(
                        cell.text_frame.paragraphs[0].runs[0].font,
                        "body"
                    )

    def _add_metric_boxes(self, slide: Slide, metrics: list) -> None:
        """Add metric KPI boxes to the slide."""
        if not metrics:
            return

        num_metrics = min(len(metrics), 5)
        box_width = (self.SLIDE_WIDTH - Inches(1 + 0.25 * (num_metrics - 1))) / num_metrics
        box_height = Inches(2)

        top = Inches(2.5)

        for i, metric in enumerate(metrics[:5]):
            left = Inches(0.5) + i * (box_width + Inches(0.25))

            # Add box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height
            )

            # Style box
            shape.fill.solid()
            primary_color = self._parse_color(self.colors.get("primary", "#3C96B4"))
            shape.fill.fore_color.rgb = primary_color
            shape.line.fill.background()

            # Add metric text
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Value
            p = tf.paragraphs[0]
            p.text = str(metric.get("value", ""))
            if p.runs:
                p.runs[0].font.size = Pt(36)
                p.runs[0].font.bold = True
                p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Label
            p2 = tf.add_paragraph()
            p2.text = metric.get("label", "")
            p2.alignment = PP_ALIGN.CENTER
            if p2.runs:
                p2.runs[0].font.size = Pt(14)
                p2.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def _add_image(self, slide: Slide, image_path: str) -> None:
        """Add an image to the slide."""
        try:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(11)
            height = Inches(5)

            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            logger.error(f"Error adding image: {e}")

    def _add_caption(self, slide: Slide, text: str) -> None:
        """Add a caption to the slide."""
        left = Inches(0.5)
        top = Inches(6.5)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(0.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER

        self._apply_font_style(p.runs[0].font, "caption")

    def _add_narrative(self, slide: Slide, text: str) -> None:
        """Add narrative text below a chart."""
        left = Inches(0.5)
        top = Inches(6)
        width = self.SLIDE_WIDTH - Inches(1)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        self._apply_font_style(p.runs[0].font, "body")

    def _apply_font_style(
        self,
        font,
        style_type: str,
        size_override: Optional[Pt] = None
    ) -> None:
        """Apply font styling from the style guide."""
        style = self.fonts.get(style_type, {})

        font.name = style.get("name", "Arial")
        font.size = size_override or Pt(style.get("size_pt", 18))
        font.bold = style.get("bold", False)

        # Apply text color
        text_colors = self.colors.get("text", {})
        color_hex = text_colors.get(style_type, text_colors.get("body", "#000000"))
        font.color.rgb = self._parse_color(color_hex)

    def _parse_color(self, color_str: str) -> RGBColor:
        """Parse a color string to RGBColor."""
        if not color_str:
            return RGBColor(0, 0, 0)

        color_str = color_str.lstrip("#")
        try:
            r = int(color_str[0:2], 16)
            g = int(color_str[2:4], 16)
            b = int(color_str[4:6], 16)
            return RGBColor(r, g, b)
        except (ValueError, IndexError):
            return RGBColor(0, 0, 0)

    def apply_style(self, shape, style_type: str) -> None:
        """Apply standard formatting to a shape."""
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    self._apply_font_style(run.font, style_type)


def main():
    """Test the slide renderer."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Slide Renderer")
    parser.add_argument(
        "--style-guide",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/config/style_guide.json",
        help="Path to style guide"
    )
    parser.add_argument(
        "--catalog",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/config/slide_catalog.json",
        help="Path to slide catalog"
    )
    parser.add_argument(
        "--output",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/output/test_render.pptx",
        help="Output file path"
    )

    args = parser.parse_args()

    # Load configs
    with open(args.style_guide, "r") as f:
        style_guide = json.load(f)
    with open(args.catalog, "r") as f:
        slide_catalog = json.load(f)

    # Create renderer
    renderer = SlideRenderer(style_guide, slide_catalog)
    prs = renderer.create_presentation()

    # Create test slides
    renderer.create_slide(prs, "title_slide", {
        "title": "Investment Pitch",
        "subtitle": "Q4 2025 Investor Presentation"
    })

    renderer.create_slide(prs, "section_divider", {
        "title": "Executive Summary"
    })

    renderer.create_slide(prs, "title_content", {
        "title": "Investment Thesis",
        "bullets": [
            "Strong market fundamentals",
            "Experienced management team",
            "Proven track record",
            "Attractive risk-adjusted returns"
        ]
    })

    renderer.create_slide(prs, "two_column", {
        "title": "Strategy Comparison",
        "left_column": {
            "header": "Strategy A",
            "bullets": ["Lower risk", "Stable returns", "Long-term focus"]
        },
        "right_column": {
            "header": "Strategy B",
            "bullets": ["Higher potential", "Growth focus", "Active management"]
        }
    })

    renderer.create_slide(prs, "key_metrics", {
        "title": "Key Performance Indicators",
        "metrics": [
            {"label": "AUM", "value": "$2.5B"},
            {"label": "IRR", "value": "18.5%"},
            {"label": "Deals", "value": "42"},
            {"label": "Team", "value": "25"}
        ]
    })

    renderer.create_slide(prs, "table_slide", {
        "title": "Historical Returns",
        "headers": ["Year", "Return", "Benchmark", "Alpha"],
        "data": [
            ["2022", "12.5%", "8.2%", "+4.3%"],
            ["2023", "15.8%", "10.1%", "+5.7%"],
            ["2024", "14.2%", "9.5%", "+4.7%"]
        ]
    })

    # Save
    Path(args.output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"Saved test presentation to: {args.output}")


if __name__ == "__main__":
    main()
