"""
Component Library Integration Module

Provides access to the extracted PowerPoint component library from the pptx_generator project.
Enables searching and reusing:
- Basic components: charts, tables, images, diagrams, shapes
- Styles: color palettes, typography, effects, gradients, shadows
- Chart styles: formatting profiles for charts
- Layout blueprints: grid systems and zone definitions
- Diagram templates: reusable shape combinations
- Text patterns: bullet structures, title styles
- Slide sequences: deck templates and flow patterns
"""

import json
import logging
from pathlib import Path
from typing import Optional, Union, List, Dict, Any
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Default path to component library (relative to this project)
DEFAULT_LIBRARY_PATH = Path(__file__).parent.parent.parent / "pptx_component_library"


@dataclass
class ComponentMatch:
    """Represents a matched component from the library."""
    id: str
    type: str
    category: str
    filename: str
    width_inches: Optional[float] = None
    height_inches: Optional[float] = None
    metadata: dict = None

    @property
    def file_path(self) -> Path:
        """Get the full path to the component file."""
        return None  # Set by ComponentLibrary


class ComponentLibrary:
    """
    Interface to the PowerPoint component library.

    Provides methods to search, retrieve, and use extracted components
    (charts, tables, images, shapes, diagrams) in new presentations.
    """

    def __init__(self, library_path: Optional[Path] = None):
        """
        Initialize the component library.

        Args:
            library_path: Path to the component library directory.
                         Defaults to ../component_library relative to pptx_generator.
        """
        self.library_path = Path(library_path) if library_path else DEFAULT_LIBRARY_PATH
        self.index: dict = {}
        self._load_index()

    def _load_index(self) -> None:
        """Load all library indexes."""
        # Main component index
        index_path = self.library_path / "library_index.json"

        if not index_path.exists():
            logger.warning(f"Component library index not found at: {index_path}")
            logger.info("Run 'python -m pptx_extractor.library_cli extract-all' to create the library first.")
            self.index = {"components": {}, "templates": {}}
        else:
            with open(index_path, 'r', encoding='utf-8') as f:
                self.index = json.load(f)

        # Load additional indexes
        self._style_index = self._load_json_index("styles/style_index.json")
        self._chart_style_index = self._load_json_index("styles/chart_style_index.json")
        self._layout_index = self._load_json_index("layouts/layout_index.json")
        self._diagram_index = self._load_json_index("diagrams/diagram_template_index.json")
        self._text_index = self._load_json_index("text_templates/text_template_index.json")
        self._sequence_index = self._load_json_index("sequences/sequence_index.json")

        total = sum(len(v) for v in self.index.get('components', {}).values())
        logger.info(f"Loaded component library with {total} components")

    def _load_json_index(self, relative_path: str) -> dict:
        """Load a JSON index file if it exists."""
        path = self.library_path / relative_path
        if path.exists():
            with open(path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {}

    @property
    def is_available(self) -> bool:
        """Check if the library is available and loaded."""
        return bool(self.index.get('components'))

    def get_stats(self) -> dict:
        """Get library statistics."""
        stats = {
            'templates': len(self.index.get('templates', {})),
            'components': {}
        }

        for comp_type, items in self.index.get('components', {}).items():
            stats['components'][comp_type] = len(items)

        return stats

    # ==================== Search Methods ====================

    def search(
        self,
        component_type: Optional[str] = None,
        category: Optional[str] = None,
        tags: Optional[list] = None,
        template: Optional[str] = None,
        min_width: Optional[float] = None,
        max_width: Optional[float] = None,
        limit: int = 20
    ) -> list[dict]:
        """
        Search for components in the library.

        Args:
            component_type: Filter by type ('images', 'charts', 'tables', etc.)
            category: Filter by category (e.g., 'column_charts', 'comparison_matrix')
            tags: Filter by tags (any match)
            template: Filter by source template name
            min_width: Minimum width in inches
            max_width: Maximum width in inches
            limit: Maximum results to return

        Returns:
            List of matching component metadata dicts
        """
        results = []

        types_to_search = [component_type] if component_type else self.index.get('components', {}).keys()

        for ctype in types_to_search:
            items = self.index.get('components', {}).get(ctype, [])

            for item in items:
                # Apply filters
                if category and item.get('category') != category:
                    continue

                if tags:
                    item_tags = [t.lower() for t in item.get('tags', [])]
                    if not any(t.lower() in item_tags for t in tags):
                        continue

                if template:
                    refs = item.get('references', [])
                    if not any(template.lower() in r.get('template', '').lower() for r in refs):
                        continue

                if min_width and item.get('width_inches', 0) < min_width:
                    continue

                if max_width and item.get('width_inches', float('inf')) > max_width:
                    continue

                # Add type info
                item_copy = item.copy()
                item_copy['_type'] = ctype
                results.append(item_copy)

        return results[:limit]

    def search_charts(
        self,
        chart_type: Optional[str] = None,
        category: Optional[str] = None,
        min_series: Optional[int] = None,
        limit: int = 20
    ) -> list[dict]:
        """
        Search for charts specifically.

        Args:
            chart_type: Filter by chart type (e.g., 'COLUMN_CLUSTERED', 'LINE')
            category: Filter by category (e.g., 'column_charts', 'line_charts')
            min_series: Minimum number of data series
            limit: Maximum results

        Returns:
            List of matching chart metadata
        """
        results = []

        for item in self.index.get('components', {}).get('charts', []):
            if chart_type and chart_type.lower() not in item.get('chart_type', '').lower():
                continue

            if category and item.get('category') != category:
                continue

            if min_series and item.get('series_count', 0) < min_series:
                continue

            results.append(item)

        return results[:limit]

    def search_tables(
        self,
        category: Optional[str] = None,
        min_rows: Optional[int] = None,
        max_rows: Optional[int] = None,
        min_cols: Optional[int] = None,
        max_cols: Optional[int] = None,
        limit: int = 20
    ) -> list[dict]:
        """
        Search for tables specifically.

        Args:
            category: Filter by category ('data_table', 'comparison_matrix', etc.)
            min_rows, max_rows: Row count filters
            min_cols, max_cols: Column count filters
            limit: Maximum results

        Returns:
            List of matching table metadata
        """
        results = []

        for item in self.index.get('components', {}).get('tables', []):
            if category and item.get('category') != category:
                continue

            rows = item.get('rows', 0)
            cols = item.get('cols', 0)

            if min_rows and rows < min_rows:
                continue
            if max_rows and rows > max_rows:
                continue
            if min_cols and cols < min_cols:
                continue
            if max_cols and cols > max_cols:
                continue

            results.append(item)

        return results[:limit]

    def search_images(
        self,
        format: Optional[str] = None,
        min_size_kb: Optional[float] = None,
        max_size_kb: Optional[float] = None,
        limit: int = 20
    ) -> list[dict]:
        """
        Search for images specifically.

        Args:
            format: Filter by format ('png', 'jpg', 'tiff')
            min_size_kb, max_size_kb: File size filters in KB
            limit: Maximum results

        Returns:
            List of matching image metadata
        """
        results = []

        for item in self.index.get('components', {}).get('images', []):
            if format and item.get('format', '').lower() != format.lower():
                continue

            size_kb = item.get('size_bytes', 0) / 1024

            if min_size_kb and size_kb < min_size_kb:
                continue
            if max_size_kb and size_kb > max_size_kb:
                continue

            results.append(item)

        return results[:limit]

    # ==================== Retrieval Methods ====================

    def get_component(self, component_id: str) -> Optional[dict]:
        """
        Get a specific component by ID.

        Args:
            component_id: The unique component ID

        Returns:
            Component metadata dict or None if not found
        """
        for comp_type, items in self.index.get('components', {}).items():
            for item in items:
                if item['id'] == component_id:
                    item_copy = item.copy()
                    item_copy['_type'] = comp_type
                    return item_copy
        return None

    def get_component_file(self, component_id: str) -> Optional[Path]:
        """
        Get the file path for a component.

        Args:
            component_id: The unique component ID

        Returns:
            Path to the component file or None
        """
        component = self.get_component(component_id)
        if not component:
            return None

        comp_type = component['_type']
        filename = component.get('filename')

        if filename:
            return self.library_path / comp_type / filename
        return None

    def get_chart_data(self, component_id: str) -> Optional[dict]:
        """
        Get chart data for a chart component.

        Args:
            component_id: The chart component ID

        Returns:
            Dict with 'categories' and 'series' data, or None
        """
        file_path = self.get_component_file(component_id)
        if not file_path or not file_path.exists():
            return None

        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)

    def get_table_data(self, component_id: str) -> Optional[dict]:
        """
        Get table data for a table component.

        Args:
            component_id: The table component ID

        Returns:
            Dict with table structure and data, or None
        """
        file_path = self.get_component_file(component_id)
        if not file_path or not file_path.exists():
            return None

        with open(file_path, 'r', encoding='utf-8') as f:
            return json.load(f)

    def get_image_bytes(self, component_id: str) -> Optional[bytes]:
        """
        Get image data for an image component.

        Args:
            component_id: The image component ID

        Returns:
            Image bytes or None
        """
        file_path = self.get_component_file(component_id)
        if not file_path or not file_path.exists():
            return None

        with open(file_path, 'rb') as f:
            return f.read()

    # ==================== Creation Methods ====================

    def add_chart_from_library(
        self,
        slide,
        component_id: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        custom_data: Optional[dict] = None
    ):
        """
        Add a chart to a slide using library component as template.

        Args:
            slide: Target slide object
            component_id: Chart component ID from library
            left, top: Position in inches
            width, height: Optional size in inches (uses original if not specified)
            custom_data: Optional dict with 'categories' and 'series' to override data

        Returns:
            The created chart shape or None
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        component = self.get_component(component_id)
        if not component or component['_type'] != 'charts':
            logger.error(f"Chart component not found: {component_id}")
            return None

        # Get chart data
        chart_data = custom_data or self.get_chart_data(component_id)
        if not chart_data:
            logger.error(f"Could not load chart data for: {component_id}")
            return None

        # Determine chart type
        chart_type_str = component.get('chart_type', 'COLUMN_CLUSTERED')
        chart_type = self._get_chart_type_enum(chart_type_str)

        # Create chart data object
        data = CategoryChartData()
        data.categories = chart_data.get('categories', [])

        for series in chart_data.get('series', []):
            data.add_series(series.get('name', 'Series'), series.get('values', []))

        # Dimensions
        w = Inches(width) if width else Inches(component.get('width_inches', 6))
        h = Inches(height) if height else Inches(component.get('height_inches', 4))

        # Add chart
        try:
            chart_shape = slide.shapes.add_chart(
                chart_type,
                Inches(left), Inches(top),
                w, h,
                data
            )
            return chart_shape
        except Exception as e:
            logger.error(f"Failed to add chart: {e}")
            return None

    def add_table_from_library(
        self,
        slide,
        component_id: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        custom_data: Optional[list] = None
    ):
        """
        Add a table to a slide using library component as template.

        Args:
            slide: Target slide object
            component_id: Table component ID from library
            left, top: Position in inches
            width, height: Optional size in inches
            custom_data: Optional 2D list to override table data

        Returns:
            The created table shape or None
        """
        component = self.get_component(component_id)
        if not component or component['_type'] != 'tables':
            logger.error(f"Table component not found: {component_id}")
            return None

        # Get table structure
        table_info = self.get_table_data(component_id)
        if not table_info:
            logger.error(f"Could not load table data for: {component_id}")
            return None

        rows = table_info.get('rows', 2)
        cols = table_info.get('cols', 2)
        data = custom_data or table_info.get('data', [])

        # Dimensions
        w = Inches(width) if width else Inches(component.get('width_inches', 8))
        h = Inches(height) if height else Inches(component.get('height_inches', 3))

        try:
            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(left), Inches(top),
                w, h
            )
            table = table_shape.table

            # Populate data
            for i, row_data in enumerate(data):
                if i >= rows:
                    break
                for j, cell_value in enumerate(row_data):
                    if j >= cols:
                        break
                    table.cell(i, j).text = str(cell_value)

            return table_shape
        except Exception as e:
            logger.error(f"Failed to add table: {e}")
            return None

    def add_image_from_library(
        self,
        slide,
        component_id: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None
    ):
        """
        Add an image to a slide from the library.

        Args:
            slide: Target slide object
            component_id: Image component ID from library
            left, top: Position in inches
            width, height: Optional size in inches

        Returns:
            The created picture shape or None
        """
        component = self.get_component(component_id)
        if not component or component['_type'] != 'images':
            logger.error(f"Image component not found: {component_id}")
            return None

        file_path = self.get_component_file(component_id)
        if not file_path or not file_path.exists():
            logger.error(f"Image file not found: {component_id}")
            return None

        # Dimensions
        w = Inches(width) if width else Inches(component.get('width_inches', 2))
        h = Inches(height) if height else Inches(component.get('height_inches', 2))

        try:
            picture = slide.shapes.add_picture(
                str(file_path),
                Inches(left), Inches(top),
                w, h
            )
            return picture
        except Exception as e:
            logger.error(f"Failed to add image: {e}")
            return None

    # ==================== Style Methods ====================

    def get_color_palettes(self, template: Optional[str] = None) -> List[dict]:
        """Get color palettes from the library."""
        palettes = self._style_index.get('color_palettes', [])
        if template:
            palettes = [p for p in palettes if p.get('template') == template]
        return palettes

    def get_color_palette(self, palette_id: str) -> Optional[dict]:
        """Get a specific color palette by ID."""
        for palette in self._style_index.get('color_palettes', []):
            if palette.get('id') == palette_id:
                return palette
        return None

    def get_typography_presets(self, preset_type: Optional[str] = None,
                                template: Optional[str] = None) -> List[dict]:
        """Get typography presets."""
        presets = self._style_index.get('typography_presets', [])
        if preset_type:
            presets = [p for p in presets if p.get('preset_type') == preset_type]
        if template:
            presets = [p for p in presets if p.get('template') == template]
        return presets

    def get_gradients(self, template: Optional[str] = None) -> List[dict]:
        """Get gradient presets."""
        gradients = self._style_index.get('gradient_presets', [])
        if template:
            gradients = [g for g in gradients if g.get('template') == template]
        return gradients

    def get_shadows(self, shadow_type: Optional[str] = None) -> List[dict]:
        """Get shadow presets."""
        shadows = self._style_index.get('shadow_presets', [])
        if shadow_type:
            shadows = [s for s in shadows if s.get('type') == shadow_type]
        return shadows

    def apply_color_palette(self, shape, palette_id: str, color_role: str = 'primary') -> bool:
        """
        Apply a color from a palette to a shape.

        Args:
            shape: The shape to apply color to
            palette_id: ID of the color palette
            color_role: Which color to use ('primary', 'accent1', 'accent2', etc.)

        Returns:
            True if successful
        """
        palette = self.get_color_palette(palette_id)
        if not palette:
            return False

        color_hex = palette.get(color_role, palette.get('primary', '#000000'))
        if color_hex:
            try:
                # Remove # if present
                color_hex = color_hex.lstrip('#')
                rgb = RGBColor(
                    int(color_hex[0:2], 16),
                    int(color_hex[2:4], 16),
                    int(color_hex[4:6], 16)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb
                return True
            except Exception as e:
                logger.error(f"Failed to apply color: {e}")
        return False

    # ==================== Chart Style Methods ====================

    def get_chart_styles(self, chart_type: Optional[str] = None,
                         template: Optional[str] = None) -> List[dict]:
        """Get chart style presets."""
        styles = self._chart_style_index.get('chart_styles', [])
        if chart_type:
            styles = [s for s in styles if chart_type.lower() in s.get('chart_type', '').lower()]
        if template:
            styles = [s for s in styles if s.get('template') == template]
        return styles

    def get_chart_style(self, style_id: str) -> Optional[dict]:
        """Get a specific chart style by ID."""
        for style in self._chart_style_index.get('chart_styles', []):
            if style.get('id') == style_id:
                return style
        return None

    def apply_chart_style(self, chart, style_id: str) -> bool:
        """
        Apply a chart style to an existing chart.

        Args:
            chart: The chart object to style
            style_id: ID of the chart style

        Returns:
            True if successful
        """
        style = self.get_chart_style(style_id)
        if not style:
            return False

        try:
            # Apply legend settings
            legend_config = style.get('legend', {})
            if legend_config.get('visible'):
                chart.has_legend = True
                # Position would require additional handling
            else:
                chart.has_legend = False

            # Apply data label settings
            labels_config = style.get('data_labels', {})
            if labels_config.get('visible'):
                for plot in chart.plots:
                    plot.has_data_labels = True

            return True
        except Exception as e:
            logger.error(f"Failed to apply chart style: {e}")
            return False

    # ==================== Layout Methods ====================

    def get_layouts(self, category: Optional[str] = None,
                    template: Optional[str] = None) -> List[dict]:
        """Get layout blueprints."""
        layouts = self._layout_index.get('blueprints', [])
        if category:
            layouts = [l for l in layouts if l.get('category') == category]
        if template:
            layouts = [l for l in layouts if l.get('template') == template]
        return layouts

    def get_grids(self, columns: Optional[int] = None,
                  rows: Optional[int] = None) -> List[dict]:
        """Get grid systems."""
        grids = self._layout_index.get('grids', [])
        if columns:
            grids = [g for g in grids if g.get('columns') == columns]
        if rows:
            grids = [g for g in grids if g.get('rows') == rows]
        return grids

    def find_layout_for_content(self, content_requirements: Dict[str, int]) -> List[dict]:
        """
        Find layouts that can accommodate given content.

        Args:
            content_requirements: Dict like {'chart': 2, 'table': 1}

        Returns:
            List of matching layouts sorted by fit score
        """
        layouts = self._layout_index.get('blueprints', [])
        matches = []

        for layout in layouts:
            content_types = layout.get('content_types', {})
            can_fit = True
            excess = 0

            for content_type, count in content_requirements.items():
                available = content_types.get(content_type, 0)
                if available < count:
                    can_fit = False
                    break
                excess += available - count

            if can_fit:
                matches.append((layout, excess))

        # Sort by closest match
        matches.sort(key=lambda x: x[1])
        return [m[0] for m in matches]

    # ==================== Diagram Template Methods ====================

    def get_diagram_templates(self, category: Optional[str] = None,
                               template: Optional[str] = None,
                               min_shapes: Optional[int] = None) -> List[dict]:
        """Get diagram templates."""
        diagrams = self._diagram_index.get('templates', [])
        if category:
            diagrams = [d for d in diagrams if d.get('category') == category]
        if template:
            diagrams = [d for d in diagrams if d.get('template') == template]
        if min_shapes:
            diagrams = [d for d in diagrams if d.get('shape_count', 0) >= min_shapes]
        return diagrams

    def get_diagram_template(self, template_id: str) -> Optional[dict]:
        """Get a specific diagram template by ID."""
        for diagram in self._diagram_index.get('templates', []):
            if diagram.get('id') == template_id:
                return diagram
        return None

    def get_diagram_categories(self) -> Dict[str, int]:
        """Get diagram template categories with counts."""
        categories = {}
        for cat, ids in self._diagram_index.get('categories', {}).items():
            categories[cat] = len(ids)
        return categories

    # ==================== Text Pattern Methods ====================

    def get_bullet_patterns(self, pattern: Optional[str] = None,
                            template: Optional[str] = None) -> List[dict]:
        """Get bullet list patterns."""
        bullets = self._text_index.get('bullet_patterns', [])
        if pattern:
            bullets = [b for b in bullets if b.get('bullet_style', {}).get('pattern') == pattern]
        if template:
            bullets = [b for b in bullets if b.get('template') == template]
        return bullets

    def get_title_patterns(self, format_type: Optional[str] = None,
                           template: Optional[str] = None) -> List[dict]:
        """Get title patterns."""
        titles = self._text_index.get('title_patterns', [])
        if format_type:
            titles = [t for t in titles if t.get('title_style', {}).get('format') == format_type]
        if template:
            titles = [t for t in titles if t.get('template') == template]
        return titles

    def get_callout_patterns(self, callout_type: Optional[str] = None) -> List[dict]:
        """Get callout patterns."""
        callouts = self._text_index.get('callouts', [])
        if callout_type:
            callouts = [c for c in callouts if c.get('callout_type') == callout_type]
        return callouts

    # ==================== Sequence Methods ====================

    def get_deck_templates(self, template: Optional[str] = None) -> List[dict]:
        """Get deck templates."""
        decks = self._sequence_index.get('deck_templates', [])
        if template:
            decks = [d for d in decks if d.get('template') == template]
        return decks

    def get_sequences(self, sequence_type: Optional[str] = None,
                      template: Optional[str] = None) -> List[dict]:
        """Get slide sequences."""
        sequences = self._sequence_index.get('sequences', [])
        if sequence_type:
            sequences = [s for s in sequences if s.get('type') == sequence_type]
        if template:
            sequences = [s for s in sequences if s.get('template') == template]
        return sequences

    def generate_deck_outline(self, structure_type: str, topic: str = 'Presentation') -> dict:
        """
        Generate a deck outline based on structure type.

        Args:
            structure_type: One of 'executive_presentation', 'data_heavy', 'brief', 'multi_section'
            topic: Topic/title for the deck

        Returns:
            Deck outline with slide definitions
        """
        outlines = {
            'executive_presentation': [
                {'type': 'title', 'title': topic},
                {'type': 'agenda', 'title': 'Agenda'},
                {'type': 'summary', 'title': 'Executive Summary'},
                {'type': 'section', 'title': 'Background'},
                {'type': 'content', 'title': 'Current Situation'},
                {'type': 'section', 'title': 'Analysis'},
                {'type': 'data', 'title': 'Key Metrics'},
                {'type': 'comparison', 'title': 'Comparison'},
                {'type': 'section', 'title': 'Recommendations'},
                {'type': 'content', 'title': 'Proposed Solution'},
                {'type': 'timeline', 'title': 'Implementation Roadmap'},
                {'type': 'summary', 'title': 'Key Takeaways'},
                {'type': 'contact', 'title': 'Questions & Next Steps'},
            ],
            'data_heavy': [
                {'type': 'title', 'title': topic},
                {'type': 'summary', 'title': 'Key Findings'},
                {'type': 'data', 'title': 'Overview Metrics'},
                {'type': 'data', 'title': 'Trend Analysis'},
                {'type': 'data', 'title': 'Segment Analysis'},
                {'type': 'comparison', 'title': 'Benchmark Comparison'},
                {'type': 'data', 'title': 'Detailed Breakdown'},
                {'type': 'summary', 'title': 'Conclusions'},
                {'type': 'appendix', 'title': 'Appendix: Data Sources'},
            ],
            'brief': [
                {'type': 'title', 'title': topic},
                {'type': 'content', 'title': 'Overview'},
                {'type': 'content', 'title': 'Key Points'},
                {'type': 'data', 'title': 'Supporting Data'},
                {'type': 'summary', 'title': 'Summary'},
            ],
            'multi_section': [
                {'type': 'title', 'title': topic},
                {'type': 'agenda', 'title': 'Contents'},
                {'type': 'section', 'title': 'Section 1'},
                {'type': 'content', 'title': 'Content 1'},
                {'type': 'section', 'title': 'Section 2'},
                {'type': 'content', 'title': 'Content 2'},
                {'type': 'section', 'title': 'Section 3'},
                {'type': 'content', 'title': 'Content 3'},
                {'type': 'summary', 'title': 'Summary'},
            ],
        }

        outline = outlines.get(structure_type, outlines['brief'])

        return {
            'structure_type': structure_type,
            'topic': topic,
            'slides': outline,
            'slide_count': len(outline),
        }

    # ==================== Full Statistics ====================

    def get_full_stats(self) -> dict:
        """Get comprehensive statistics for all library components."""
        stats = {
            'templates': len(self.index.get('templates', {})),
            'components': {},
            'styles': {},
            'layouts': {},
            'diagrams': {},
            'text_patterns': {},
            'sequences': {},
        }

        # Basic components
        for comp_type, items in self.index.get('components', {}).items():
            stats['components'][comp_type] = len(items)

        # Styles
        for key in ['color_palettes', 'typography_presets', 'gradient_presets', 'shadow_presets']:
            stats['styles'][key] = len(self._style_index.get(key, []))

        # Chart styles
        stats['chart_styles'] = len(self._chart_style_index.get('chart_styles', []))

        # Layouts
        stats['layouts']['blueprints'] = len(self._layout_index.get('blueprints', []))
        stats['layouts']['grids'] = len(self._layout_index.get('grids', []))

        # Diagrams
        stats['diagrams']['templates'] = len(self._diagram_index.get('templates', []))

        # Text patterns
        for key in ['bullet_patterns', 'title_patterns', 'text_blocks', 'callouts']:
            stats['text_patterns'][key] = len(self._text_index.get(key, []))

        # Sequences
        stats['sequences']['deck_templates'] = len(self._sequence_index.get('deck_templates', []))
        stats['sequences']['sequences'] = len(self._sequence_index.get('sequences', []))

        return stats

    # ==================== Helper Methods ====================

    def _get_chart_type_enum(self, chart_type_str: str):
        """Convert chart type string to XL_CHART_TYPE enum."""
        from pptx.enum.chart import XL_CHART_TYPE

        # Extract just the type name from strings like "COLUMN_CLUSTERED (51)"
        type_name = chart_type_str.split('(')[0].strip().upper()

        chart_map = {
            'COLUMN_CLUSTERED': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'COLUMN_STACKED': XL_CHART_TYPE.COLUMN_STACKED,
            'COLUMN_STACKED_100': XL_CHART_TYPE.COLUMN_STACKED_100,
            'BAR_CLUSTERED': XL_CHART_TYPE.BAR_CLUSTERED,
            'BAR_STACKED': XL_CHART_TYPE.BAR_STACKED,
            'BAR_STACKED_100': XL_CHART_TYPE.BAR_STACKED_100,
            'LINE': XL_CHART_TYPE.LINE,
            'LINE_MARKERS': XL_CHART_TYPE.LINE_MARKERS,
            'LINE_STACKED': XL_CHART_TYPE.LINE_STACKED,
            'PIE': XL_CHART_TYPE.PIE,
            'PIE_EXPLODED': XL_CHART_TYPE.PIE_EXPLODED,
            'DOUGHNUT': XL_CHART_TYPE.DOUGHNUT,
            'AREA': XL_CHART_TYPE.AREA,
            'AREA_STACKED': XL_CHART_TYPE.AREA_STACKED,
            'SCATTER': XL_CHART_TYPE.XY_SCATTER,
            'XY_SCATTER': XL_CHART_TYPE.XY_SCATTER,
            'RADAR': XL_CHART_TYPE.RADAR,
        }

        return chart_map.get(type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def list_categories(self, component_type: str) -> dict:
        """
        List all categories for a component type with counts.

        Args:
            component_type: Type of component ('charts', 'tables', etc.)

        Returns:
            Dict mapping category names to counts
        """
        categories = {}

        for item in self.index.get('components', {}).get(component_type, []):
            cat = item.get('category', 'uncategorized')
            categories[cat] = categories.get(cat, 0) + 1

        return categories

    def list_templates(self) -> list[dict]:
        """List all source templates in the library."""
        return [
            {
                'name': name,
                'slide_count': info.get('slide_count', 0),
                'components': info.get('components', {})
            }
            for name, info in self.index.get('templates', {}).items()
        ]


# Convenience function for quick access
def get_library(library_path: Optional[Path] = None) -> ComponentLibrary:
    """Get a ComponentLibrary instance."""
    return ComponentLibrary(library_path)


# CLI for testing
def main():
    """Test the component library integration."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Component Library Integration")
    parser.add_argument("--stats", action="store_true", help="Show library statistics")
    parser.add_argument("--search", type=str, help="Search query")
    parser.add_argument("--type", type=str, help="Component type filter")
    parser.add_argument("--category", type=str, help="Category filter")
    parser.add_argument("--list-categories", type=str, help="List categories for a type")

    args = parser.parse_args()

    library = ComponentLibrary()

    if not library.is_available:
        print("Component library not available. Run the extractor first.")
        return

    if args.stats:
        stats = library.get_stats()
        print("\nLibrary Statistics:")
        print(f"  Templates: {stats['templates']}")
        print("  Components:")
        for comp_type, count in stats['components'].items():
            print(f"    {comp_type}: {count}")

    if args.list_categories:
        categories = library.list_categories(args.list_categories)
        print(f"\nCategories for {args.list_categories}:")
        for cat, count in sorted(categories.items(), key=lambda x: -x[1]):
            print(f"  {cat}: {count}")

    if args.search:
        results = library.search(
            component_type=args.type,
            category=args.category,
            limit=10
        )
        print(f"\nSearch results ({len(results)} found):")
        for r in results:
            print(f"  [{r['_type']}] {r['id']} - {r.get('category', 'N/A')}")


if __name__ == "__main__":
    main()
