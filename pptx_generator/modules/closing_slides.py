"""
Closing Slides Module

Appends standard closing slides (Contact, Disclosures, End) to presentations.
Also handles front page logo insertion.

Usage:
    from pptx_generator.modules.closing_slides import append_closing_slides, add_front_page_logo

    # Add closing slides
    append_closing_slides("output.pptx", background_image="path/to/image.png")

    # Add logo to front page
    add_front_page_logo("output.pptx")

    # For BTR presentations, use section image as end slide background
    append_closing_slides_btr("output.pptx")
"""

import json
from io import BytesIO
from pathlib import Path
from typing import Optional, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree


# Default source for closing slides
DEFAULT_SOURCE = Path(__file__).parent.parent.parent / "pptx_generator/output/light_industrial/Light_Industrial_Thesis_vFinal.pptx"

# Layout indices in the template
LAYOUT_CONTACT = 6
LAYOUT_DISCLAIMERS = 7
LAYOUT_END = 13

# White PCCP logo file
WHITE_LOGO_PATH = Path(__file__).parent.parent.parent / "logos/pccp_logo_white.png"

# Closing slide indices (0-based) in the source presentation
CLOSING_SLIDE_INDICES = [43, 44, 45]  # Contact, Disclosures, End

# Logo position on front page (matches style guide)
LOGO_LEFT = Inches(0.4)
LOGO_TOP = Inches(0.4)
LOGO_WIDTH = Inches(1.8)
LOGO_ASPECT_RATIO = 2.382  # Width/Height for PCCP logo

# Logo position on end slide (centered)
END_LOGO_LEFT = Inches(4.25)
END_LOGO_TOP = Inches(3.35)
END_LOGO_WIDTH = Inches(2.5)


def get_style_guide_config() -> dict:
    """Load closing slides config from style guide if available."""
    style_guide_path = Path(__file__).parent.parent.parent / "style_guides/pccp_cs_style_guide_2026.01.01.json"
    if style_guide_path.exists():
        with open(style_guide_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return data.get("closing_slides", {})
    return {}


def extract_logo_from_presentation(source_path: Path, slide_idx: int = 0) -> Optional[bytes]:
    """
    Extract the white PCCP logo from the front page of a presentation.

    Args:
        source_path: Path to source presentation
        slide_idx: Index of slide containing logo (default: 0 for front page)

    Returns:
        Logo image bytes or None if not found
    """
    prs = Presentation(str(source_path))

    if slide_idx >= len(prs.slides):
        return None

    slide = prs.slides[slide_idx]

    # Look for logo image in top-left area
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Logo is typically small and in top-left
            if shape.left < Inches(3) and shape.top < Inches(2) and shape.width < Inches(4):
                return shape.image.blob

    return None


def add_front_page_logo(
    target_path: Path,
    logo_path: Optional[Path] = None,
    output_path: Optional[Path] = None
) -> Path:
    """
    Add the white PCCP logo to the front page of a presentation.

    Args:
        target_path: Path to presentation to modify
        logo_path: Path to logo image file (default: logos/pccp_logo_white.png)
        output_path: Path to save modified presentation (default: overwrites target)

    Returns:
        Path to saved presentation
    """
    target_path = Path(target_path)
    if logo_path is None:
        logo_path = WHITE_LOGO_PATH
    logo_path = Path(logo_path)

    if output_path is None:
        output_path = target_path
    output_path = Path(output_path)

    if not logo_path.exists():
        print(f"Warning: Logo file not found: {logo_path}")
        return target_path

    # Open target and add logo to first slide
    prs = Presentation(str(target_path))

    if len(prs.slides) == 0:
        print("Warning: Target presentation has no slides")
        return target_path

    slide = prs.slides[0]

    # Check if logo already exists in top-left area
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.left < Inches(3) and shape.top < Inches(2) and shape.width < Inches(4):
                # Small image in top-left - likely logo already present
                print("Logo already present on front page")
                return target_path

    # Calculate logo height based on aspect ratio
    logo_height = LOGO_WIDTH / LOGO_ASPECT_RATIO

    # Add logo
    pic = slide.shapes.add_picture(str(logo_path), LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, logo_height)

    # Move logo to FRONT (top of z-order) so it's above background image
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.append(sp)  # Append to end = top of z-order

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Added logo to front page: {output_path.name}")

    return output_path


def copy_slide(source_prs: Presentation, target_prs: Presentation, slide_idx: int) -> None:
    """
    Copy a slide from source to target presentation.

    Args:
        source_prs: Source presentation
        target_prs: Target presentation
        slide_idx: Index of slide to copy from source
    """
    if slide_idx >= len(source_prs.slides):
        print(f"Warning: Slide index {slide_idx} out of range")
        return

    source_slide = source_prs.slides[slide_idx]

    # Get the layout - try to find matching layout in target
    source_layout_name = source_slide.slide_layout.name if source_slide.slide_layout else ""

    # Find or use blank layout
    target_layout = None
    for layout in target_prs.slide_layouts:
        if layout.name and layout.name.lower() == source_layout_name.lower():
            target_layout = layout
            break

    if target_layout is None:
        # Use last layout (often blank or section)
        target_layout = target_prs.slide_layouts[-1]

    # Add new slide
    new_slide = target_prs.slides.add_slide(target_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        _copy_shape(shape, new_slide, source_prs, target_prs)


def _copy_shape(shape, target_slide, source_prs, target_prs):
    """Copy a shape to the target slide."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Copy picture
            image_blob = shape.image.blob
            image_stream = BytesIO(image_blob)
            target_slide.shapes.add_picture(
                image_stream,
                shape.left, shape.top,
                shape.width, shape.height
            )
        elif hasattr(shape, 'text_frame'):
            # Copy text box
            new_shape = target_slide.shapes.add_textbox(
                shape.left, shape.top,
                shape.width, shape.height
            )
            # Copy text content
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                if para_idx == 0:
                    new_para = new_shape.text_frame.paragraphs[0]
                else:
                    new_para = new_shape.text_frame.add_paragraph()

                for run in para.runs:
                    new_run = new_para.add_run()
                    new_run.text = run.text
                    # Copy font properties
                    if run.font.name:
                        new_run.font.name = run.font.name
                    if run.font.size:
                        new_run.font.size = run.font.size
                    if run.font.bold is not None:
                        new_run.font.bold = run.font.bold
                    try:
                        if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            new_run.font.color.rgb = run.font.color.rgb
                    except AttributeError:
                        pass  # Color not set or not accessible
    except Exception as e:
        print(f"Warning: Could not copy shape: {e}")


def replace_background_image(slide, new_image_path: Path) -> bool:
    """
    Replace the background image on a slide.

    Args:
        slide: Slide to modify
        new_image_path: Path to new background image

    Returns:
        True if replacement was successful
    """
    if not new_image_path.exists():
        print(f"Warning: Background image not found: {new_image_path}")
        return False

    # Find and remove existing background image
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.width > Inches(8) and shape.height > Inches(6):
                shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Add new background image
    pic = slide.shapes.add_picture(
        str(new_image_path),
        Inches(0), Inches(0),
        Inches(11), Inches(8.5)
    )

    # Move to back
    spTree = slide.shapes._spTree
    sp = spTree[-1]
    spTree.remove(sp)
    spTree.insert(2, sp)

    return True


def append_closing_slides(
    target_path: Path,
    source_path: Optional[Path] = None,
    output_path: Optional[Path] = None,
    slide_indices: Optional[List[int]] = None,
    end_slide_background: Optional[Path] = None
) -> Path:
    """
    Append closing slides (Contact, Disclosures, End) to a presentation.

    Args:
        target_path: Path to presentation to modify
        source_path: Path to source presentation with closing slides
        output_path: Path to save modified presentation (default: overwrites target)
        slide_indices: Indices of slides to copy (default: [43, 44, 45])
        end_slide_background: Optional path to background image for End slide

    Returns:
        Path to saved presentation
    """
    target_path = Path(target_path)
    if source_path is None:
        source_path = DEFAULT_SOURCE
    source_path = Path(source_path)

    if output_path is None:
        output_path = target_path
    output_path = Path(output_path)

    if slide_indices is None:
        slide_indices = CLOSING_SLIDE_INDICES

    if not source_path.exists():
        print(f"Error: Source presentation not found: {source_path}")
        return target_path

    # Load presentations
    source_prs = Presentation(str(source_path))
    target_prs = Presentation(str(target_path))

    # Copy each closing slide
    for idx in slide_indices:
        if idx < len(source_prs.slides):
            copy_slide(source_prs, target_prs, idx)
            print(f"  Copied slide {idx + 1} from source")
        else:
            print(f"  Warning: Slide {idx + 1} not found in source")

    # Replace End slide background if specified
    if end_slide_background:
        # End slide is the last one we added
        end_slide = target_prs.slides[-1]
        if replace_background_image(end_slide, Path(end_slide_background)):
            print(f"  Replaced End slide background with: {end_slide_background}")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    target_prs.save(str(output_path))
    print(f"Saved with closing slides: {output_path}")

    return output_path


def add_closing_slides_to_presentation(
    pptx_path: Path,
    background_image: Optional[Path] = None,
    add_logo: bool = True
) -> Path:
    """
    Complete workflow: Add logo and closing slides to a presentation.

    Args:
        pptx_path: Path to presentation
        background_image: Optional background image for End slide
        add_logo: Whether to add logo to front page

    Returns:
        Path to modified presentation
    """
    pptx_path = Path(pptx_path)

    print(f"\nProcessing: {pptx_path.name}")

    # Add logo first
    if add_logo:
        add_front_page_logo(pptx_path)

    # Add closing slides
    append_closing_slides(pptx_path, end_slide_background=background_image)

    return pptx_path


def get_section_background_image(prs: Presentation, section_slide_idx: int = 3) -> Optional[bytes]:
    """
    Extract background image from a section slide.

    Args:
        prs: Presentation object
        section_slide_idx: Index of section slide to get image from (0-based)

    Returns:
        Image bytes or None if not found
    """
    slides = list(prs.slides)
    if section_slide_idx >= len(slides):
        return None

    slide = slides[section_slide_idx]
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.width > Inches(8):  # Background image
                return shape.image.blob
    return None


def append_closing_slides_with_layouts(
    target_path: Path,
    source_path: Optional[Path] = None,
    output_path: Optional[Path] = None,
    end_slide_background: Optional[bytes] = None,
    add_end_logo: bool = True
) -> Path:
    """
    Append closing slides using proper master layouts (preserves formatting).

    Args:
        target_path: Path to presentation to modify
        source_path: Path to source presentation for content
        output_path: Path to save modified presentation
        end_slide_background: Background image bytes for End slide
        add_end_logo: Whether to add white logo to End slide

    Returns:
        Path to saved presentation
    """
    target_path = Path(target_path)
    if source_path is None:
        source_path = DEFAULT_SOURCE
    source_path = Path(source_path)
    if output_path is None:
        output_path = target_path
    output_path = Path(output_path)

    if not source_path.exists():
        print(f"Error: Source presentation not found: {source_path}")
        return target_path

    # Load presentations
    source_prs = Presentation(str(source_path))
    target_prs = Presentation(str(target_path))

    # Get content from source slides
    source_contact = source_prs.slides[43]
    source_disclaim = source_prs.slides[44]

    def get_placeholder_text(slide, idx):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == idx:
                if shape.has_text_frame:
                    return shape.text_frame.text
        return ''

    # Contact content (idx 0=title, 21-26=offices)
    contact_content = {
        0: get_placeholder_text(source_contact, 0),
        21: get_placeholder_text(source_contact, 21),
        22: get_placeholder_text(source_contact, 22),
        23: get_placeholder_text(source_contact, 23),
        24: get_placeholder_text(source_contact, 24),
        25: get_placeholder_text(source_contact, 25),
        26: get_placeholder_text(source_contact, 26),
    }

    # Disclaimers content
    disclaim_content = {
        0: get_placeholder_text(source_disclaim, 0),
        21: get_placeholder_text(source_disclaim, 21),
    }

    # Add Contact slide using layout
    contact_layout = target_prs.slide_layouts[LAYOUT_CONTACT]
    contact_slide = target_prs.slides.add_slide(contact_layout)
    for idx, text in contact_content.items():
        for shape in contact_slide.placeholders:
            if shape.placeholder_format.idx == idx and shape.has_text_frame:
                shape.text_frame.paragraphs[0].text = text
    print("  Added Contact slide")

    # Add Disclaimers slide using layout
    disclaim_layout = target_prs.slide_layouts[LAYOUT_DISCLAIMERS]
    disclaim_slide = target_prs.slides.add_slide(disclaim_layout)
    for idx, text in disclaim_content.items():
        for shape in disclaim_slide.placeholders:
            if shape.placeholder_format.idx == idx and shape.has_text_frame:
                shape.text_frame.paragraphs[0].text = text
    print("  Added Disclosures slide")

    # Add End slide using layout
    end_layout = target_prs.slide_layouts[LAYOUT_END]
    end_slide = target_prs.slides.add_slide(end_layout)

    # Add background image
    if end_slide_background:
        pic = end_slide.shapes.add_picture(
            BytesIO(end_slide_background),
            Inches(0), Inches(0),
            Inches(11), Inches(8.5)
        )
        # Apply 16:9 to letter crop (13.6% from each side)
        blipFill = pic._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
        if blipFill is not None:
            srcRect = etree.SubElement(blipFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
            srcRect.set('l', '13600')
            srcRect.set('r', '13600')
            srcRect.set('t', '0')
            srcRect.set('b', '0')
        # Move to back
        spTree = end_slide.shapes._spTree
        sp = spTree[-1]
        spTree.remove(sp)
        spTree.insert(2, sp)
        print("  Added End slide background")

    # Add white logo to End slide
    if add_end_logo and WHITE_LOGO_PATH.exists():
        logo = end_slide.shapes.add_picture(
            str(WHITE_LOGO_PATH),
            END_LOGO_LEFT, END_LOGO_TOP,
            END_LOGO_WIDTH
        )
        print("  Added End slide logo")

    print("  Added End slide")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    target_prs.save(str(output_path))
    print(f"Saved with closing slides: {output_path}")

    return output_path


def append_closing_slides_btr(
    target_path: Path,
    output_path: Optional[Path] = None,
    section_slide_idx: int = 3
) -> Path:
    """
    Append closing slides for BTR presentations.
    Uses a section slide background image for the End slide.

    Args:
        target_path: Path to BTR presentation
        output_path: Path to save (default: overwrites target)
        section_slide_idx: Index of section slide to get background from

    Returns:
        Path to saved presentation
    """
    target_path = Path(target_path)
    if output_path is None:
        output_path = target_path
    output_path = Path(output_path)

    # Load presentation to get section background
    prs = Presentation(str(target_path))
    bg_image = get_section_background_image(prs, section_slide_idx)

    if bg_image:
        print(f"  Using section slide {section_slide_idx + 1} background for End slide")
    else:
        print("  Warning: No section background found, using default")

    return append_closing_slides_with_layouts(
        target_path,
        output_path=output_path,
        end_slide_background=bg_image,
        add_end_logo=True
    )


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python closing_slides.py <pptx_file> [background_image]")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    background = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    add_closing_slides_to_presentation(pptx_path, background)
