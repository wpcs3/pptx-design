"""
Template Analyzer Module

Analyzes PowerPoint templates to extract:
- Style guide (colors, fonts, spacing)
- Slide catalog (slide types, layouts, element patterns)
"""

import json
import logging
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)


class TemplateAnalyzer:
    """Analyzes PowerPoint templates to extract styles and patterns."""

    def __init__(self, templates_dir: str):
        """
        Initialize the analyzer.

        Args:
            templates_dir: Path to directory containing PPTX templates
        """
        self.templates_dir = Path(templates_dir)
        self.templates: list[Path] = []
        self.style_data: dict = {}
        self.slide_data: list[dict] = []

        self._discover_templates()

    def _discover_templates(self) -> None:
        """Find all PPTX files in templates directory."""
        self.templates = list(self.templates_dir.rglob("*.pptx"))
        logger.info(f"Found {len(self.templates)} templates")

    def analyze_all(self) -> tuple[dict, dict]:
        """
        Analyze all templates and extract style guide and slide catalog.

        Returns:
            Tuple of (style_guide, slide_catalog)
        """
        all_colors = []
        all_fonts = []
        all_slides = []
        master_layouts = defaultdict(list)

        for template_path in self.templates:
            logger.info(f"Analyzing: {template_path.name}")
            try:
                prs = Presentation(str(template_path))

                # Extract colors and fonts from this template
                colors, fonts = self._extract_styles(prs)
                all_colors.extend(colors)
                all_fonts.extend(fonts)

                # Extract slide information
                slides = self._extract_slides(prs, template_path.name)
                all_slides.extend(slides)

                # Extract master layouts
                for layout in prs.slide_layouts:
                    master_layouts[layout.name].append(template_path.name)

            except Exception as e:
                logger.error(f"Error analyzing {template_path}: {e}")

        # Build style guide from aggregated data
        style_guide = self._build_style_guide(all_colors, all_fonts, master_layouts)

        # Build slide catalog from extracted slides
        slide_catalog = self._build_slide_catalog(all_slides)

        return style_guide, slide_catalog

    def _extract_styles(self, prs: Presentation) -> tuple[list, list]:
        """Extract colors and fonts from a presentation."""
        colors = []
        fonts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                # Extract colors from shapes
                shape_colors = self._extract_shape_colors(shape)
                colors.extend(shape_colors)

                # Extract fonts from text frames
                if shape.has_text_frame:
                    shape_fonts = self._extract_text_fonts(shape.text_frame)
                    fonts.extend(shape_fonts)

        return colors, fonts

    def _extract_shape_colors(self, shape) -> list[dict]:
        """Extract colors from a shape."""
        colors = []

        try:
            # Fill color
            if hasattr(shape, 'fill') and shape.fill:
                fill = shape.fill
                if fill.type is not None:
                    color = self._get_fill_color(fill)
                    if color:
                        colors.append({
                            "hex": color,
                            "context": "fill",
                            "shape_type": str(shape.shape_type) if hasattr(shape, 'shape_type') else "unknown"
                        })

            # Line color
            if hasattr(shape, 'line') and shape.line:
                line = shape.line
                if line.fill and line.fill.type is not None:
                    color = self._get_fill_color(line.fill)
                    if color:
                        colors.append({
                            "hex": color,
                            "context": "line"
                        })
        except Exception:
            pass

        return colors

    def _get_fill_color(self, fill) -> Optional[str]:
        """Get hex color from a fill."""
        try:
            if fill.type == 1:  # Solid fill
                fore_color = fill.fore_color
                if fore_color.type == 1:  # RGB
                    rgb = fore_color.rgb
                    return f"#{rgb}"
                elif fore_color.type == 2:  # Theme color
                    # Try to get the actual RGB value
                    try:
                        rgb = fore_color.rgb
                        if rgb:
                            return f"#{rgb}"
                    except:
                        pass
        except Exception:
            pass
        return None

    def _extract_text_fonts(self, text_frame) -> list[dict]:
        """Extract font information from a text frame."""
        fonts = []

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    font = run.font
                    font_info = {
                        "name": font.name,
                        "size_pt": font.size.pt if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                        "color": None
                    }

                    # Get font color
                    if font.color and font.color.type == 1:
                        try:
                            font_info["color"] = f"#{font.color.rgb}"
                        except:
                            pass

                    if font_info["name"] or font_info["size_pt"]:
                        fonts.append(font_info)
                except Exception:
                    pass

        return fonts

    def _extract_slides(self, prs: Presentation, template_name: str) -> list[dict]:
        """Extract slide information from a presentation."""
        slides = []

        for idx, slide in enumerate(prs.slides):
            slide_info = {
                "template": template_name,
                "slide_index": idx,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "unknown",
                "shape_count": len(slide.shapes),
                "shapes": [],
                "has_title": False,
                "has_body": False,
                "has_chart": False,
                "has_table": False,
                "has_image": False,
                "text_preview": ""
            }

            text_parts = []

            for shape in slide.shapes:
                shape_info = self._analyze_shape(shape)
                slide_info["shapes"].append(shape_info)

                # Track content types
                if shape_info["type"] == "PLACEHOLDER" and shape_info.get("placeholder_type") == "TITLE":
                    slide_info["has_title"] = True
                if shape_info["type"] == "PLACEHOLDER" and shape_info.get("placeholder_type") in ["BODY", "OBJECT"]:
                    slide_info["has_body"] = True
                if shape_info["type"] == "CHART":
                    slide_info["has_chart"] = True
                if shape_info["type"] == "TABLE":
                    slide_info["has_table"] = True
                if shape_info["type"] == "PICTURE":
                    slide_info["has_image"] = True

                # Collect text for preview
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        text_parts.append(text[:100])

            slide_info["text_preview"] = " | ".join(text_parts[:3])[:200]
            slides.append(slide_info)

        return slides

    def _analyze_shape(self, shape) -> dict:
        """Analyze a single shape."""
        shape_info = {
            "type": str(shape.shape_type.name) if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
            "left": shape.left / Emu(914400) if shape.left else 0,  # Convert to inches
            "top": shape.top / Emu(914400) if shape.top else 0,
            "width": shape.width / Emu(914400) if shape.width else 0,
            "height": shape.height / Emu(914400) if shape.height else 0,
            "has_text": shape.has_text_frame,
        }

        # Check for placeholder type
        if shape.is_placeholder:
            try:
                shape_info["placeholder_type"] = str(shape.placeholder_format.type.name)
            except:
                shape_info["placeholder_type"] = "unknown"

        return shape_info

    def _build_style_guide(self, colors: list, fonts: list, master_layouts: dict) -> dict:
        """Build style guide from extracted data."""
        # Count color frequencies
        color_counts = Counter()
        fill_colors = []
        text_colors = []
        line_colors = []

        for color_info in colors:
            if color_info.get("hex"):
                hex_val = color_info["hex"]
                color_counts[hex_val] += 1

                if color_info.get("context") == "fill":
                    fill_colors.append(hex_val)
                elif color_info.get("context") == "line":
                    line_colors.append(hex_val)

        # Count font frequencies
        font_name_counts = Counter()
        font_sizes = defaultdict(list)

        for font_info in fonts:
            if font_info.get("name"):
                font_name_counts[font_info["name"]] += 1
            if font_info.get("size_pt"):
                if font_info.get("bold"):
                    font_sizes["title"].append(font_info["size_pt"])
                else:
                    font_sizes["body"].append(font_info["size_pt"])
            if font_info.get("color"):
                text_colors.append(font_info["color"])

        # Get most common values
        top_colors = [c[0] for c in color_counts.most_common(10)]
        top_fonts = [f[0] for f in font_name_counts.most_common(5)]

        # Determine primary/secondary colors
        fill_color_counts = Counter(fill_colors)
        top_fills = [c[0] for c in fill_color_counts.most_common(5)]

        text_color_counts = Counter(text_colors)
        top_text_colors = [c[0] for c in text_color_counts.most_common(3)]

        # Calculate average font sizes
        avg_title_size = sum(font_sizes["title"]) / len(font_sizes["title"]) if font_sizes["title"] else 28
        avg_body_size = sum(font_sizes["body"]) / len(font_sizes["body"]) if font_sizes["body"] else 18

        style_guide = {
            "colors": {
                "primary": top_fills[0] if top_fills else "#333333",
                "secondary": top_fills[1] if len(top_fills) > 1 else "#666666",
                "accent": top_fills[2:5] if len(top_fills) > 2 else ["#0066CC"],
                "text": {
                    "title": top_text_colors[0] if top_text_colors else "#333333",
                    "body": top_text_colors[1] if len(top_text_colors) > 1 else "#666666",
                    "subtle": top_text_colors[2] if len(top_text_colors) > 2 else "#999999"
                },
                "backgrounds": [c for c in top_colors if self._is_light_color(c)][:3] or ["#FFFFFF"],
                "all_colors": top_colors
            },
            "fonts": {
                "title": {
                    "name": top_fonts[0] if top_fonts else "Arial",
                    "size_pt": round(avg_title_size),
                    "bold": True
                },
                "subtitle": {
                    "name": top_fonts[0] if top_fonts else "Arial",
                    "size_pt": round(avg_title_size * 0.7),
                    "bold": False
                },
                "body": {
                    "name": top_fonts[0] if top_fonts else "Arial",
                    "size_pt": round(avg_body_size)
                },
                "caption": {
                    "name": top_fonts[0] if top_fonts else "Arial",
                    "size_pt": round(avg_body_size * 0.7)
                },
                "all_fonts": list(font_name_counts.keys())
            },
            "spacing": {
                "margins": {"left": 0.5, "right": 0.5, "top": 0.5, "bottom": 0.5},
                "line_spacing": 1.15,
                "paragraph_spacing_pt": 12
            },
            "master_slides": {
                layout: {
                    "usage_count": len(templates),
                    "templates": templates
                }
                for layout, templates in master_layouts.items()
            }
        }

        return style_guide

    def _is_light_color(self, hex_color: str) -> bool:
        """Check if a color is light (suitable for background)."""
        try:
            hex_val = hex_color.lstrip("#")
            r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
            # Calculate relative luminance
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return luminance > 0.7
        except:
            return False

    def _build_slide_catalog(self, slides: list[dict]) -> dict:
        """Build slide catalog from extracted slides."""
        # Cluster slides by characteristics
        slide_types = defaultdict(list)

        for slide in slides:
            # Create a type signature based on content
            type_key = self._get_slide_type_signature(slide)
            slide_types[type_key].append(slide)

        # Build catalog entries
        catalog_entries = []
        type_id_counter = defaultdict(int)

        for type_sig, slide_list in slide_types.items():
            # Determine semantic name based on signature
            type_name, type_id_base = self._get_semantic_type_name(type_sig)
            type_id_counter[type_id_base] += 1

            # Use first slide as example
            example = slide_list[0]

            entry = {
                "id": f"{type_id_base}_{type_id_counter[type_id_base]}" if type_id_counter[type_id_base] > 1 else type_id_base,
                "name": type_name,
                "description": self._generate_type_description(type_sig, slide_list),
                "master_layout": example["layout_name"],
                "signature": type_sig,
                "elements": self._extract_element_types(example),
                "usage": f"Found in {len(slide_list)} slides across templates",
                "examples": [
                    {
                        "template": s["template"],
                        "slide_index": s["slide_index"],
                        "text_preview": s["text_preview"]
                    }
                    for s in slide_list[:5]  # First 5 examples
                ],
                "shape_count": example["shape_count"],
                "occurrence_count": len(slide_list)
            }
            catalog_entries.append(entry)

        # Sort by occurrence count
        catalog_entries.sort(key=lambda x: x["occurrence_count"], reverse=True)

        return {
            "slide_types": catalog_entries,
            "layout_mapping": self._build_layout_mapping(slides),
            "statistics": {
                "total_slides_analyzed": len(slides),
                "unique_type_signatures": len(slide_types),
                "templates_analyzed": len(set(s["template"] for s in slides))
            }
        }

    def _get_slide_type_signature(self, slide: dict) -> str:
        """Create a type signature for slide clustering."""
        parts = [
            f"layout:{slide['layout_name']}",
            f"shapes:{slide['shape_count']}",
            f"title:{slide['has_title']}",
            f"body:{slide['has_body']}",
            f"chart:{slide['has_chart']}",
            f"table:{slide['has_table']}",
            f"image:{slide['has_image']}"
        ]
        return "|".join(parts)

    def _get_semantic_type_name(self, signature: str) -> tuple[str, str]:
        """Get semantic name and ID base from signature."""
        # Parse signature
        parts = dict(p.split(":") for p in signature.split("|"))

        layout = parts.get("layout", "").lower()
        has_chart = parts.get("chart") == "True"
        has_table = parts.get("table") == "True"
        has_image = parts.get("image") == "True"
        has_title = parts.get("title") == "True"
        has_body = parts.get("body") == "True"
        shape_count = int(parts.get("shapes", 0))

        # Determine type based on characteristics
        if "title slide" in layout or (has_title and not has_body and shape_count < 5):
            return "Title Slide", "title_slide"
        elif "section" in layout:
            return "Section Divider", "section_divider"
        elif has_chart and has_title:
            return "Data Visualization Slide", "data_chart"
        elif has_table:
            return "Table Slide", "table_slide"
        elif has_image and has_title:
            return "Image Slide", "image_slide"
        elif "two" in layout.lower() or "comparison" in layout.lower():
            return "Two-Column Layout", "two_column"
        elif has_title and has_body:
            return "Title and Content", "title_content"
        elif shape_count > 10:
            return "Complex Layout", "complex_layout"
        else:
            return "General Slide", "general_slide"

    def _generate_type_description(self, signature: str, slides: list) -> str:
        """Generate description for a slide type."""
        parts = dict(p.split(":") for p in signature.split("|"))

        desc_parts = []
        if parts.get("chart") == "True":
            desc_parts.append("contains data visualization")
        if parts.get("table") == "True":
            desc_parts.append("includes tabular data")
        if parts.get("image") == "True":
            desc_parts.append("features images")
        if parts.get("title") == "True" and parts.get("body") == "True":
            desc_parts.append("has title and body content")

        layout = parts.get("layout", "Unknown")
        base = f"Slide using '{layout}' layout"

        if desc_parts:
            return f"{base}; {'; '.join(desc_parts)}"
        return base

    def _extract_element_types(self, slide: dict) -> list[dict]:
        """Extract element types from slide shapes."""
        elements = []

        for shape in slide["shapes"]:
            element = {
                "type": shape["type"].lower().replace("_", " "),
                "position": {
                    "left": round(shape["left"], 2),
                    "top": round(shape["top"], 2),
                    "width": round(shape["width"], 2),
                    "height": round(shape["height"], 2)
                }
            }

            if shape.get("placeholder_type"):
                element["purpose"] = shape["placeholder_type"].lower()

            elements.append(element)

        return elements

    def _build_layout_mapping(self, slides: list) -> dict:
        """Build mapping of layout names to usage statistics."""
        layout_stats = defaultdict(lambda: {"count": 0, "templates": set()})

        for slide in slides:
            layout = slide["layout_name"]
            layout_stats[layout]["count"] += 1
            layout_stats[layout]["templates"].add(slide["template"])

        return {
            layout: {
                "count": stats["count"],
                "templates": list(stats["templates"])
            }
            for layout, stats in layout_stats.items()
        }

    def save_results(self, output_dir: str) -> tuple[Path, Path]:
        """
        Analyze templates and save results to config files.

        Args:
            output_dir: Directory to save config files

        Returns:
            Tuple of (style_guide_path, slide_catalog_path)
        """
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        style_guide, slide_catalog = self.analyze_all()

        style_guide_path = output_path / "style_guide.json"
        slide_catalog_path = output_path / "slide_catalog.json"

        with open(style_guide_path, "w", encoding="utf-8") as f:
            json.dump(style_guide, f, indent=2)

        with open(slide_catalog_path, "w", encoding="utf-8") as f:
            json.dump(slide_catalog, f, indent=2)

        logger.info(f"Saved style guide to: {style_guide_path}")
        logger.info(f"Saved slide catalog to: {slide_catalog_path}")

        return style_guide_path, slide_catalog_path


def main():
    """Run template analysis from command line."""
    import argparse

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    parser = argparse.ArgumentParser(description="Analyze PowerPoint templates")
    parser.add_argument(
        "--templates-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_templates",
        help="Directory containing PPTX templates"
    )
    parser.add_argument(
        "--output-dir",
        default="C:/Users/wpcol/claudecode/pptx-design/pptx_generator/config",
        help="Directory to save config files"
    )

    args = parser.parse_args()

    analyzer = TemplateAnalyzer(args.templates_dir)
    style_path, catalog_path = analyzer.save_results(args.output_dir)

    print(f"\nStyle guide saved to: {style_path}")
    print(f"Slide catalog saved to: {catalog_path}")


if __name__ == "__main__":
    main()
