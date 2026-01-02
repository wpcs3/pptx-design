"""
End Module Cloner for PCCP Style Guide Presentations

This module provides functionality to append the standard PCCP end module
(Contact, Disclosures, End slides) to any presentation using the pccp_cs_style_guide.

The end module is cloned exactly from a template file to ensure consistent
formatting across all presentations.

Usage:
    from pptx_generator.modules.end_module import append_end_module

    # Append to an existing presentation
    append_end_module(presentation_object)

    # Or append to a file
    append_end_module_to_file('my_presentation.pptx')
"""

import copy
from pathlib import Path
from typing import Optional, Union

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree


# Path to the end module template
END_MODULE_TEMPLATE = Path(__file__).parent.parent / 'assets' / 'end_module' / 'pccp_end_module_template.pptx'

# Expected slide layouts in the end module
END_MODULE_LAYOUTS = ['Contact', 'Disclaimers', 'End']


def get_end_module_template_path() -> Path:
    """Get the path to the end module template file."""
    if not END_MODULE_TEMPLATE.exists():
        raise FileNotFoundError(
            f"End module template not found at: {END_MODULE_TEMPLATE}\n"
            "Please ensure the template file exists."
        )
    return END_MODULE_TEMPLATE


def clone_slide(source_prs: Presentation, source_slide_idx: int,
                target_prs: Presentation) -> None:
    """
    Clone a slide from source presentation to target presentation.

    This performs a deep copy of the slide including all shapes, images,
    and formatting. The slide layout is matched by name in the target.

    Args:
        source_prs: Source presentation object
        source_slide_idx: Index of slide to clone (0-based)
        target_prs: Target presentation object
    """
    source_slide = source_prs.slides[source_slide_idx]
    source_layout = source_slide.slide_layout
    source_layout_name = source_layout.name if source_layout else None

    # Find matching layout in target presentation
    target_layout = None
    for layout in target_prs.slide_masters[0].slide_layouts:
        if layout.name == source_layout_name:
            target_layout = layout
            break

    if target_layout is None:
        # Fall back to first layout if no match
        print(f"  Warning: Layout '{source_layout_name}' not found in target, using first layout")
        target_layout = target_prs.slide_masters[0].slide_layouts[0]

    # Add new slide with matching layout
    target_slide = target_prs.slides.add_slide(target_layout)

    # Copy all shapes from source to target
    _copy_slide_content(source_slide, target_slide, source_prs, target_prs)


def _copy_slide_content(source_slide, target_slide, source_prs, target_prs):
    """
    Copy all content from source slide to target slide.

    Handles placeholders, shapes, images, and text with formatting.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.dml.color import RGBColor

    # Build a map of placeholder indices to target placeholders
    target_placeholders = {}
    for shape in target_slide.shapes:
        if shape.is_placeholder:
            ph_idx = shape.placeholder_format.idx
            target_placeholders[ph_idx] = shape

    # Track which placeholders we've filled
    filled_placeholders = set()

    # First pass: handle placeholders
    for source_shape in source_slide.shapes:
        if source_shape.is_placeholder:
            ph_idx = source_shape.placeholder_format.idx

            if ph_idx in target_placeholders:
                target_shape = target_placeholders[ph_idx]
                filled_placeholders.add(ph_idx)

                # Copy text content with formatting
                if source_shape.has_text_frame and target_shape.has_text_frame:
                    _copy_text_frame(source_shape.text_frame, target_shape.text_frame)

                # Handle picture placeholders
                if source_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    try:
                        ph_type = source_shape.placeholder_format.type
                        if ph_type == PP_PLACEHOLDER.PICTURE:
                            # Picture placeholder - check if it has an image
                            if hasattr(source_shape, 'image') and source_shape.image:
                                # Insert the image into the placeholder
                                image_blob = source_shape.image.blob
                                target_shape.insert_picture(image_blob)
                    except Exception:
                        pass

    # Second pass: handle non-placeholder shapes (images, custom shapes)
    for source_shape in source_slide.shapes:
        if not source_shape.is_placeholder:
            if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Clone picture shape
                _clone_picture(source_shape, target_slide, source_prs, target_prs)
            elif source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Clone text box
                _clone_textbox(source_shape, target_slide)
            elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Clone auto shape
                _clone_autoshape(source_shape, target_slide)


def _copy_text_frame(source_tf, target_tf):
    """Copy text frame content with formatting."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    # Clear existing paragraphs in target (keep first one)
    while len(target_tf.paragraphs) > 1:
        p = target_tf.paragraphs[-1]._p
        p.getparent().remove(p)

    # Copy each paragraph
    for para_idx, source_para in enumerate(source_tf.paragraphs):
        if para_idx == 0:
            target_para = target_tf.paragraphs[0]
        else:
            target_para = target_tf.add_paragraph()

        # Copy paragraph properties
        try:
            target_para.alignment = source_para.alignment
            target_para.level = source_para.level
            if source_para.line_spacing:
                target_para.line_spacing = source_para.line_spacing
            if source_para.space_before:
                target_para.space_before = source_para.space_before
            if source_para.space_after:
                target_para.space_after = source_para.space_after
        except Exception:
            pass

        # Clear existing runs in target paragraph
        for run in list(target_para.runs):
            run._r.getparent().remove(run._r)

        # Copy runs with formatting
        for source_run in source_para.runs:
            target_run = target_para.add_run()
            target_run.text = source_run.text

            # Copy font properties
            try:
                if source_run.font.name:
                    target_run.font.name = source_run.font.name
                if source_run.font.size:
                    target_run.font.size = source_run.font.size
                if source_run.font.bold is not None:
                    target_run.font.bold = source_run.font.bold
                if source_run.font.italic is not None:
                    target_run.font.italic = source_run.font.italic
                if source_run.font.underline is not None:
                    target_run.font.underline = source_run.font.underline
                if source_run.font.color and source_run.font.color.rgb:
                    target_run.font.color.rgb = source_run.font.color.rgb
            except Exception:
                pass

        # If no runs, copy text directly
        if not source_para.runs and source_para.text:
            target_para.text = source_para.text


def _clone_picture(source_shape, target_slide, source_prs, target_prs):
    """Clone a picture shape to the target slide."""
    import io

    try:
        # Get image data as bytes and wrap in BytesIO
        image_blob = source_shape.image.blob
        image_stream = io.BytesIO(image_blob)

        # Add picture at same position and size
        picture = target_slide.shapes.add_picture(
            image_stream,
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )

        # Move to back if it's a background image (full slide size)
        slide_width = target_prs.slide_width
        slide_height = target_prs.slide_height

        if source_shape.width >= slide_width * 0.9 and source_shape.height >= slide_height * 0.9:
            # This is a background image - move to back
            spTree = target_slide.shapes._spTree
            pic_elem = picture._element
            spTree.remove(pic_elem)
            spTree.insert(2, pic_elem)  # Insert after nvGrpSpPr and grpSpPr

    except Exception as e:
        print(f"  Warning: Could not clone picture: {e}")


def _clone_textbox(source_shape, target_slide):
    """Clone a text box shape to the target slide."""
    try:
        # Add text box at same position and size
        textbox = target_slide.shapes.add_textbox(
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )

        # Copy text content
        if source_shape.has_text_frame:
            _copy_text_frame(source_shape.text_frame, textbox.text_frame)

    except Exception as e:
        print(f"  Warning: Could not clone text box: {e}")


def _clone_autoshape(source_shape, target_slide):
    """Clone an auto shape to the target slide."""
    try:
        # Get shape type
        from pptx.enum.shapes import MSO_SHAPE

        # Add shape at same position and size
        shape = target_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # Default, may need adjustment
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )

        # Copy text if present
        if source_shape.has_text_frame and shape.has_text_frame:
            _copy_text_frame(source_shape.text_frame, shape.text_frame)

    except Exception as e:
        print(f"  Warning: Could not clone auto shape: {e}")


def append_end_module(prs: Presentation, verbose: bool = True) -> int:
    """
    Append the PCCP end module to a presentation.

    Clones the Contact, Disclosures, and End slides from the template
    to the end of the target presentation.

    Args:
        prs: Target presentation object
        verbose: Print progress messages

    Returns:
        Number of slides appended
    """
    template_path = get_end_module_template_path()
    template_prs = Presentation(str(template_path))

    if verbose:
        print(f"  Appending end module from template...")

    slides_added = 0
    for i in range(len(template_prs.slides)):
        layout_name = template_prs.slides[i].slide_layout.name
        if verbose:
            print(f"    Cloning slide {i+1}: {layout_name}")
        clone_slide(template_prs, i, prs)
        slides_added += 1

    if verbose:
        print(f"  Added {slides_added} end module slides")

    return slides_added


def append_end_module_to_file(pptx_path: Union[str, Path],
                               output_path: Optional[Union[str, Path]] = None,
                               verbose: bool = True) -> Path:
    """
    Append the PCCP end module to a presentation file.

    Args:
        pptx_path: Path to the presentation file
        output_path: Output path (defaults to overwriting input)
        verbose: Print progress messages

    Returns:
        Path to the output file
    """
    pptx_path = Path(pptx_path)
    if output_path is None:
        output_path = pptx_path
    else:
        output_path = Path(output_path)

    prs = Presentation(str(pptx_path))
    append_end_module(prs, verbose)
    prs.save(str(output_path))

    if verbose:
        print(f"  Saved: {output_path}")

    return output_path


def verify_end_module(prs: Presentation) -> dict:
    """
    Verify that a presentation has the correct end module.

    Args:
        prs: Presentation object to check

    Returns:
        Dictionary with verification results
    """
    total_slides = len(prs.slides)

    if total_slides < 3:
        return {
            'has_end_module': False,
            'missing': END_MODULE_LAYOUTS,
            'message': 'Presentation has fewer than 3 slides'
        }

    # Check last 3 slides
    results = {
        'has_end_module': True,
        'found': [],
        'missing': [],
        'message': ''
    }

    for i, expected_layout in enumerate(END_MODULE_LAYOUTS):
        slide_idx = total_slides - 3 + i
        slide = prs.slides[slide_idx]
        actual_layout = slide.slide_layout.name if slide.slide_layout else 'Unknown'

        if actual_layout == expected_layout:
            results['found'].append(expected_layout)
        else:
            results['missing'].append(expected_layout)
            results['has_end_module'] = False

    if results['has_end_module']:
        results['message'] = 'End module verified: Contact, Disclosures, End slides present'
    else:
        results['message'] = f"End module incomplete. Missing: {', '.join(results['missing'])}"

    return results


def remove_existing_end_module(prs: Presentation, verbose: bool = True) -> int:
    """
    Remove any existing end module slides before appending a fresh one.

    Checks last slides for Contact, Disclosures, End layouts and removes them.

    Args:
        prs: Presentation object
        verbose: Print progress messages

    Returns:
        Number of slides removed
    """
    removed = 0

    # Check from the end backwards
    while len(prs.slides) > 0:
        last_slide = prs.slides[-1]
        layout_name = last_slide.slide_layout.name if last_slide.slide_layout else ''

        if layout_name in END_MODULE_LAYOUTS:
            # Remove this slide
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
            removed += 1

            if verbose:
                print(f"    Removed existing {layout_name} slide")
        else:
            break

    return removed


if __name__ == "__main__":
    # Test the end module functionality
    import sys

    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
        print(f"Appending end module to: {pptx_path}")
        append_end_module_to_file(pptx_path)
    else:
        print("Usage: python end_module.py <pptx_file>")
        print("\nVerifying end module template exists...")
        try:
            path = get_end_module_template_path()
            print(f"Template found: {path}")

            # Open and verify
            prs = Presentation(str(path))
            print(f"Template has {len(prs.slides)} slides:")
            for i, slide in enumerate(prs.slides):
                layout = slide.slide_layout.name if slide.slide_layout else 'Unknown'
                print(f"  {i+1}. {layout}")
        except Exception as e:
            print(f"Error: {e}")
