"""
Presentation Review Module

Automated style guide compliance checking, gap analysis, and correction system.
Runs after presentation generation to ensure adherence to PCCP style specifications.

Usage:
    from pptx_generator.modules.presentation_review import PresentationReviewer

    reviewer = PresentationReviewer()
    gap_analysis = reviewer.analyze(presentation_path)
    reviewer.apply_corrections(presentation_path, gap_analysis)
"""

import json
import re
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class Priority(Enum):
    """Gap priority levels."""
    HIGH = "HIGH"
    MEDIUM = "MEDIUM"
    LOW = "LOW"


class GapCategory(Enum):
    """Categories of style gaps."""
    CHART_FORMATTING = "Chart Formatting"
    TABLE_FORMATTING = "Table Formatting"
    TYPOGRAPHY = "Typography"
    COLOR_PALETTE = "Color Palette"
    LAYOUT = "Layout"
    CONTENT_FORMAT = "Content Format"
    STRUCTURE = "Structure"
    FOOTER = "Footer"
    DIMENSIONS = "Dimensions"
    IMAGE_FORMATTING = "Image Formatting"
    MISSING_ELEMENT = "Missing Element"
    PRESENTATION_STRUCTURE = "Presentation Structure"


@dataclass
class StyleGap:
    """Represents a single style gap finding."""
    slide_number: int
    shape_name: str
    category: GapCategory
    priority: Priority
    description: str
    expected: str
    actual: str
    auto_fixable: bool = True
    fix_action: Optional[str] = None


@dataclass
class GapAnalysis:
    """Complete gap analysis report."""
    presentation_path: str
    style_guide_version: str
    analysis_date: str
    total_slides: int
    gaps: List[StyleGap] = field(default_factory=list)
    summary: Dict[str, int] = field(default_factory=dict)
    compliance_score: float = 0.0

    def to_markdown(self) -> str:
        """Generate markdown report."""
        lines = [
            f"# Presentation Style Gap Analysis",
            f"",
            f"**Presentation:** `{Path(self.presentation_path).name}`",
            f"**Style Guide:** `{self.style_guide_version}`",
            f"**Analysis Date:** {self.analysis_date}",
            f"**Total Slides:** {self.total_slides}",
            f"**Compliance Score:** {self.compliance_score:.1f}%",
            f"",
            f"---",
            f"",
            f"## Executive Summary",
            f"",
            f"| Category | Gaps Found | Priority |",
            f"|----------|------------|----------|",
        ]

        # Group by category
        by_category = {}
        for gap in self.gaps:
            cat = gap.category.value
            if cat not in by_category:
                by_category[cat] = {"count": 0, "high": 0, "medium": 0, "low": 0}
            by_category[cat]["count"] += 1
            by_category[cat][gap.priority.value.lower()] += 1

        for cat, stats in sorted(by_category.items()):
            priority_str = []
            if stats["high"] > 0:
                priority_str.append(f"{stats['high']} HIGH")
            if stats["medium"] > 0:
                priority_str.append(f"{stats['medium']} MEDIUM")
            if stats["low"] > 0:
                priority_str.append(f"{stats['low']} LOW")
            lines.append(f"| {cat} | {stats['count']} | {', '.join(priority_str)} |")

        lines.extend([
            f"",
            f"---",
            f"",
            f"## Slide-by-Slide Analysis",
            f"",
        ])

        # Group gaps by slide
        by_slide = {}
        for gap in self.gaps:
            if gap.slide_number not in by_slide:
                by_slide[gap.slide_number] = []
            by_slide[gap.slide_number].append(gap)

        for slide_num in sorted(by_slide.keys()):
            slide_gaps = by_slide[slide_num]
            lines.extend([
                f"### Slide {slide_num}",
                f"",
                f"| Element | Category | Issue | Expected | Actual | Priority | Auto-Fix |",
                f"|---------|----------|-------|----------|--------|----------|----------|",
            ])

            for gap in slide_gaps:
                auto_fix = "Yes" if gap.auto_fixable else "No"
                expected = gap.expected[:30] + "..." if len(gap.expected) > 30 else gap.expected
                actual = gap.actual[:30] + "..." if len(gap.actual) > 30 else gap.actual
                lines.append(
                    f"| {gap.shape_name} | {gap.category.value} | {gap.description} | "
                    f"`{expected}` | `{actual}` | {gap.priority.value} | {auto_fix} |"
                )

            lines.append("")

        lines.extend([
            f"---",
            f"",
            f"## Correction Actions",
            f"",
        ])

        # List auto-fixable issues
        fixable = [g for g in self.gaps if g.auto_fixable]
        non_fixable = [g for g in self.gaps if not g.auto_fixable]

        lines.append(f"### Auto-Fixable Issues ({len(fixable)})")
        lines.append("")
        for gap in fixable:
            lines.append(f"- [ ] Slide {gap.slide_number}: {gap.description}")

        lines.append("")
        lines.append(f"### Manual Review Required ({len(non_fixable)})")
        lines.append("")
        for gap in non_fixable:
            lines.append(f"- [ ] Slide {gap.slide_number}: {gap.description}")

        lines.extend([
            f"",
            f"---",
            f"",
            f"*Generated by pptx_generator Presentation Review System*",
        ])

        return "\n".join(lines)


@dataclass
class TypographySpec:
    """Typography specification for a text element."""
    font_name: str
    font_size_pt: float
    font_bold: bool
    font_color: str
    placeholder_idx: Optional[int] = None
    placeholder_types: Optional[List[str]] = None
    horizontal_alignment: Optional[str] = None  # LEFT, CENTER, RIGHT
    vertical_anchor: Optional[str] = None  # TOP, MIDDLE, BOTTOM


@dataclass
class ChartSpec:
    """Chart formatting specification."""
    title_font_name: str = "Arial"
    title_font_size_pt: float = 18.0
    title_font_bold: bool = True
    title_color: str = "#000000"
    data_label_font_name: str = "Arial"
    data_label_font_size_pt: float = 14.0
    data_label_color: str = "#000000"
    axis_label_font_name: str = "Arial"
    axis_label_font_size_pt: float = 14.0
    axis_label_color: str = "#000000"
    legend_font_name: str = "Arial"
    legend_font_size_pt: float = 14.0
    legend_color: str = "#000000"
    gridline_enabled: bool = True
    gridline_width_pt: float = 0.5
    gridline_color: str = "#E2E8F0"  # Light gray per PCCP style guide
    tick_marks: str = "NONE"
    bar_fill_color: str = "#309CE7"
    primary_series_color: str = "#309CE7"
    secondary_series_color: str = "#051C2C"
    tertiary_series_color: str = "#A6A6A6"
    line_width_pt: float = 3.0


@dataclass
class TableSpec:
    """Table formatting specification."""
    header_color: str = "#051C2C"
    header_text_color: str = "#FFFFFF"
    header_font_name: str = "Arial"
    header_font_bold: bool = True
    header_font_size_pt: float = 16.0
    row_odd_color: str = "#FFFFFF"
    row_even_color: str = "#F5F5F5"
    row_text_color: str = "#000000"
    row_font_name: str = "Arial"
    row_font_size_pt: float = 14.0
    cell_margin_lr_inches: float = 0.1
    cell_margin_tb_inches: float = 0.05


@dataclass
class StyleGuideSpec:
    """Parsed style guide specifications."""
    version: str
    slide_width: float  # inches
    slide_height: float  # inches
    # Typography
    title_content: TypographySpec  # Title on content slides (32pt bold black)
    title_section: TypographySpec  # Title on section/front page slides (44pt bold white)
    subtitle: TypographySpec  # Subheaders on content slides (20pt bold black)
    subtitle_frontpage: TypographySpec  # Subtitle on front page ONLY (18pt bold white)
    body: TypographySpec
    content_header: TypographySpec  # Headers within side-by-side bullet lists only (18pt bold black)
    section_name: TypographySpec
    footnote: TypographySpec
    primary_font: str
    # Charts (new structured spec)
    chart: ChartSpec
    # Tables (new structured spec)
    table: TableSpec
    # Legacy chart fields (for backwards compatibility)
    chart_gridline_width: float  # points
    chart_gridline_color: str  # hex
    chart_tick_marks: str  # "NONE"
    chart_data_label_font_size: float  # points
    chart_axis_label_font_size: float  # points
    # Legacy table fields (for backwards compatibility)
    table_header_color: str  # hex
    table_header_font_size: float  # points
    table_row_odd_color: str  # hex
    table_row_even_color: str  # hex
    table_row_font_size: float  # points
    table_cell_margin_lr: float  # inches
    table_cell_margin_tb: float  # inches
    # Metric boxes
    metric_box_color: str  # hex
    metric_value_font_size: float  # points
    metric_label_font_size: float  # points


class StyleGuideLoader:
    """Loads and parses style guide specifications."""

    def __init__(self, style_guides_dir: Optional[Path] = None):
        if style_guides_dir is None:
            style_guides_dir = Path(__file__).parent.parent.parent / "style_guides"
        self.style_guides_dir = style_guides_dir

    def get_latest_guide(self) -> Path:
        """Get the most recent style guide file (JSON preferred, then MD)."""
        # Try JSON first
        json_guides = list(self.style_guides_dir.glob("pccp_cs_style_guide_*.json"))
        if json_guides:
            json_guides.sort(key=lambda p: p.stem, reverse=True)
            return json_guides[0]
        # Fall back to markdown
        md_guides = list(self.style_guides_dir.glob("pccp_cs_style_guide_*.md"))
        if md_guides:
            md_guides.sort(key=lambda p: p.stem, reverse=True)
            return md_guides[0]
        raise FileNotFoundError(f"No style guides found in {self.style_guides_dir}")

    def get_guide_by_version(self, version: str) -> Path:
        """Get a specific version of the style guide."""
        # Try JSON first
        json_path = self.style_guides_dir / f"pccp_cs_style_guide_{version}.json"
        if json_path.exists():
            return json_path
        # Fall back to markdown
        md_path = self.style_guides_dir / f"pccp_cs_style_guide_{version}.md"
        if md_path.exists():
            return md_path
        raise FileNotFoundError(f"Style guide version {version} not found")

    def _parse_typography(self, data: Dict[str, Any]) -> TypographySpec:
        """Parse typography spec from dict."""
        return TypographySpec(
            font_name=data.get("font_name", "Arial"),
            font_size_pt=data.get("font_size_pt", 12.0),
            font_bold=data.get("font_bold", False),
            font_color=data.get("font_color", "#061F32"),
            placeholder_idx=data.get("placeholder_idx"),
            placeholder_types=data.get("placeholder_types"),
            horizontal_alignment=data.get("horizontal_alignment"),
            vertical_anchor=data.get("vertical_anchor"),
        )

    def _parse_chart_spec(self, data: Dict[str, Any]) -> ChartSpec:
        """Parse chart spec from dict."""
        return ChartSpec(
            title_font_name=data.get("title_font_name", "Arial"),
            title_font_size_pt=data.get("title_font_size_pt", 18.0),
            title_font_bold=data.get("title_font_bold", True),
            title_color=data.get("title_color", "#000000"),
            data_label_font_name=data.get("data_label_font_name", "Arial"),
            data_label_font_size_pt=data.get("data_label_font_size_pt", 14.0),
            data_label_color=data.get("data_label_color", "#000000"),
            axis_label_font_name=data.get("axis_label_font_name", "Arial"),
            axis_label_font_size_pt=data.get("axis_label_font_size_pt", 14.0),
            axis_label_color=data.get("axis_label_color", "#000000"),
            legend_font_name=data.get("legend_font_name", "Arial"),
            legend_font_size_pt=data.get("legend_font_size_pt", 14.0),
            legend_color=data.get("legend_color", "#000000"),
            gridline_enabled=data.get("gridline_enabled", True),
            gridline_width_pt=data.get("gridline_width_pt", 0.5),
            gridline_color=data.get("gridline_color", "#E2E8F0"),
            tick_marks=data.get("value_axis_tick_mark", "NONE"),
            bar_fill_color=data.get("bar_fill_color", "#309CE7"),
            primary_series_color=data.get("primary_series_color", "#309CE7"),
            secondary_series_color=data.get("secondary_series_color", "#051C2C"),
            tertiary_series_color=data.get("tertiary_series_color", "#A6A6A6"),
            line_width_pt=data.get("line_width_pt", 3.0),
        )

    def _parse_table_spec(self, data: Dict[str, Any]) -> TableSpec:
        """Parse table spec from dict."""
        return TableSpec(
            header_color=data.get("header_color", "#051C2C"),
            header_text_color=data.get("header_text_color", "#FFFFFF"),
            header_font_name=data.get("header_font_name", "Arial"),
            header_font_bold=data.get("header_font_bold", True),
            header_font_size_pt=data.get("header_font_size_pt", 16.0),
            row_odd_color=data.get("row_odd_color", "#FFFFFF"),
            row_even_color=data.get("row_even_color", "#F5F5F5"),
            row_text_color=data.get("row_text_color", "#000000"),
            row_font_name=data.get("row_font_name", "Arial"),
            row_font_size_pt=data.get("row_font_size_pt", 14.0),
            cell_margin_lr_inches=data.get("cell_margin_lr_inches", 0.1),
            cell_margin_tb_inches=data.get("cell_margin_tb_inches", 0.05),
        )

    def parse_guide(self, guide_path: Optional[Path] = None) -> StyleGuideSpec:
        """Parse style guide into specifications."""
        if guide_path is None:
            guide_path = self.get_latest_guide()

        # Extract version from filename
        version = guide_path.stem.replace("pccp_cs_style_guide_", "")

        # Load from JSON if available
        if guide_path.suffix == ".json":
            with open(guide_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            typo = data.get("typography", {})
            chart = data.get("chart", {})
            table = data.get("table", {})
            metric = data.get("metric_box", {})
            layout = data.get("layout", {})

            # Parse structured specs
            chart_spec = self._parse_chart_spec(chart)
            table_spec = self._parse_table_spec(table)

            return StyleGuideSpec(
                version=data.get("version", version),
                slide_width=layout.get("width_inches", 11.0),
                slide_height=layout.get("height_inches", 8.5),
                # Typography - support both old and new format
                title_content=self._parse_typography(typo.get("title_content", typo.get("title", {}))),
                title_section=self._parse_typography(typo.get("title_section", typo.get("title_frontpage", {}))),
                subtitle=self._parse_typography(typo.get("subtitle", {})),
                subtitle_frontpage=self._parse_typography(typo.get("subtitle_frontpage", {"font_name": "Arial", "font_size_pt": 18.0, "font_bold": True, "font_color": "#FFFFFF"})),
                body=self._parse_typography(typo.get("body", {})),
                content_header=self._parse_typography(typo.get("content_header", {"font_name": "Arial", "font_size_pt": 18.0, "font_bold": True, "font_color": "#000000"})),
                section_name=self._parse_typography(typo.get("section_name", {})),
                footnote=self._parse_typography(typo.get("footnote", {})),
                primary_font=typo.get("primary_font", "Arial"),
                # Structured specs
                chart=chart_spec,
                table=table_spec,
                # Legacy chart fields (for backwards compatibility)
                chart_gridline_width=chart_spec.gridline_width_pt,
                chart_gridline_color=chart_spec.gridline_color,
                chart_tick_marks=chart_spec.tick_marks,
                chart_data_label_font_size=chart_spec.data_label_font_size_pt,
                chart_axis_label_font_size=chart_spec.axis_label_font_size_pt,
                # Legacy table fields (for backwards compatibility)
                table_header_color=table_spec.header_color,
                table_header_font_size=table_spec.header_font_size_pt,
                table_row_odd_color=table_spec.row_odd_color,
                table_row_even_color=table_spec.row_even_color,
                table_row_font_size=table_spec.row_font_size_pt,
                table_cell_margin_lr=table_spec.cell_margin_lr_inches,
                table_cell_margin_tb=table_spec.cell_margin_tb_inches,
                # Metric boxes
                metric_box_color=metric.get("background_color", "#051C2C"),
                metric_value_font_size=metric.get("value_font_size_pt", 28.0),
                metric_label_font_size=metric.get("label_font_size_pt", 14.0),
            )

        # Fallback: Return default specifications per PCCP CS Style Guide
        default_chart = ChartSpec()
        default_table = TableSpec()
        # PCCP style uses #000000 (black) for content slide titles/subtitles
        return StyleGuideSpec(
            version=version,
            slide_width=11.0,
            slide_height=8.5,
            title_content=TypographySpec("Arial", 32.0, True, "#000000", 0),  # Black for content slides
            title_section=TypographySpec("Arial", 44.0, True, "#FFFFFF"),
            subtitle=TypographySpec("Arial", 20.0, True, "#000000", 1),  # Black for content slides
            subtitle_frontpage=TypographySpec("Arial", 18.0, True, "#FFFFFF"),
            body=TypographySpec("Arial", 14.0, False, "#2D3748", placeholder_types=["BODY", "OBJECT"]),  # Slate for body
            content_header=TypographySpec("Arial", 18.0, True, "#000000"),  # Black
            section_name=TypographySpec("Arial", 9.0, False, "#718096", 17),  # Gray/secondary
            footnote=TypographySpec("Arial", 8.0, False, "#A6A6A6", 20, horizontal_alignment="LEFT", vertical_anchor="BOTTOM"),  # Medium gray, bottom-left
            primary_font="Arial",
            chart=default_chart,
            table=default_table,
            chart_gridline_width=default_chart.gridline_width_pt,
            chart_gridline_color=default_chart.gridline_color,
            chart_tick_marks=default_chart.tick_marks,
            chart_data_label_font_size=default_chart.data_label_font_size_pt,
            chart_axis_label_font_size=default_chart.axis_label_font_size_pt,
            table_header_color=default_table.header_color,
            table_header_font_size=default_table.header_font_size_pt,
            table_row_odd_color=default_table.row_odd_color,
            table_row_even_color=default_table.row_even_color,
            table_row_font_size=default_table.row_font_size_pt,
            table_cell_margin_lr=default_table.cell_margin_lr_inches,
            table_cell_margin_tb=default_table.cell_margin_tb_inches,
            metric_box_color="#051C2C",
            metric_value_font_size=28.0,
            metric_label_font_size=14.0,
        )


class PresentationAnalyzer:
    """Analyzes presentations for style guide compliance."""

    def __init__(self, spec: StyleGuideSpec):
        self.spec = spec
        self.gaps: List[StyleGap] = []

    def _emu_to_inches(self, emu: int) -> float:
        """Convert EMU to inches."""
        return round(emu / 914400, 3) if emu else 0

    def _emu_to_pt(self, emu: int) -> float:
        """Convert EMU to points."""
        return round(emu / 12700, 1) if emu else 0

    def _color_to_hex(self, color) -> Optional[str]:
        """Convert color to hex string."""
        try:
            if hasattr(color, 'rgb') and color.rgb:
                return f"#{color.rgb}"
        except:
            pass
        return None

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _colors_match(self, actual: str, expected: str, tolerance: int = 5) -> bool:
        """Check if two colors match within tolerance."""
        if actual is None or expected is None:
            return False
        try:
            actual_rgb = self._hex_to_rgb(actual)
            expected_rgb = self._hex_to_rgb(expected)
            return all(abs(a - e) <= tolerance for a, e in zip(actual_rgb, expected_rgb))
        except:
            return False

    def analyze(self, pptx_path: Path) -> GapAnalysis:
        """Analyze presentation for style compliance."""
        self.gaps = []
        prs = Presentation(str(pptx_path))

        # Check slide dimensions
        self._check_dimensions(prs)

        # Analyze each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            self._analyze_slide(slide, slide_num)

        # Check presentation structure (end module, required slides)
        self._check_presentation_structure(prs)

        # Calculate compliance score
        total_checks = max(len(self.gaps) + 10, 1)  # Assume at least 10 items checked
        compliance = ((total_checks - len(self.gaps)) / total_checks) * 100

        # Build summary
        summary = {}
        for gap in self.gaps:
            cat = gap.category.value
            summary[cat] = summary.get(cat, 0) + 1

        return GapAnalysis(
            presentation_path=str(pptx_path),
            style_guide_version=self.spec.version,
            analysis_date=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            total_slides=len(prs.slides),
            gaps=self.gaps,
            summary=summary,
            compliance_score=compliance,
        )

    def _check_dimensions(self, prs: Presentation):
        """Check slide dimensions."""
        width = self._emu_to_inches(prs.slide_width)
        height = self._emu_to_inches(prs.slide_height)

        if abs(width - self.spec.slide_width) > 0.1:
            self.gaps.append(StyleGap(
                slide_number=0,
                shape_name="Presentation",
                category=GapCategory.DIMENSIONS,
                priority=Priority.HIGH,
                description="Slide width mismatch",
                expected=f"{self.spec.slide_width}\"",
                actual=f"{width}\"",
                auto_fixable=False,
            ))

        if abs(height - self.spec.slide_height) > 0.1:
            self.gaps.append(StyleGap(
                slide_number=0,
                shape_name="Presentation",
                category=GapCategory.DIMENSIONS,
                priority=Priority.HIGH,
                description="Slide height mismatch",
                expected=f"{self.spec.slide_height}\"",
                actual=f"{height}\"",
                auto_fixable=False,
            ))

    def _check_presentation_structure(self, prs: Presentation):
        """Check presentation structure for required elements.

        Uses both layout name matching (exact) and keyword fallback for
        verifying the end module (Contact, Disclosures/Disclaimers, End).
        """
        total_slides = len(prs.slides)

        if total_slides < 3:
            self.gaps.append(StyleGap(
                slide_number=0,
                shape_name="Presentation Structure",
                category=GapCategory.PRESENTATION_STRUCTURE,
                priority=Priority.HIGH,
                description="Presentation has fewer than 3 slides",
                expected="At least 3 slides including end module",
                actual=f"{total_slides} slides",
                auto_fixable=False,
            ))
            return

        # Expected end module layout names (exact match)
        end_module_layouts = ['Contact', 'Disclaimers', 'End']

        # Fallback keywords for content-based matching
        end_module_keywords = {
            'Contact': ['contact', 'reach out', 'get in touch', 'inquiry@', 'pccpllc.com'],
            'Disclaimers': ['disclosure', 'disclaimer', 'legal', 'confidential', 'informational purposes'],
            'End': [],  # End slide may have no text, matched by layout only
        }

        found_end_slides = {layout: False for layout in end_module_layouts}

        # Check last 5 slides for end module content
        for slide_idx in range(max(0, total_slides - 5), total_slides):
            slide = prs.slides[slide_idx]

            # First check by layout name (most reliable)
            layout_name = slide.slide_layout.name if slide.slide_layout else ''
            if layout_name in end_module_layouts:
                found_end_slides[layout_name] = True
                continue

            # Fallback: check by content keywords
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        slide_text += para.text.lower() + " "

            for module_type, keywords in end_module_keywords.items():
                if not found_end_slides[module_type] and keywords:
                    if any(kw in slide_text for kw in keywords):
                        found_end_slides[module_type] = True

        # Report missing end module slides
        if not found_end_slides['Contact']:
            self.gaps.append(StyleGap(
                slide_number=total_slides,
                shape_name="Presentation Structure",
                category=GapCategory.PRESENTATION_STRUCTURE,
                priority=Priority.HIGH,
                description="Missing Contact slide in end module",
                expected="Contact slide with PCCP office information (layout name: 'Contact')",
                actual="Not found in last 5 slides",
                auto_fixable=False,
            ))

        if not found_end_slides['Disclaimers']:
            self.gaps.append(StyleGap(
                slide_number=total_slides,
                shape_name="Presentation Structure",
                category=GapCategory.PRESENTATION_STRUCTURE,
                priority=Priority.HIGH,
                description="Missing Disclosures slide in end module",
                expected="Disclosures/legal disclaimer slide (layout name: 'Disclaimers')",
                actual="Not found in last 5 slides",
                auto_fixable=False,
            ))

        if not found_end_slides['End']:
            self.gaps.append(StyleGap(
                slide_number=total_slides,
                shape_name="Presentation Structure",
                category=GapCategory.PRESENTATION_STRUCTURE,
                priority=Priority.MEDIUM,
                description="Missing End slide in end module",
                expected="Closing slide with logo (layout name: 'End')",
                actual="Not found in last 5 slides",
                auto_fixable=False,
            ))

    def _check_section_slide_background(self, slide, slide_num: int):
        """Check that section/title slides have background images."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Count picture shapes that could be backgrounds
        has_background_image = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Check if image covers significant portion of slide (>50%)
                if shape.width > Inches(5.5) and shape.height > Inches(4.25):
                    has_background_image = True
                    break

        if not has_background_image:
            self.gaps.append(StyleGap(
                slide_number=slide_num,
                shape_name="Slide Background",
                category=GapCategory.MISSING_ELEMENT,
                priority=Priority.MEDIUM,
                description="Section/title slide missing background image",
                expected="Background image covering slide",
                actual="No background image found",
                auto_fixable=False,
            ))

    def _check_source_footnote(self, slide, slide_num: int):
        """Check that content slides have source footnotes where expected."""
        has_footnote = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.lower().strip()
                    # Check for source indicators
                    if text.startswith('source:') or text.startswith('sources:'):
                        has_footnote = True
                        break
                    # Also check placeholder idx 20 (footnote placeholder)
                    ph_idx = self._get_placeholder_idx(shape)
                    if ph_idx == 20 and text:
                        has_footnote = True
                        break

        # Only flag missing footnotes on content slides with data (charts/tables)
        has_data_content = any(
            s.has_chart or s.has_table for s in slide.shapes
        )

        if has_data_content and not has_footnote:
            self.gaps.append(StyleGap(
                slide_number=slide_num,
                shape_name="Source Footnote",
                category=GapCategory.MISSING_ELEMENT,
                priority=Priority.LOW,
                description="Data slide missing source footnote",
                expected="Source: [data source]",
                actual="No source footnote found",
                auto_fixable=False,
            ))

    def _get_placeholder_type(self, shape) -> Optional[str]:
        """Get the placeholder type of a shape."""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = shape.placeholder_format.type
                if ph_type:
                    return str(ph_type).split('(')[0].strip()
        except:
            pass
        return None

    def _get_placeholder_idx(self, shape) -> Optional[int]:
        """Get the placeholder index of a shape."""
        try:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                return shape.placeholder_format.idx
        except:
            pass
        return None

    def _is_section_slide(self, slide) -> bool:
        """
        Determine if a slide is a section/front page slide vs a content slide.
        Section slides typically have:
        - Layout names containing 'section', 'divider', 'front', 'title'
        - Background images covering the slide
        - Fewer shapes (typically just title)
        """
        layout_name = ""
        try:
            if slide.slide_layout and slide.slide_layout.name:
                layout_name = slide.slide_layout.name.lower()
        except:
            pass

        # Check layout name
        section_keywords = ['section', 'divider', 'front', 'title slide', 'frontpage', 'end']
        if any(kw in layout_name for kw in section_keywords):
            return True

        # Check for full-page background image (section slides often have these)
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Large image covering most of slide
                if shape.width > Inches(8) and shape.height > Inches(6):
                    return True

        # Check shape count - section slides typically have very few shapes
        shape_count = len(list(slide.shapes))
        if shape_count <= 3:
            return True

        return False

    def _requires_background_image(self, slide) -> bool:
        """
        Determine if a slide should have a background image.
        Only slides with explicit section/title/front/end layout names require backgrounds.
        This is separate from _is_section_slide which also uses heuristics.
        """
        layout_name = ""
        try:
            if slide.slide_layout and slide.slide_layout.name:
                layout_name = slide.slide_layout.name.lower()
        except:
            pass

        # Only explicit section/title slide layouts require background images
        # Exclude 'contact', 'disclaimers', 'agenda' which may have few shapes but don't need backgrounds
        background_required_keywords = ['section', 'divider', 'frontpage', 'title slide', 'end']
        excluded_keywords = ['contact', 'disclaimer', 'agenda']

        # Check if excluded first
        if any(excl in layout_name for excl in excluded_keywords):
            return False

        return any(kw in layout_name for kw in background_required_keywords)

    def _analyze_slide(self, slide, slide_num: int):
        """Analyze a single slide."""
        is_section = self._is_section_slide(slide)
        is_front_page = (slide_num == 1)

        # Check section/title slides for background images (only for explicit layouts)
        if is_front_page or self._requires_background_image(slide):
            self._check_section_slide_background(slide, slide_num)

        # Check content slides for source footnotes
        if not is_section:
            self._check_source_footnote(slide, slide_num)

        # Detect side-by-side layout by checking layout name
        is_side_by_side = False
        if hasattr(slide, 'slide_layout') and slide.slide_layout:
            layout_name = slide.slide_layout.name.lower() if slide.slide_layout.name else ""
            is_side_by_side = "side by side" in layout_name or "two column" in layout_name

        for shape in slide.shapes:
            # Check charts
            if shape.has_chart:
                self._check_chart(shape, slide_num)

            # Check tables
            if shape.has_table:
                self._check_table(shape, slide_num)

            # Check background images for aspect ratio issues
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Large images (>80% of slide) are background images
                if shape.width > Inches(8.8) and shape.height > Inches(6.8):
                    self._check_background_image(shape, slide_num)

            # Check text formatting based on placeholder type
            if shape.has_text_frame:
                ph_type = self._get_placeholder_type(shape)
                ph_idx = self._get_placeholder_idx(shape)
                self._check_text(shape, slide_num, ph_type, ph_idx, is_section,
                               is_front_page=is_front_page, is_side_by_side=is_side_by_side)

    def _check_chart(self, shape, slide_num: int):
        """Check chart formatting compliance."""
        chart = shape.chart

        # Check value axis gridlines
        if hasattr(chart, 'value_axis') and chart.value_axis:
            va = chart.value_axis

            # Check if major gridlines exist
            if not va.has_major_gridlines:
                self.gaps.append(StyleGap(
                    slide_number=slide_num,
                    shape_name=shape.name,
                    category=GapCategory.CHART_FORMATTING,
                    priority=Priority.HIGH,
                    description="Missing major horizontal gridlines",
                    expected="ON",
                    actual="OFF",
                    auto_fixable=True,
                    fix_action="enable_major_gridlines",
                ))
            else:
                # Check gridline formatting - always ensure proper format
                try:
                    gridlines = va.major_gridlines
                    line = gridlines.format.line

                    # Check width - flag if not set (0/None) or wrong value
                    width_pt = self._emu_to_pt(line.width) if line.width else 0
                    if abs(width_pt - self.spec.chart_gridline_width) > 0.1:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Gridline width incorrect",
                            expected=f"{self.spec.chart_gridline_width}pt",
                            actual=f"{width_pt}pt" if width_pt else "not set",
                            auto_fixable=True,
                            fix_action="set_gridline_format",
                        ))

                    # Check color - flag if not set or wrong
                    color = self._color_to_hex(line.color) if hasattr(line, 'color') and line.color else None
                    if not color or not self._colors_match(color, self.spec.chart_gridline_color):
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Gridline color incorrect",
                            expected=self.spec.chart_gridline_color,
                            actual=color if color else "not set",
                            auto_fixable=True,
                            fix_action="set_gridline_format",
                        ))
                except Exception as e:
                    pass

            # Check tick marks
            try:
                if va.major_tick_mark != XL_TICK_MARK.NONE:
                    self.gaps.append(StyleGap(
                        slide_number=slide_num,
                        shape_name=shape.name,
                        category=GapCategory.CHART_FORMATTING,
                        priority=Priority.HIGH,
                        description="Value axis tick marks should be NONE",
                        expected="NONE",
                        actual=str(va.major_tick_mark),
                        auto_fixable=True,
                        fix_action="set_value_tick_none",
                    ))
            except:
                pass

        # Check category axis tick marks and gridlines
        if hasattr(chart, 'category_axis') and chart.category_axis:
            ca = chart.category_axis
            try:
                if ca.major_tick_mark != XL_TICK_MARK.NONE:
                    self.gaps.append(StyleGap(
                        slide_number=slide_num,
                        shape_name=shape.name,
                        category=GapCategory.CHART_FORMATTING,
                        priority=Priority.HIGH,
                        description="Category axis tick marks should be NONE",
                        expected="NONE",
                        actual=str(ca.major_tick_mark),
                        auto_fixable=True,
                        fix_action="set_category_tick_none",
                    ))
            except:
                pass

            # Check category axis gridlines (per PCCP style guide, both axes should have gridlines)
            try:
                if not ca.has_major_gridlines:
                    self.gaps.append(StyleGap(
                        slide_number=slide_num,
                        shape_name=shape.name,
                        category=GapCategory.CHART_FORMATTING,
                        priority=Priority.MEDIUM,
                        description="Missing major vertical gridlines on category axis",
                        expected="ON",
                        actual="OFF",
                        auto_fixable=True,
                        fix_action="enable_category_gridlines",
                    ))
                else:
                    # Check gridline formatting
                    gridlines = ca.major_gridlines
                    line = gridlines.format.line
                    width_pt = self._emu_to_pt(line.width) if line.width else 0
                    if abs(width_pt - self.spec.chart_gridline_width) > 0.1:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Category axis gridline width incorrect",
                            expected=f"{self.spec.chart_gridline_width}pt",
                            actual=f"{width_pt}pt" if width_pt else "not set",
                            auto_fixable=True,
                            fix_action="set_category_gridline_format",
                        ))
            except:
                pass

        # Check chart title font
        if hasattr(chart, 'chart_title') and chart.chart_title:
            try:
                title = chart.chart_title
                if title.has_text_frame:
                    for para in title.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                size_pt = self._emu_to_pt(run.font.size)
                                if abs(size_pt - self.spec.chart.title_font_size_pt) > 0.5:
                                    self.gaps.append(StyleGap(
                                        slide_number=slide_num,
                                        shape_name=shape.name,
                                        category=GapCategory.CHART_FORMATTING,
                                        priority=Priority.HIGH,
                                        description="Chart title font size incorrect",
                                        expected=f"{self.spec.chart.title_font_size_pt}pt",
                                        actual=f"{size_pt}pt",
                                        auto_fixable=True,
                                        fix_action="set_chart_title_font_size",
                                    ))
                            if run.font.bold is not None and run.font.bold != self.spec.chart.title_font_bold:
                                self.gaps.append(StyleGap(
                                    slide_number=slide_num,
                                    shape_name=shape.name,
                                    category=GapCategory.CHART_FORMATTING,
                                    priority=Priority.MEDIUM,
                                    description="Chart title bold setting incorrect",
                                    expected="Bold" if self.spec.chart.title_font_bold else "Not Bold",
                                    actual="Bold" if run.font.bold else "Not Bold",
                                    auto_fixable=True,
                                    fix_action="set_chart_title_font_bold",
                                ))
                            break
                        break
            except:
                pass

        # Check axis label fonts (value axis) - flag if not set or wrong
        if hasattr(chart, 'value_axis') and chart.value_axis:
            try:
                va = chart.value_axis
                if hasattr(va, 'tick_labels') and va.tick_labels:
                    font = va.tick_labels.font
                    size_pt = self._emu_to_pt(font.size) if font.size else None
                    if size_pt is None or abs(size_pt - self.spec.chart.axis_label_font_size_pt) > 0.5:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Value axis label font size incorrect",
                            expected=f"{self.spec.chart.axis_label_font_size_pt}pt",
                            actual=f"{size_pt}pt" if size_pt else "not set",
                            auto_fixable=True,
                            fix_action="set_chart_fonts",
                        ))
            except:
                pass

        # Check axis label fonts (category axis) - flag if not set or wrong
        if hasattr(chart, 'category_axis') and chart.category_axis:
            try:
                ca = chart.category_axis
                if hasattr(ca, 'tick_labels') and ca.tick_labels:
                    font = ca.tick_labels.font
                    size_pt = self._emu_to_pt(font.size) if font.size else None
                    if size_pt is None or abs(size_pt - self.spec.chart.axis_label_font_size_pt) > 0.5:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Category axis label font size incorrect",
                            expected=f"{self.spec.chart.axis_label_font_size_pt}pt",
                            actual=f"{size_pt}pt" if size_pt else "not set",
                            auto_fixable=True,
                            fix_action="set_chart_fonts",
                        ))
            except:
                pass

        # Check legend font - flag if not set or wrong
        if hasattr(chart, 'legend') and chart.legend:
            try:
                legend = chart.legend
                if hasattr(legend, 'font') and legend.font:
                    font = legend.font
                    size_pt = self._emu_to_pt(font.size) if font.size else None
                    if size_pt is None or abs(size_pt - self.spec.chart.legend_font_size_pt) > 0.5:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.CHART_FORMATTING,
                            priority=Priority.HIGH,
                            description="Legend font size incorrect",
                            expected=f"{self.spec.chart.legend_font_size_pt}pt",
                            actual=f"{size_pt}pt" if size_pt else "not set",
                            auto_fixable=True,
                            fix_action="set_chart_fonts",
                        ))
            except:
                pass

        # Check bar fill color (for bar/column charts)
        try:
            from pptx.enum.chart import XL_CHART_TYPE
            chart_type = chart.chart_type
            # Check if this is a bar or column chart
            if chart_type in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED,
                              XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED,
                              XL_CHART_TYPE.BAR_STACKED_100, XL_CHART_TYPE.COLUMN_STACKED_100):
                # Check the first series fill color
                if chart.series and len(chart.series) > 0:
                    series = chart.series[0]
                    if hasattr(series, 'format') and series.format.fill:
                        fill = series.format.fill
                        if hasattr(fill, 'fore_color') and fill.fore_color:
                            color = self._color_to_hex(fill.fore_color)
                            if color and not self._colors_match(color, self.spec.chart.bar_fill_color):
                                self.gaps.append(StyleGap(
                                    slide_number=slide_num,
                                    shape_name=shape.name,
                                    category=GapCategory.CHART_FORMATTING,
                                    priority=Priority.HIGH,
                                    description="Bar/column fill color incorrect",
                                    expected=self.spec.chart.bar_fill_color,
                                    actual=color,
                                    auto_fixable=True,
                                    fix_action="set_bar_fill_color",
                                ))
        except:
            pass

        # Check data label font
        try:
            for series in chart.series:
                if hasattr(series, 'data_labels') and series.data_labels:
                    dl = series.data_labels
                    if hasattr(dl, 'font') and dl.font:
                        font = dl.font
                        if font.size:
                            size_pt = self._emu_to_pt(font.size)
                            if abs(size_pt - self.spec.chart.data_label_font_size_pt) > 0.5:
                                self.gaps.append(StyleGap(
                                    slide_number=slide_num,
                                    shape_name=shape.name,
                                    category=GapCategory.CHART_FORMATTING,
                                    priority=Priority.HIGH,
                                    description="Data label font size incorrect",
                                    expected=f"{self.spec.chart.data_label_font_size_pt}pt",
                                    actual=f"{size_pt}pt",
                                    auto_fixable=True,
                                    fix_action="set_data_label_font_size",
                                ))
                                break
                    break
        except:
            pass

    def _check_table(self, shape, slide_num: int):
        """Check table formatting compliance."""
        table = shape.table

        # Check header row color
        if table.rows:
            header_row = table.rows[0]
            for cell in header_row.cells:
                try:
                    if cell.fill and cell.fill.type:
                        color = self._color_to_hex(cell.fill.fore_color)
                        if color and not self._colors_match(color, self.spec.table_header_color):
                            self.gaps.append(StyleGap(
                                slide_number=slide_num,
                                shape_name=shape.name,
                                category=GapCategory.TABLE_FORMATTING,
                                priority=Priority.HIGH,
                                description="Table header color incorrect",
                                expected=self.spec.table_header_color,
                                actual=color,
                                auto_fixable=True,
                                fix_action="set_header_color",
                            ))
                        break  # Only check first cell
                except:
                    pass

        # Check header row text formatting
        if table.rows:
            header_row = table.rows[0]
            for cell in header_row.cells:
                try:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    # Check font size
                                    if run.font.size:
                                        size_pt = self._emu_to_pt(run.font.size)
                                        if abs(size_pt - self.spec.table.header_font_size_pt) > 0.5:
                                            self.gaps.append(StyleGap(
                                                slide_number=slide_num,
                                                shape_name=shape.name,
                                                category=GapCategory.TABLE_FORMATTING,
                                                priority=Priority.HIGH,
                                                description="Table header font size incorrect",
                                                expected=f"{self.spec.table.header_font_size_pt}pt",
                                                actual=f"{size_pt}pt",
                                                auto_fixable=True,
                                                fix_action="set_header_font_size",
                                            ))
                                    # Check bold
                                    if run.font.bold is not None and run.font.bold != self.spec.table.header_font_bold:
                                        self.gaps.append(StyleGap(
                                            slide_number=slide_num,
                                            shape_name=shape.name,
                                            category=GapCategory.TABLE_FORMATTING,
                                            priority=Priority.MEDIUM,
                                            description="Table header bold setting incorrect",
                                            expected="Bold" if self.spec.table.header_font_bold else "Not Bold",
                                            actual="Bold" if run.font.bold else "Not Bold",
                                            auto_fixable=True,
                                            fix_action="set_header_font_bold",
                                        ))
                                    # Check text color
                                    color = self._color_to_hex(run.font.color)
                                    if color and not self._colors_match(color, self.spec.table.header_text_color, tolerance=10):
                                        self.gaps.append(StyleGap(
                                            slide_number=slide_num,
                                            shape_name=shape.name,
                                            category=GapCategory.TABLE_FORMATTING,
                                            priority=Priority.HIGH,
                                            description="Table header text color incorrect",
                                            expected=self.spec.table.header_text_color,
                                            actual=color,
                                            auto_fixable=True,
                                            fix_action="set_header_text_color",
                                        ))
                                    break
                            break
                    break
                except:
                    pass

        # Check row text formatting (non-header rows)
        for row_idx in range(1, min(len(table.rows), 5)):  # Check first few rows
            row = table.rows[row_idx]
            for cell in row.cells:
                try:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    # Check font size
                                    if run.font.size:
                                        size_pt = self._emu_to_pt(run.font.size)
                                        if abs(size_pt - self.spec.table.row_font_size_pt) > 0.5:
                                            self.gaps.append(StyleGap(
                                                slide_number=slide_num,
                                                shape_name=shape.name,
                                                category=GapCategory.TABLE_FORMATTING,
                                                priority=Priority.HIGH,
                                                description=f"Table row {row_idx} font size incorrect",
                                                expected=f"{self.spec.table.row_font_size_pt}pt",
                                                actual=f"{size_pt}pt",
                                                auto_fixable=True,
                                                fix_action="set_row_font_size",
                                            ))
                                    # Check text color
                                    color = self._color_to_hex(run.font.color)
                                    if color and not self._colors_match(color, self.spec.table.row_text_color, tolerance=10):
                                        self.gaps.append(StyleGap(
                                            slide_number=slide_num,
                                            shape_name=shape.name,
                                            category=GapCategory.TABLE_FORMATTING,
                                            priority=Priority.MEDIUM,
                                            description=f"Table row {row_idx} text color incorrect",
                                            expected=self.spec.table.row_text_color,
                                            actual=color,
                                            auto_fixable=True,
                                            fix_action="set_row_text_color",
                                        ))
                                    break
                            break
                    break
                except:
                    pass
            break  # Only check one row per table

        # Check alternating row colors
        for row_idx in range(1, min(len(table.rows), 10)):
            row = table.rows[row_idx]
            expected_color = self.spec.table_row_odd_color if row_idx % 2 == 1 else self.spec.table_row_even_color

            for cell in row.cells:
                try:
                    if cell.fill and cell.fill.type:
                        color = self._color_to_hex(cell.fill.fore_color)
                        if color and not self._colors_match(color, expected_color, tolerance=10):
                            self.gaps.append(StyleGap(
                                slide_number=slide_num,
                                shape_name=shape.name,
                                category=GapCategory.TABLE_FORMATTING,
                                priority=Priority.MEDIUM,
                                description=f"Row {row_idx} color incorrect",
                                expected=expected_color,
                                actual=color,
                                auto_fixable=True,
                                fix_action="set_row_color",
                            ))
                        break  # Only check first cell
                except:
                    pass

        # Check cell margins
        if table.rows:
            cell = table.rows[0].cells[0]
            margin_lr = self._emu_to_inches(cell.margin_left)
            margin_tb = self._emu_to_inches(cell.margin_top)

            if abs(margin_lr - self.spec.table_cell_margin_lr) > 0.02:
                self.gaps.append(StyleGap(
                    slide_number=slide_num,
                    shape_name=shape.name,
                    category=GapCategory.TABLE_FORMATTING,
                    priority=Priority.LOW,
                    description="Cell left/right margin incorrect",
                    expected=f"{self.spec.table_cell_margin_lr}\"",
                    actual=f"{margin_lr}\"",
                    auto_fixable=True,
                    fix_action="set_cell_margins",
                ))

    def _get_spec_for_placeholder(self, ph_type: Optional[str], ph_idx: Optional[int],
                                    is_section: bool = False) -> Optional[TypographySpec]:
        """Get the typography spec for a placeholder type/index."""
        # Check by placeholder index first (more specific)
        if ph_idx is not None:
            if ph_idx == self.spec.subtitle.placeholder_idx:
                return self.spec.subtitle
            if ph_idx == self.spec.title_content.placeholder_idx:
                # Use section or content title spec based on slide type
                return self.spec.title_section if is_section else self.spec.title_content
            if ph_idx == self.spec.section_name.placeholder_idx:
                return self.spec.section_name
            if ph_idx == self.spec.footnote.placeholder_idx:
                return self.spec.footnote

        # Check by placeholder type - IMPORTANT: Check SUBTITLE before TITLE since "TITLE" is substring of "SUBTITLE"
        if ph_type:
            if "SUBTITLE" in ph_type:
                return self.spec.subtitle
            if "TITLE" in ph_type and "SUBTITLE" not in ph_type:
                return self.spec.title_section if is_section else self.spec.title_content
            if self.spec.body.placeholder_types:
                for pt in self.spec.body.placeholder_types:
                    if pt in ph_type:
                        return self.spec.body

        return None

    def _get_element_name(self, ph_type: Optional[str], ph_idx: Optional[int],
                          is_section: bool = False) -> str:
        """Get human-readable name for a placeholder."""
        # IMPORTANT: Check SUBTITLE before TITLE since "TITLE" is substring of "SUBTITLE"
        if ph_idx == self.spec.subtitle.placeholder_idx or (ph_type and "SUBTITLE" in ph_type):
            return "Subtitle"
        if ph_idx == self.spec.title_content.placeholder_idx or (ph_type and "TITLE" in ph_type and "SUBTITLE" not in ph_type):
            return "Section Title" if is_section else "Title"
        if ph_idx == self.spec.section_name.placeholder_idx:
            return "Section Name"
        if ph_idx == self.spec.footnote.placeholder_idx:
            return "Footnote"
        if ph_type and "BODY" in ph_type:
            return "Body Text"
        return "Text"

    def _get_effective_font_size(self, run, para) -> Optional[float]:
        """Get effective font size, checking run then paragraph default."""
        if run.font.size:
            return self._emu_to_pt(run.font.size)
        # Check paragraph default font
        try:
            if para.font and para.font.size:
                return self._emu_to_pt(para.font.size)
        except:
            pass
        return None

    def _is_bullet_list_textbox(self, tf) -> bool:
        """Check if a text frame contains a bulleted list (multiple paragraphs with content)."""
        non_empty_paras = [p for p in tf.paragraphs if p.text.strip()]
        # Consider it a bullet list if it has 2+ non-empty paragraphs
        return len(non_empty_paras) >= 2

    def _check_background_image(self, shape, slide_num: int):
        """Check background image for aspect ratio stretching.

        For 16:9 images on letter-size slides, checks if srcRect cropping
        is applied to prevent vertical stretching.
        """
        try:
            from PIL import Image
            import io

            # Get the image blob and calculate aspect ratio
            image_blob = shape.image.blob
            img = Image.open(io.BytesIO(image_blob))
            orig_width, orig_height = img.size
            image_aspect = orig_width / orig_height
            slide_aspect = 11.0 / 8.5  # Letter size = 1.294

            # If image is 16:9 (1.778) or wider, it should have srcRect cropping
            if image_aspect > slide_aspect + 0.05:
                # Check if srcRect cropping is applied
                pic_elem = shape._element
                has_crop = False

                # Look for srcRect in blipFill
                for child in pic_elem:
                    if child.tag.endswith('}blipFill'):
                        for subchild in child:
                            if subchild.tag.endswith('}srcRect'):
                                # Check if left/right crop values are set
                                l_crop = subchild.get('l', '0')
                                r_crop = subchild.get('r', '0')
                                if int(l_crop) > 0 or int(r_crop) > 0:
                                    has_crop = True
                                break
                        break

                if not has_crop:
                    expected_crop = (1 - (slide_aspect / image_aspect)) / 2 * 100
                    self.gaps.append(StyleGap(
                        slide_number=slide_num,
                        shape_name=shape.name or "Background Image",
                        category=GapCategory.IMAGE_FORMATTING,
                        priority=Priority.MEDIUM,
                        description="16:9 background image stretched to letter-size without aspect ratio correction",
                        expected=f"srcRect cropping of ~{expected_crop:.1f}% on each side",
                        actual="No srcRect cropping applied",
                        auto_fixable=True,
                        fix_action="fix_background_image_aspect_ratio"
                    ))
        except Exception as e:
            pass  # Silently skip if we can't check the image

    def _check_text(self, shape, slide_num: int, ph_type: Optional[str] = None,
                    ph_idx: Optional[int] = None, is_section: bool = False,
                    is_front_page: bool = False, is_side_by_side: bool = False):
        """Check text formatting compliance.

        Args:
            shape: The shape to check
            slide_num: Slide number (1-indexed)
            ph_type: Placeholder type string
            ph_idx: Placeholder index
            is_section: Whether this is a section divider slide
            is_front_page: Whether this is the front page (slide 1)
            is_side_by_side: Whether this is a side-by-side layout slide
        """
        tf = shape.text_frame

        # Get the expected typography spec for this placeholder
        typo_spec = self._get_spec_for_placeholder(ph_type, ph_idx, is_section)
        element_name = self._get_element_name(ph_type, ph_idx, is_section)

        # Front page / section slide subtitle uses special styling (18pt bold white)
        # Note: Front page subtitle is typically BODY placeholder, not SUBTITLE placeholder
        # This also applies to section slides (e.g., "Thank You" slides with Frontpage layout)
        if (is_front_page or is_section) and ph_type and ("SUBTITLE" in ph_type or "BODY" in ph_type):
            # Check if this is the subtitle area (not the title)
            if "TITLE" not in element_name:
                typo_spec = self.spec.subtitle_frontpage
                element_name = "Section/Front Page Subtitle"

        # Special handling for side-by-side bullet lists ONLY
        # First paragraph gets content_header style (18pt bold, no bullet indent)
        is_bullet_list = self._is_bullet_list_textbox(tf)
        first_para_checked = False

        if typo_spec:
            # Check typography for each paragraph
            for para_idx, para in enumerate(tf.paragraphs):
                if not para.text.strip():
                    continue

                # ONLY for side-by-side layouts with bullet lists, first paragraph uses content_header
                if is_side_by_side and is_bullet_list and ph_type and any(pt in ph_type for pt in ["BODY", "OBJECT", "CONTENT"]):
                    if not first_para_checked:
                        # First paragraph should use content_header style (18pt bold)
                        current_spec = self.spec.content_header
                        current_element = "Side-by-Side Header"
                        first_para_checked = True
                    else:
                        # Remaining paragraphs use body style (14pt regular)
                        current_spec = self.spec.body
                        current_element = "Bullet Text"
                else:
                    current_spec = typo_spec
                    current_element = element_name

                # Determine fix action based on context
                if is_front_page and "Front Page Subtitle" in current_element:
                    fix_action_base = "set_frontpage_subtitle"
                elif is_side_by_side and is_bullet_list:
                    fix_action_base = "set_side_by_side_fonts"
                else:
                    fix_action_base = f"set_font:{element_name.lower().replace(' ', '_')}"

                # Check each run for font properties
                for run in para.runs:
                    if not run.text.strip():
                        continue

                    # Check font name
                    if run.font.name and run.font.name != current_spec.font_name:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.TYPOGRAPHY,
                            priority=Priority.HIGH,
                            description=f"{current_element} font name incorrect",
                            expected=current_spec.font_name,
                            actual=run.font.name,
                            auto_fixable=True,
                            fix_action=fix_action_base,
                        ))
                        break  # One gap per paragraph

                    # Check font size - use effective size (run or paragraph level)
                    size_pt = self._get_effective_font_size(run, para)
                    if size_pt is not None:
                        if abs(size_pt - current_spec.font_size_pt) > 0.5:
                            self.gaps.append(StyleGap(
                                slide_number=slide_num,
                                shape_name=shape.name,
                                category=GapCategory.TYPOGRAPHY,
                                priority=Priority.HIGH,
                                description=f"{current_element} font size incorrect",
                                expected=f"{current_spec.font_size_pt}pt",
                                actual=f"{size_pt}pt",
                                auto_fixable=True,
                                fix_action=fix_action_base,
                            ))
                            break

                    # Check font bold
                    if run.font.bold is not None and run.font.bold != current_spec.font_bold:
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.TYPOGRAPHY,
                            priority=Priority.MEDIUM,
                            description=f"{current_element} bold setting incorrect",
                            expected="Bold" if current_spec.font_bold else "Not Bold",
                            actual="Bold" if run.font.bold else "Not Bold",
                            auto_fixable=True,
                            fix_action=fix_action_base,
                        ))
                        break

                    # Check font color
                    color = self._color_to_hex(run.font.color)
                    if color and not self._colors_match(color, current_spec.font_color, tolerance=10):
                        self.gaps.append(StyleGap(
                            slide_number=slide_num,
                            shape_name=shape.name,
                            category=GapCategory.TYPOGRAPHY,
                            priority=Priority.MEDIUM,
                            description=f"{current_element} font color incorrect",
                            expected=current_spec.font_color,
                            actual=color,
                            auto_fixable=True,
                            fix_action=fix_action_base,
                        ))
                        break

        # Check alignment for elements with alignment requirements (e.g., footnotes)
        if typo_spec and typo_spec.horizontal_alignment:
            # Check horizontal alignment on first paragraph
            for para in tf.paragraphs:
                if para.text.strip():
                    try:
                        actual_align = para.alignment
                        if actual_align is not None:
                            expected_align = getattr(PP_ALIGN, typo_spec.horizontal_alignment, None)
                            if expected_align and actual_align != expected_align:
                                actual_name = actual_align.name if hasattr(actual_align, 'name') else str(actual_align)
                                self.gaps.append(StyleGap(
                                    slide_number=slide_num,
                                    shape_name=shape.name,
                                    category=GapCategory.TYPOGRAPHY,
                                    priority=Priority.MEDIUM,
                                    description=f"{element_name} horizontal alignment incorrect",
                                    expected=typo_spec.horizontal_alignment,
                                    actual=actual_name,
                                    auto_fixable=True,
                                    fix_action="set_footnote_alignment",
                                ))
                    except:
                        pass
                    break  # Only check first non-empty paragraph

        if typo_spec and typo_spec.vertical_anchor:
            # Check vertical anchor on text frame
            try:
                # python-pptx exposes text frame's vertical anchor via .anchor property
                # but only on shapes that support it
                if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, '_txBody'):
                    # Try to get the vertical anchor from the XML
                    tx_body = shape.text_frame._txBody
                    body_pr = tx_body.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
                    if body_pr is not None:
                        anchor_attr = body_pr.get('anchor')
                        if anchor_attr:
                            # Map XML anchor values to expected values
                            anchor_map = {'t': 'TOP', 'ctr': 'MIDDLE', 'b': 'BOTTOM'}
                            actual_anchor = anchor_map.get(anchor_attr, anchor_attr.upper())
                            if actual_anchor != typo_spec.vertical_anchor:
                                self.gaps.append(StyleGap(
                                    slide_number=slide_num,
                                    shape_name=shape.name,
                                    category=GapCategory.TYPOGRAPHY,
                                    priority=Priority.MEDIUM,
                                    description=f"{element_name} vertical anchor incorrect",
                                    expected=typo_spec.vertical_anchor,
                                    actual=actual_anchor,
                                    auto_fixable=True,
                                    fix_action="set_footnote_alignment",
                                ))
            except:
                pass

        # Check for bullet formatting (ALL CAPS + em-dash pattern)
        for para in tf.paragraphs:
            text = para.text.strip()
            if text.startswith('•') or text.startswith('-'):
                # Check for ALL CAPS header + em-dash pattern
                if '—' not in text and ' - ' not in text:
                    # Might be missing the proper format
                    if not re.match(r'^[•\-]\s*[A-Z][A-Z\s]+\s*[—–-]', text):
                        # Check if this looks like a bullet that should have format
                        if len(text) > 20 and not text.startswith('• Source'):
                            self.gaps.append(StyleGap(
                                slide_number=slide_num,
                                shape_name=shape.name,
                                category=GapCategory.CONTENT_FORMAT,
                                priority=Priority.MEDIUM,
                                description="Bullet missing ALL CAPS + em-dash format",
                                expected="• CATEGORY — Description",
                                actual=text[:50] + "..." if len(text) > 50 else text,
                                auto_fixable=False,
                            ))


class PresentationCorrector:
    """Applies corrections to fix style gaps."""

    def __init__(self, spec: StyleGuideSpec):
        self.spec = spec

    def _hex_to_rgb_color(self, hex_color: str) -> RGBColor:
        """Convert hex to RGBColor."""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _get_typography_spec(self, element_type: str) -> Optional[TypographySpec]:
        """Get typography spec by element type."""
        specs = {
            "title": self.spec.title_content,
            "title_content": self.spec.title_content,
            "section_title": self.spec.title_section,
            "title_section": self.spec.title_section,
            "subtitle": self.spec.subtitle,
            "subtitle_frontpage": self.spec.subtitle_frontpage,
            "front_page_subtitle": self.spec.subtitle_frontpage,
            "body_text": self.spec.body,
            "text": self.spec.body,
            "content_header": self.spec.content_header,
            "side-by-side_header": self.spec.content_header,
            "section_name": self.spec.section_name,
            "footnote": self.spec.footnote,
        }
        return specs.get(element_type)

    def _apply_typography_to_shape(self, shape, typo_spec: TypographySpec,
                                   apply_name: bool = False, apply_size: bool = False,
                                   apply_bold: bool = False, apply_color: bool = False):
        """Apply typography settings to all text in a shape."""
        if not shape.has_text_frame:
            return

        for para in shape.text_frame.paragraphs:
            # Apply to paragraph default font as well
            if apply_size and typo_spec.font_size_pt:
                try:
                    para.font.size = Pt(typo_spec.font_size_pt)
                except:
                    pass
            for run in para.runs:
                if apply_name and typo_spec.font_name:
                    run.font.name = typo_spec.font_name
                if apply_size and typo_spec.font_size_pt:
                    run.font.size = Pt(typo_spec.font_size_pt)
                if apply_bold:
                    run.font.bold = typo_spec.font_bold
                if apply_color and typo_spec.font_color:
                    run.font.color.rgb = self._hex_to_rgb_color(typo_spec.font_color)

    def apply_corrections(self, pptx_path: Path, gap_analysis: GapAnalysis) -> Tuple[int, int]:
        """
        Apply automatic corrections to the presentation.

        Returns:
            Tuple of (corrections_applied, corrections_failed)
        """
        prs = Presentation(str(pptx_path))
        applied = 0
        failed = 0

        # Group gaps by slide for efficiency
        by_slide = {}
        for gap in gap_analysis.gaps:
            if gap.auto_fixable:
                if gap.slide_number not in by_slide:
                    by_slide[gap.slide_number] = []
                by_slide[gap.slide_number].append(gap)

        for slide_num, gaps in by_slide.items():
            if slide_num == 0:
                continue  # Skip presentation-level gaps

            slide = prs.slides[slide_num - 1]

            for gap in gaps:
                try:
                    # Find the shape
                    shape = None
                    for s in slide.shapes:
                        if s.name == gap.shape_name:
                            shape = s
                            break

                    if shape is None:
                        failed += 1
                        continue

                    # Apply the correction
                    if gap.fix_action == "enable_major_gridlines":
                        if shape.has_chart:
                            shape.chart.value_axis.has_major_gridlines = True
                            applied += 1

                    elif gap.fix_action == "set_gridline_width":
                        if shape.has_chart and shape.chart.value_axis.has_major_gridlines:
                            gridlines = shape.chart.value_axis.major_gridlines
                            gridlines.format.line.width = Pt(self.spec.chart_gridline_width)
                            applied += 1

                    elif gap.fix_action == "set_gridline_color":
                        if shape.has_chart and shape.chart.value_axis.has_major_gridlines:
                            gridlines = shape.chart.value_axis.major_gridlines
                            gridlines.format.line.color.rgb = self._hex_to_rgb_color(
                                self.spec.chart_gridline_color
                            )
                            applied += 1

                    elif gap.fix_action == "set_gridline_format":
                        # Combined gridline format - set both width and color
                        if shape.has_chart and shape.chart.value_axis:
                            va = shape.chart.value_axis
                            if va.has_major_gridlines:
                                gridlines = va.major_gridlines
                                gridlines.format.line.width = Pt(self.spec.chart.gridline_width_pt)
                                gridlines.format.line.color.rgb = self._hex_to_rgb_color(
                                    self.spec.chart.gridline_color
                                )
                            applied += 1

                    elif gap.fix_action == "set_value_tick_none":
                        if shape.has_chart:
                            shape.chart.value_axis.major_tick_mark = XL_TICK_MARK.NONE
                            shape.chart.value_axis.minor_tick_mark = XL_TICK_MARK.NONE
                            applied += 1

                    elif gap.fix_action == "set_category_tick_none":
                        if shape.has_chart:
                            shape.chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
                            shape.chart.category_axis.minor_tick_mark = XL_TICK_MARK.NONE
                            applied += 1

                    elif gap.fix_action == "set_header_color":
                        if shape.has_table:
                            for cell in shape.table.rows[0].cells:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = self._hex_to_rgb_color(
                                    self.spec.table_header_color
                                )
                            applied += 1

                    elif gap.fix_action == "set_row_color":
                        if shape.has_table:
                            for row_idx in range(1, len(shape.table.rows)):
                                color = (self.spec.table_row_odd_color if row_idx % 2 == 1
                                        else self.spec.table_row_even_color)
                                for cell in shape.table.rows[row_idx].cells:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = self._hex_to_rgb_color(color)
                            applied += 1

                    elif gap.fix_action == "set_cell_margins":
                        if shape.has_table:
                            margin_lr = Inches(self.spec.table_cell_margin_lr)
                            margin_tb = Inches(self.spec.table_cell_margin_tb)
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    cell.margin_left = margin_lr
                                    cell.margin_right = margin_lr
                                    cell.margin_top = margin_tb
                                    cell.margin_bottom = margin_tb
                            applied += 1

                    # Typography corrections - apply ALL typography properties (size, bold, color) at once
                    elif gap.fix_action and gap.fix_action.startswith("set_font_"):
                        element_type = gap.fix_action.split(":")[1] if ":" in gap.fix_action else None
                        if element_type:
                            typo_spec = self._get_typography_spec(element_type)
                            if typo_spec and shape.has_text_frame:
                                # Apply all typography properties for consistency
                                self._apply_typography_to_shape(
                                    shape, typo_spec,
                                    apply_name=True,
                                    apply_size=True,
                                    apply_bold=True,
                                    apply_color=True
                                )
                                applied += 1

                    # Chart font corrections
                    elif gap.fix_action == "set_chart_title_font_size":
                        if shape.has_chart and shape.chart.chart_title:
                            title = shape.chart.chart_title
                            if title.has_text_frame:
                                for para in title.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.size = Pt(self.spec.chart.title_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_chart_title_font_bold":
                        if shape.has_chart and shape.chart.chart_title:
                            title = shape.chart.chart_title
                            if title.has_text_frame:
                                for para in title.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.bold = self.spec.chart.title_font_bold
                            applied += 1

                    elif gap.fix_action == "set_value_axis_font_size":
                        if shape.has_chart and shape.chart.value_axis:
                            va = shape.chart.value_axis
                            if hasattr(va, 'tick_labels') and va.tick_labels:
                                va.tick_labels.font.size = Pt(self.spec.chart.axis_label_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_category_axis_font_size":
                        if shape.has_chart and shape.chart.category_axis:
                            ca = shape.chart.category_axis
                            if hasattr(ca, 'tick_labels') and ca.tick_labels:
                                ca.tick_labels.font.size = Pt(self.spec.chart.axis_label_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_legend_font_size":
                        if shape.has_chart and shape.chart.legend:
                            legend = shape.chart.legend
                            if hasattr(legend, 'font'):
                                legend.font.size = Pt(self.spec.chart.legend_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_data_label_font_size":
                        if shape.has_chart:
                            for series in shape.chart.series:
                                if hasattr(series, 'data_labels') and series.data_labels:
                                    if hasattr(series.data_labels, 'font'):
                                        series.data_labels.font.size = Pt(self.spec.chart.data_label_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_chart_fonts":
                        # Combined chart fonts - set all fonts at once
                        if shape.has_chart:
                            chart = shape.chart
                            font_size = Pt(self.spec.chart.axis_label_font_size_pt)

                            # Set value axis font
                            if chart.value_axis and hasattr(chart.value_axis, 'tick_labels'):
                                chart.value_axis.tick_labels.font.size = font_size

                            # Set category axis font
                            if chart.category_axis and hasattr(chart.category_axis, 'tick_labels'):
                                chart.category_axis.tick_labels.font.size = font_size

                            # Set legend font
                            if chart.legend and hasattr(chart.legend, 'font'):
                                chart.legend.font.size = Pt(self.spec.chart.legend_font_size_pt)

                            # Set data labels font
                            for series in chart.series:
                                if hasattr(series, 'data_labels') and series.data_labels:
                                    if hasattr(series.data_labels, 'font'):
                                        series.data_labels.font.size = Pt(self.spec.chart.data_label_font_size_pt)

                            applied += 1

                    elif gap.fix_action == "set_bar_fill_color":
                        if shape.has_chart:
                            # Set fill color for first series (primary bar color)
                            if shape.chart.series and len(shape.chart.series) > 0:
                                series = shape.chart.series[0]
                                if hasattr(series, 'format'):
                                    series.format.fill.solid()
                                    series.format.fill.fore_color.rgb = self._hex_to_rgb_color(
                                        self.spec.chart.bar_fill_color
                                    )
                            applied += 1

                    elif gap.fix_action == "set_frontpage_subtitle":
                        # Apply front page subtitle style (18pt bold white)
                        if shape.has_text_frame:
                            spec = self.spec.subtitle_frontpage
                            for para in shape.text_frame.paragraphs:
                                try:
                                    para.font.size = Pt(spec.font_size_pt)
                                except:
                                    pass
                                for run in para.runs:
                                    run.font.name = spec.font_name
                                    run.font.size = Pt(spec.font_size_pt)
                                    run.font.bold = spec.font_bold
                                    run.font.color.rgb = self._hex_to_rgb_color(spec.font_color)
                            applied += 1

                    elif gap.fix_action == "set_side_by_side_fonts":
                        # Apply content_header to first paragraph (18pt bold), body to rest (14pt regular)
                        # For side-by-side layouts ONLY
                        if shape.has_text_frame:
                            tf = shape.text_frame
                            first_para_done = False
                            for para in tf.paragraphs:
                                if not para.text.strip():
                                    continue
                                if not first_para_done:
                                    # First paragraph: content_header style (18pt bold)
                                    spec = self.spec.content_header
                                    first_para_done = True
                                else:
                                    # Remaining paragraphs: body style (14pt regular)
                                    spec = self.spec.body
                                # Apply to paragraph default and all runs
                                try:
                                    para.font.size = Pt(spec.font_size_pt)
                                except:
                                    pass
                                for run in para.runs:
                                    run.font.name = spec.font_name
                                    run.font.size = Pt(spec.font_size_pt)
                                    run.font.bold = spec.font_bold
                                    run.font.color.rgb = self._hex_to_rgb_color(spec.font_color)
                            applied += 1

                    # Table font corrections
                    elif gap.fix_action == "set_header_font_size":
                        if shape.has_table:
                            for cell in shape.table.rows[0].cells:
                                if cell.text_frame:
                                    for para in cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            run.font.size = Pt(self.spec.table.header_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_header_font_bold":
                        if shape.has_table:
                            for cell in shape.table.rows[0].cells:
                                if cell.text_frame:
                                    for para in cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            run.font.bold = self.spec.table.header_font_bold
                            applied += 1

                    elif gap.fix_action == "set_header_text_color":
                        if shape.has_table:
                            for cell in shape.table.rows[0].cells:
                                if cell.text_frame:
                                    for para in cell.text_frame.paragraphs:
                                        for run in para.runs:
                                            run.font.color.rgb = self._hex_to_rgb_color(
                                                self.spec.table.header_text_color
                                            )
                            applied += 1

                    elif gap.fix_action == "set_row_font_size":
                        if shape.has_table:
                            for row_idx in range(1, len(shape.table.rows)):
                                for cell in shape.table.rows[row_idx].cells:
                                    if cell.text_frame:
                                        for para in cell.text_frame.paragraphs:
                                            for run in para.runs:
                                                run.font.size = Pt(self.spec.table.row_font_size_pt)
                            applied += 1

                    elif gap.fix_action == "set_row_text_color":
                        if shape.has_table:
                            for row_idx in range(1, len(shape.table.rows)):
                                for cell in shape.table.rows[row_idx].cells:
                                    if cell.text_frame:
                                        for para in cell.text_frame.paragraphs:
                                            for run in para.runs:
                                                run.font.color.rgb = self._hex_to_rgb_color(
                                                    self.spec.table.row_text_color
                                                )
                            applied += 1

                    elif gap.fix_action == "fix_background_image_aspect_ratio":
                        # Fix 16:9 image stretched to letter-size slide
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                from PIL import Image
                                from lxml import etree
                                import io

                                # Get image aspect ratio
                                image_blob = shape.image.blob
                                img = Image.open(io.BytesIO(image_blob))
                                orig_width, orig_height = img.size
                                image_aspect = orig_width / orig_height
                                slide_aspect = 11.0 / 8.5  # Letter size

                                if image_aspect > slide_aspect + 0.05:
                                    # Calculate crop percentage
                                    crop_ratio = 1 - (slide_aspect / image_aspect)
                                    crop_pct = int((crop_ratio / 2) * 100000)

                                    # Find blipFill element
                                    pic_elem = shape._element
                                    blipFill = None
                                    for child in pic_elem:
                                        if child.tag.endswith('}blipFill'):
                                            blipFill = child
                                            break

                                    if blipFill is not None:
                                        # Remove existing srcRect
                                        for subchild in list(blipFill):
                                            if subchild.tag.endswith('}srcRect'):
                                                blipFill.remove(subchild)

                                        # Create srcRect with cropping
                                        srcRect = etree.SubElement(
                                            blipFill,
                                            '{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect'
                                        )
                                        srcRect.set('l', str(crop_pct))
                                        srcRect.set('r', str(crop_pct))
                                        srcRect.set('t', '0')
                                        srcRect.set('b', '0')
                                        applied += 1
                                    else:
                                        failed += 1
                                else:
                                    applied += 1  # No fix needed
                            except Exception as e:
                                failed += 1
                        else:
                            failed += 1

                    else:
                        failed += 1

                except Exception as e:
                    failed += 1

        # Save corrected presentation
        corrected_path = pptx_path.parent / f"{pptx_path.stem}_corrected{pptx_path.suffix}"
        prs.save(str(corrected_path))

        return applied, failed


class PresentationReviewer:
    """
    Main class for presentation review workflow.

    Combines style guide loading, analysis, gap reporting, and corrections.
    """

    def __init__(self, style_guide_version: Optional[str] = None):
        """
        Initialize the reviewer.

        Args:
            style_guide_version: Specific version to use (e.g., "2026.01.01").
                               If None, uses the latest available.
        """
        self.loader = StyleGuideLoader()

        if style_guide_version:
            guide_path = self.loader.get_guide_by_version(style_guide_version)
        else:
            guide_path = self.loader.get_latest_guide()

        self.spec = self.loader.parse_guide(guide_path)
        self.analyzer = PresentationAnalyzer(self.spec)
        self.corrector = PresentationCorrector(self.spec)

    def analyze(self, pptx_path: Path) -> GapAnalysis:
        """
        Analyze a presentation for style guide compliance.

        Args:
            pptx_path: Path to the PPTX file.

        Returns:
            GapAnalysis with all findings.
        """
        return self.analyzer.analyze(pptx_path)

    def generate_report(self, gap_analysis: GapAnalysis, output_path: Optional[Path] = None) -> str:
        """
        Generate a markdown gap analysis report.

        Args:
            gap_analysis: The analysis results.
            output_path: Optional path to save the report.

        Returns:
            Markdown report as string.
        """
        report = gap_analysis.to_markdown()

        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(report)

        return report

    def apply_corrections(self, pptx_path: Path, gap_analysis: GapAnalysis) -> Tuple[int, int]:
        """
        Apply automatic corrections based on gap analysis.

        Args:
            pptx_path: Path to the PPTX file.
            gap_analysis: The analysis results.

        Returns:
            Tuple of (corrections_applied, corrections_failed).
        """
        return self.corrector.apply_corrections(pptx_path, gap_analysis)

    def review_and_correct(self, pptx_path: Path, output_dir: Optional[Path] = None) -> Dict[str, Any]:
        """
        Complete review workflow: analyze, report, and correct.

        Args:
            pptx_path: Path to the PPTX file.
            output_dir: Directory for output files. Defaults to same as input.

        Returns:
            Dictionary with analysis results and correction stats.
        """
        pptx_path = Path(pptx_path)
        if output_dir is None:
            output_dir = pptx_path.parent

        # Step 1: Analyze
        gap_analysis = self.analyze(pptx_path)

        # Step 2: Generate report
        report_path = output_dir / f"{pptx_path.stem}_gap_analysis.md"
        self.generate_report(gap_analysis, report_path)

        # Step 3: Apply corrections
        applied, failed = self.apply_corrections(pptx_path, gap_analysis)

        corrected_path = output_dir / f"{pptx_path.stem}_corrected.pptx"

        return {
            "presentation": str(pptx_path),
            "style_guide_version": self.spec.version,
            "total_gaps": len(gap_analysis.gaps),
            "compliance_score": gap_analysis.compliance_score,
            "corrections_applied": applied,
            "corrections_failed": failed,
            "report_path": str(report_path),
            "corrected_path": str(corrected_path) if applied > 0 else None,
        }


def review_presentation(pptx_path: str, style_guide_version: Optional[str] = None) -> Dict[str, Any]:
    """
    Convenience function to review a presentation.

    Args:
        pptx_path: Path to the PPTX file.
        style_guide_version: Optional specific version of style guide.

    Returns:
        Review results dictionary.
    """
    reviewer = PresentationReviewer(style_guide_version)
    return reviewer.review_and_correct(Path(pptx_path))


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python presentation_review.py <pptx_file> [style_guide_version]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    version = sys.argv[2] if len(sys.argv) > 2 else None

    results = review_presentation(pptx_path, version)

    print("\nPresentation Review Results")
    print("=" * 50)
    print(f"Presentation: {results['presentation']}")
    print(f"Style Guide: {results['style_guide_version']}")
    print(f"Total Gaps: {results['total_gaps']}")
    print(f"Compliance Score: {results['compliance_score']:.1f}%")
    print(f"Corrections Applied: {results['corrections_applied']}")
    print(f"Corrections Failed: {results['corrections_failed']}")
    print(f"Report: {results['report_path']}")
    if results['corrected_path']:
        print(f"Corrected: {results['corrected_path']}")
