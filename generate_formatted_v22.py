"""Generate complete presentation using user's manually formatted template.

This script:
1. Uses the user's edited PPTX as template (preserves slide master & layouts)
2. Reads the original outline JSON for all content
3. Applies the extracted formatting specifications to all slides
"""

import json
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================================
# FORMATTING SPECIFICATIONS (extracted from user's edited PPTX)
# ============================================================================

COLORS = {
    'dark_blue': RGBColor(0x05, 0x1C, 0x2C),      # #051C2C - titles, table headers
    'accent_blue': RGBColor(0x30, 0x9C, 0xE7),    # #309CE7 - metric boxes
    'white': RGBColor(0xFF, 0xFF, 0xFF),          # #FFFFFF - text on dark
    'dark_text': RGBColor(0x06, 0x1F, 0x32),      # #061F32 - body text
}

# Layout indices in slide master
LAYOUTS = {
    'frontpage': 0,       # Title slide
    'default': 1,         # Standard content
    '1_default': 2,       # Two content areas
    'top_left_title': 3,
    'section_breaker': 4, # Section dividers
    '1_4_grey': 5,
    '1_3_grey': 6,
    '1_2_grey': 7,        # Two-column comparison
    '3_4_grey': 8,
    '1_1_2_grey': 9,
    'end': 10,
}

# Positioning (in inches)
POS = {
    'title': {'left': 0.40, 'top': 0.40, 'width': 10.20, 'height': 1.00},
    'subtitle': {'left': 0.40, 'top': 1.70, 'width': 10.20, 'height': 0.80},
    'content': {'left': 0.40, 'top': 2.80, 'width': 10.20, 'height': 4.70},
    'footer': {'left': 0.40, 'top': 7.50, 'width': 10.20, 'height': 0.40},
    'table': {'left': 0.40, 'top': 2.80, 'width': 10.20},
    'chart': {'left': 0.40, 'top': 2.80, 'width': 10.20, 'height': 4.00},
}

# Font specifications
FONTS = {
    'title_slide_title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': 'white'},
    'title_slide_subtitle': {'name': 'Arial', 'size': 20, 'bold': True, 'color': 'white'},
    'title_slide_date': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'white'},
    'content_title': {'name': 'Arial', 'size': 36, 'bold': True, 'color': 'dark_text'},
    'content_title_long': {'name': 'Arial', 'size': 32, 'bold': True, 'color': 'dark_text'},
    'subtitle': {'name': 'Arial', 'size': 20, 'bold': True, 'color': 'dark_text'},
    'body': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'dark_text'},
    'body_header': {'name': 'Arial', 'size': 14, 'bold': True, 'color': 'dark_text'},
    'metric_value': {'name': 'Arial', 'size': 28, 'bold': True, 'color': 'white'},
    'metric_label': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'white'},
    'table_header': {'name': 'Arial', 'size': 16, 'bold': True, 'color': 'white'},
    'table_data': {'name': 'Arial', 'size': 14, 'bold': False, 'color': 'dark_text'},
    'section_title': {'name': 'Arial', 'size': 44, 'bold': True, 'color': 'white'},
    'footnote': {'name': 'Arial', 'size': 10, 'bold': False, 'color': 'dark_text'},
}

# Table formatting
TABLE_FORMAT = {
    'header_row_height': 0.50,
    'data_row_height': 0.40,
    'header_fill': 'dark_blue',
    'data_fill_odd': None,  # White/transparent
    'data_fill_even': RGBColor(0xF5, 0xF5, 0xF5),  # Light gray
}


def apply_font(run, font_spec):
    """Apply font specification to a run."""
    font = run.font
    font.name = font_spec['name']
    font.size = Pt(font_spec['size'])
    font.bold = font_spec['bold']
    if font_spec['color'] in COLORS:
        font.color.rgb = COLORS[font_spec['color']]


def format_bullet_text(text):
    """Convert bullet text to header-dash format.

    Input: "HEADER: Some description text"
    Output: ("HEADER", "Some description text")
    """
    # Check for colon separator
    if ': ' in text:
        parts = text.split(': ', 1)
        return parts[0], parts[1] if len(parts) > 1 else ''
    # Check for dash separator
    if ' - ' in text:
        parts = text.split(' - ', 1)
        return parts[0], parts[1] if len(parts) > 1 else ''
    return None, text


def add_formatted_paragraph(text_frame, text, font_spec, level=0, is_bullet=False):
    """Add a paragraph with proper formatting."""
    # Check if we need to use existing paragraph or add new
    if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
        para = text_frame.paragraphs[0]
    else:
        para = text_frame.add_paragraph()

    para.level = level

    if is_bullet:
        # Format as "HEADER - rest of text"
        header, body = format_bullet_text(text)
        if header:
            # Add bold header
            run = para.add_run()
            run.text = header + ' - '
            apply_font(run, {**font_spec, 'bold': True})
            # Add normal body
            if body:
                run = para.add_run()
                run.text = body
                apply_font(run, {**font_spec, 'bold': False})
        else:
            run = para.add_run()
            run.text = text
            apply_font(run, font_spec)
    else:
        run = para.add_run()
        run.text = text
        apply_font(run, font_spec)

    # Remove bullet formatting via XML
    try:
        pPr = para._p.get_or_add_pPr()
        # Add buNone to disable bullets
        buNone = etree.SubElement(pPr, qn('a:buNone'))
    except:
        pass

    return para


def format_table(table, has_header=True):
    """Apply standard table formatting."""
    for row_idx, row in enumerate(table.rows):
        # Set row height
        if row_idx == 0 and has_header:
            row.height = Inches(TABLE_FORMAT['header_row_height'])
        else:
            row.height = Inches(TABLE_FORMAT['data_row_height'])

        for cell in row.cells:
            # Apply cell fill
            if row_idx == 0 and has_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['dark_blue']
            elif row_idx % 2 == 1:
                # Odd rows (1, 3, 5...) - white
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                # Even rows (2, 4, 6...) - light gray
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

            # Format text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if row_idx == 0 and has_header:
                        apply_font(run, FONTS['table_header'])
                    else:
                        apply_font(run, FONTS['table_data'])


def create_metric_box(slide, left, top, value, label):
    """Create a metric box (KPI card)."""
    shape = slide.shapes.add_shape(
        1,  # Rounded rectangle
        Inches(left), Inches(top),
        Inches(2.40), Inches(1.20)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['accent_blue']
    shape.line.fill.background()  # No border

    # Add text
    tf = shape.text_frame
    tf.word_wrap = True

    # Value paragraph
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = value
    apply_font(run, FONTS['metric_value'])

    # Label paragraph
    para = tf.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    apply_font(run, FONTS['metric_label'])

    return shape


def get_slide_content_from_outline(outline, slide_idx):
    """Extract slide content from outline by flattening sections."""
    slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            slides.append({
                'section': section.get('name', ''),
                **slide
            })

    if slide_idx < len(slides):
        return slides[slide_idx]
    return None


def main():
    """Generate the formatted presentation."""
    print("Generating formatted presentation v22...")
    print("=" * 60)

    # Load the user's edited PPTX as template
    template_path = Path("pptx_generator/output/Light_Industrial_Thesis_v21_CS_edits_v2.pptx")
    if not template_path.exists():
        print(f"Error: Template not found: {template_path}")
        return

    # Load the original outline for content
    outline_path = Path("pptx_generator/output/light_industrial_thesis_v18.json")
    if not outline_path.exists():
        print(f"Error: Outline not found: {outline_path}")
        return

    with open(outline_path) as f:
        outline = json.load(f)

    # Open the template
    prs = Presentation(str(template_path))

    print(f"Template loaded: {len(prs.slides)} slides")
    print(f"Outline has {sum(len(s.get('slides', [])) for s in outline.get('sections', []))} slides")

    # Get all slides from outline
    all_slides = []
    for section in outline.get('sections', []):
        for slide in section.get('slides', []):
            all_slides.append({
                'section': section.get('name', ''),
                **slide
            })

    # The user's template has 20 slides
    # We need to add slides 21-45 from the outline (indices 20-44)

    existing_slide_count = len(prs.slides)
    print(f"\nAdding slides {existing_slide_count + 1} to {len(all_slides)}...")

    # Get the default layout for new slides
    default_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['default']]
    section_layout = prs.slide_masters[0].slide_layouts[LAYOUTS['section_breaker']]

    # Add remaining slides
    for idx in range(existing_slide_count, len(all_slides)):
        slide_data = all_slides[idx]
        slide_type = slide_data.get('slide_type', 'title_content')
        content = slide_data.get('content', {})

        print(f"  Adding slide {idx + 1}: {slide_type} - {content.get('title', 'N/A')[:40]}...")

        # Choose layout based on slide type
        if slide_type == 'section_divider':
            layout = section_layout
        else:
            layout = default_layout

        slide = prs.slides.add_slide(layout)

        # Find and populate placeholders
        title_shape = None
        subtitle_shape = None
        content_shape = None

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # TITLE
                    title_shape = shape
                elif ph_type == 4:  # SUBTITLE
                    subtitle_shape = shape
                elif ph_type == 7:  # OBJECT/CONTENT
                    content_shape = shape

        # Set title
        if title_shape and content.get('title'):
            title_text = content['title']
            tf = title_shape.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = title_text
            # Use smaller font for long titles
            if len(title_text) > 40:
                apply_font(run, FONTS['content_title_long'])
            else:
                apply_font(run, FONTS['content_title'])

        # Set subtitle/takeaway
        if subtitle_shape:
            subtitle_text = content.get('takeaway') or content.get('subtitle', '')
            if subtitle_text:
                tf = subtitle_shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = subtitle_text
                apply_font(run, FONTS['subtitle'])

        # Set content based on slide type
        if slide_type in ['title_content', 'content']:
            bullets = content.get('bullets', [])
            if bullets and content_shape:
                tf = content_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()

                    para.level = 0

                    # Format as "HEADER - body text"
                    header, body = format_bullet_text(bullet)
                    if header:
                        run = para.add_run()
                        run.text = header + ' - '
                        apply_font(run, FONTS['body_header'])
                        if body:
                            run = para.add_run()
                            run.text = body
                            apply_font(run, FONTS['body'])
                    else:
                        run = para.add_run()
                        run.text = bullet
                        apply_font(run, FONTS['body'])

                    # Remove bullet formatting
                    try:
                        pPr = para._p.get_or_add_pPr()
                        buNone = etree.SubElement(pPr, qn('a:buNone'))
                    except:
                        pass

        elif slide_type == 'table_slide':
            # Check for table data in both formats
            table_data = content.get('table', {})
            if table_data:
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
            else:
                # Alternative format: headers and data directly in content
                headers = content.get('headers', [])
                rows = content.get('data', [])

            if headers and rows:
                # Calculate table dimensions
                num_cols = len(headers)
                num_rows = len(rows) + 1  # +1 for header

                # Add table
                table = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(POS['table']['left']),
                    Inches(POS['table']['top']),
                    Inches(POS['table']['width']),
                    Inches(0.5 + 0.4 * len(rows))  # Height based on rows
                ).table

                # Populate header row
                for c_idx, header in enumerate(headers):
                    cell = table.cell(0, c_idx)
                    cell.text = header

                # Populate data rows
                for r_idx, row in enumerate(rows):
                    for c_idx, cell_val in enumerate(row):
                        if c_idx < num_cols:
                            cell = table.cell(r_idx + 1, c_idx)
                            cell.text = str(cell_val) if cell_val else ''

                # Apply formatting
                format_table(table)

        elif slide_type == 'data_chart':
            chart_data = content.get('chart_data', {})
            # Charts would need python-pptx chart functionality
            # For now, add a placeholder note
            if content_shape:
                tf = content_shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = f"[Chart: {chart_data.get('type', 'N/A')}]"
                apply_font(run, FONTS['body'])

        elif slide_type == 'key_metrics':
            metrics = content.get('metrics', [])
            # Add metric boxes
            start_left = 0.40
            for m_idx, metric in enumerate(metrics[:4]):  # Max 4 metrics
                left = start_left + (m_idx * 2.60)
                create_metric_box(
                    slide, left, 4.20,
                    metric.get('value', 'N/A'),
                    metric.get('label', '')
                )

    # Save the presentation
    output_path = Path("pptx_generator/output/Light_Industrial_Thesis_v22.pptx")
    prs.save(str(output_path))

    print(f"\n{'=' * 60}")
    print(f"Generated presentation: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Convert to PDF
    print("\nConverting to PDF...")
    try:
        import subprocess
        pdf_path = output_path.with_suffix('.pdf')
        cmd = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_path.parent),
            str(output_path)
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        print(f"PDF saved: {pdf_path}")
    except Exception as e:
        print(f"PDF conversion failed: {e}")

    return output_path


if __name__ == "__main__":
    main()
