"""
Create a presentation using extracted components from the library.

This script demonstrates how to use the component library to generate
presentations with consistent styling and reusable components.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Add project to path
sys.path.insert(0, str(Path(__file__).parent))

from pptx_generator.modules.component_library import ComponentLibrary


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


class PresentationBuilder:
    """Build presentations using the component library."""

    def __init__(self):
        self.library = ComponentLibrary()
        self.prs = Presentation()

        # Set slide dimensions (16:9 widescreen)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # Get color palette
        palettes = self.library.get_color_palettes()
        if palettes:
            self.palette = palettes[0]
            print(f"Using color palette from: {self.palette.get('template', 'unknown')}")
        else:
            self.palette = {
                'primary': '#051C2C',
                'accent1': '#22A3DF',
                'accent2': '#00A86B',
                'background': '#FFFFFF',
                'text': '#333333'
            }

        # Get typography presets
        self.typography = {}
        for preset in self.library.get_typography_presets():
            preset_type = preset.get('preset_type', 'body')
            if preset_type not in self.typography:
                self.typography[preset_type] = preset

        print(f"Loaded typography presets: {list(self.typography.keys())}")

    def _get_blank_layout(self):
        """Get or create a blank slide layout."""
        return self.prs.slide_layouts[6]  # Usually the blank layout

    def _apply_typography(self, text_frame, preset_type: str = 'body'):
        """Apply typography preset to text frame."""
        preset = self.typography.get(preset_type, {})
        font_name = preset.get('font_name', 'Arial')
        font_size = preset.get('font_size_pt', 12)
        is_bold = preset.get('bold', False)

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = is_bold

    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add a title slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))
        background.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.333), Inches(0.8)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        return slide

    def add_section_slide(self, section_title: str):
        """Add a section divider slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Accent bar on left
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(0.3), self.prs.slide_height
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(self.palette.get('accent1', '#22A3DF'))
        accent_bar.line.fill.background()

        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(11), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        return slide

    def add_content_slide(self, title: str, bullets: list):
        """Add a content slide with bullets."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        # Accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.15),
            Inches(2), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.palette.get('accent1', '#22A3DF'))
        line.line.fill.background()

        # Bullet points
        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.5),
            Inches(11.5), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb('#333333')
            p.level = 0
            p.space_before = Pt(12)

        return slide

    def add_data_slide(self, title: str, chart_id: str = None):
        """Add a data/chart slide using library chart."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        # Try to add chart from library
        if chart_id:
            chart_shape = self.library.add_chart_from_library(
                slide, chart_id,
                left=1, top=1.8,
                width=11, height=5
            )
            if chart_shape:
                print(f"  Added chart from library: {chart_id}")
                return slide

        # Find a chart from library
        charts = self.library.search(component_type='charts')
        if charts:
            chart_id = charts[0].get('id')
            chart_shape = self.library.add_chart_from_library(
                slide, chart_id,
                left=1, top=1.8,
                width=11, height=5
            )
            if chart_shape:
                print(f"  Added chart from library: {chart_id}")
                return slide

        # Fallback: Create a simple bar chart with sample data
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
        chart_data.add_series('Revenue', (4.2, 4.8, 5.1, 5.9))
        chart_data.add_series('Profit', (1.2, 1.4, 1.6, 1.9))

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(1.8),
            Inches(11), Inches(5),
            chart_data
        )

        return slide

    def add_comparison_slide(self, title: str, items: list):
        """Add a comparison slide with multiple columns."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        # Create comparison boxes
        num_items = min(len(items), 3)
        box_width = 3.8
        spacing = 0.3
        start_x = (13.333 - (num_items * box_width + (num_items - 1) * spacing)) / 2

        colors = [
            self.palette.get('accent1', '#22A3DF'),
            self.palette.get('accent2', '#00A86B'),
            self.palette.get('primary', '#051C2C')
        ]

        for i, item in enumerate(items[:num_items]):
            x = start_x + i * (box_width + spacing)

            # Box header
            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.5),
                Inches(box_width), Inches(0.8)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
            header.line.fill.background()

            # Header text
            tf = header.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get('title', f'Option {i+1}')
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(10)

            # Box content
            content_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(2.4),
                Inches(box_width), Inches(4.5)
            )
            content_box.fill.solid()
            content_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
            content_box.line.color.rgb = hex_to_rgb(colors[i % len(colors)])

            # Content text
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].space_before = Pt(15)

            for j, point in enumerate(item.get('points', [])):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"  {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = hex_to_rgb('#333333')
                p.space_before = Pt(8)

        return slide

    def add_timeline_slide(self, title: str, milestones: list):
        """Add a timeline slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        # Timeline base line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(4),
            Inches(11.333), Inches(0.1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))
        line.line.fill.background()

        # Add milestones
        num_milestones = min(len(milestones), 5)
        spacing = 11.333 / (num_milestones)

        for i, milestone in enumerate(milestones[:num_milestones]):
            x = 1 + spacing * i + spacing / 2 - 0.5

            # Circle marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 0.35), Inches(3.85),
                Inches(0.3), Inches(0.3)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = hex_to_rgb(self.palette.get('accent1', '#22A3DF'))
            marker.line.fill.background()

            # Date/Phase label
            date_box = slide.shapes.add_textbox(
                Inches(x), Inches(4.3),
                Inches(1), Inches(0.5)
            )
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = milestone.get('date', f'Phase {i+1}')
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self.palette.get('accent1', '#22A3DF'))
            p.alignment = PP_ALIGN.CENTER

            # Description
            desc_box = slide.shapes.add_textbox(
                Inches(x - 0.5), Inches(2) if i % 2 == 0 else Inches(5),
                Inches(2), Inches(1.5)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = milestone.get('description', '')
            p.font.size = Pt(14)
            p.font.color.rgb = hex_to_rgb('#333333')
            p.alignment = PP_ALIGN.CENTER

        return slide

    def add_summary_slide(self, title: str, key_points: list):
        """Add a summary/takeaways slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))

        # Key points with numbered icons
        num_points = min(len(key_points), 4)

        for i, point in enumerate(key_points[:num_points]):
            y = 1.5 + i * 1.4

            # Number circle
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(y),
                Inches(0.6), Inches(0.6)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = hex_to_rgb(self.palette.get('accent1', '#22A3DF'))
            num_circle.line.fill.background()

            # Number text
            tf = num_circle.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE

            # Point text
            point_box = slide.shapes.add_textbox(
                Inches(1.7), Inches(y + 0.1),
                Inches(10.5), Inches(0.6)
            )
            tf = point_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb('#333333')

        return slide

    def add_contact_slide(self, title: str = "Questions & Next Steps",
                          contact_info: dict = None):
        """Add a closing/contact slide."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())

        # Background with gradient effect (using two shapes)
        bg1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, Inches(4)
        )
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = hex_to_rgb(self.palette.get('primary', '#051C2C'))
        bg1.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(11.333), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Contact info if provided
        if contact_info:
            info_text = []
            if contact_info.get('name'):
                info_text.append(contact_info['name'])
            if contact_info.get('email'):
                info_text.append(contact_info['email'])
            if contact_info.get('phone'):
                info_text.append(contact_info['phone'])

            info_box = slide.shapes.add_textbox(
                Inches(1), Inches(5),
                Inches(11.333), Inches(1.5)
            )
            tf = info_box.text_frame
            for i, text in enumerate(info_text):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(18)
                p.font.color.rgb = hex_to_rgb('#333333')
                p.alignment = PP_ALIGN.CENTER

        return slide

    def save(self, output_path: str):
        """Save the presentation."""
        self.prs.save(output_path)
        print(f"\nPresentation saved to: {output_path}")


def main():
    """Create a sample presentation using extracted components."""
    print("=" * 60)
    print("Creating Presentation Using Extracted Components")
    print("=" * 60)

    builder = PresentationBuilder()

    # 1. Title Slide
    print("\nAdding slides...")
    builder.add_title_slide(
        "Strategic Business Review",
        "Q4 2024 Performance & 2025 Outlook"
    )
    print("  1. Title slide")

    # 2. Agenda
    builder.add_content_slide(
        "Agenda",
        [
            "Executive Summary",
            "Market Analysis & Performance",
            "Key Metrics & Trends",
            "Strategic Recommendations",
            "Implementation Roadmap",
            "Next Steps"
        ]
    )
    print("  2. Agenda slide")

    # 3. Executive Summary
    builder.add_summary_slide(
        "Executive Summary",
        [
            "Revenue grew 18% YoY, exceeding targets by 3 percentage points",
            "Market share increased to 24% in core segments",
            "Customer satisfaction scores improved to 4.6/5.0",
            "Operating margin expanded 2.1 percentage points"
        ]
    )
    print("  3. Executive Summary slide")

    # 4. Section: Analysis
    builder.add_section_slide("Market Analysis")
    print("  4. Section divider")

    # 5. Current Situation
    builder.add_content_slide(
        "Current Market Position",
        [
            "Strong brand recognition in enterprise segment (78% awareness)",
            "Growing adoption of cloud-based solutions (+42% YoY)",
            "Competitive pressure from new market entrants",
            "Expanding addressable market ($12.4B projected by 2026)",
            "Digital transformation driving increased demand"
        ]
    )
    print("  5. Current situation slide")

    # 6. Key Metrics (with chart)
    builder.add_data_slide("Key Performance Metrics")
    print("  6. Data/chart slide")

    # 7. Comparison
    builder.add_comparison_slide(
        "Strategic Options Comparison",
        [
            {
                "title": "Option A: Organic Growth",
                "points": [
                    "Lower risk profile",
                    "Predictable investment",
                    "2-3 year timeline",
                    "15% projected ROI"
                ]
            },
            {
                "title": "Option B: Acquisition",
                "points": [
                    "Rapid market entry",
                    "Higher upfront cost",
                    "Integration challenges",
                    "25% projected ROI"
                ]
            },
            {
                "title": "Option C: Partnership",
                "points": [
                    "Shared investment",
                    "Access to new tech",
                    "Flexible structure",
                    "20% projected ROI"
                ]
            }
        ]
    )
    print("  7. Comparison slide")

    # 8. Section: Recommendations
    builder.add_section_slide("Recommendations")
    print("  8. Section divider")

    # 9. Proposed Solution
    builder.add_content_slide(
        "Recommended Approach",
        [
            "Pursue hybrid strategy combining organic growth with strategic partnership",
            "Invest $24M over 18 months in product development",
            "Establish joint venture with technology partner for AI capabilities",
            "Expand sales team by 35% in key growth markets",
            "Launch customer success program to drive retention"
        ]
    )
    print("  9. Proposed solution slide")

    # 10. Timeline
    builder.add_timeline_slide(
        "Implementation Roadmap",
        [
            {"date": "Q1 2025", "description": "Strategic planning & partner selection"},
            {"date": "Q2 2025", "description": "Partnership finalization & team expansion"},
            {"date": "Q3 2025", "description": "Product development & pilot launch"},
            {"date": "Q4 2025", "description": "Full market rollout & optimization"}
        ]
    )
    print("  10. Timeline slide")

    # 11. Key Takeaways
    builder.add_summary_slide(
        "Key Takeaways",
        [
            "Market opportunity is substantial and growing",
            "Hybrid approach balances risk and reward effectively",
            "Investment required: $24M over 18 months",
            "Expected ROI: 22% with payback in 2.5 years"
        ]
    )
    print("  11. Key takeaways slide")

    # 12. Contact/Next Steps
    builder.add_contact_slide(
        "Questions & Next Steps",
        {
            "name": "Strategic Planning Team",
            "email": "strategy@company.com"
        }
    )
    print("  12. Contact slide")

    # Save the presentation
    output_path = Path(__file__).parent / "outputs" / "strategic_business_review.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)
    builder.save(str(output_path))

    print("\n" + "=" * 60)
    print("PRESENTATION COMPLETE")
    print("=" * 60)
    print(f"Slides: 12")
    print(f"Output: {output_path}")


if __name__ == "__main__":
    main()
