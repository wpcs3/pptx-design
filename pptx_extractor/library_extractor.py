"""
Library Extractor - Extract and catalog components from PowerPoint templates.

Creates a searchable library of:
- Images
- Charts (with data)
- Tables (with structure)
- Shapes/Diagrams
- Slide layouts

Each component is saved individually with metadata for easy lookup.
"""

import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


class LibraryExtractor:
    """Extract and catalog components from PowerPoint templates."""

    def __init__(self, output_dir: Path):
        """
        Initialize the library extractor.

        Args:
            output_dir: Base directory for the component library
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories for each component type
        self.dirs = {
            'images': self.output_dir / 'images',
            'charts': self.output_dir / 'charts',
            'tables': self.output_dir / 'tables',
            'shapes': self.output_dir / 'shapes',
            'diagrams': self.output_dir / 'diagrams',
            'layouts': self.output_dir / 'layouts',
            'slides': self.output_dir / 'slides',
        }
        for d in self.dirs.values():
            d.mkdir(parents=True, exist_ok=True)

        # Index file path
        self.index_path = self.output_dir / 'library_index.json'

        # Load existing index or create new
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'templates': {},
            'components': {
                'images': [],
                'charts': [],
                'tables': [],
                'shapes': [],
                'diagrams': [],
                'layouts': [],
                'slides': [],
            },
            'categories': defaultdict(list),
            'tags': defaultdict(list),
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()

        # Convert defaultdicts to regular dicts for JSON serialization
        index_copy = self.index.copy()
        index_copy['categories'] = dict(self.index['categories'])
        index_copy['tags'] = dict(self.index['tags'])

        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(index_copy, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: bytes) -> str:
        """Generate unique ID from content hash."""
        return hashlib.md5(content).hexdigest()[:12]

    def _get_shape_type_name(self, shape) -> str:
        """Get human-readable shape type name."""
        try:
            return shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type)
        except:
            return 'UNKNOWN'

    def _extract_colors_from_shape(self, shape) -> list:
        """Extract colors used in a shape."""
        colors = []
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                    try:
                        colors.append(str(shape.fill.fore_color.rgb))
                    except:
                        pass
            if hasattr(shape, 'line') and shape.line.fill.type is not None:
                try:
                    colors.append(str(shape.line.color.rgb))
                except:
                    pass
        except:
            pass
        return colors

    def extract_template(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all components from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template (defaults to filename)

        Returns:
            Summary of extracted components
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting components from: {template_name}")

        prs = Presentation(pptx_path)

        # Template metadata
        template_info = {
            'name': template_name,
            'source_file': str(pptx_path),
            'extracted_date': datetime.now().isoformat(),
            'slide_count': len(prs.slides),
            'slide_width_inches': prs.slide_width.inches,
            'slide_height_inches': prs.slide_height.inches,
            'components': {
                'images': 0,
                'charts': 0,
                'tables': 0,
                'shapes': 0,
                'diagrams': 0,
            }
        }

        # Extract components from each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'

            # Extract slide thumbnail/layout info
            self._extract_slide_layout(slide, slide_num, template_name, layout_name)

            for shape in slide.shapes:
                shape_type = self._get_shape_type_name(shape)

                if shape_type == 'PICTURE':
                    self._extract_image(shape, slide_num, template_name)
                    template_info['components']['images'] += 1

                elif shape_type == 'CHART':
                    self._extract_chart(shape, slide_num, template_name)
                    template_info['components']['charts'] += 1

                elif shape_type == 'TABLE':
                    self._extract_table(shape, slide_num, template_name)
                    template_info['components']['tables'] += 1

                elif shape_type == 'GROUP':
                    self._extract_diagram(shape, slide_num, template_name)
                    template_info['components']['diagrams'] += 1

                elif shape_type in ('AUTO_SHAPE', 'FREEFORM'):
                    self._extract_shape(shape, slide_num, template_name)
                    template_info['components']['shapes'] += 1

        # Save template info to index
        self.index['templates'][template_name] = template_info
        self._save_index()

        logger.info(f"Extracted: {template_info['components']}")
        return template_info

    def _extract_image(self, shape, slide_num: int, template_name: str):
        """Extract and save an image."""
        try:
            image = shape.image
            image_bytes = image.blob

            # Generate unique ID
            img_id = self._generate_id(image_bytes)

            # Check if already extracted
            existing = [c for c in self.index['components']['images'] if c['id'] == img_id]
            if existing:
                # Add reference to existing image
                existing[0]['references'].append({
                    'template': template_name,
                    'slide': slide_num,
                })
                return

            # Determine file extension
            ext = image.ext
            filename = f"{img_id}.{ext}"
            filepath = self.dirs['images'] / filename

            # Save image
            with open(filepath, 'wb') as f:
                f.write(image_bytes)

            # Create metadata
            metadata = {
                'id': img_id,
                'type': 'image',
                'filename': filename,
                'format': ext,
                'size_bytes': len(image_bytes),
                'width_inches': shape.width.inches if shape.width else None,
                'height_inches': shape.height.inches if shape.height else None,
                'position': {
                    'left_inches': shape.left.inches if shape.left else None,
                    'top_inches': shape.top.inches if shape.top else None,
                },
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': [],
                'category': 'uncategorized',
            }

            self.index['components']['images'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract image from slide {slide_num}: {e}")

    def _extract_chart(self, shape, slide_num: int, template_name: str):
        """Extract and save chart data and metadata."""
        try:
            chart = shape.chart

            # Get chart type
            chart_type = 'unknown'
            try:
                chart_type = str(chart.chart_type).split('.')[-1] if chart.chart_type else 'unknown'
            except:
                pass

            # Extract chart data
            chart_data = {
                'categories': [],
                'series': [],
            }

            try:
                # Get categories
                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories') and plot.categories:
                        chart_data['categories'] = [str(c) for c in plot.categories]

                # Get series data
                for series in chart.series:
                    series_data = {
                        'name': str(series.name) if series.name else 'Series',
                        'values': [],
                    }
                    try:
                        series_data['values'] = [v for v in series.values]
                    except:
                        pass
                    chart_data['series'].append(series_data)
            except Exception as e:
                logger.debug(f"Could not extract full chart data: {e}")

            # Generate ID from chart structure
            chart_str = json.dumps(chart_data, sort_keys=True)
            chart_id = self._generate_id(chart_str.encode())

            # Check if similar chart exists
            existing = [c for c in self.index['components']['charts'] if c['id'] == chart_id]
            if existing:
                existing[0]['references'].append({
                    'template': template_name,
                    'slide': slide_num,
                })
                return

            # Save chart data
            filename = f"{chart_id}.json"
            filepath = self.dirs['charts'] / filename

            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(chart_data, f, indent=2)

            # Create metadata
            metadata = {
                'id': chart_id,
                'type': 'chart',
                'chart_type': chart_type,
                'filename': filename,
                'has_title': chart.has_title if hasattr(chart, 'has_title') else False,
                'title': chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and chart.has_title else None,
                'series_count': len(chart_data['series']),
                'category_count': len(chart_data['categories']),
                'width_inches': shape.width.inches if shape.width else None,
                'height_inches': shape.height.inches if shape.height else None,
                'position': {
                    'left_inches': shape.left.inches if shape.left else None,
                    'top_inches': shape.top.inches if shape.top else None,
                },
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': [chart_type],
                'category': self._categorize_chart(chart_type),
            }

            self.index['components']['charts'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract chart from slide {slide_num}: {e}")

    def _categorize_chart(self, chart_type: str) -> str:
        """Categorize chart by type."""
        chart_type_lower = chart_type.lower()
        if 'bar' in chart_type_lower:
            return 'bar_charts'
        elif 'column' in chart_type_lower:
            return 'column_charts'
        elif 'line' in chart_type_lower:
            return 'line_charts'
        elif 'pie' in chart_type_lower or 'doughnut' in chart_type_lower:
            return 'pie_charts'
        elif 'area' in chart_type_lower:
            return 'area_charts'
        elif 'scatter' in chart_type_lower or 'xy' in chart_type_lower:
            return 'scatter_charts'
        elif 'radar' in chart_type_lower:
            return 'radar_charts'
        else:
            return 'other_charts'

    def _extract_table(self, shape, slide_num: int, template_name: str):
        """Extract and save table structure and data."""
        try:
            table = shape.table

            # Extract table data
            table_data = {
                'rows': len(table.rows),
                'cols': len(table.columns),
                'data': [],
                'column_widths': [],
                'row_heights': [],
            }

            # Get column widths
            for col in table.columns:
                table_data['column_widths'].append(col.width.inches if col.width else None)

            # Get row heights and cell data
            for row in table.rows:
                table_data['row_heights'].append(row.height.inches if row.height else None)
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text if cell.text else ''
                    row_data.append(cell_text)
                table_data['data'].append(row_data)

            # Generate ID
            table_str = json.dumps(table_data, sort_keys=True)
            table_id = self._generate_id(table_str.encode())

            # Check if similar table exists
            existing = [c for c in self.index['components']['tables'] if c['id'] == table_id]
            if existing:
                existing[0]['references'].append({
                    'template': template_name,
                    'slide': slide_num,
                })
                return

            # Save table data
            filename = f"{table_id}.json"
            filepath = self.dirs['tables'] / filename

            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(table_data, f, indent=2)

            # Determine table category
            category = 'data_table'
            if table_data['rows'] <= 2:
                category = 'header_table'
            elif table_data['cols'] >= 5:
                category = 'comparison_matrix'
            elif all(len(row[0]) > 20 for row in table_data['data'] if row and row[0]):
                category = 'text_table'

            # Create metadata
            metadata = {
                'id': table_id,
                'type': 'table',
                'filename': filename,
                'rows': table_data['rows'],
                'cols': table_data['cols'],
                'width_inches': shape.width.inches if shape.width else None,
                'height_inches': shape.height.inches if shape.height else None,
                'position': {
                    'left_inches': shape.left.inches if shape.left else None,
                    'top_inches': shape.top.inches if shape.top else None,
                },
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': [f"{table_data['rows']}x{table_data['cols']}"],
                'category': category,
            }

            self.index['components']['tables'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract table from slide {slide_num}: {e}")

    def _extract_shape(self, shape, slide_num: int, template_name: str):
        """Extract and save shape metadata."""
        try:
            # Get shape details
            shape_info = {
                'auto_shape_type': None,
                'has_text': False,
                'text': '',
                'fill_type': None,
                'colors': [],
            }

            try:
                if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type:
                    shape_info['auto_shape_type'] = str(shape.auto_shape_type).split('.')[-1]
            except:
                pass

            if hasattr(shape, 'text') and shape.text:
                shape_info['has_text'] = True
                shape_info['text'] = shape.text[:100]  # Truncate long text

            shape_info['colors'] = self._extract_colors_from_shape(shape)

            # Generate ID
            shape_str = f"{shape_info['auto_shape_type']}_{shape.width}_{shape.height}"
            shape_id = self._generate_id(shape_str.encode())

            # Save shape metadata
            filename = f"{shape_id}.json"
            filepath = self.dirs['shapes'] / filename

            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(shape_info, f, indent=2)

            # Create metadata
            metadata = {
                'id': shape_id,
                'type': 'shape',
                'shape_type': shape_info['auto_shape_type'],
                'filename': filename,
                'has_text': shape_info['has_text'],
                'width_inches': shape.width.inches if shape.width else None,
                'height_inches': shape.height.inches if shape.height else None,
                'position': {
                    'left_inches': shape.left.inches if shape.left else None,
                    'top_inches': shape.top.inches if shape.top else None,
                },
                'colors': shape_info['colors'],
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': [shape_info['auto_shape_type']] if shape_info['auto_shape_type'] else [],
                'category': 'basic_shapes',
            }

            self.index['components']['shapes'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract shape from slide {slide_num}: {e}")

    def _extract_diagram(self, group_shape, slide_num: int, template_name: str):
        """Extract and save diagram (group of shapes) metadata."""
        try:
            # Analyze group contents
            diagram_info = {
                'shape_count': 0,
                'shape_types': [],
                'has_text': False,
                'text_content': [],
                'colors': [],
            }

            def analyze_shapes(shapes):
                for shape in shapes:
                    diagram_info['shape_count'] += 1
                    shape_type = self._get_shape_type_name(shape)
                    diagram_info['shape_types'].append(shape_type)

                    if hasattr(shape, 'text') and shape.text:
                        diagram_info['has_text'] = True
                        diagram_info['text_content'].append(shape.text[:50])

                    diagram_info['colors'].extend(self._extract_colors_from_shape(shape))

                    # Recurse into nested groups
                    if shape_type == 'GROUP' and hasattr(shape, 'shapes'):
                        analyze_shapes(shape.shapes)

            analyze_shapes(group_shape.shapes)

            # Determine diagram category based on content
            category = self._categorize_diagram(diagram_info)

            # Generate ID
            diagram_str = json.dumps(diagram_info, sort_keys=True)
            diagram_id = self._generate_id(diagram_str.encode())

            # Save diagram info
            filename = f"{diagram_id}.json"
            filepath = self.dirs['diagrams'] / filename

            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(diagram_info, f, indent=2)

            # Create metadata
            metadata = {
                'id': diagram_id,
                'type': 'diagram',
                'filename': filename,
                'shape_count': diagram_info['shape_count'],
                'has_text': diagram_info['has_text'],
                'width_inches': group_shape.width.inches if group_shape.width else None,
                'height_inches': group_shape.height.inches if group_shape.height else None,
                'position': {
                    'left_inches': group_shape.left.inches if group_shape.left else None,
                    'top_inches': group_shape.top.inches if group_shape.top else None,
                },
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': list(set(diagram_info['shape_types'])),
                'category': category,
            }

            self.index['components']['diagrams'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract diagram from slide {slide_num}: {e}")

    def _categorize_diagram(self, diagram_info: dict) -> str:
        """Categorize diagram based on its content."""
        shape_types = diagram_info['shape_types']
        text_content = ' '.join(diagram_info['text_content']).lower()

        # Check for common patterns
        if shape_types.count('AUTO_SHAPE') >= 3:
            if 'arrow' in str(shape_types).lower():
                return 'process_flow'
            elif diagram_info['shape_count'] >= 6:
                return 'framework'

        if 'AUTO_SHAPE' in shape_types and diagram_info['has_text']:
            return 'labeled_diagram'

        return 'custom_diagram'

    def _extract_slide_layout(self, slide, slide_num: int, template_name: str, layout_name: str):
        """Extract slide layout information."""
        try:
            # Count shapes by type
            shape_counts = defaultdict(int)
            for shape in slide.shapes:
                shape_type = self._get_shape_type_name(shape)
                shape_counts[shape_type] += 1

            # Get slide title if available
            title = ''
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    title = shape.text[:80].replace('\n', ' ').strip()
                    break

            # Generate layout ID
            layout_str = f"{template_name}_{slide_num}_{layout_name}"
            layout_id = self._generate_id(layout_str.encode())

            # Create metadata
            metadata = {
                'id': layout_id,
                'type': 'layout',
                'template': template_name,
                'slide_num': slide_num,
                'layout_name': layout_name,
                'title': title,
                'shape_counts': dict(shape_counts),
                'total_shapes': sum(shape_counts.values()),
                'references': [{
                    'template': template_name,
                    'slide': slide_num,
                }],
                'tags': [layout_name],
                'category': self._categorize_layout(layout_name, shape_counts),
            }

            self.index['components']['layouts'].append(metadata)

        except Exception as e:
            logger.warning(f"Failed to extract layout from slide {slide_num}: {e}")

    def _categorize_layout(self, layout_name: str, shape_counts: dict) -> str:
        """Categorize slide layout."""
        layout_lower = layout_name.lower()

        if 'blank' in layout_lower:
            return 'blank'
        elif 'title' in layout_lower or 'front' in layout_lower:
            return 'title_slide'
        elif 'agenda' in layout_lower or 'content' in layout_lower:
            return 'agenda'
        elif shape_counts.get('CHART', 0) > 0:
            return 'chart_slide'
        elif shape_counts.get('TABLE', 0) > 0:
            return 'table_slide'
        elif shape_counts.get('GROUP', 0) > 0:
            return 'diagram_slide'
        else:
            return 'content_slide'

    def get_summary(self) -> dict:
        """Get summary of the library contents."""
        summary = {
            'total_templates': len(self.index['templates']),
            'total_components': {},
            'categories': {},
        }

        for component_type, components in self.index['components'].items():
            summary['total_components'][component_type] = len(components)

            # Count by category
            categories = defaultdict(int)
            for comp in components:
                cat = comp.get('category', 'uncategorized')
                categories[cat] += 1
            summary['categories'][component_type] = dict(categories)

        return summary

    def search(self,
               component_type: Optional[str] = None,
               category: Optional[str] = None,
               tags: Optional[list] = None,
               template: Optional[str] = None) -> list:
        """
        Search the library for components.

        Args:
            component_type: Filter by type (images, charts, tables, etc.)
            category: Filter by category
            tags: Filter by tags (any match)
            template: Filter by source template

        Returns:
            List of matching components
        """
        results = []

        types_to_search = [component_type] if component_type else self.index['components'].keys()

        for ctype in types_to_search:
            if ctype not in self.index['components']:
                continue

            for component in self.index['components'][ctype]:
                # Apply filters
                if category and component.get('category') != category:
                    continue

                if tags:
                    comp_tags = component.get('tags', [])
                    if not any(t in comp_tags for t in tags):
                        continue

                if template:
                    refs = component.get('references', [])
                    if not any(r['template'] == template for r in refs):
                        continue

                results.append(component)

        return results


def extract_all_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract components from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = LibraryExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting: {pptx_path.name}")
        try:
            extractor.extract_template(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    # Save final index
    extractor._save_index()

    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python library_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_all_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
