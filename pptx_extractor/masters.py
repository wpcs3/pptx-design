"""
Slide Master and Layout Extraction Module

Extracts slide master and layout information from PPTX files and provides
functionality to recreate them.
"""
import json
import logging
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Add parent directory to path for config import
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))

from config import OUTPUT_DIR, DESCRIPTION_DIR

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class MasterError(Exception):
    """Exception raised when master/layout extraction fails."""
    pass


def emu_to_inches(emu: int) -> float:
    """Convert EMUs to inches."""
    return emu / 914400


def extract_color(color_obj) -> Optional[str]:
    """
    Extract hex color from a python-pptx color object.

    Args:
        color_obj: A python-pptx color object

    Returns:
        Hex color string or None
    """
    try:
        if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
            rgb = color_obj.rgb
            return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"
        if hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
            return f"theme:{color_obj.theme_color}"
    except Exception:
        pass
    return None


def extract_font_info(font) -> dict:
    """
    Extract font information from a python-pptx font object.

    Args:
        font: A python-pptx font object

    Returns:
        Dict with font properties
    """
    info = {}

    try:
        if font.name:
            info["font_family"] = font.name
        if font.size:
            info["font_size_pt"] = font.size.pt
        if font.bold is not None:
            info["bold"] = font.bold
        if font.italic is not None:
            info["italic"] = font.italic
        if font.color and font.color.rgb:
            info["font_color"] = extract_color(font.color)
    except Exception as e:
        logger.debug(f"Error extracting font info: {e}")

    return info


def extract_fill_info(fill) -> dict:
    """
    Extract fill information from a python-pptx fill object.

    Args:
        fill: A python-pptx fill object

    Returns:
        Dict with fill properties
    """
    info = {}

    try:
        fill_type = fill.type
        if fill_type is not None:
            info["type"] = str(fill_type).split('.')[-1].lower()

        if hasattr(fill, 'fore_color') and fill.fore_color:
            color = extract_color(fill.fore_color)
            if color:
                info["color"] = color

        if hasattr(fill, 'back_color') and fill.back_color:
            color = extract_color(fill.back_color)
            if color:
                info["back_color"] = color

    except Exception as e:
        logger.debug(f"Error extracting fill info: {e}")

    return info


def extract_placeholder_info(shape) -> Optional[dict]:
    """
    Extract placeholder information from a shape.

    Args:
        shape: A python-pptx shape object

    Returns:
        Dict with placeholder info or None if not a placeholder
    """
    if not shape.is_placeholder:
        return None

    ph = shape.placeholder_format

    info = {
        "type": str(ph.type).split('.')[-1] if ph.type else "UNKNOWN",
        "idx": ph.idx,
    }

    # Position and size
    info["position"] = {
        "left_inches": emu_to_inches(shape.left) if shape.left else 0,
        "top_inches": emu_to_inches(shape.top) if shape.top else 0,
        "width_inches": emu_to_inches(shape.width) if shape.width else 0,
        "height_inches": emu_to_inches(shape.height) if shape.height else 0,
    }

    # Text properties if it has a text frame
    if shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            info["text_properties"] = extract_font_info(p.font)
            if p.alignment:
                info["text_properties"]["alignment"] = str(p.alignment).split('.')[-1].lower()

    return info


def extract_shape_info(shape) -> dict:
    """
    Extract information from a shape.

    Args:
        shape: A python-pptx shape object

    Returns:
        Dict with shape information
    """
    info = {
        "name": shape.name,
        "shape_type": str(shape.shape_type).split('.')[-1] if shape.shape_type else "UNKNOWN",
        "position": {
            "left_inches": emu_to_inches(shape.left) if shape.left else 0,
            "top_inches": emu_to_inches(shape.top) if shape.top else 0,
            "width_inches": emu_to_inches(shape.width) if shape.width else 0,
            "height_inches": emu_to_inches(shape.height) if shape.height else 0,
        }
    }

    # Check if placeholder
    if shape.is_placeholder:
        info["placeholder"] = extract_placeholder_info(shape)

    # Fill info
    if hasattr(shape, 'fill'):
        fill_info = extract_fill_info(shape.fill)
        if fill_info:
            info["fill"] = fill_info

    # Line/border info
    if hasattr(shape, 'line') and shape.line:
        try:
            if shape.line.color and shape.line.color.rgb:
                info["line_color"] = extract_color(shape.line.color)
            if shape.line.width:
                info["line_width_pt"] = shape.line.width.pt
        except Exception:
            pass

    # Text info
    if shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            info["text_properties"] = extract_font_info(p.font)
            if p.text:
                info["text"] = p.text

    return info


def extract_layout_info(layout) -> dict:
    """
    Extract information from a slide layout.

    Args:
        layout: A python-pptx slide layout object

    Returns:
        Dict with layout information
    """
    info = {
        "name": layout.name,
        "placeholders": [],
        "shapes": []
    }

    for shape in layout.shapes:
        if shape.is_placeholder:
            ph_info = extract_placeholder_info(shape)
            if ph_info:
                info["placeholders"].append(ph_info)
        else:
            shape_info = extract_shape_info(shape)
            info["shapes"].append(shape_info)

    return info


def extract_slide_master_info(master) -> dict:
    """
    Extract information from a slide master.

    Args:
        master: A python-pptx slide master object

    Returns:
        Dict with slide master information
    """
    info = {
        "name": master.name if hasattr(master, 'name') else "Slide Master",
        "placeholders": [],
        "shapes": [],
        "layouts": []
    }

    # Extract background
    try:
        if master.background and master.background.fill:
            info["background"] = extract_fill_info(master.background.fill)
    except Exception as e:
        logger.debug(f"Error extracting master background: {e}")

    # Extract shapes on the master
    for shape in master.shapes:
        if shape.is_placeholder:
            ph_info = extract_placeholder_info(shape)
            if ph_info:
                info["placeholders"].append(ph_info)
        else:
            shape_info = extract_shape_info(shape)
            info["shapes"].append(shape_info)

    # Extract layouts
    for layout in master.slide_layouts:
        layout_info = extract_layout_info(layout)
        info["layouts"].append(layout_info)

    return info


def extract_all_masters(pptx_path: Path) -> dict:
    """
    Extract all slide masters and layouts from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Dict with complete master/layout information
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise MasterError(f"PPTX file not found: {pptx_path}")

    prs = Presentation(pptx_path)

    result = {
        "source_file": pptx_path.name,
        "slide_dimensions": {
            "width_inches": emu_to_inches(prs.slide_width),
            "height_inches": emu_to_inches(prs.slide_height)
        },
        "slide_masters": []
    }

    for master in prs.slide_masters:
        master_info = extract_slide_master_info(master)
        result["slide_masters"].append(master_info)

    logger.info(f"Extracted {len(result['slide_masters'])} slide masters from {pptx_path.name}")

    return result


def extract_slide_layout_usage(pptx_path: Path) -> list[dict]:
    """
    Extract which layout each slide in a presentation uses.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        List of dicts with slide index and layout name
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise MasterError(f"PPTX file not found: {pptx_path}")

    prs = Presentation(pptx_path)

    usage = []
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        usage.append({
            "slide_index": i,
            "slide_number": i + 1,
            "layout_name": layout_name
        })

    return usage


def save_master_info(
    master_info: dict,
    template_name: str,
    output_dir: Optional[Path] = None
) -> Path:
    """
    Save extracted master information to JSON.

    Args:
        master_info: Dict from extract_all_masters()
        template_name: Name for the output file
        output_dir: Output directory (defaults to DESCRIPTION_DIR)

    Returns:
        Path to saved file
    """
    if output_dir is None:
        output_dir = DESCRIPTION_DIR
    else:
        output_dir = Path(output_dir)

    output_dir.mkdir(parents=True, exist_ok=True)

    output_path = output_dir / f"{template_name}_masters.json"

    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(master_info, f, indent=2, ensure_ascii=False)

    logger.info(f"Master info saved: {output_path}")
    return output_path


def get_layout_names(pptx_path: Path) -> list[str]:
    """
    Get list of all layout names available in a PPTX.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        List of layout names
    """
    prs = Presentation(pptx_path)
    names = []

    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            names.append(layout.name)

    return names


def find_layout_by_name(prs: Presentation, layout_name: str):
    """
    Find a slide layout by name in a presentation.

    Args:
        prs: Presentation object
        layout_name: Name of the layout to find

    Returns:
        SlideLayout object or None if not found
    """
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.lower() == layout_name.lower():
                return layout
    return None


def create_slide_with_layout(
    prs: Presentation,
    layout_name: str,
    fallback_to_blank: bool = True
):
    """
    Create a new slide using a specific layout.

    Args:
        prs: Presentation object
        layout_name: Name of the layout to use
        fallback_to_blank: If True, use blank layout if named layout not found

    Returns:
        The new slide object

    Raises:
        MasterError: If layout not found and fallback is False
    """
    layout = find_layout_by_name(prs, layout_name)

    if layout is None:
        if fallback_to_blank:
            # Find blank layout
            layout = find_layout_by_name(prs, "Blank")
            if layout is None:
                # Use last layout as fallback
                layout = prs.slide_layouts[-1]
            logger.warning(f"Layout '{layout_name}' not found, using '{layout.name}'")
        else:
            raise MasterError(f"Layout not found: {layout_name}")

    return prs.slides.add_slide(layout)


if __name__ == "__main__":
    import sys
    from rich.console import Console
    from rich.table import Table

    console = Console()

    if len(sys.argv) < 2:
        console.print("Usage: python -m src.masters <pptx_file>")
        console.print("\nExtracts slide master and layout information from a PPTX file.")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])

    console.print(f"\n[bold]Extracting masters from:[/bold] {pptx_path.name}\n")

    try:
        # Extract master info
        master_info = extract_all_masters(pptx_path)

        # Display summary
        console.print(f"[bold]Slide Dimensions:[/bold] {master_info['slide_dimensions']['width_inches']:.2f}\" x {master_info['slide_dimensions']['height_inches']:.2f}\"")
        console.print(f"[bold]Slide Masters:[/bold] {len(master_info['slide_masters'])}")

        for i, master in enumerate(master_info['slide_masters']):
            console.print(f"\n[bold cyan]Master {i + 1}: {master['name']}[/bold cyan]")
            console.print(f"  Placeholders: {len(master['placeholders'])}")
            console.print(f"  Shapes: {len(master['shapes'])}")
            console.print(f"  Layouts: {len(master['layouts'])}")

            # Show layouts
            if master['layouts']:
                table = Table(title=f"Layouts in {master['name']}")
                table.add_column("Index", style="cyan")
                table.add_column("Name", style="green")
                table.add_column("Placeholders", style="dim")

                for j, layout in enumerate(master['layouts']):
                    table.add_row(
                        str(j),
                        layout['name'],
                        str(len(layout['placeholders']))
                    )

                console.print(table)

        # Show slide usage
        usage = extract_slide_layout_usage(pptx_path)
        if usage:
            console.print("\n[bold]Slide Layout Usage:[/bold]")
            for u in usage:
                console.print(f"  Slide {u['slide_number']}: {u['layout_name']}")

        # Save to file
        template_name = pptx_path.stem
        output_path = save_master_info(master_info, template_name)
        console.print(f"\n[bold green]Saved to:[/bold green] {output_path}")

    except Exception as e:
        console.print(f"[bold red]Error:[/bold red] {e}")
        sys.exit(1)
