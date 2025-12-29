"""
Master Slide Extractor

Extracts all master slides and layouts with exact positions, formatting,
and placeholder information for use as templates.
"""
import json
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
from dataclasses import dataclass, asdict
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def emu_to_inches(emu: int) -> float:
    """Convert EMUs to inches."""
    if emu is None:
        return 0.0
    return round(emu / 914400, 3)


def get_color_hex(color) -> Optional[str]:
    """Extract hex color from a color object."""
    try:
        if color is None:
            return None
        if hasattr(color, 'rgb') and color.rgb:
            return f"#{color.rgb}"
        if hasattr(color, 'theme_color') and color.theme_color:
            return f"theme:{color.theme_color}"
        return None
    except:
        return None


def get_font_info(font) -> Dict[str, Any]:
    """Extract font information."""
    info = {}
    try:
        if font.name:
            info["font_name"] = font.name
        if font.size:
            info["font_size_pt"] = font.size.pt
        if font.bold is not None:
            info["bold"] = font.bold
        if font.italic is not None:
            info["italic"] = font.italic
        color = get_color_hex(font.color)
        if color:
            info["font_color"] = color
    except:
        pass
    return info


def extract_placeholder_info(placeholder) -> Dict[str, Any]:
    """Extract detailed placeholder information."""
    info = {
        "idx": placeholder.placeholder_format.idx,
        "type": str(placeholder.placeholder_format.type).split('.')[-1].strip('()'),
        "position": {
            "left_inches": emu_to_inches(placeholder.left),
            "top_inches": emu_to_inches(placeholder.top),
            "width_inches": emu_to_inches(placeholder.width),
            "height_inches": emu_to_inches(placeholder.height),
        }
    }

    # Get text formatting if available
    if placeholder.has_text_frame:
        tf = placeholder.text_frame
        info["text_frame"] = {
            "word_wrap": tf.word_wrap,
        }

        # Get paragraph formatting from first paragraph
        if tf.paragraphs:
            p = tf.paragraphs[0]
            para_info = {}
            if p.alignment:
                para_info["alignment"] = str(p.alignment).split('.')[-1].strip('()')
            if p.font:
                para_info.update(get_font_info(p.font))
            if para_info:
                info["paragraph_format"] = para_info

    return info


def extract_shape_info(shape) -> Dict[str, Any]:
    """Extract shape information."""
    info = {
        "shape_type": str(shape.shape_type).split('.')[-1].strip('()'),
        "position": {
            "left_inches": emu_to_inches(shape.left),
            "top_inches": emu_to_inches(shape.top),
            "width_inches": emu_to_inches(shape.width),
            "height_inches": emu_to_inches(shape.height),
        }
    }

    # Get name if available
    if shape.name:
        info["name"] = shape.name

    # Get fill color
    try:
        if hasattr(shape, 'fill') and shape.fill:
            fill = shape.fill
            if fill.type is not None:
                info["fill_type"] = str(fill.type).split('.')[-1].strip('()')
            if hasattr(fill, 'fore_color') and fill.fore_color:
                color = get_color_hex(fill.fore_color)
                if color:
                    info["fill_color"] = color
    except:
        pass

    # Get line/border
    try:
        if hasattr(shape, 'line') and shape.line:
            line = shape.line
            line_info = {}
            if line.width:
                line_info["width_pt"] = line.width.pt
            if hasattr(line, 'color') and line.color:
                color = get_color_hex(line.color)
                if color:
                    line_info["color"] = color
            if line_info:
                info["line"] = line_info
    except:
        pass

    # Get text if present
    if shape.has_text_frame:
        try:
            text = shape.text_frame.text.strip()
            if text:
                info["text"] = text
        except:
            pass

    return info


def extract_layout(layout) -> Dict[str, Any]:
    """Extract a single layout with all details."""
    layout_info = {
        "name": layout.name,
        "placeholders": [],
        "shapes": [],
    }

    # Extract placeholders
    for placeholder in layout.placeholders:
        ph_info = extract_placeholder_info(placeholder)
        layout_info["placeholders"].append(ph_info)

    # Extract non-placeholder shapes
    for shape in layout.shapes:
        if not shape.is_placeholder:
            shape_info = extract_shape_info(shape)
            layout_info["shapes"].append(shape_info)

    return layout_info


def extract_master(master) -> Dict[str, Any]:
    """Extract a slide master with all layouts."""
    master_info = {
        "name": master.name if hasattr(master, 'name') else "Default",
        "placeholders": [],
        "shapes": [],
        "layouts": [],
    }

    # Extract master-level placeholders
    for placeholder in master.placeholders:
        ph_info = extract_placeholder_info(placeholder)
        master_info["placeholders"].append(ph_info)

    # Extract master-level shapes
    for shape in master.shapes:
        if not shape.is_placeholder:
            shape_info = extract_shape_info(shape)
            master_info["shapes"].append(shape_info)

    # Extract all layouts
    for layout in master.slide_layouts:
        layout_info = extract_layout(layout)
        master_info["layouts"].append(layout_info)

    return master_info


def extract_all_masters(pptx_path: Path) -> Dict[str, Any]:
    """
    Extract all master slides and layouts from a PPTX file.

    Returns a comprehensive dictionary with:
    - Slide dimensions
    - All masters with their layouts
    - Placeholder positions and formatting
    - Shape positions and styles
    """
    prs = Presentation(pptx_path)

    # Get slide dimensions
    width_inches = emu_to_inches(prs.slide_width)
    height_inches = emu_to_inches(prs.slide_height)

    result = {
        "source_file": pptx_path.name,
        "slide_dimensions": {
            "width_inches": width_inches,
            "height_inches": height_inches,
        },
        "masters": [],
        "layout_index": {},  # Quick lookup by layout name
    }

    # Extract each master
    for master in prs.slide_masters:
        master_info = extract_master(master)
        result["masters"].append(master_info)

        # Build layout index
        for i, layout in enumerate(master_info["layouts"]):
            result["layout_index"][layout["name"]] = {
                "master_idx": len(result["masters"]) - 1,
                "layout_idx": i,
            }

    # Map slide numbers to layouts
    result["slide_layout_mapping"] = {}
    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        result["slide_layout_mapping"][i + 1] = layout_name

    logger.info(f"Extracted {len(result['masters'])} masters with "
                f"{sum(len(m['layouts']) for m in result['masters'])} layouts")

    return result


def save_masters(pptx_path: Path, output_dir: Path = None) -> Path:
    """Extract and save masters to JSON file."""
    if output_dir is None:
        output_dir = Path("descriptions")
    output_dir.mkdir(parents=True, exist_ok=True)

    # Extract
    masters_data = extract_all_masters(pptx_path)

    # Save
    output_name = f"{pptx_path.stem}_masters_detailed.json"
    output_path = output_dir / output_name

    with open(output_path, 'w') as f:
        json.dump(masters_data, f, indent=2)

    logger.info(f"Saved masters to: {output_path}")
    return output_path


def get_layout_template(masters_data: Dict, layout_name: str) -> Optional[Dict]:
    """Get a layout template by name."""
    if layout_name not in masters_data.get("layout_index", {}):
        return None

    idx_info = masters_data["layout_index"][layout_name]
    master = masters_data["masters"][idx_info["master_idx"]]
    layout = master["layouts"][idx_info["layout_idx"]]

    return layout


def list_layouts(masters_data: Dict) -> List[str]:
    """List all available layout names."""
    return list(masters_data.get("layout_index", {}).keys())


if __name__ == "__main__":
    import sys

    logging.basicConfig(level=logging.INFO, format="%(message)s")

    if len(sys.argv) < 2:
        print("Usage: python master_extractor.py <template.pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        # Try in templates directory
        from config import TEMPLATE_DIR
        for template in TEMPLATE_DIR.glob("**/*.pptx"):
            if pptx_path.name in template.name:
                pptx_path = template
                break

    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    output = save_masters(pptx_path)

    # Print summary
    with open(output) as f:
        data = json.load(f)

    print(f"\nExtracted from: {data['source_file']}")
    print(f"Dimensions: {data['slide_dimensions']['width_inches']}\" x {data['slide_dimensions']['height_inches']}\"")
    print(f"\nLayouts available ({len(data['layout_index'])}):")
    for name in sorted(data['layout_index'].keys()):
        idx = data['layout_index'][name]
        layout = data['masters'][idx['master_idx']]['layouts'][idx['layout_idx']]
        print(f"  - {name}: {len(layout['placeholders'])} placeholders, {len(layout['shapes'])} shapes")
