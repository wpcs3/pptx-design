"""
PPTX Generator Module

Generates PowerPoint slides programmatically from natural language descriptions
using python-pptx.
"""
import logging
import re
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE_DASH_STYLE

# Add parent directory to path for config import
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))

from config import OUTPUT_DIR, DEFAULT_SLIDE_WIDTH_INCHES, DEFAULT_SLIDE_HEIGHT_INCHES

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class GeneratorError(Exception):
    """Exception raised when generation fails."""
    pass


# Mapping of shape type names to MSO_SHAPE constants
SHAPE_TYPE_MAP = {
    # Basic shapes
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "round_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,

    # Triangles
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "isosceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,

    # Polygons
    "pentagon": MSO_SHAPE.PENTAGON,
    "regular_pentagon": MSO_SHAPE.REGULAR_PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "heptagon": MSO_SHAPE.HEPTAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "decagon": MSO_SHAPE.DECAGON,
    "dodecagon": MSO_SHAPE.DODECAGON,

    # Block arrows
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "arrow_left": MSO_SHAPE.LEFT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "arrow_up": MSO_SHAPE.UP_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "arrow_down": MSO_SHAPE.DOWN_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "up_down_arrow": MSO_SHAPE.UP_DOWN_ARROW,
    "quad_arrow": MSO_SHAPE.QUAD_ARROW,
    "bent_arrow": MSO_SHAPE.BENT_ARROW,
    "bent_up_arrow": MSO_SHAPE.BENT_UP_ARROW,
    "curved_right_arrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
    "curved_left_arrow": MSO_SHAPE.CURVED_LEFT_ARROW,
    "curved_up_arrow": MSO_SHAPE.CURVED_UP_ARROW,
    "curved_down_arrow": MSO_SHAPE.CURVED_DOWN_ARROW,
    "circular_arrow": MSO_SHAPE.CIRCULAR_ARROW,
    "u_turn_arrow": MSO_SHAPE.U_TURN_ARROW,
    "notched_right_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "striped_right_arrow": MSO_SHAPE.STRIPED_RIGHT_ARROW,
    "swoosh_arrow": MSO_SHAPE.SWOOSH_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,

    # Stars
    "star": MSO_SHAPE.STAR_5_POINT,
    "star_4": MSO_SHAPE.STAR_4_POINT,
    "star_4_point": MSO_SHAPE.STAR_4_POINT,
    "star_5": MSO_SHAPE.STAR_5_POINT,
    "star_5_point": MSO_SHAPE.STAR_5_POINT,
    "star_6": MSO_SHAPE.STAR_6_POINT,
    "star_6_point": MSO_SHAPE.STAR_6_POINT,
    "star_7": MSO_SHAPE.STAR_7_POINT,
    "star_7_point": MSO_SHAPE.STAR_7_POINT,
    "star_8": MSO_SHAPE.STAR_8_POINT,
    "star_8_point": MSO_SHAPE.STAR_8_POINT,
    "star_10": MSO_SHAPE.STAR_10_POINT,
    "star_10_point": MSO_SHAPE.STAR_10_POINT,
    "star_12": MSO_SHAPE.STAR_12_POINT,
    "star_12_point": MSO_SHAPE.STAR_12_POINT,
    "star_16": MSO_SHAPE.STAR_16_POINT,
    "star_16_point": MSO_SHAPE.STAR_16_POINT,
    "star_24": MSO_SHAPE.STAR_24_POINT,
    "star_24_point": MSO_SHAPE.STAR_24_POINT,
    "star_32": MSO_SHAPE.STAR_32_POINT,
    "star_32_point": MSO_SHAPE.STAR_32_POINT,
    "explosion": MSO_SHAPE.EXPLOSION1,
    "explosion1": MSO_SHAPE.EXPLOSION1,
    "explosion2": MSO_SHAPE.EXPLOSION2,
    "sun": MSO_SHAPE.SUN,

    # Callouts
    "callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "rectangular_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "rounded_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    "rounded_rectangular_callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    "oval_callout": MSO_SHAPE.OVAL_CALLOUT,
    "cloud_callout": MSO_SHAPE.CLOUD_CALLOUT,
    "line_callout_1": MSO_SHAPE.LINE_CALLOUT_1,
    "line_callout_2": MSO_SHAPE.LINE_CALLOUT_2,
    "line_callout_3": MSO_SHAPE.LINE_CALLOUT_3,
    "line_callout_4": MSO_SHAPE.LINE_CALLOUT_4,

    # Flowchart shapes
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
    "flowchart_data": MSO_SHAPE.FLOWCHART_DATA,
    "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "flowchart_document": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "flowchart_multidocument": MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
    "flowchart_connector": MSO_SHAPE.FLOWCHART_CONNECTOR,
    "flowchart_offpage_connector": MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR,
    "flowchart_preparation": MSO_SHAPE.FLOWCHART_PREPARATION,
    "flowchart_manual_input": MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
    "flowchart_manual_operation": MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
    "flowchart_predefined_process": MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
    "flowchart_internal_storage": MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
    "flowchart_stored_data": MSO_SHAPE.FLOWCHART_STORED_DATA,
    "flowchart_delay": MSO_SHAPE.FLOWCHART_DELAY,
    "flowchart_alternate_process": MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS,
    "flowchart_merge": MSO_SHAPE.FLOWCHART_MERGE,
    "flowchart_extract": MSO_SHAPE.FLOWCHART_EXTRACT,
    "flowchart_sort": MSO_SHAPE.FLOWCHART_SORT,
    "flowchart_or": MSO_SHAPE.FLOWCHART_OR,
    "flowchart_summing_junction": MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
    "flowchart_display": MSO_SHAPE.FLOWCHART_DISPLAY,

    # Equation shapes
    "math_plus": MSO_SHAPE.MATH_PLUS,
    "math_minus": MSO_SHAPE.MATH_MINUS,
    "math_multiply": MSO_SHAPE.MATH_MULTIPLY,
    "math_divide": MSO_SHAPE.MATH_DIVIDE,
    "math_equal": MSO_SHAPE.MATH_EQUAL,
    "math_not_equal": MSO_SHAPE.MATH_NOT_EQUAL,

    # Special shapes
    "cross": MSO_SHAPE.CROSS,
    "plus": MSO_SHAPE.CROSS,
    "cube": MSO_SHAPE.CUBE,
    "can": MSO_SHAPE.CAN,
    "cylinder": MSO_SHAPE.CAN,
    "donut": MSO_SHAPE.DONUT,
    "ring": MSO_SHAPE.DONUT,
    "cloud": MSO_SHAPE.CLOUD,
    "heart": MSO_SHAPE.HEART,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
    "moon": MSO_SHAPE.MOON,
    "smiley_face": MSO_SHAPE.SMILEY_FACE,
    "smiley": MSO_SHAPE.SMILEY_FACE,
    "no_symbol": MSO_SHAPE.NO_SYMBOL,
    "prohibited": MSO_SHAPE.NO_SYMBOL,
    "tear": MSO_SHAPE.TEAR,
    "teardrop": MSO_SHAPE.TEAR,
    "wave": MSO_SHAPE.WAVE,
    "double_wave": MSO_SHAPE.DOUBLE_WAVE,
    "funnel": MSO_SHAPE.FUNNEL,
    "gear_6": MSO_SHAPE.GEAR_6,
    "gear_9": MSO_SHAPE.GEAR_9,
    "bevel": MSO_SHAPE.BEVEL,
    "frame": MSO_SHAPE.FRAME,
    "half_frame": MSO_SHAPE.HALF_FRAME,
    "corner": MSO_SHAPE.CORNER,
    "diagonal_stripe": MSO_SHAPE.DIAGONAL_STRIPE,
    "pie": MSO_SHAPE.PIE,
    "pie_wedge": MSO_SHAPE.PIE_WEDGE,
    "arc": MSO_SHAPE.ARC,
    "block_arc": MSO_SHAPE.BLOCK_ARC,
    "chord": MSO_SHAPE.CHORD,
    "plaque": MSO_SHAPE.PLAQUE,
    "folded_corner": MSO_SHAPE.FOLDED_CORNER,

    # Ribbons and banners
    "ribbon": MSO_SHAPE.UP_RIBBON,
    "up_ribbon": MSO_SHAPE.UP_RIBBON,
    "down_ribbon": MSO_SHAPE.DOWN_RIBBON,
    "curved_up_ribbon": MSO_SHAPE.CURVED_UP_RIBBON,
    "curved_down_ribbon": MSO_SHAPE.CURVED_DOWN_RIBBON,
    "left_right_ribbon": MSO_SHAPE.LEFT_RIGHT_RIBBON,
    "vertical_scroll": MSO_SHAPE.VERTICAL_SCROLL,
    "horizontal_scroll": MSO_SHAPE.HORIZONTAL_SCROLL,

    # Brackets and braces
    "left_brace": MSO_SHAPE.LEFT_BRACE,
    "right_brace": MSO_SHAPE.RIGHT_BRACE,
    "left_bracket": MSO_SHAPE.LEFT_BRACKET,
    "right_bracket": MSO_SHAPE.RIGHT_BRACKET,
    "double_brace": MSO_SHAPE.DOUBLE_BRACE,
    "double_bracket": MSO_SHAPE.DOUBLE_BRACKET,

    # Specialty rectangles
    "snip_rectangle": MSO_SHAPE.SNIP_1_RECTANGLE,
    "snip_1_rectangle": MSO_SHAPE.SNIP_1_RECTANGLE,
    "snip_2_rectangle": MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
    "snip_2_same_rectangle": MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
    "snip_2_diag_rectangle": MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
    "snip_round_rectangle": MSO_SHAPE.SNIP_ROUND_RECTANGLE,
    "round_1_rectangle": MSO_SHAPE.ROUND_1_RECTANGLE,
    "round_2_rectangle": MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
    "round_2_same_rectangle": MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
    "round_2_diag_rectangle": MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,

    # Action buttons
    "action_button_home": MSO_SHAPE.ACTION_BUTTON_HOME,
    "action_button_help": MSO_SHAPE.ACTION_BUTTON_HELP,
    "action_button_back": MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS,
    "action_button_forward": MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT,
    "action_button_beginning": MSO_SHAPE.ACTION_BUTTON_BEGINNING,
    "action_button_end": MSO_SHAPE.ACTION_BUTTON_END,
    "action_button_return": MSO_SHAPE.ACTION_BUTTON_RETURN,
    "action_button_document": MSO_SHAPE.ACTION_BUTTON_DOCUMENT,
    "action_button_sound": MSO_SHAPE.ACTION_BUTTON_SOUND,
    "action_button_movie": MSO_SHAPE.ACTION_BUTTON_MOVIE,
    "action_button_info": MSO_SHAPE.ACTION_BUTTON_INFORMATION,
    "action_button_custom": MSO_SHAPE.ACTION_BUTTON_CUSTOM,

    # Speech balloon
    "balloon": MSO_SHAPE.BALLOON,
    "speech_bubble": MSO_SHAPE.BALLOON,
}

# Mapping of alignment names to PP_ALIGN constants
ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "centre": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

# Mapping of vertical alignment names to MSO_ANCHOR constants
VERTICAL_ALIGNMENT_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "center": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

# Mapping of auto-size names to MSO_AUTO_SIZE constants
AUTO_SIZE_MAP = {
    "none": MSO_AUTO_SIZE.NONE,
    "shape_to_fit_text": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
    "text_to_fit_shape": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    "shrink_to_fit": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,  # Alias
}

# Font weight suffixes for common fonts
# Maps (base_font, weight) to the actual font name
FONT_WEIGHT_SUFFIXES = {
    "thin": " Thin",
    "extralight": " ExtraLight",
    "light": " Light",
    "regular": "",
    "medium": " Medium",
    "semibold": " SemiBold",
    "bold": " Bold",
    "extrabold": " ExtraBold",
    "black": " Black",
}


def get_font_with_weight(font_family: str, font_weight: str = None) -> tuple[str, bool]:
    """
    Get the actual font name based on font family and weight.

    Args:
        font_family: Base font family name (e.g., "Calibri", "Arial")
        font_weight: Weight name (thin, light, regular, medium, semibold, bold, black)

    Returns:
        Tuple of (actual_font_name, use_bold_flag)
        - For weights that map to font variants, returns (variant_name, False)
        - For "bold" weight without variant, returns (font_family, True)
    """
    if not font_weight:
        return font_family, False

    weight_lower = font_weight.lower().strip()

    # If bold is requested, check if font has Bold variant
    # Otherwise just use the bold flag
    if weight_lower == "bold":
        return font_family, True

    # For other weights, try to use font variant
    if weight_lower in FONT_WEIGHT_SUFFIXES:
        suffix = FONT_WEIGHT_SUFFIXES[weight_lower]
        if suffix:
            # Return the variant name, no bold flag
            return f"{font_family}{suffix}", False
        else:
            # Regular weight
            return font_family, False

    # Unknown weight, return as-is
    return font_family, False

# Mapping of underline style names to MSO_TEXT_UNDERLINE_TYPE constants
UNDERLINE_STYLE_MAP = {
    "none": MSO_TEXT_UNDERLINE_TYPE.NONE,
    "single": MSO_TEXT_UNDERLINE_TYPE.SINGLE_LINE,
    "double": MSO_TEXT_UNDERLINE_TYPE.DOUBLE_LINE,
    "heavy": MSO_TEXT_UNDERLINE_TYPE.HEAVY_LINE,
    "dotted": MSO_TEXT_UNDERLINE_TYPE.DOTTED_LINE,
    "dashed": MSO_TEXT_UNDERLINE_TYPE.DASH_LINE,
    "dot_dash": MSO_TEXT_UNDERLINE_TYPE.DOT_DASH_LINE,
    "dot_dot_dash": MSO_TEXT_UNDERLINE_TYPE.DOT_DOT_DASH_LINE,
    "wavy": MSO_TEXT_UNDERLINE_TYPE.WAVY_LINE,
    "heavy_wavy": MSO_TEXT_UNDERLINE_TYPE.WAVY_HEAVY_LINE,
    "wavy_double": MSO_TEXT_UNDERLINE_TYPE.WAVY_DOUBLE_LINE,
}


def parse_formatted_text(text: str) -> list[tuple[str, dict]]:
    """
    Parse text with inline formatting markers into segments.

    Supports:
    - **bold** or __bold__ for bold text
    - *italic* or _italic_ for italic text (single * or _ not adjacent to word chars)
    - ***bold italic*** for both
    - ^superscript^ for superscript
    - ~subscript~ for subscript

    Args:
        text: Text string potentially containing format markers

    Returns:
        List of (text_segment, format_dict) tuples where format_dict has
        keys 'bold', 'italic', 'superscript', 'subscript' with boolean values
    """
    if not text:
        return [("", {"bold": False, "italic": False, "superscript": False, "subscript": False})]

    segments = []
    # Pattern to match formatting markers
    # Order: bold+italic first, then bold, then italic, then superscript, then subscript
    pattern = r'(\*\*\*(.+?)\*\*\*|__(.+?)__|_([^_]+?)_|\*\*(.+?)\*\*|\*([^*]+?)\*|\^([^^]+?)\^|~([^~]+?)~)'

    last_end = 0
    for match in re.finditer(pattern, text):
        # Add any text before this match as plain text
        if match.start() > last_end:
            plain_text = text[last_end:match.start()]
            if plain_text:
                segments.append((plain_text, {"bold": False, "italic": False, "superscript": False, "subscript": False}))

        # Determine which group matched and its formatting
        full_match = match.group(0)
        base_format = {"bold": False, "italic": False, "superscript": False, "subscript": False}

        if full_match.startswith('***'):
            # Bold + italic
            content = match.group(2)
            segments.append((content, {"bold": True, "italic": True, "superscript": False, "subscript": False}))
        elif full_match.startswith('**'):
            # Bold
            content = match.group(5)
            segments.append((content or full_match, {"bold": True, "italic": False, "superscript": False, "subscript": False}))
        elif full_match.startswith('__'):
            # Bold (alternative)
            content = match.group(3)
            segments.append((content or full_match, {"bold": True, "italic": False, "superscript": False, "subscript": False}))
        elif full_match.startswith('*'):
            # Italic
            content = match.group(6)
            segments.append((content or full_match, {"bold": False, "italic": True, "superscript": False, "subscript": False}))
        elif full_match.startswith('_'):
            # Italic (alternative)
            content = match.group(4)
            segments.append((content or full_match, {"bold": False, "italic": True, "superscript": False, "subscript": False}))
        elif full_match.startswith('^'):
            # Superscript
            content = match.group(7)
            segments.append((content or full_match, {"bold": False, "italic": False, "superscript": True, "subscript": False}))
        elif full_match.startswith('~'):
            # Subscript
            content = match.group(8)
            segments.append((content or full_match, {"bold": False, "italic": False, "superscript": False, "subscript": True}))

        last_end = match.end()

    # Add any remaining text after the last match
    if last_end < len(text):
        remaining = text[last_end:]
        if remaining:
            segments.append((remaining, {"bold": False, "italic": False, "superscript": False, "subscript": False}))

    # If no matches found, return the original text as plain
    if not segments:
        segments.append((text, {"bold": False, "italic": False, "superscript": False, "subscript": False}))

    return segments


def apply_formatted_text_to_paragraph(
    paragraph,
    text: str,
    base_font_name: str = None,
    base_font_size_pt: float = None,
    base_font_color: RGBColor = None,
    base_bold: bool = False,
    base_italic: bool = False,
    base_underline: bool = False,
    base_underline_style: str = None,
    base_strikethrough: bool = False,
    base_double_strikethrough: bool = False,
    character_spacing_pt: float = None
) -> None:
    """
    Apply text with inline formatting to a paragraph using multiple runs.

    Args:
        paragraph: PowerPoint paragraph object
        text: Text string with optional **bold** and *italic* markers
        base_font_name: Default font family
        base_font_size_pt: Default font size in points
        base_font_color: Default font color
        base_bold: Whether base text should be bold
        base_italic: Whether base text should be italic
        base_underline: Whether base text should be underlined
        base_underline_style: Underline style (single, double, wavy, dotted, dashed, etc.)
        base_strikethrough: Whether base text should have strikethrough
        base_double_strikethrough: Whether base text should have double strikethrough
        character_spacing_pt: Character spacing in points (positive = expanded, negative = condensed)
    """
    # Clear existing text
    paragraph.clear()

    # Parse the formatted text
    segments = parse_formatted_text(text)

    for segment_text, formatting in segments:
        if not segment_text:
            continue

        # Add a run for this segment
        run = paragraph.add_run()
        run.text = segment_text

        # Apply base formatting
        if base_font_name:
            run.font.name = base_font_name
        if base_font_size_pt:
            run.font.size = Pt(base_font_size_pt)
        if base_font_color:
            run.font.color.rgb = base_font_color

        # Apply segment-specific formatting (combine with base)
        run.font.bold = base_bold or formatting.get("bold", False)
        run.font.italic = base_italic or formatting.get("italic", False)

        # Apply underline
        if base_underline or base_underline_style:
            if base_underline_style and base_underline_style.lower() in UNDERLINE_STYLE_MAP:
                run.font.underline = UNDERLINE_STYLE_MAP[base_underline_style.lower()]
            else:
                run.font.underline = True

        # Apply strikethrough
        if base_double_strikethrough:
            # Access via XML for double strikethrough
            from pptx.oxml.ns import qn
            rPr = run.font._element
            rPr.set(qn('a:dblStrike'), 'true')
        elif base_strikethrough:
            # Use the font._element to set strikethrough via XML
            from pptx.oxml.ns import qn
            rPr = run.font._element
            rPr.set(qn('a:strike'), 'sngStrike')

        # Apply superscript/subscript from inline formatting
        if formatting.get("superscript", False):
            # Superscript uses positive baseline (30% = 30000)
            from pptx.oxml.ns import qn
            rPr = run.font._element
            rPr.set(qn('a:baseline'), '30000')
        elif formatting.get("subscript", False):
            # Subscript uses negative baseline (-25% = -25000)
            from pptx.oxml.ns import qn
            rPr = run.font._element
            rPr.set(qn('a:baseline'), '-25000')

        # Apply character spacing (kerning)
        if character_spacing_pt is not None:
            from pptx.oxml.ns import qn
            # Character spacing in 100ths of a point
            spacing_hundredths = int(character_spacing_pt * 100)
            rPr = run.font._element
            rPr.set(qn('a:spc'), str(spacing_hundredths))


# Theme color name to MSO_THEME_COLOR mapping
THEME_COLOR_MAP = {
    "accent1": MSO_THEME_COLOR.ACCENT_1,
    "accent2": MSO_THEME_COLOR.ACCENT_2,
    "accent3": MSO_THEME_COLOR.ACCENT_3,
    "accent4": MSO_THEME_COLOR.ACCENT_4,
    "accent5": MSO_THEME_COLOR.ACCENT_5,
    "accent6": MSO_THEME_COLOR.ACCENT_6,
    "dark1": MSO_THEME_COLOR.DARK_1,
    "dark2": MSO_THEME_COLOR.DARK_2,
    "light1": MSO_THEME_COLOR.LIGHT_1,
    "light2": MSO_THEME_COLOR.LIGHT_2,
    "background1": MSO_THEME_COLOR.BACKGROUND_1,
    "background2": MSO_THEME_COLOR.BACKGROUND_2,
    "text1": MSO_THEME_COLOR.TEXT_1,
    "text2": MSO_THEME_COLOR.TEXT_2,
    "hyperlink": MSO_THEME_COLOR.HYPERLINK,
    "followed_hyperlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
}


class ThemeColor:
    """Represents a theme color reference that can be applied to elements."""

    def __init__(self, theme_color_name: str, luminance_mod: float = None, luminance_off: float = None):
        self.name = theme_color_name
        self.mso_value = THEME_COLOR_MAP.get(theme_color_name.lower())
        self.luminance_mod = luminance_mod  # e.g., 0.5 for 50% brightness
        self.luminance_off = luminance_off  # e.g., 0.4 for 40% lighter

    def __repr__(self):
        return f"ThemeColor({self.name})"


def is_theme_color(color_str: str) -> bool:
    """Check if a color string represents a theme color."""
    if not isinstance(color_str, str):
        return False
    color_lower = color_str.strip().lower()
    # Check for direct theme color names
    if color_lower in THEME_COLOR_MAP:
        return True
    # Check for theme: prefix
    if color_lower.startswith("theme:"):
        return True
    return False


def parse_theme_color(color_str: str) -> ThemeColor:
    """Parse a theme color string into a ThemeColor object."""
    color_str = str(color_str).strip().lower()

    # Handle theme: prefix (e.g., "theme:accent1", "theme:accent1:50%")
    if color_str.startswith("theme:"):
        parts = color_str[6:].split(":")
        theme_name = parts[0]
        luminance_mod = None
        luminance_off = None

        if len(parts) > 1:
            mod_str = parts[1]
            if mod_str.endswith("%"):
                luminance_mod = float(mod_str[:-1]) / 100.0

        return ThemeColor(theme_name, luminance_mod, luminance_off)

    # Direct theme color name
    return ThemeColor(color_str)


def apply_theme_color_to_fill(shape, theme_color: ThemeColor) -> None:
    """Apply a theme color to a shape's fill."""
    if theme_color.mso_value is None:
        logger.warning(f"Unknown theme color: {theme_color.name}")
        return

    shape.fill.solid()
    shape.fill.fore_color.theme_color = theme_color.mso_value

    # Apply luminance modifiers if specified
    if theme_color.luminance_mod is not None:
        # Brightness/luminance requires XML manipulation
        # This is a simplified version - full implementation would modify XML
        pass


def parse_color(color_str: str) -> RGBColor:
    """
    Convert various color formats to python-pptx RGBColor.

    Supports:
    - Hex format: "#RRGGBB", "RRGGBB", "#RGB"
    - RGB format: "rgb(R, G, B)"
    - Named colors: "red", "blue", etc.
    - Theme colors: "accent1", "theme:accent2", etc. (returns None - use is_theme_color first)

    Args:
        color_str: Color string in any supported format

    Returns:
        RGBColor object, or None for theme colors (check with is_theme_color first)

    Raises:
        GeneratorError: If color format is not recognized
    """
    if isinstance(color_str, RGBColor):
        return color_str

    if isinstance(color_str, ThemeColor):
        return None  # Theme colors handled separately

    color_str = str(color_str).strip().lower()

    # Handle special "none" / "transparent" values - return None to indicate no fill
    if color_str in ("none", "transparent", "null", ""):
        return None

    # Check for theme colors - return None, caller should use is_theme_color check
    if is_theme_color(color_str):
        return None

    # Named colors
    named_colors = {
        "white": (255, 255, 255),
        "black": (0, 0, 0),
        "red": (255, 0, 0),
        "green": (0, 128, 0),
        "blue": (0, 0, 255),
        "yellow": (255, 255, 0),
        "cyan": (0, 255, 255),
        "magenta": (255, 0, 255),
        "gray": (128, 128, 128),
        "grey": (128, 128, 128),
        "orange": (255, 165, 0),
        "purple": (128, 0, 128),
        "navy": (0, 0, 128),
        "teal": (0, 128, 128),
    }

    if color_str in named_colors:
        r, g, b = named_colors[color_str]
        return RGBColor(r, g, b)

    # Hex format: #RRGGBB or RRGGBB
    hex_match = re.match(r'^#?([0-9a-f]{6})$', color_str)
    if hex_match:
        hex_val = hex_match.group(1)
        r = int(hex_val[0:2], 16)
        g = int(hex_val[2:4], 16)
        b = int(hex_val[4:6], 16)
        return RGBColor(r, g, b)

    # Short hex format: #RGB
    short_hex_match = re.match(r'^#?([0-9a-f]{3})$', color_str)
    if short_hex_match:
        hex_val = short_hex_match.group(1)
        r = int(hex_val[0] * 2, 16)
        g = int(hex_val[1] * 2, 16)
        b = int(hex_val[2] * 2, 16)
        return RGBColor(r, g, b)

    # RGB format: rgb(R, G, B)
    rgb_match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color_str)
    if rgb_match:
        r = int(rgb_match.group(1))
        g = int(rgb_match.group(2))
        b = int(rgb_match.group(3))
        return RGBColor(min(r, 255), min(g, 255), min(b, 255))

    raise GeneratorError(f"Unrecognized color format: {color_str}")


def parse_measurement(
    measurement: Any,
    reference_size: int = None,
    default_unit: str = "inches"
) -> int:
    """
    Convert various measurement formats to EMUs (English Metric Units).

    Supports:
    - Inches: "2 inches", "2in", "2.5"
    - Points: "24pt", "24 pt"
    - Centimeters: "5cm", "5 cm"
    - Percentage: "50%"
    - EMU: "914400"
    - Numeric (treated as inches by default): 2.5

    Args:
        measurement: Measurement string or number
        reference_size: Reference size in EMUs for percentage calculations
        default_unit: Default unit if none specified (inches, pt, cm)

    Returns:
        Measurement in EMUs (integer)

    Raises:
        GeneratorError: If measurement format is not recognized
    """
    # Already an integer (assume EMUs if large, inches if small)
    if isinstance(measurement, int):
        if measurement > 10000:  # Probably EMUs
            return measurement
        return Inches(measurement)

    if isinstance(measurement, float):
        return Inches(measurement)

    measurement_str = str(measurement).strip().lower()

    # Inches: "2 inches", "2in", "2.5in"
    inches_match = re.match(r'^([\d.]+)\s*(inches?|in)?$', measurement_str)
    if inches_match:
        value = float(inches_match.group(1))
        return Inches(value)

    # Points: "24pt", "24 pt"
    pt_match = re.match(r'^([\d.]+)\s*pt$', measurement_str)
    if pt_match:
        value = float(pt_match.group(1))
        return Pt(value)

    # Centimeters: "5cm", "5 cm"
    cm_match = re.match(r'^([\d.]+)\s*cm$', measurement_str)
    if cm_match:
        value = float(cm_match.group(1))
        # 1 cm = 0.393701 inches
        return Inches(value * 0.393701)

    # Percentage: "50%"
    pct_match = re.match(r'^([\d.]+)\s*%$', measurement_str)
    if pct_match:
        if reference_size is None:
            raise GeneratorError(
                f"Cannot parse percentage '{measurement}' without reference size"
            )
        value = float(pct_match.group(1)) / 100
        return int(reference_size * value)

    # EMU: large number
    emu_match = re.match(r'^(\d+)$', measurement_str)
    if emu_match:
        value = int(emu_match.group(1))
        if value > 10000:  # Probably EMUs
            return value
        return Inches(value)

    raise GeneratorError(f"Unrecognized measurement format: {measurement}")


def create_presentation(
    width_inches: float = None,
    height_inches: float = None
) -> Presentation:
    """
    Create a new presentation with specified dimensions.

    Args:
        width_inches: Slide width in inches (default: 13.333 for 16:9)
        height_inches: Slide height in inches (default: 7.5 for 16:9)

    Returns:
        New Presentation object
    """
    if width_inches is None:
        width_inches = DEFAULT_SLIDE_WIDTH_INCHES
    if height_inches is None:
        height_inches = DEFAULT_SLIDE_HEIGHT_INCHES

    prs = Presentation()
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)

    return prs


def add_blank_slide(prs: Presentation) -> Any:
    """
    Add a blank slide to the presentation.

    Args:
        prs: Presentation object

    Returns:
        The new slide object
    """
    # Use blank layout (usually index 6)
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            blank_layout = layout
            break

    if blank_layout is None:
        # Fallback to last layout (usually blank)
        blank_layout = prs.slide_layouts[-1]

    return prs.slides.add_slide(blank_layout)


def apply_background(slide, background: dict) -> None:
    """
    Apply background to a slide.

    Args:
        slide: Slide object
        background: Background specification dict with keys:
            - type: "solid" | "gradient" | "image"
            - color: Hex color for solid
            - gradient_start, gradient_end: Colors for gradient
            - image_path: Path for image background
    """
    bg_type = background.get("type", "solid").lower()
    fill = slide.background.fill

    if bg_type == "solid":
        fill.solid()
        color = background.get("color", "#FFFFFF")
        fill.fore_color.rgb = parse_color(color)

    elif bg_type == "gradient":
        fill.gradient()
        fill.gradient_stops[0].color.rgb = parse_color(
            background.get("gradient_start", "#FFFFFF")
        )
        fill.gradient_stops[1].color.rgb = parse_color(
            background.get("gradient_end", "#000000")
        )

    # Note: Image backgrounds are more complex and would need additional handling


def apply_textbox(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a text box to a slide.

    Args:
        slide: Slide object
        element: Element specification with position and text properties
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    position = element.get("position", {})

    # Parse position
    left = parse_measurement(
        position.get("left_inches", 0),
        slide_width
    )
    top = parse_measurement(
        position.get("top_inches", 0),
        slide_height
    )
    width = parse_measurement(
        position.get("width_inches", 2),
        slide_width
    )
    height = parse_measurement(
        position.get("height_inches", 1),
        slide_height
    )

    # Create text box
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame

    # Get text properties - support both "text_properties" and "text_content" formats
    text_props = element.get("text_properties", {})
    text_content = element.get("text_content", {})

    # Merge text_content into text_props if text_props is empty (for LLM extraction compatibility)
    if not text_props and text_content:
        text_props = {
            "placeholder_text": text_content.get("text", ""),
            "font_family": text_content.get("font_name"),
            "font_size_pt": text_content.get("font_size_pt"),
            "font_color": text_content.get("font_color"),
            "bold": text_content.get("bold", False),
            "italic": text_content.get("italic", False),
            "alignment": text_content.get("alignment"),
        }

    # Text frame properties - word wrap
    if text_props.get("wrap_text") is not None:
        tf.word_wrap = text_props["wrap_text"]
    else:
        tf.word_wrap = True  # Default to True

    # Auto-size behavior
    auto_size_value = text_props.get("auto_size", "none").lower()
    if auto_size_value in AUTO_SIZE_MAP:
        tf.auto_size = AUTO_SIZE_MAP[auto_size_value]

    # Text frame margins/padding
    if text_props.get("margin_left_inches") is not None:
        tf.margin_left = Inches(text_props["margin_left_inches"])
    if text_props.get("margin_right_inches") is not None:
        tf.margin_right = Inches(text_props["margin_right_inches"])
    if text_props.get("margin_top_inches") is not None:
        tf.margin_top = Inches(text_props["margin_top_inches"])
    if text_props.get("margin_bottom_inches") is not None:
        tf.margin_bottom = Inches(text_props["margin_bottom_inches"])

    # Vertical alignment
    v_align = text_props.get("vertical_alignment", "top").lower()
    if v_align in VERTICAL_ALIGNMENT_MAP:
        tf.anchor = VERTICAL_ALIGNMENT_MAP[v_align]

    # Get or create paragraph
    p = tf.paragraphs[0]

    # Get text and formatting properties
    placeholder_text = text_props.get("placeholder_text", "")
    font_family = text_props.get("font_family")
    font_weight = text_props.get("font_weight")
    font_size_pt = text_props.get("font_size_pt")
    font_color = parse_color(text_props["font_color"]) if text_props.get("font_color") else None
    base_bold = text_props.get("bold", False)
    base_italic = text_props.get("italic", False)

    # Apply font weight to get actual font name and bold flag
    if font_family and font_weight:
        font_family, weight_bold = get_font_with_weight(font_family, font_weight)
        base_bold = base_bold or weight_bold
    base_underline = text_props.get("underline", False)
    base_underline_style = text_props.get("underline_style")
    base_strikethrough = text_props.get("strikethrough", False)
    base_double_strikethrough = text_props.get("double_strikethrough", False)
    character_spacing_pt = text_props.get("character_spacing_pt")

    # Alignment
    alignment = text_props.get("alignment", "left").lower()
    if alignment in ALIGNMENT_MAP:
        p.alignment = ALIGNMENT_MAP[alignment]

    # Paragraph spacing
    if text_props.get("space_before_pt"):
        p.space_before = Pt(text_props["space_before_pt"])
    if text_props.get("space_after_pt"):
        p.space_after = Pt(text_props["space_after_pt"])

    # Line spacing
    if text_props.get("line_spacing_pt"):
        # Fixed line spacing in points
        p.line_spacing = Pt(text_props["line_spacing_pt"])
    elif text_props.get("line_spacing_multiple"):
        # Relative line spacing (1.0 = single, 1.5 = 1.5x, 2.0 = double)
        p.line_spacing = text_props["line_spacing_multiple"]

    # Check if text contains inline formatting markers
    has_inline_formatting = any(marker in placeholder_text for marker in ['**', '__', '*', '_', '^', '~'])

    # Use formatted text parser if inline formatting or underline/strikethrough or character spacing
    needs_runs = has_inline_formatting or base_underline or base_underline_style or base_strikethrough or base_double_strikethrough or character_spacing_pt

    if needs_runs:
        # Use formatted text parser for inline bold/italic and underline/strikethrough
        apply_formatted_text_to_paragraph(
            p,
            placeholder_text,
            base_font_name=font_family,
            base_font_size_pt=font_size_pt,
            base_font_color=font_color,
            base_bold=base_bold,
            base_italic=base_italic,
            base_underline=base_underline,
            base_underline_style=base_underline_style,
            base_strikethrough=base_strikethrough,
            base_double_strikethrough=base_double_strikethrough,
            character_spacing_pt=character_spacing_pt
        )
    else:
        # Simple text without inline formatting
        p.text = placeholder_text

        # Font properties applied to entire paragraph
        font = p.font
        if font_family:
            font.name = font_family
        if font_size_pt:
            font.size = Pt(font_size_pt)
        if font_color:
            font.color.rgb = font_color
        if base_bold:
            font.bold = True
        if base_italic:
            font.italic = True

    # Shape fill (background of text box)
    shape_props = element.get("shape_properties", {})
    if shape_props.get("fill_color"):
        fill_color = parse_color(shape_props["fill_color"])
        if fill_color is not None:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = fill_color
        else:
            textbox.fill.background()  # No fill / transparent

        # Handle transparency
        if "fill_transparency" in shape_props:
            transparency = shape_props["fill_transparency"]
            if transparency > 0:
                # Note: python-pptx has limited transparency support
                pass

    # Border
    if shape_props.get("border_color"):
        textbox.line.color.rgb = parse_color(shape_props["border_color"])
    if shape_props.get("border_width_pt"):
        textbox.line.width = Pt(shape_props["border_width_pt"])

    # Rotation
    rotation_degrees = position.get("rotation_degrees", text_props.get("rotation_degrees"))
    if rotation_degrees is not None:
        textbox.rotation = float(rotation_degrees)

    # Shadow effects
    shadow_config = shape_props.get("shadow", text_props.get("shadow"))
    if shadow_config:
        apply_shadow(textbox, shadow_config)


def apply_bullet_list(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a bullet list to a slide.

    Args:
        slide: Slide object
        element: Element specification with:
            - position: left_inches, top_inches, width_inches, height_inches
            - items: List of {"text": str, "level": int} dicts
            - bullet_properties: bullet_type, bullet_char, bullet_color, indent_inches
            - text_properties: font settings applied to all items
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    position = element.get("position", {})

    # Parse position
    left = parse_measurement(position.get("left_inches", 0), slide_width)
    top = parse_measurement(position.get("top_inches", 0), slide_height)
    width = parse_measurement(position.get("width_inches", 2), slide_width)
    height = parse_measurement(position.get("height_inches", 1), slide_height)

    # Create text box
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    # Get properties
    items = element.get("items", [])
    bullet_props = element.get("bullet_properties", {})
    text_props = element.get("text_properties", {})

    # Text formatting
    font_family = text_props.get("font_family")
    font_size_pt = text_props.get("font_size_pt", 14)
    font_color = parse_color(text_props["font_color"]) if text_props.get("font_color") else None

    # Bullet properties
    bullet_type = bullet_props.get("bullet_type", "bullet")  # bullet, numbered, letter, roman, none
    bullet_char = bullet_props.get("bullet_char", "•")
    bullet_color = parse_color(bullet_props["bullet_color"]) if bullet_props.get("bullet_color") else None
    indent_inches = bullet_props.get("indent_inches", 0.25)
    number_format = bullet_props.get("number_format", "arabic_period")  # arabic_period, arabic_paren, roman_uc_period, etc.
    start_at = bullet_props.get("start_at", 1)

    # Margins
    if text_props.get("margin_left_inches") is not None:
        tf.margin_left = Inches(text_props["margin_left_inches"])
    if text_props.get("margin_top_inches") is not None:
        tf.margin_top = Inches(text_props["margin_top_inches"])

    # Add each item as a paragraph
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Item text and level
        item_text = item.get("text", "") if isinstance(item, dict) else str(item)
        item_level = item.get("level", 0) if isinstance(item, dict) else 0

        # Set paragraph level for indentation
        p.level = item_level

        # Set text
        p.text = item_text

        # Apply bullet/numbering
        if bullet_type == "bullet":
            p.bullet = True
            # Custom bullet character
            if bullet_char and bullet_char != "•":
                # Access XML to set custom bullet character
                from pptx.oxml.ns import qn
                from pptx.oxml import parse_xml
                pPr = p._pPr
                # Remove existing buNone if present
                buNone = pPr.find(qn('a:buNone'))
                if buNone is not None:
                    pPr.remove(buNone)
                # Add buChar element
                buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>')
                pPr.append(buChar)
        elif bullet_type in ("numbered", "number", "arabic"):
            # Numbered list - arabic numerals
            from pptx.oxml.ns import qn
            from pptx.oxml import parse_xml
            pPr = p._pPr
            # Remove existing bullet elements
            for elem_name in ['a:buNone', 'a:buChar', 'a:buAutoNum']:
                elem = pPr.find(qn(elem_name))
                if elem is not None:
                    pPr.remove(elem)
            # Map number_format to type attribute
            format_map = {
                "arabic_period": "arabicPeriod",      # 1. 2. 3.
                "arabic_paren": "arabicParenR",       # 1) 2) 3)
                "arabic_plain": "arabicPlain",       # 1 2 3
                "roman_uc_period": "romanUcPeriod",   # I. II. III.
                "roman_lc_period": "romanLcPeriod",   # i. ii. iii.
                "alpha_uc_period": "alphaUcPeriod",   # A. B. C.
                "alpha_lc_period": "alphaLcPeriod",   # a. b. c.
                "alpha_uc_paren": "alphaUcParenR",    # A) B) C)
                "alpha_lc_paren": "alphaLcParenR",    # a) b) c)
            }
            autonumType = format_map.get(number_format, "arabicPeriod")
            buAutoNum = parse_xml(f'<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="{autonumType}" startAt="{start_at}"/>')
            pPr.append(buAutoNum)
        elif bullet_type in ("letter", "alpha"):
            # Letter list - a. b. c.
            from pptx.oxml.ns import qn
            from pptx.oxml import parse_xml
            pPr = p._pPr
            for elem_name in ['a:buNone', 'a:buChar', 'a:buAutoNum']:
                elem = pPr.find(qn(elem_name))
                if elem is not None:
                    pPr.remove(elem)
            autonumType = "alphaLcPeriod" if number_format.startswith("alpha_lc") else "alphaUcPeriod"
            buAutoNum = parse_xml(f'<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="{autonumType}" startAt="{start_at}"/>')
            pPr.append(buAutoNum)
        elif bullet_type == "roman":
            # Roman numeral list - I. II. III.
            from pptx.oxml.ns import qn
            from pptx.oxml import parse_xml
            pPr = p._pPr
            for elem_name in ['a:buNone', 'a:buChar', 'a:buAutoNum']:
                elem = pPr.find(qn(elem_name))
                if elem is not None:
                    pPr.remove(elem)
            autonumType = "romanLcPeriod" if number_format.startswith("roman_lc") else "romanUcPeriod"
            buAutoNum = parse_xml(f'<a:buAutoNum xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="{autonumType}" startAt="{start_at}"/>')
            pPr.append(buAutoNum)
        elif bullet_type == "none":
            p.bullet = False

        # Set font properties
        font = p.font
        if font_family:
            font.name = font_family
        if font_size_pt:
            font.size = Pt(font_size_pt)
        if font_color:
            font.color.rgb = font_color

        # Item-specific formatting
        if isinstance(item, dict):
            if item.get("bold"):
                font.bold = True
            if item.get("italic"):
                font.italic = True

    # Shape fill (background)
    shape_props = element.get("shape_properties", {})
    if shape_props.get("fill_color"):
        fill_color = parse_color(shape_props["fill_color"])
        if fill_color is not None:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = fill_color
        else:
            textbox.fill.background()


def apply_table(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a table to a slide.

    Args:
        slide: Slide object
        element: Element specification with:
            - position: left_inches, top_inches, width_inches, height_inches
            - rows: Number of rows
            - cols: Number of columns
            - data: 2D list of cell values
            - table_style: Styling options
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    position = element.get("position", {})

    # Parse position
    left = parse_measurement(position.get("left_inches", 0), slide_width)
    top = parse_measurement(position.get("top_inches", 0), slide_height)
    width = parse_measurement(position.get("width_inches", 6), slide_width)
    height = parse_measurement(position.get("height_inches", 2), slide_height)

    # Get table dimensions
    data = element.get("data", [])
    rows = element.get("rows", len(data) if data else 2)
    cols = element.get("cols", len(data[0]) if data and data[0] else 3)

    # Create table
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Get styling options
    table_style = element.get("table_style", {})
    first_row_header = table_style.get("first_row_header", False)
    banded_rows = table_style.get("banded_rows", False)
    header_fill_color = parse_color(table_style["header_fill_color"]) if table_style.get("header_fill_color") else None
    header_font_color = parse_color(table_style["header_font_color"]) if table_style.get("header_font_color") else None
    row_fill_colors = table_style.get("row_fill_colors", [])  # List of colors for alternating rows
    border_color = parse_color(table_style["border_color"]) if table_style.get("border_color") else None
    border_width_pt = table_style.get("border_width_pt", 1)
    font_family = table_style.get("font_family")
    font_size_pt = table_style.get("font_size_pt", 11)
    font_color = parse_color(table_style["font_color"]) if table_style.get("font_color") else None

    # Populate table with data
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)

            # Get cell text
            if data and row_idx < len(data) and col_idx < len(data[row_idx]):
                cell_text = str(data[row_idx][col_idx])
            else:
                cell_text = ""

            cell.text = cell_text

            # Apply cell formatting
            # Cell fill
            is_header = first_row_header and row_idx == 0

            if is_header and header_fill_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill_color
            elif banded_rows and row_fill_colors:
                # Alternating row colors
                data_row_idx = row_idx - (1 if first_row_header else 0)
                if data_row_idx >= 0:
                    color_idx = data_row_idx % len(row_fill_colors)
                    row_color = parse_color(row_fill_colors[color_idx])
                    if row_color:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = row_color

            # Text formatting
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if font_family:
                        run.font.name = font_family
                    if font_size_pt:
                        run.font.size = Pt(font_size_pt)

                    # Header-specific formatting
                    if is_header:
                        if header_font_color:
                            run.font.color.rgb = header_font_color
                        run.font.bold = True
                    elif font_color:
                        run.font.color.rgb = font_color

    # Apply cell-specific formatting
    cells_config = element.get("cells", [])
    for cell_config in cells_config:
        row_idx = cell_config.get("row", 0)
        col_idx = cell_config.get("col", 0)

        if row_idx >= rows or col_idx >= cols:
            continue

        cell = table.cell(row_idx, col_idx)

        # Merge cells
        if cell_config.get("merge_right"):
            merge_cols = cell_config["merge_right"]
            end_col = min(col_idx + merge_cols, cols - 1)
            if end_col > col_idx:
                cell.merge(table.cell(row_idx, end_col))

        if cell_config.get("merge_down"):
            merge_rows = cell_config["merge_down"]
            end_row = min(row_idx + merge_rows, rows - 1)
            if end_row > row_idx:
                cell.merge(table.cell(end_row, col_idx))

        # Cell-specific fill color
        if cell_config.get("fill_color"):
            cell_fill = parse_color(cell_config["fill_color"])
            if cell_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fill

        # Cell-specific text formatting
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                if cell_config.get("bold"):
                    run.font.bold = True
                if cell_config.get("italic"):
                    run.font.italic = True
                if cell_config.get("font_color"):
                    cell_font_color = parse_color(cell_config["font_color"])
                    if cell_font_color:
                        run.font.color.rgb = cell_font_color
                if cell_config.get("font_size_pt"):
                    run.font.size = Pt(cell_config["font_size_pt"])

        # Cell text alignment
        if cell_config.get("alignment"):
            align_str = cell_config["alignment"].lower()
            if align_str in ALIGNMENT_MAP:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = ALIGNMENT_MAP[align_str]

        # Vertical alignment
        if cell_config.get("vertical_alignment"):
            v_align = cell_config["vertical_alignment"].lower()
            if v_align in VERTICAL_ALIGNMENT_MAP:
                cell.text_frame.anchor = VERTICAL_ALIGNMENT_MAP[v_align]


def apply_image(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add an image to a slide.

    Args:
        slide: Slide object
        element: Element specification with:
            - position: left_inches, top_inches, width_inches, height_inches
            - image_properties:
                - source: Local file path
                - source_base64: Base64-encoded image data
                - maintain_aspect_ratio: Whether to preserve aspect ratio
                - crop_left, crop_right, crop_top, crop_bottom: Crop percentages (0-1)
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    import base64
    import tempfile
    import os

    position = element.get("position", {})
    image_props = element.get("image_properties", {})

    # Parse position
    left = parse_measurement(position.get("left_inches", 0), slide_width)
    top = parse_measurement(position.get("top_inches", 0), slide_height)
    width = parse_measurement(position.get("width_inches"), slide_width) if position.get("width_inches") else None
    height = parse_measurement(position.get("height_inches"), slide_height) if position.get("height_inches") else None

    # Get image source
    source = image_props.get("source")
    source_base64 = image_props.get("source_base64")
    maintain_aspect_ratio = image_props.get("maintain_aspect_ratio", True)

    image_path = None
    temp_file = None

    try:
        if source:
            # Local file path
            image_path = Path(source)
            if not image_path.exists():
                logger.warning(f"Image not found: {source}")
                return

        elif source_base64:
            # Base64 encoded image
            # Strip data URL prefix if present
            if "," in source_base64:
                source_base64 = source_base64.split(",", 1)[1]

            # Decode and save to temp file
            image_data = base64.b64decode(source_base64)

            # Detect image type from header
            if image_data[:8] == b'\x89PNG\r\n\x1a\n':
                ext = ".png"
            elif image_data[:2] == b'\xff\xd8':
                ext = ".jpg"
            elif image_data[:4] == b'GIF8':
                ext = ".gif"
            else:
                ext = ".png"  # Default

            temp_file = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            temp_file.write(image_data)
            temp_file.close()
            image_path = Path(temp_file.name)

        else:
            logger.warning("No image source provided")
            return

        # Add picture
        if maintain_aspect_ratio:
            # Add with one dimension, let python-pptx maintain aspect ratio
            if width and not height:
                picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
            elif height and not width:
                picture = slide.shapes.add_picture(str(image_path), left, top, height=height)
            elif width and height:
                # Add with both, then adjust
                picture = slide.shapes.add_picture(str(image_path), left, top)
                # Scale to fit within bounds while maintaining aspect ratio
                img_width = picture.width
                img_height = picture.height
                scale_w = width / img_width
                scale_h = height / img_height
                scale = min(scale_w, scale_h)
                picture.width = int(img_width * scale)
                picture.height = int(img_height * scale)
            else:
                picture = slide.shapes.add_picture(str(image_path), left, top)
        else:
            # Stretch to exact dimensions
            if width and height:
                picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
            elif width:
                picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
            elif height:
                picture = slide.shapes.add_picture(str(image_path), left, top, height=height)
            else:
                picture = slide.shapes.add_picture(str(image_path), left, top)

        # Apply cropping if specified
        crop_left = image_props.get("crop_left", 0)
        crop_right = image_props.get("crop_right", 0)
        crop_top = image_props.get("crop_top", 0)
        crop_bottom = image_props.get("crop_bottom", 0)

        if any([crop_left, crop_right, crop_top, crop_bottom]):
            # Access crop via XML
            from pptx.oxml.ns import qn
            spPr = picture._element.spPr
            blipFill = picture._element.find(qn('p:blipFill'))
            if blipFill is not None:
                srcRect = blipFill.find(qn('a:srcRect'))
                if srcRect is None:
                    from pptx.oxml import parse_xml
                    srcRect = parse_xml('<a:srcRect xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
                    blipFill.append(srcRect)
                # Crop values are in percentages * 100000
                if crop_left:
                    srcRect.set('l', str(int(crop_left * 100000)))
                if crop_right:
                    srcRect.set('r', str(int(crop_right * 100000)))
                if crop_top:
                    srcRect.set('t', str(int(crop_top * 100000)))
                if crop_bottom:
                    srcRect.set('b', str(int(crop_bottom * 100000)))

        # Apply rotation
        rotation_degrees = position.get("rotation_degrees", image_props.get("rotation_degrees"))
        if rotation_degrees is not None:
            picture.rotation = float(rotation_degrees)

    finally:
        # Clean up temp file
        if temp_file:
            try:
                os.unlink(temp_file.name)
            except:
                pass


def apply_shape(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a shape to a slide.

    Args:
        slide: Slide object
        element: Element specification
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    position = element.get("position", {})
    shape_props = element.get("shape_properties", {})

    # Parse position
    left = parse_measurement(position.get("left_inches", 0), slide_width)
    top = parse_measurement(position.get("top_inches", 0), slide_height)
    width = parse_measurement(position.get("width_inches", 2), slide_width)
    height = parse_measurement(position.get("height_inches", 1), slide_height)

    # Determine shape type
    shape_type_str = shape_props.get("shape_type", "rectangle").lower()
    shape_type = SHAPE_TYPE_MAP.get(shape_type_str, MSO_SHAPE.RECTANGLE)

    # Create shape
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # Corner radius for rounded rectangles
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        corner_radius = shape_props.get("corner_radius_percent")
        if corner_radius is not None:
            # Corner radius is stored as a percentage (0-50 in user input, 0-0.5 for adjustments)
            # Convert to adjustment value (0-0.5 range)
            adj_value = min(0.5, max(0, corner_radius / 100))
            try:
                shape.adjustments[0] = adj_value
            except Exception as e:
                logger.warning(f"Failed to set corner radius: {e}")

    # Fill - check for different fill types
    fill_type = shape_props.get("fill_type", "solid")

    if fill_type == "gradient" or shape_props.get("gradient_stops") or \
       (shape_props.get("gradient_start") and shape_props.get("gradient_end")):
        # Multi-stop gradient fill
        gradient_config = shape_props.get("gradient", {})
        if not gradient_config:
            # Build config from individual properties
            gradient_config = {
                "stops": shape_props.get("gradient_stops", []),
                "angle_degrees": shape_props.get("gradient_angle", 0),
                "type": shape_props.get("gradient_type", "linear"),
                "start_color": shape_props.get("gradient_start"),
                "end_color": shape_props.get("gradient_end"),
            }
        apply_gradient_fill(shape, gradient_config)

    elif fill_type == "pattern" or shape_props.get("pattern"):
        # Pattern fill
        pattern_config = shape_props.get("pattern_fill", {})
        if not pattern_config:
            # Build config from individual properties
            pattern_config = {
                "pattern": shape_props.get("pattern", "horizontal"),
                "foreground_color": shape_props.get("pattern_fg_color", "#000000"),
                "background_color": shape_props.get("pattern_bg_color", "#FFFFFF"),
            }
        apply_pattern_fill(shape, pattern_config)

    elif fill_type == "picture" and shape_props.get("fill_image"):
        # Picture fill
        import base64
        import tempfile
        import os

        fill_image = shape_props["fill_image"]
        fill_image_base64 = shape_props.get("fill_image_base64")

        try:
            if fill_image_base64:
                # Base64 encoded image
                if "," in fill_image_base64:
                    fill_image_base64 = fill_image_base64.split(",", 1)[1]
                image_data = base64.b64decode(fill_image_base64)
                temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                temp_file.write(image_data)
                temp_file.close()
                fill_image = temp_file.name

            image_path = Path(fill_image)
            if image_path.exists():
                # Use blip fill for picture
                from pptx.oxml.ns import qn
                from pptx.oxml import parse_xml
                from pptx.parts.image import Image

                # Add image to the part
                image_part, rId = shape.part.relate_to(
                    Image.from_file(str(image_path)),
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                )

                # Create blipFill XML
                blipFill = parse_xml(f'''
                    <a:blipFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                        <a:blip r:embed="{rId}"/>
                        <a:stretch>
                            <a:fillRect/>
                        </a:stretch>
                    </a:blipFill>
                ''')

                # Replace existing fill with blipFill
                spPr = shape._element.spPr
                # Remove existing fill
                for fill_elem in spPr.findall(qn('a:solidFill')):
                    spPr.remove(fill_elem)
                for fill_elem in spPr.findall(qn('a:noFill')):
                    spPr.remove(fill_elem)
                for fill_elem in spPr.findall(qn('a:blipFill')):
                    spPr.remove(fill_elem)
                spPr.insert(0, blipFill)

            if fill_image_base64 and 'temp_file' in locals():
                os.unlink(temp_file.name)
        except Exception as e:
            logger.warning(f"Failed to apply picture fill: {e}")
            # Fall back to solid fill
            if shape_props.get("fill_color"):
                fill_color = parse_color(shape_props["fill_color"])
                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill_color

    elif shape_props.get("fill_color"):
        fill_color_str = shape_props["fill_color"]
        # Check if it's a theme color
        if is_theme_color(fill_color_str):
            theme_color = parse_theme_color(fill_color_str)
            apply_theme_color_to_fill(shape, theme_color)
        else:
            fill_color = parse_color(fill_color_str)
            if fill_color is not None:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
            else:
                shape.fill.background()  # No fill / transparent
    elif shape_props.get("no_fill"):
        shape.fill.background()

    # Border/line
    if shape_props.get("border_color"):
        border_color = parse_color(shape_props["border_color"])
        if border_color is not None:
            shape.line.color.rgb = border_color
    if shape_props.get("border_width_pt"):
        shape.line.width = Pt(shape_props["border_width_pt"])
    if shape_props.get("no_border"):
        shape.line.fill.background()

    # Transparency support
    fill_transparency = shape_props.get("fill_transparency", shape_props.get("transparency"))
    if fill_transparency:
        apply_fill_transparency(shape, float(fill_transparency))

    line_transparency = shape_props.get("line_transparency", shape_props.get("border_transparency"))
    if line_transparency:
        apply_line_transparency(shape, float(line_transparency))

    # Rotation
    rotation_degrees = position.get("rotation_degrees", shape_props.get("rotation_degrees"))
    if rotation_degrees is not None:
        shape.rotation = float(rotation_degrees)

    # Shadow effects
    shadow_config = shape_props.get("shadow")
    if shadow_config:
        apply_shadow(shape, shadow_config)

    # Hyperlink
    hyperlink_config = shape_props.get("hyperlink")
    if hyperlink_config:
        if isinstance(hyperlink_config, str):
            # Simple URL string
            hyperlink_config = {"url": hyperlink_config}
        apply_hyperlink(shape, hyperlink_config)

    # Shape action (alternative to hyperlink for navigation actions)
    action_config = shape_props.get("action", shape_props.get("click_action"))
    if action_config:
        if isinstance(action_config, str):
            action_config = {"type": action_config}
        apply_shape_action(shape, action_config)


def apply_shadow(shape, shadow_config: dict) -> None:
    """
    Apply shadow effect to a shape.

    Args:
        shape: PowerPoint shape object
        shadow_config: Shadow configuration dict with:
            - type: "outer" (default), "inner", or "perspective"
            - color: Shadow color (hex or named)
            - transparency: 0-1 (0 = opaque, 1 = fully transparent)
            - blur_radius_pt: Blur radius in points
            - distance_pt: Distance from shape in points
            - angle_degrees: Angle of shadow direction
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml

    shadow_type = shadow_config.get("type", "outer")
    color_str = shadow_config.get("color", "#000000")
    transparency = shadow_config.get("transparency", 0.6)
    blur_radius_pt = shadow_config.get("blur_radius_pt", 4)
    distance_pt = shadow_config.get("distance_pt", 3)
    angle_degrees = shadow_config.get("angle_degrees", 45)

    # Parse color
    color = parse_color(color_str)
    if color is None:
        color = RGBColor(0, 0, 0)

    # Convert measurements to EMUs (1 point = 12700 EMUs)
    blur_rad = int(blur_radius_pt * 12700)
    dist = int(distance_pt * 12700)

    # Convert angle to 60000ths of a degree (PowerPoint units)
    angle = int(angle_degrees * 60000)

    # Calculate alpha (transparency to alpha: 0% transparency = 100% alpha)
    alpha_val = int((1 - transparency) * 100000)

    # Format color as hex
    color_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

    try:
        spPr = shape._element.spPr

        # Remove existing effect list if present
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is not None:
            spPr.remove(effectLst)

        if shadow_type == "outer":
            # Outer shadow (drop shadow)
            effect_xml = f'''
                <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:outerShdw blurRad="{blur_rad}" dist="{dist}" dir="{angle}" algn="bl" rotWithShape="0">
                        <a:srgbClr val="{color_hex}">
                            <a:alpha val="{alpha_val}"/>
                        </a:srgbClr>
                    </a:outerShdw>
                </a:effectLst>
            '''
        elif shadow_type == "inner":
            # Inner shadow
            effect_xml = f'''
                <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:innerShdw blurRad="{blur_rad}" dist="{dist}" dir="{angle}">
                        <a:srgbClr val="{color_hex}">
                            <a:alpha val="{alpha_val}"/>
                        </a:srgbClr>
                    </a:innerShdw>
                </a:effectLst>
            '''
        else:
            # Default to outer shadow
            effect_xml = f'''
                <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                    <a:outerShdw blurRad="{blur_rad}" dist="{dist}" dir="{angle}" algn="bl" rotWithShape="0">
                        <a:srgbClr val="{color_hex}">
                            <a:alpha val="{alpha_val}"/>
                        </a:srgbClr>
                    </a:outerShdw>
                </a:effectLst>
            '''

        effectLst = parse_xml(effect_xml)
        spPr.append(effectLst)

    except Exception as e:
        logger.warning(f"Failed to apply shadow: {e}")


def apply_gradient_fill(shape, gradient_config: dict) -> None:
    """
    Apply a multi-stop gradient fill to a shape using XML manipulation.

    Args:
        shape: The shape to apply gradient to
        gradient_config: Gradient configuration dict containing:
            - stops: List of dicts with 'color' and 'position' (0.0 to 1.0)
            - angle_degrees: Gradient angle (0=right, 90=up, 180=left, 270=down)
            - type: 'linear' (default), 'radial', 'rectangular', 'path'
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml

    stops = gradient_config.get("stops", [])
    if not stops:
        # Fallback: check for simple start/end colors
        start_color = gradient_config.get("start_color", gradient_config.get("gradient_start"))
        end_color = gradient_config.get("end_color", gradient_config.get("gradient_end"))
        if start_color and end_color:
            stops = [
                {"color": start_color, "position": 0.0},
                {"color": end_color, "position": 1.0}
            ]
        else:
            logger.warning("No gradient stops provided")
            return

    angle_degrees = gradient_config.get("angle_degrees", 0)
    gradient_type = gradient_config.get("type", "linear").lower()

    # Convert angle to PowerPoint units (60000ths of a degree)
    # PowerPoint uses 90 degrees offset from standard
    angle_ppt = int(angle_degrees * 60000)

    try:
        spPr = shape._element.spPr

        # Remove existing fills
        for fill_elem in spPr.findall(qn('a:solidFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:noFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:gradFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:blipFill')):
            spPr.remove(fill_elem)

        # Build gradient stops XML
        stops_xml = ""
        for stop in stops:
            color = parse_color(stop.get("color", "#000000"))
            if color is None:
                color = RGBColor(0, 0, 0)
            color_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
            position = int(stop.get("position", 0) * 100000)  # Convert to 100ths of percent
            transparency = stop.get("transparency", 0)
            alpha_val = int((1 - transparency) * 100000)

            if transparency > 0:
                stops_xml += f'''
                    <a:gs pos="{position}">
                        <a:srgbClr val="{color_hex}">
                            <a:alpha val="{alpha_val}"/>
                        </a:srgbClr>
                    </a:gs>'''
            else:
                stops_xml += f'''
                    <a:gs pos="{position}">
                        <a:srgbClr val="{color_hex}"/>
                    </a:gs>'''

        # Build full gradient XML based on type
        if gradient_type == "radial":
            grad_xml = f'''
                <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
                    <a:gsLst>{stops_xml}
                    </a:gsLst>
                    <a:path path="circle">
                        <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
                    </a:path>
                </a:gradFill>
            '''
        elif gradient_type == "rectangular":
            grad_xml = f'''
                <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
                    <a:gsLst>{stops_xml}
                    </a:gsLst>
                    <a:path path="rect">
                        <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
                    </a:path>
                </a:gradFill>
            '''
        elif gradient_type == "path":
            grad_xml = f'''
                <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
                    <a:gsLst>{stops_xml}
                    </a:gsLst>
                    <a:path path="shape">
                        <a:fillToRect l="50000" t="50000" r="50000" b="50000"/>
                    </a:path>
                </a:gradFill>
            '''
        else:  # linear (default)
            grad_xml = f'''
                <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
                    <a:gsLst>{stops_xml}
                    </a:gsLst>
                    <a:lin ang="{angle_ppt}" scaled="1"/>
                </a:gradFill>
            '''

        gradFill = parse_xml(grad_xml)
        spPr.insert(0, gradFill)

    except Exception as e:
        logger.warning(f"Failed to apply gradient fill: {e}")


# Pattern fill type mapping (OOXML preset patterns)
PATTERN_TYPE_MAP = {
    # Horizontal patterns
    "horizontal": "horz",
    "horizontal_lines": "horz",
    "light_horizontal": "ltHorz",
    "dark_horizontal": "dkHorz",
    "narrow_horizontal": "narHorz",
    "wide_horizontal": "wdUpDiag",  # Using alternative

    # Vertical patterns
    "vertical": "vert",
    "vertical_lines": "vert",
    "light_vertical": "ltVert",
    "dark_vertical": "dkVert",
    "narrow_vertical": "narVert",

    # Diagonal patterns
    "diagonal_up": "upDiag",
    "diagonal_down": "dnDiag",
    "light_diagonal_up": "ltUpDiag",
    "light_diagonal_down": "ltDnDiag",
    "dark_diagonal_up": "dkUpDiag",
    "dark_diagonal_down": "dkDnDiag",
    "wide_diagonal_up": "wdUpDiag",
    "wide_diagonal_down": "wdDnDiag",

    # Grid patterns
    "small_grid": "smGrid",
    "large_grid": "lgGrid",
    "dotted_grid": "dotGrid",
    "cross": "dashHorz",  # Similar to cross

    # Diagonal crosshatch
    "diagonal_cross": "diagCross",

    # Dot patterns
    "dots": "pct5",
    "small_dots": "pct5",
    "large_dots": "pct50",
    "pct5": "pct5",
    "pct10": "pct10",
    "pct20": "pct20",
    "pct25": "pct25",
    "pct30": "pct30",
    "pct40": "pct40",
    "pct50": "pct50",
    "pct60": "pct60",
    "pct70": "pct70",
    "pct75": "pct75",
    "pct80": "pct80",
    "pct90": "pct90",

    # Brick patterns
    "horizontal_brick": "horzBrick",
    "diagonal_brick": "diagBrick",

    # Check/plaid patterns
    "small_check": "smCheck",
    "large_check": "lgCheck",

    # Confetti/scatter patterns
    "small_confetti": "smConfetti",
    "large_confetti": "lgConfetti",

    # Special patterns
    "trellis": "trellis",
    "zig_zag": "zigZag",
    "wave": "wave",
    "weave": "weave",
    "shingle": "shingle",
    "plaid": "plaid",
    "divot": "divot",
    "sphere": "sphere",
    "solid_diamond": "solidDmnd",
    "open_diamond": "openDmnd",
    "dotted_diamond": "dotDmnd",
}


def apply_pattern_fill(shape, pattern_config: dict) -> None:
    """
    Apply a pattern fill to a shape using XML manipulation.

    Args:
        shape: The shape to apply pattern to
        pattern_config: Pattern configuration dict containing:
            - pattern: Pattern name (e.g., 'horizontal', 'diagonal_up', 'dots')
            - foreground_color: Color for the pattern lines/dots
            - background_color: Background color behind the pattern
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml

    pattern_name = pattern_config.get("pattern", pattern_config.get("type", "horizontal"))
    fg_color = pattern_config.get("foreground_color", pattern_config.get("fg_color", "#000000"))
    bg_color = pattern_config.get("background_color", pattern_config.get("bg_color", "#FFFFFF"))

    # Map pattern name to OOXML preset
    pattern_preset = PATTERN_TYPE_MAP.get(pattern_name.lower(), "horz")

    # Parse colors
    fg = parse_color(fg_color)
    bg = parse_color(bg_color)
    if fg is None:
        fg = RGBColor(0, 0, 0)
    if bg is None:
        bg = RGBColor(255, 255, 255)

    fg_hex = f"{fg[0]:02X}{fg[1]:02X}{fg[2]:02X}"
    bg_hex = f"{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}"

    try:
        spPr = shape._element.spPr

        # Remove existing fills
        for fill_elem in spPr.findall(qn('a:solidFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:noFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:gradFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:blipFill')):
            spPr.remove(fill_elem)
        for fill_elem in spPr.findall(qn('a:pattFill')):
            spPr.remove(fill_elem)

        patt_xml = f'''
            <a:pattFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="{pattern_preset}">
                <a:fgClr>
                    <a:srgbClr val="{fg_hex}"/>
                </a:fgClr>
                <a:bgClr>
                    <a:srgbClr val="{bg_hex}"/>
                </a:bgClr>
            </a:pattFill>
        '''

        pattFill = parse_xml(patt_xml)
        spPr.insert(0, pattFill)

    except Exception as e:
        logger.warning(f"Failed to apply pattern fill: {e}")


def apply_hyperlink(shape, hyperlink_config: dict) -> None:
    """
    Apply a hyperlink to a shape.

    Args:
        shape: The shape to add hyperlink to
        hyperlink_config: Hyperlink configuration dict containing:
            - url: Web URL to link to
            - slide: Slide number to link to (internal link)
            - tooltip: Tooltip text to show on hover
    """
    url = hyperlink_config.get("url", hyperlink_config.get("href"))
    slide_num = hyperlink_config.get("slide", hyperlink_config.get("slide_number"))
    tooltip = hyperlink_config.get("tooltip")

    try:
        if url:
            # External URL hyperlink
            shape.click_action.hyperlink.address = url
            if tooltip:
                shape.click_action.hyperlink.tooltip = tooltip
        elif slide_num is not None:
            # Internal slide link
            # This requires additional handling as python-pptx has limited support
            # For now, we'll set a placeholder action
            logger.info(f"Internal slide links require additional implementation for slide {slide_num}")

    except Exception as e:
        logger.warning(f"Failed to apply hyperlink: {e}")


# Action type mapping for shape click actions
ACTION_TYPE_MAP = {
    "next_slide": "ppActionNextSlide",
    "previous_slide": "ppActionPreviousSlide",
    "first_slide": "ppActionFirstSlide",
    "last_slide": "ppActionLastSlide",
    "end_show": "ppActionEndShow",
    "hyperlink": "ppActionHyperlink",
    "run_program": "ppActionRunProgram",
    "play_sound": "ppActionOLEVerb",
    "none": "ppActionNone",
}


def apply_text_columns(textbox, num_columns: int, spacing_inches: float = 0.3) -> None:
    """
    Configure text columns for a text box.

    Args:
        textbox: TextBox shape
        num_columns: Number of columns (1-9)
        spacing_inches: Space between columns in inches
    """
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    if num_columns < 1 or num_columns > 9:
        logger.warning(f"Invalid column count: {num_columns}. Must be 1-9.")
        return

    try:
        # Access the bodyPr element
        txBody = textbox._element.txBody
        if txBody is None:
            return

        bodyPr = txBody.bodyPr
        if bodyPr is None:
            bodyPr = OxmlElement("a:bodyPr")
            txBody.insert(0, bodyPr)

        # Set number of columns
        bodyPr.set("numCol", str(num_columns))

        # Set column spacing (in EMUs)
        if num_columns > 1:
            spacing_emu = int(spacing_inches * 914400)
            bodyPr.set("spcCol", str(spacing_emu))

    except Exception as e:
        logger.warning(f"Failed to set text columns: {e}")


def apply_callout_pointer(shape, pointer_config: dict) -> None:
    """
    Configure the pointer/tail of a callout shape.

    Args:
        shape: Callout shape
        pointer_config: Pointer configuration with:
            - x_percent: X position as percentage of shape width (0-100)
            - y_percent: Y position as percentage of shape height (0-100)
    """
    try:
        # Callout adjustments depend on the specific callout type
        # The adjustments control the pointer position
        x_pct = pointer_config.get("x_percent", 50)
        y_pct = pointer_config.get("y_percent", 150)  # Default below shape

        # Adjustment values are normalized to shape size
        # The exact adjustment indices depend on the callout type
        if hasattr(shape, "adjustments") and len(shape.adjustments) >= 2:
            # Typical callout has at least 2 adjustments for pointer x,y
            shape.adjustments[0] = x_pct / 100.0 * 0.5  # Normalize
            shape.adjustments[1] = y_pct / 100.0 * 0.5

    except Exception as e:
        logger.warning(f"Failed to configure callout pointer: {e}")


def apply_math_equation(slide, element: dict, slide_width: int, slide_height: int) -> None:
    """
    Add a mathematical equation to a slide.

    Note: PowerPoint's native equation editor uses OMML (Office Math Markup Language).
    This function provides basic equation support using text with math fonts.

    Args:
        slide: Slide object
        element: Element specification with:
            - equation: The equation text (e.g., "x = (-b ± √(b² - 4ac)) / 2a")
            - position: Position dict
            - text_properties: Font settings
    """
    position = element.get("position", {})
    equation_text = element.get("equation", element.get("text", ""))
    text_props = element.get("text_properties", {})

    # Use math-friendly font
    text_props.setdefault("font_family", "Cambria Math")
    text_props.setdefault("font_size_pt", 24)
    text_props.setdefault("alignment", "center")

    # Create as a text box with the equation
    modified_element = {
        "type": "textbox",
        "position": position,
        "text_properties": {
            **text_props,
            "placeholder_text": equation_text,
        }
    }

    apply_textbox(slide, modified_element, slide_width, slide_height)

    logger.info(f"Added math equation: {equation_text[:50]}...")


def apply_shape_action(shape, action_config: dict) -> None:
    """
    Apply a click action to a shape.

    Args:
        shape: The shape to add action to
        action_config: Action configuration dict containing:
            - type: Action type (next_slide, previous_slide, hyperlink, etc.)
            - url: URL for hyperlink action
            - slide: Target slide number
            - program: Path to program for run_program action
            - sound: Path to sound file for play_sound action
    """
    from pptx.enum.action import PP_ACTION

    action_type = action_config.get("type", "hyperlink")
    url = action_config.get("url")
    slide_num = action_config.get("slide")
    tooltip = action_config.get("tooltip")

    try:
        click_action = shape.click_action

        if action_type == "next_slide":
            click_action.action = PP_ACTION.NEXT_SLIDE
        elif action_type == "previous_slide":
            click_action.action = PP_ACTION.PREVIOUS_SLIDE
        elif action_type == "first_slide":
            click_action.action = PP_ACTION.FIRST_SLIDE
        elif action_type == "last_slide":
            click_action.action = PP_ACTION.LAST_SLIDE
        elif action_type == "end_show":
            click_action.action = PP_ACTION.END_SHOW
        elif action_type == "hyperlink" and url:
            click_action.hyperlink.address = url
            if tooltip:
                click_action.hyperlink.tooltip = tooltip
        elif action_type == "none":
            click_action.action = PP_ACTION.NONE
        else:
            logger.warning(f"Unknown or unsupported action type: {action_type}")

    except Exception as e:
        logger.warning(f"Failed to apply shape action: {e}")


def apply_text_hyperlink(run, url: str, tooltip: str = None) -> None:
    """
    Apply a hyperlink to a text run.

    Args:
        run: The text run to add hyperlink to
        url: URL to link to
        tooltip: Optional tooltip text
    """
    try:
        run.hyperlink.address = url
        if tooltip:
            # Tooltip requires XML manipulation
            pass
    except Exception as e:
        logger.warning(f"Failed to apply text hyperlink: {e}")


def apply_fill_transparency(shape, transparency: float) -> None:
    """
    Apply transparency to a shape's fill.

    Args:
        shape: The shape to modify
        transparency: Transparency value from 0.0 (opaque) to 1.0 (fully transparent)
    """
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from pptx.oxml.xmlchemy import OxmlElement

    if transparency <= 0:
        return  # No transparency needed

    # Clamp to valid range
    transparency = min(1.0, max(0.0, transparency))
    alpha_val = int((1 - transparency) * 100000)  # Convert to OOXML units

    try:
        spPr = shape._element.spPr

        # Find the fill element
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            # Add alpha modifier to the color
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                # Remove any existing alpha
                for alpha in srgbClr.findall(qn('a:alpha')):
                    srgbClr.remove(alpha)
                # Add new alpha
                alpha_elem = OxmlElement('a:alpha')
                alpha_elem.set('val', str(alpha_val))
                srgbClr.append(alpha_elem)

    except Exception as e:
        logger.warning(f"Failed to apply fill transparency: {e}")


def apply_line_transparency(shape, transparency: float) -> None:
    """
    Apply transparency to a shape's line/border.

    Args:
        shape: The shape to modify
        transparency: Transparency value from 0.0 (opaque) to 1.0 (fully transparent)
    """
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    if transparency <= 0:
        return

    transparency = min(1.0, max(0.0, transparency))
    alpha_val = int((1 - transparency) * 100000)

    try:
        spPr = shape._element.spPr
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            solidFill = ln.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    for alpha in srgbClr.findall(qn('a:alpha')):
                        srgbClr.remove(alpha)
                    alpha_elem = OxmlElement('a:alpha')
                    alpha_elem.set('val', str(alpha_val))
                    srgbClr.append(alpha_elem)

    except Exception as e:
        logger.warning(f"Failed to apply line transparency: {e}")


def _get_dash_style(dash_name: str) -> Optional[MSO_LINE_DASH_STYLE]:
    """
    Convert dash style name to MSO_LINE_DASH_STYLE enum.

    Args:
        dash_name: Name of the dash style (e.g., "dash", "dot", "dashed")

    Returns:
        MSO_LINE_DASH_STYLE enum value or None for solid
    """
    dash_map = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dots": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "dotted": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
        "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
        "dash": MSO_LINE_DASH_STYLE.DASH,
        "dashed": MSO_LINE_DASH_STYLE.DASH,
        "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
        "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
        "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
        "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
    }
    return dash_map.get(dash_name.lower().strip(), MSO_LINE_DASH_STYLE.SOLID)


# Arrow head type mapping
ARROW_TYPE_MAP = {
    "none": None,
    "triangle": "triangle",
    "stealth": "stealth",
    "diamond": "diamond",
    "oval": "oval",
    "open": "arrow",
    "arrow": "arrow",
}

# Arrow size mapping (small, medium, large to width/length values)
ARROW_SIZE_MAP = {
    "small": ("sm", "sm"),
    "medium": ("med", "med"),
    "large": ("lg", "lg"),
}


def _apply_arrow_heads(connector, shape_props: dict) -> None:
    """
    Apply arrow head properties to a connector line using XML manipulation.

    Args:
        connector: The connector shape
        shape_props: Shape properties containing arrow settings

    Supports two formats:
    1. String format: arrow_head_start="triangle", arrow_size_start="large"
    2. Dict format: start_arrow={"type": "triangle", "size": "large"}
    """
    # Get arrow head settings (support both string and dict formats)
    start_config = shape_props.get("arrow_head_start", shape_props.get("start_arrow", "none"))
    end_config = shape_props.get("arrow_head_end", shape_props.get("end_arrow", "none"))

    # Extract type and size from config (can be string or dict)
    if isinstance(start_config, dict):
        start_type = start_config.get("type", "none")
        start_size = start_config.get("size", "medium")
    else:
        start_type = start_config
        start_size = shape_props.get("arrow_size_start", shape_props.get("start_arrow_size", "medium"))

    if isinstance(end_config, dict):
        end_type = end_config.get("type", "none")
        end_size = end_config.get("size", "medium")
    else:
        end_type = end_config
        end_size = shape_props.get("arrow_size_end", shape_props.get("end_arrow_size", "medium"))

    # Map arrow types
    start_arrow = ARROW_TYPE_MAP.get(str(start_type).lower(), None) if start_type else None
    end_arrow = ARROW_TYPE_MAP.get(str(end_type).lower(), None) if end_type else None

    if not start_arrow and not end_arrow:
        return

    # Import XML utilities
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    # Get the line element from the connector
    spPr = connector._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = OxmlElement("a:ln")
        spPr.append(ln)

    # Apply start arrow (headEnd in OOXML terminology)
    if start_arrow:
        size_w, size_len = ARROW_SIZE_MAP.get(str(start_size).lower(), ("med", "med"))
        head_end = ln.find(qn("a:headEnd"))
        if head_end is None:
            head_end = OxmlElement("a:headEnd")
            ln.append(head_end)
        head_end.set("type", start_arrow)
        head_end.set("w", size_w)
        head_end.set("len", size_len)

    # Apply end arrow (tailEnd in OOXML terminology)
    if end_arrow:
        size_w, size_len = ARROW_SIZE_MAP.get(str(end_size).lower(), ("med", "med"))
        tail_end = ln.find(qn("a:tailEnd"))
        if tail_end is None:
            tail_end = OxmlElement("a:tailEnd")
            ln.append(tail_end)
        tail_end.set("type", end_arrow)
        tail_end.set("w", size_w)
        tail_end.set("len", size_len)


# Connector type mapping
CONNECTOR_TYPE_MAP = {
    "straight": MSO_CONNECTOR.STRAIGHT,
    "elbow": MSO_CONNECTOR.ELBOW,
    "curved": MSO_CONNECTOR.CURVE,
}


# Chart type mapping
CHART_TYPE_MAP = None  # Will be initialized on first use to avoid import overhead


def _get_chart_type_map():
    """Lazily initialize chart type mapping."""
    global CHART_TYPE_MAP
    if CHART_TYPE_MAP is None:
        from pptx.enum.chart import XL_CHART_TYPE
        CHART_TYPE_MAP = {
            # Bar charts
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "bar_100": XL_CHART_TYPE.BAR_STACKED_100,

            # Column charts
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "column_100": XL_CHART_TYPE.COLUMN_STACKED_100,

            # Line charts
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "line_stacked": XL_CHART_TYPE.LINE_STACKED,

            # Pie charts
            "pie": XL_CHART_TYPE.PIE,
            "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "doughnut_exploded": XL_CHART_TYPE.DOUGHNUT_EXPLODED,

            # Area charts
            "area": XL_CHART_TYPE.AREA,
            "area_stacked": XL_CHART_TYPE.AREA_STACKED,
            "area_100": XL_CHART_TYPE.AREA_STACKED_100,

            # Scatter/XY charts
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
            "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,

            # Radar charts
            "radar": XL_CHART_TYPE.RADAR,
            "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
            "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
        }
    return CHART_TYPE_MAP


def apply_chart(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a chart to a slide.

    Args:
        slide: Slide object
        element: Element specification dict with:
            - chart_type: Type of chart (bar, line, pie, etc.)
            - categories: List of category labels
            - series: List of series dicts with 'name' and 'values'
            - title: Optional chart title
            - position: Chart position and size
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    position = element.get("position", {})

    # Get chart properties (support both nested and flat formats)
    chart_props = element.get("chart_properties", {})
    chart_type_name = chart_props.get("chart_type", element.get("chart_type", "column"))
    categories = chart_props.get("categories", element.get("categories", []))
    series_list = chart_props.get("series", element.get("series", []))
    chart_title = chart_props.get("title", element.get("title", element.get("chart_title")))

    # Get chart type
    chart_type_map = _get_chart_type_map()
    chart_type = chart_type_map.get(chart_type_name.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Position and size
    left = parse_measurement(position.get("left_inches", 1), slide_width)
    top = parse_measurement(position.get("top_inches", 1), slide_height)
    width = parse_measurement(position.get("width_inches", 6), slide_width)
    height = parse_measurement(position.get("height_inches", 4), slide_height)

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series in series_list:
        series_name = series.get("name", "Series")
        series_values = series.get("values", [])
        chart_data.add_series(series_name, series_values)

    # Add chart to slide
    chart = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    ).chart

    # Set title if provided
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_title

    # Apply styling
    chart_style = chart_props.get("style", element.get("chart_style", element.get("style", {})))

    # Legend positioning
    if chart_style.get("legend_position"):
        from pptx.enum.chart import XL_LEGEND_POSITION
        legend_map = {
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "left": XL_LEGEND_POSITION.LEFT,
            "right": XL_LEGEND_POSITION.RIGHT,
            "top": XL_LEGEND_POSITION.TOP,
            "corner": XL_LEGEND_POSITION.CORNER,
        }
        pos = legend_map.get(chart_style["legend_position"].lower())
        if pos:
            chart.has_legend = True
            chart.legend.position = pos
    elif chart_style.get("show_legend") == False:
        chart.has_legend = False

    logger.info(f"Created {chart_type_name} chart with {len(series_list)} series")


def apply_group(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Create a group of shapes on a slide.

    Args:
        slide: Slide object
        element: Element specification dict with:
            - elements: List of child element definitions to group
            - position: Optional position offset for the group
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs

    Note: python-pptx has limited group shape support. This creates
    individual shapes and positions them as if grouped.
    """
    from pptx.shapes.group import GroupShape

    child_elements = element.get("elements", element.get("children", []))
    position = element.get("position", {})

    # Get group offset
    group_left = parse_measurement(position.get("left_inches", 0), slide_width)
    group_top = parse_measurement(position.get("top_inches", 0), slide_height)

    # Create each child element with offset
    created_shapes = []
    for child in child_elements:
        child_pos = child.get("position", {})

        # Apply group offset to child position
        if "left_inches" in child_pos:
            child_pos["left_inches"] = child_pos["left_inches"] + position.get("left_inches", 0)
        if "top_inches" in child_pos:
            child_pos["top_inches"] = child_pos["top_inches"] + position.get("top_inches", 0)

        child["position"] = child_pos

        # Track shapes before adding
        shapes_before = set(s.shape_id for s in slide.shapes)

        # Add the child element
        apply_element(slide, child, slide_width, slide_height)

        # Find newly added shape(s)
        shapes_after = set(s.shape_id for s in slide.shapes)
        new_shape_ids = shapes_after - shapes_before
        for shape_id in new_shape_ids:
            for shape in slide.shapes:
                if shape.shape_id == shape_id:
                    created_shapes.append(shape)
                    break

    # Note: Actual grouping requires XML manipulation that's not well-supported
    # in python-pptx. The shapes are positioned correctly but not technically grouped.
    # For true grouping, you would need to:
    # 1. Create a group shape container
    # 2. Move child shapes into the group
    # 3. Adjust coordinates to be relative to the group

    if len(created_shapes) > 1:
        logger.info(f"Created group with {len(created_shapes)} shapes (positioned but not formally grouped)")


def apply_connector(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int,
    shape_registry: dict = None
) -> None:
    """
    Add a connector line between two shapes on a slide.

    The connector can be attached to specific shapes using their element IDs,
    or positioned using explicit coordinates.

    Args:
        slide: Slide object
        element: Element specification dict with:
            - start_shape_id: ID of the starting shape
            - end_shape_id: ID of the ending shape
            - start_point/end_point: Fallback coordinates if shapes not found
            - connector_type: 'straight', 'elbow', or 'curved'
            - shape_properties: Line color, width, arrows, etc.
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        shape_registry: Dict mapping element IDs to shape objects
    """
    position = element.get("position", {})
    shape_props = element.get("shape_properties", {})

    # Get connector type
    connector_type_name = element.get("connector_type", "straight")
    connector_type = CONNECTOR_TYPE_MAP.get(connector_type_name.lower(), MSO_CONNECTOR.STRAIGHT)

    # Get start and end shape IDs
    start_shape_id = element.get("start_shape_id", element.get("from_shape"))
    end_shape_id = element.get("end_shape_id", element.get("to_shape"))

    # Get shapes from registry if available
    start_shape = shape_registry.get(start_shape_id) if shape_registry and start_shape_id else None
    end_shape = shape_registry.get(end_shape_id) if shape_registry and end_shape_id else None

    # Determine start and end points
    if start_shape is not None:
        # Use center of start shape
        start_x = start_shape.left + start_shape.width // 2
        start_y = start_shape.top + start_shape.height // 2
    else:
        start_x = parse_measurement(position.get("start_x_inches", 0), slide_width)
        start_y = parse_measurement(position.get("start_y_inches", 0), slide_height)

    if end_shape is not None:
        # Use center of end shape
        end_x = end_shape.left + end_shape.width // 2
        end_y = end_shape.top + end_shape.height // 2
    else:
        end_x = parse_measurement(position.get("end_x_inches", 1), slide_width)
        end_y = parse_measurement(position.get("end_y_inches", 1), slide_height)

    # Create connector
    connector = slide.shapes.add_connector(
        connector_type,
        start_x, start_y,
        end_x, end_y
    )

    # Style the connector
    line_color = parse_color(shape_props.get("line_color", shape_props.get("border_color", "#000000")))
    if line_color:
        connector.line.color.rgb = line_color

    line_width = shape_props.get("line_width_pt", shape_props.get("border_width_pt", 1.0))
    connector.line.width = Pt(line_width)

    # Dash style
    dash_style_name = shape_props.get("dash_style", shape_props.get("line_style"))
    if dash_style_name:
        dash_style = _get_dash_style(str(dash_style_name))
        if dash_style:
            connector.line.dash_style = dash_style

    # Apply arrow heads
    _apply_arrow_heads(connector, shape_props)

    # Connect to shapes if available
    # Note: python-pptx doesn't fully support shape connections, but we position correctly
    if start_shape is not None:
        # Position at edge of start shape (right side by default)
        start_x = start_shape.left + start_shape.width
        start_y = start_shape.top + start_shape.height // 2

    if end_shape is not None:
        # Position at edge of end shape (left side by default)
        end_x = end_shape.left
        end_y = end_shape.top + end_shape.height // 2


def apply_line(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add a line to a slide.

    Supports two position formats:
    1. Start/end points: start_x, start_y, end_x, end_y
    2. Box format: left_inches, top_inches, width_inches, height_inches

    Supports line styles:
    - solid (default): Uses filled rectangle (unless arrows are specified)
    - dashed/dotted: Uses connector with dash style

    Supports arrow heads:
    - arrow_head_start/arrow_head_end: none, triangle, stealth, diamond, oval, arrow
    - arrow_size_start/arrow_size_end: small, medium, large

    Args:
        slide: Slide object
        element: Element specification
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    position = element.get("position", {})
    shape_props = element.get("shape_properties", {})

    # Default line thickness in points (visible but thin)
    default_thickness_pt = 1.5

    # Check which position format is being used
    if "start_x" in position or "end_x" in position:
        # Format 1: Start/end points
        start_x = parse_measurement(position.get("start_x", 0), slide_width)
        start_y = parse_measurement(position.get("start_y", 0), slide_height)
        end_x = parse_measurement(position.get("end_x", 1), slide_width)
        end_y = parse_measurement(position.get("end_y", 0), slide_height)
    else:
        # Format 2: Box format (left/top/width/height) - convert to start/end
        left = parse_measurement(position.get("left_inches", 0), slide_width)
        top = parse_measurement(position.get("top_inches", 0), slide_height)
        width = parse_measurement(position.get("width_inches", 1), slide_width)
        height = parse_measurement(position.get("height_inches", 0), slide_height)

        # Convert box to start/end points
        start_x = left
        start_y = top
        end_x = left + width
        end_y = top + height

    # Get line thickness from shape properties
    thickness_pt = shape_props.get("border_width_pt", default_thickness_pt)
    # Ensure minimum visible thickness
    thickness_pt = max(thickness_pt, 0.75)
    thickness_emu = Pt(thickness_pt)

    # Get line color
    line_color = None
    if shape_props.get("fill_color"):
        line_color = parse_color(shape_props["fill_color"])
    if line_color is None and shape_props.get("border_color"):
        line_color = parse_color(shape_props["border_color"])
    if line_color is None:
        line_color = RGBColor(0, 0, 0)  # Default black

    # Check for dash style
    dash_style_name = shape_props.get("dash_style", shape_props.get("line_style", "solid"))
    dash_style = _get_dash_style(str(dash_style_name)) if dash_style_name else None

    # Check if arrows are specified (support both string and dict format)
    start_arrow_config = shape_props.get("arrow_head_start", shape_props.get("start_arrow", "none"))
    end_arrow_config = shape_props.get("arrow_head_end", shape_props.get("end_arrow", "none"))

    # Extract arrow type from config (can be string or dict with 'type' key)
    def get_arrow_type(config):
        if isinstance(config, dict):
            return config.get("type", "none")
        return config if config else "none"

    has_start_arrow = get_arrow_type(start_arrow_config)
    has_end_arrow = get_arrow_type(end_arrow_config)
    needs_arrows = (has_start_arrow and has_start_arrow.lower() != "none") or \
                   (has_end_arrow and has_end_arrow.lower() != "none")

    # Use connector for dashed/dotted lines OR when arrows are needed
    use_connector = (dash_style and dash_style != MSO_LINE_DASH_STYLE.SOLID) or needs_arrows

    if use_connector:
        # Use connector for dashed lines or lines with arrows
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            start_x, start_y,
            end_x, end_y
        )
        connector.line.color.rgb = line_color
        connector.line.width = thickness_emu
        if dash_style and dash_style != MSO_LINE_DASH_STYLE.SOLID:
            connector.line.dash_style = dash_style

        # Apply arrow heads if specified
        if needs_arrows:
            _apply_arrow_heads(connector, shape_props)
    else:
        # Use rectangle for solid lines without arrows (better visual appearance)
        left = min(start_x, end_x)
        top = min(start_y, end_y)
        width = abs(end_x - start_x)
        height = abs(end_y - start_y)

        # Ensure minimum dimensions for visibility
        if width > 0 and height <= thickness_emu:
            # Horizontal line
            height = thickness_emu
        elif height > 0 and width <= thickness_emu:
            # Vertical line
            width = thickness_emu
        else:
            # Diagonal or both dimensions specified
            if width < thickness_emu:
                width = thickness_emu
            if height < thickness_emu:
                height = thickness_emu

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = line_color
        shape.line.fill.background()  # No border for rectangle-as-line


def apply_element(
    slide,
    element: dict,
    slide_width: int,
    slide_height: int
) -> None:
    """
    Add an element to a slide based on its type.

    Args:
        slide: Slide object
        element: Element specification dict
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
    """
    element_type = element.get("type", "textbox").lower()

    if element_type in ("textbox", "text_box", "text"):
        apply_textbox(slide, element, slide_width, slide_height)
    elif element_type in ("bullet_list", "bullets", "list"):
        apply_bullet_list(slide, element, slide_width, slide_height)
    elif element_type == "table":
        apply_table(slide, element, slide_width, slide_height)
    elif element_type in ("shape", "rectangle", "oval", "circle"):
        apply_shape(slide, element, slide_width, slide_height)
    elif element_type == "line":
        apply_line(slide, element, slide_width, slide_height)
    elif element_type == "connector":
        apply_connector(slide, element, slide_width, slide_height)
    elif element_type == "group":
        apply_group(slide, element, slide_width, slide_height)
    elif element_type == "chart":
        apply_chart(slide, element, slide_width, slide_height)
    elif element_type == "image":
        apply_image(slide, element, slide_width, slide_height)
    elif element_type == "image_placeholder":
        # Check if image source is provided
        image_props = element.get("image_properties", {})
        if image_props.get("source") or image_props.get("source_base64"):
            apply_image(slide, element, slide_width, slide_height)
        else:
            # No image source, treat as a shape placeholder
            element["shape_properties"] = element.get("shape_properties", {})
            element["shape_properties"]["fill_color"] = element["shape_properties"].get(
                "fill_color", "#CCCCCC"
            )
            apply_shape(slide, element, slide_width, slide_height)
    elif element_type in ("equation", "math", "formula"):
        apply_math_equation(slide, element, slide_width, slide_height)
    else:
        logger.warning(f"Unknown element type: {element_type}")


def batch_create_elements(
    slide,
    elements: list,
    slide_width: int,
    slide_height: int,
    shape_registry: dict = None
) -> dict:
    """
    Create multiple elements in batch with optional ID tracking.

    Args:
        slide: Slide object
        elements: List of element specifications
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        shape_registry: Optional dict to store created shapes by ID

    Returns:
        Dict mapping element IDs to created shape objects
    """
    created_shapes = {}

    for element in elements:
        element_id = element.get("id")

        # Track shapes before adding
        shapes_before = set(s.shape_id for s in slide.shapes)

        # Create the element
        try:
            apply_element(slide, element, slide_width, slide_height)
        except Exception as e:
            logger.error(f"Failed to create element {element_id}: {e}")
            continue

        # Find newly added shape(s) and track by ID
        if element_id:
            shapes_after = set(s.shape_id for s in slide.shapes)
            new_shape_ids = shapes_after - shapes_before
            for shape_id in new_shape_ids:
                for shape in slide.shapes:
                    if shape.shape_id == shape_id:
                        created_shapes[element_id] = shape
                        if shape_registry is not None:
                            shape_registry[element_id] = shape
                        break

    return created_shapes


def validate_description(description: dict) -> list:
    """
    Validate a slide description for common errors.

    Args:
        description: The description dict to validate

    Returns:
        List of validation error/warning messages
    """
    errors = []

    # Check required fields
    if "elements" not in description:
        errors.append("Warning: No 'elements' field found")

    # Check slide dimensions
    dimensions = description.get("slide_dimensions", {})
    width = dimensions.get("width_inches")
    height = dimensions.get("height_inches")
    if width and (width < 1 or width > 100):
        errors.append(f"Warning: Unusual slide width: {width} inches")
    if height and (height < 1 or height > 100):
        errors.append(f"Warning: Unusual slide height: {height} inches")

    # Validate elements
    elements = description.get("elements", [])
    element_ids = set()

    for i, element in enumerate(elements):
        element_id = element.get("id", f"element_{i}")

        # Check for duplicate IDs
        if element_id in element_ids:
            errors.append(f"Error: Duplicate element ID '{element_id}'")
        element_ids.add(element_id)

        # Check element type
        element_type = element.get("type")
        if not element_type:
            errors.append(f"Warning: Element {element_id} has no type specified")

        # Check position
        position = element.get("position", {})
        if not position:
            errors.append(f"Warning: Element {element_id} has no position specified")
        else:
            # Check for required position fields
            if "left_inches" not in position and "left" not in position:
                errors.append(f"Warning: Element {element_id} missing left position")
            if "top_inches" not in position and "top" not in position:
                errors.append(f"Warning: Element {element_id} missing top position")

        # Check for color format issues
        shape_props = element.get("shape_properties", {})
        text_props = element.get("text_properties", {})

        for color_key in ["fill_color", "border_color", "font_color"]:
            color_val = shape_props.get(color_key) or text_props.get(color_key)
            if color_val and isinstance(color_val, str):
                try:
                    parse_color(color_val)
                except GeneratorError:
                    errors.append(f"Error: Invalid color '{color_val}' in element {element_id}")

    return errors


def get_slide_layout(prs, layout_name: str = None, layout_index: int = None):
    """
    Get a slide layout from the presentation.

    Args:
        prs: Presentation object
        layout_name: Name of the layout (e.g., 'Title Slide', 'Blank')
        layout_index: Index of the layout in the master

    Returns:
        SlideLayout object, or None if not found
    """
    if layout_index is not None and prs.slide_layouts:
        if 0 <= layout_index < len(prs.slide_layouts):
            return prs.slide_layouts[layout_index]

    if layout_name:
        layout_name_lower = layout_name.lower().strip()
        for layout in prs.slide_layouts:
            if layout.name.lower().strip() == layout_name_lower:
                return layout

    return None


def get_available_layouts(prs) -> list:
    """
    Get list of available layout names from a presentation.

    Args:
        prs: Presentation object

    Returns:
        List of layout name strings
    """
    return [layout.name for layout in prs.slide_layouts]


def add_slide_with_layout(prs, layout_name: str = None, layout_index: int = None):
    """
    Add a new slide with a specific layout.

    Args:
        prs: Presentation object
        layout_name: Name of the layout to use
        layout_index: Index of the layout (0 = first layout)

    Returns:
        Slide object
    """
    layout = get_slide_layout(prs, layout_name, layout_index)
    if layout is None:
        # Fall back to blank layout (usually index 6) or first available
        layout = get_slide_layout(prs, "Blank") or prs.slide_layouts[0]

    return prs.slides.add_slide(layout)


def generate_slide_from_description(
    description: dict,
    output_path: Optional[Path] = None
) -> Path:
    """
    Generate a PPTX file from a structured description.

    Args:
        description: Structured description dict (from descriptor.py)
            - Can include 'layout' field with 'name' or 'index' for master layout selection
        output_path: Path for the output PPTX file

    Returns:
        Path to the generated PPTX file

    Raises:
        GeneratorError: If generation fails
    """
    # Get slide dimensions
    dimensions = description.get("slide_dimensions", {})
    width = dimensions.get("width_inches", DEFAULT_SLIDE_WIDTH_INCHES)
    height = dimensions.get("height_inches", DEFAULT_SLIDE_HEIGHT_INCHES)

    # Create presentation
    prs = create_presentation(width, height)

    # Check for layout specification
    layout_config = description.get("layout", {})
    layout_name = layout_config.get("name") if isinstance(layout_config, dict) else None
    layout_index = layout_config.get("index") if isinstance(layout_config, dict) else None

    if layout_name or layout_index is not None:
        slide = add_slide_with_layout(prs, layout_name, layout_index)
    else:
        slide = add_blank_slide(prs)

    # Get dimensions in EMUs for calculations
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Apply background
    if "background" in description:
        apply_background(slide, description["background"])

    # Apply elements - sort by z_order (lower values = behind, higher = in front)
    elements = description.get("elements", [])
    # Sort by z_order if present, default to 0 (maintain original order for same z_order)
    sorted_elements = sorted(
        enumerate(elements),
        key=lambda x: (x[1].get("z_order", 0), x[0])
    )

    for _, element in sorted_elements:
        try:
            apply_element(slide, element, slide_width, slide_height)
        except Exception as e:
            element_id = element.get("id", "unknown")
            logger.error(f"Failed to apply element {element_id}: {e}")

    # Determine output path
    if output_path is None:
        output_path = OUTPUT_DIR / "generated_slide.pptx"
    else:
        output_path = Path(output_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(output_path)
    logger.info(f"Generated PPTX saved: {output_path}")

    return output_path


def generate_from_json(json_path: Path, output_path: Optional[Path] = None) -> Path:
    """
    Generate a PPTX file from a JSON description file.

    Args:
        json_path: Path to the JSON description file
        output_path: Path for the output PPTX file

    Returns:
        Path to the generated PPTX file
    """
    import json

    with open(json_path, 'r', encoding='utf-8') as f:
        description = json.load(f)

    if output_path is None:
        output_path = OUTPUT_DIR / f"{Path(json_path).stem}_generated.pptx"

    return generate_slide_from_description(description, output_path)


def add_slide_to_presentation(
    prs: Presentation,
    description: dict
) -> Any:
    """
    Add a slide to an existing presentation from a description.

    Args:
        prs: Existing Presentation object
        description: Structured description dict for the slide

    Returns:
        The newly added slide object
    """
    slide = add_blank_slide(prs)

    # Get dimensions in EMUs for calculations
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Apply background
    if "background" in description:
        apply_background(slide, description["background"])

    # Apply elements - sort by z_order (lower values = behind, higher = in front)
    elements = description.get("elements", [])
    sorted_elements = sorted(
        enumerate(elements),
        key=lambda x: (x[1].get("z_order", 0), x[0])
    )

    for _, element in sorted_elements:
        try:
            apply_element(slide, element, slide_width, slide_height)
        except Exception as e:
            element_id = element.get("id", "unknown")
            logger.error(f"Failed to apply element {element_id}: {e}")

    return slide


def generate_multi_slide_presentation(
    descriptions: list[dict],
    output_path: Optional[Path] = None,
    width_inches: float = None,
    height_inches: float = None
) -> Path:
    """
    Generate a multi-slide PPTX from a list of descriptions.

    Args:
        descriptions: List of structured description dicts, one per slide
        output_path: Path for the output PPTX file
        width_inches: Slide width (uses first description's dimensions if None)
        height_inches: Slide height (uses first description's dimensions if None)

    Returns:
        Path to the generated PPTX file

    Raises:
        GeneratorError: If generation fails or no descriptions provided
    """
    if not descriptions:
        raise GeneratorError("No descriptions provided for multi-slide generation")

    # Get dimensions from first description if not specified
    if width_inches is None or height_inches is None:
        first_dims = descriptions[0].get("slide_dimensions", {})
        if width_inches is None:
            width_inches = first_dims.get("width_inches", DEFAULT_SLIDE_WIDTH_INCHES)
        if height_inches is None:
            height_inches = first_dims.get("height_inches", DEFAULT_SLIDE_HEIGHT_INCHES)

    # Create presentation
    prs = create_presentation(width_inches, height_inches)

    # Add each slide
    for i, description in enumerate(descriptions):
        logger.info(f"Adding slide {i + 1}/{len(descriptions)}")
        add_slide_to_presentation(prs, description)

    # Determine output path
    if output_path is None:
        output_path = OUTPUT_DIR / "combined_presentation.pptx"
    else:
        output_path = Path(output_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Save presentation
    prs.save(output_path)
    logger.info(f"Generated multi-slide PPTX saved: {output_path} ({len(descriptions)} slides)")

    return output_path


def combine_json_descriptions(
    json_paths: list[Path],
    output_path: Optional[Path] = None
) -> Path:
    """
    Combine multiple JSON description files into a single multi-slide PPTX.

    Args:
        json_paths: List of paths to JSON description files (in slide order)
        output_path: Path for the output PPTX file

    Returns:
        Path to the generated PPTX file
    """
    import json

    descriptions = []
    for json_path in json_paths:
        with open(json_path, 'r', encoding='utf-8') as f:
            descriptions.append(json.load(f))

    return generate_multi_slide_presentation(descriptions, output_path)


if __name__ == "__main__":
    # Test the generator with a simple description
    import sys

    if len(sys.argv) > 1:
        # Generate from JSON file
        json_path = Path(sys.argv[1])
        output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else None
        result = generate_from_json(json_path, output_path)
        print(f"Generated: {result}")
    else:
        # Create a test slide
        test_description = {
            "slide_dimensions": {
                "width_inches": 13.333,
                "height_inches": 7.5
            },
            "background": {
                "type": "solid",
                "color": "#1A365D"
            },
            "elements": [
                {
                    "id": "title",
                    "type": "textbox",
                    "position": {
                        "left_inches": 1,
                        "top_inches": 2.5,
                        "width_inches": 11.333,
                        "height_inches": 1.5
                    },
                    "text_properties": {
                        "placeholder_text": "PRESENTATION TITLE",
                        "font_family": "Calibri Light",
                        "font_size_pt": 44,
                        "font_color": "#FFFFFF",
                        "bold": True,
                        "alignment": "center",
                        "vertical_alignment": "middle"
                    }
                },
                {
                    "id": "subtitle",
                    "type": "textbox",
                    "position": {
                        "left_inches": 1,
                        "top_inches": 4.5,
                        "width_inches": 11.333,
                        "height_inches": 1
                    },
                    "text_properties": {
                        "placeholder_text": "Subtitle goes here",
                        "font_family": "Calibri",
                        "font_size_pt": 24,
                        "font_color": "#A0AEC0",
                        "alignment": "center"
                    }
                },
                {
                    "id": "accent_line",
                    "type": "shape",
                    "position": {
                        "left_inches": 5.5,
                        "top_inches": 4.2,
                        "width_inches": 2.5,
                        "height_inches": 0.05
                    },
                    "shape_properties": {
                        "shape_type": "rectangle",
                        "fill_color": "#4299E1",
                        "no_border": True
                    }
                }
            ]
        }

        output = generate_slide_from_description(
            test_description,
            OUTPUT_DIR / "test_generated.pptx"
        )
        print(f"Test slide generated: {output}")
