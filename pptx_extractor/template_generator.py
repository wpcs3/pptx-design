"""
Template-Based Slide Generator

Generates slides using actual master layouts from templates,
eliminating the need for LLM to approximate positions.

Usage:
    generator = TemplateGenerator("template.pptx")
    generator.create_slide("Default", {
        "title": "My Title",
        "subtitle": "My Subtitle"
    })
    generator.save("output.pptx")
"""
import json
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

logger = logging.getLogger(__name__)


# Placeholder type mapping
PLACEHOLDER_CONTENT_MAP = {
    "TITLE": ["title", "heading", "header"],
    "SUBTITLE": ["subtitle", "subheading", "sub_title"],
    "BODY": ["body", "content", "text", "bullet", "bullets"],
    "FOOTER": ["footer", "company", "company_name"],
    "SLIDE_NUMBER": ["slide_number", "page", "page_number"],
    "DATE": ["date"],
}


def parse_color(color_str: str) -> Optional[RGBColor]:
    """Parse color string to RGBColor."""
    if not color_str or color_str.lower() in ['none', 'transparent']:
        return None
    if color_str.startswith('#'):
        color_str = color_str[1:]
    if len(color_str) == 3:
        color_str = ''.join(c * 2 for c in color_str)
    try:
        return RGBColor.from_string(color_str)
    except:
        return None


class TemplateGenerator:
    """
    Generate slides using master layouts from a template.

    This approach:
    - Uses exact placeholder positions from the template
    - Preserves all formatting from the original
    - Only requires content to fill placeholders
    """

    def __init__(self, template_path: Union[str, Path], clear_slides: bool = True):
        """
        Initialize with a template PPTX file.

        Args:
            template_path: Path to the template PPTX
            clear_slides: If True, remove all existing slides (keep only layouts)
        """
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            # Try templates directory
            from config import TEMPLATE_DIR
            for t in TEMPLATE_DIR.glob("**/*.pptx"):
                if self.template_path.name in t.name:
                    self.template_path = t
                    break

        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        # Load template
        self.prs = Presentation(self.template_path)
        self._build_layout_index()

        # Remove existing slides if requested
        if clear_slides:
            self._clear_slides()

        logger.info(f"Loaded template: {self.template_path.name} "
                    f"({len(self.layout_index)} layouts)")

    def _clear_slides(self):
        """Remove all existing slides from the presentation."""
        # python-pptx doesn't have a direct delete method, so we access the XML
        slide_ids = list(self.prs.slides._sldIdLst)
        for slide_id in slide_ids:
            rId = slide_id.rId
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(slide_id)

    def _build_layout_index(self):
        """Build index of available layouts."""
        self.layout_index = {}
        self.layouts_by_name = {}

        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                self.layout_index[layout.name] = layout
                self.layouts_by_name[layout.name.lower()] = layout

                # Also index by common variations
                name_lower = layout.name.lower()
                if "title" in name_lower:
                    self.layouts_by_name["title"] = layout
                if "default" in name_lower:
                    self.layouts_by_name["content"] = layout
                if "blank" in name_lower and "dark" not in name_lower:
                    self.layouts_by_name["blank"] = layout

    def list_layouts(self) -> List[Dict[str, Any]]:
        """List all available layouts with their placeholders."""
        layouts = []
        for name, layout in self.layout_index.items():
            placeholders = []
            for ph in layout.placeholders:
                ph_type = str(ph.placeholder_format.type).split('.')[-1].strip('()')
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": ph_type,
                })
            layouts.append({
                "name": name,
                "placeholders": placeholders,
            })
        return layouts

    def get_layout(self, layout_name: str):
        """Get a layout by name (case-insensitive, partial match)."""
        # Exact match
        if layout_name in self.layout_index:
            return self.layout_index[layout_name]

        # Case-insensitive match
        name_lower = layout_name.lower()
        if name_lower in self.layouts_by_name:
            return self.layouts_by_name[name_lower]

        # Partial match
        for name, layout in self.layout_index.items():
            if name_lower in name.lower():
                return layout

        raise ValueError(f"Layout not found: {layout_name}. "
                         f"Available: {list(self.layout_index.keys())}")

    def create_slide(
        self,
        layout_name: str,
        content: Dict[str, Any],
        apply_formatting: Optional[Dict[str, Any]] = None
    ) -> int:
        """
        Create a new slide using a layout template.

        Args:
            layout_name: Name of the layout to use
            content: Dictionary mapping placeholder names to content
                     e.g., {"title": "My Title", "body": "Content here"}
            apply_formatting: Optional formatting overrides

        Returns:
            Index of the created slide
        """
        layout = self.get_layout(layout_name)
        slide = self.prs.slides.add_slide(layout)

        # Fill placeholders
        for ph in slide.placeholders:
            ph_type = str(ph.placeholder_format.type).split('.')[-1].strip('()')
            ph_idx = ph.placeholder_format.idx

            # Find matching content
            text_content = None

            # Try direct idx match
            if str(ph_idx) in content:
                text_content = content[str(ph_idx)]

            # Try type-based match
            if text_content is None:
                for type_key, content_keys in PLACEHOLDER_CONTENT_MAP.items():
                    if type_key in ph_type.upper():
                        for key in content_keys:
                            if key in content:
                                text_content = content[key]
                                break
                    if text_content:
                        break

            # Apply content
            if text_content and ph.has_text_frame:
                self._set_placeholder_text(ph, text_content, apply_formatting)

        # Handle special content types (metrics, tables, charts, two-column)
        if "metrics" in content:
            self._add_metrics(slide, content["metrics"])

        if "headers" in content and "data" in content:
            self._add_table(slide, content["headers"], content["data"])

        if "chart_data" in content:
            self._add_chart(slide, content["chart_data"])

        if "left_column" in content and "right_column" in content:
            self._add_two_column_content(slide, content["left_column"], content["right_column"])

        logger.info(f"Created slide with layout: {layout_name}")
        return len(self.prs.slides) - 1

    def _set_placeholder_text(
        self,
        placeholder,
        content: Union[str, List[str], Dict],
        formatting: Optional[Dict] = None
    ):
        """Set text in a placeholder with optional formatting."""
        tf = placeholder.text_frame

        # Handle different content types
        if isinstance(content, str):
            # Simple text
            tf.paragraphs[0].text = content
        elif isinstance(content, list):
            # Bullet points
            for i, item in enumerate(content):
                if i == 0:
                    tf.paragraphs[0].text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
        elif isinstance(content, dict):
            # Detailed content with formatting
            text = content.get("text", "")
            tf.paragraphs[0].text = text

            # Apply formatting from content dict
            if "font_size_pt" in content:
                for p in tf.paragraphs:
                    p.font.size = Pt(content["font_size_pt"])
            if "font_color" in content:
                color = parse_color(content["font_color"])
                if color:
                    for p in tf.paragraphs:
                        p.font.color.rgb = color
            if "bold" in content:
                for p in tf.paragraphs:
                    p.font.bold = content["bold"]

        # Apply global formatting overrides
        if formatting:
            for p in tf.paragraphs:
                if "font_size_pt" in formatting:
                    p.font.size = Pt(formatting["font_size_pt"])
                if "font_color" in formatting:
                    color = parse_color(formatting["font_color"])
                    if color:
                        p.font.color.rgb = color

    def _add_metrics(self, slide, metrics: List[Dict[str, str]]):
        """Add KPI metric boxes to a slide."""
        if not metrics:
            return

        # Calculate positions for metric boxes
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height

        num_metrics = len(metrics)
        box_width = Inches(2.0)
        box_height = Inches(1.2)
        spacing = Inches(0.3)

        # Calculate total width and start position for centering
        total_width = num_metrics * box_width + (num_metrics - 1) * spacing
        start_x = (slide_width - total_width) // 2
        start_y = Inches(2.5)  # Below title area

        for i, metric in enumerate(metrics):
            x = start_x + i * (box_width + spacing)

            # Create box shape
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                x, start_y, box_width, box_height
            )

            # Style the box
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light grey
            shape.line.color.rgb = RGBColor(200, 200, 200)

            # Add value text (large)
            tf = shape.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            value_para = tf.paragraphs[0]
            value_para.text = metric.get("value", "")
            value_para.font.size = Pt(28)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(51, 51, 51)

            # Add label text (small)
            label_para = tf.add_paragraph()
            label_para.text = metric.get("label", "")
            label_para.font.size = Pt(12)
            label_para.font.color.rgb = RGBColor(102, 102, 102)
            label_para.alignment = PP_ALIGN.CENTER

    def _add_table(self, slide, headers: List[str], data: List[List[str]]):
        """Add a table to a slide."""
        if not headers or not data:
            return

        rows = len(data) + 1  # +1 for header
        cols = len(headers)

        # Position and size
        left = Inches(0.75)
        top = Inches(2.0)
        width = Inches(11.5)
        height = Inches(0.4) * rows

        # Create table
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Set column widths
        col_width = width // cols
        for i in range(cols):
            table.columns[i].width = col_width

        # Fill header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(51, 51, 51)  # Dark header

            # Style header text
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(11)
            para.font.color.rgb = RGBColor(255, 255, 255)  # White text
            para.alignment = PP_ALIGN.CENTER

        # Fill data rows
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < cols:  # Prevent index out of range
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_value)

                    # Alternate row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(248, 248, 248)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(10)
                    para.alignment = PP_ALIGN.CENTER

    def _add_chart(self, slide, chart_data: Dict[str, Any]):
        """Add a chart to a slide."""
        chart_type_str = chart_data.get("type", "column").lower()
        categories = chart_data.get("categories", [])
        series_list = chart_data.get("series", [])

        if not categories or not series_list:
            return

        # Map chart type string to enum
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA,
        }
        chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # Create chart data
        data = CategoryChartData()
        data.categories = categories

        for series in series_list:
            series_name = series.get("name", "Series")
            series_values = series.get("values", [])
            # Ensure values are numeric
            numeric_values = []
            for v in series_values:
                try:
                    numeric_values.append(float(v))
                except (ValueError, TypeError):
                    numeric_values.append(0)
            data.add_series(series_name, numeric_values)

        # Position and size
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(10.0)
        height = Inches(4.5)

        # Add chart to slide
        chart = slide.shapes.add_chart(
            chart_type, left, top, width, height, data
        ).chart

        # Style the chart
        chart.has_legend = True
        if len(series_list) > 1:
            chart.legend.include_in_layout = False

    def _add_two_column_content(
        self,
        slide,
        left: Dict[str, Any],
        right: Dict[str, Any]
    ):
        """Add two-column content with headers and bullets."""
        left_header = left.get("header", "")
        left_bullets = left.get("bullets", [])
        right_header = right.get("header", "")
        right_bullets = right.get("bullets", [])

        # Left column text box
        left_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(2.0),
            Inches(5.5), Inches(4.5)
        )
        left_tf = left_box.text_frame
        left_tf.word_wrap = True

        # Left header
        if left_header:
            p = left_tf.paragraphs[0]
            p.text = left_header
            p.font.bold = True
            p.font.size = Pt(18)
            p.space_after = Pt(12)

        # Left bullets
        for bullet in left_bullets:
            p = left_tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(6)

        # Right column text box
        right_box = slide.shapes.add_textbox(
            Inches(6.75), Inches(2.0),
            Inches(5.5), Inches(4.5)
        )
        right_tf = right_box.text_frame
        right_tf.word_wrap = True

        # Right header
        if right_header:
            p = right_tf.paragraphs[0]
            p.text = right_header
            p.font.bold = True
            p.font.size = Pt(18)
            p.space_after = Pt(12)

        # Right bullets
        for bullet in right_bullets:
            p = right_tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(6)

    def duplicate_slide(self, slide_index: int) -> int:
        """Duplicate an existing slide."""
        source = self.prs.slides[slide_index]
        layout = source.slide_layout

        # Create new slide with same layout
        new_slide = self.prs.slides.add_slide(layout)

        # Copy shapes (this is simplified - full copy is complex)
        # For now, just copy placeholder content
        for src_ph in source.placeholders:
            for dst_ph in new_slide.placeholders:
                if src_ph.placeholder_format.idx == dst_ph.placeholder_format.idx:
                    if src_ph.has_text_frame and dst_ph.has_text_frame:
                        dst_ph.text_frame.paragraphs[0].text = src_ph.text_frame.text

        return len(self.prs.slides) - 1

    def remove_template_slides(self):
        """Remove all slides from the template (keep only new ones)."""
        # Note: python-pptx doesn't support slide deletion directly
        # This is a workaround - create new presentation
        pass

    def save(self, output_path: Union[str, Path]):
        """Save the presentation."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        logger.info(f"Saved: {output_path}")
        return output_path

    def create_presentation(
        self,
        slides: List[Dict[str, Any]],
        output_path: Union[str, Path]
    ) -> Path:
        """
        Create a complete presentation from slide specifications.

        Args:
            slides: List of slide specs, each with "layout" and "content"
            output_path: Where to save the presentation

        Example:
            slides = [
                {"layout": "Frontpage", "content": {"title": "My Deck", "subtitle": "2024"}},
                {"layout": "Default", "content": {"title": "Intro", "body": "Content here"}},
            ]
        """
        for slide_spec in slides:
            layout = slide_spec.get("layout", "Default")
            content = slide_spec.get("content", {})
            formatting = slide_spec.get("formatting")
            self.create_slide(layout, content, formatting)

        return self.save(output_path)


def create_from_template(
    template_name: str,
    slides: List[Dict[str, Any]],
    output_path: str
) -> Path:
    """
    Convenience function to create a presentation from a template.

    Args:
        template_name: Name of template file
        slides: List of slide specifications
        output_path: Output file path

    Returns:
        Path to created presentation
    """
    generator = TemplateGenerator(template_name)
    return generator.create_presentation(slides, output_path)


if __name__ == "__main__":
    import sys

    logging.basicConfig(level=logging.INFO, format="%(message)s")

    if len(sys.argv) < 2:
        print("Usage: python template_generator.py <template.pptx>")
        print("\nThis will list available layouts and create a test presentation.")
        sys.exit(1)

    template = sys.argv[1]
    generator = TemplateGenerator(template)

    print(f"\nTemplate: {generator.template_path.name}")
    print(f"Available layouts ({len(generator.layout_index)}):")
    for layout_info in generator.list_layouts():
        ph_types = [p["type"] for p in layout_info["placeholders"]]
        print(f"  - {layout_info['name']}: {ph_types}")

    # Create test presentation
    print("\nCreating test presentation...")
    test_slides = [
        {
            "layout": "Frontpage",
            "content": {
                "title": "Test Presentation",
                "subtitle": "Generated from Template"
            }
        },
        {
            "layout": "Default",
            "content": {
                "title": "Introduction",
                "subtitle": "This slide uses the Default layout"
            }
        },
        {
            "layout": "Agenda",
            "content": {
                "title": "Agenda",
                "body": ["Item 1", "Item 2", "Item 3"]
            }
        }
    ]

    output = Path("outputs/template_test_output.pptx")
    generator.create_presentation(test_slides, output)
    print(f"Created: {output}")
