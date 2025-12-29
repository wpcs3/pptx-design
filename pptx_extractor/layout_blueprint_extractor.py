"""
Layout Blueprint Extractor - Extract grid systems and layout patterns from PowerPoint templates.

This module extracts:
- Grid systems (detected from element positions)
- Zone definitions (title, content, sidebar, etc.)
- Spatial relationships between elements
- Content flow patterns
- Margin and gutter calculations
"""

import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any, Tuple
from collections import defaultdict
import math

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class LayoutBlueprintExtractor:
    """Extract layout blueprints and grid systems from PowerPoint templates."""

    def __init__(self, output_dir: Path):
        """Initialize the layout blueprint extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectory for layouts
        self.layouts_dir = self.output_dir / 'layouts' / 'blueprints'
        self.grids_dir = self.output_dir / 'layouts' / 'grids'
        self.layouts_dir.mkdir(parents=True, exist_ok=True)
        self.grids_dir.mkdir(parents=True, exist_ok=True)

        # Layout index
        self.index_path = self.output_dir / 'layouts' / 'layout_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'blueprints': [],
            'grids': [],
            'patterns': [],
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def extract_template_layouts(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all layout blueprints from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted layouts
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting layout blueprints from: {template_name}")

        prs = Presentation(pptx_path)

        # Get slide dimensions
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches

        results = {
            'template': template_name,
            'slide_dimensions': {
                'width': slide_width,
                'height': slide_height,
            },
            'blueprints': [],
            'grids': [],
            'patterns': [],
        }

        # Track layout patterns
        pattern_map = {}

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1

            # Extract blueprint for this slide
            blueprint = self._extract_slide_blueprint(slide, slide_num, template_name,
                                                      slide_width, slide_height)
            if blueprint:
                results['blueprints'].append(blueprint)

                # Detect grid system
                grid = self._detect_grid_system(blueprint, slide_width, slide_height)
                if grid:
                    grid['slide_num'] = slide_num
                    grid['template'] = template_name

                    # Deduplicate grids
                    grid_key = self._create_grid_key(grid)
                    if grid_key not in pattern_map:
                        pattern_map[grid_key] = grid
                        results['grids'].append(grid)

        # Analyze patterns across slides
        results['patterns'] = self._analyze_layout_patterns(results['blueprints'])

        # Save results
        self._save_results(results, template_name)

        logger.info(f"Extracted {len(results['blueprints'])} blueprints, "
                   f"{len(results['grids'])} grids, {len(results['patterns'])} patterns")

        return results

    def _extract_slide_blueprint(self, slide, slide_num: int, template_name: str,
                                 slide_width: float, slide_height: float) -> Optional[dict]:
        """Extract a detailed layout blueprint from a slide."""
        try:
            layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'

            # Collect all shape positions and types
            elements = []
            for shape in slide.shapes:
                element = self._extract_element_info(shape, slide_width, slide_height)
                if element:
                    elements.append(element)

            if not elements:
                return None

            # Determine zones
            zones = self._identify_zones(elements, slide_width, slide_height)

            # Calculate margins
            margins = self._calculate_margins(elements, slide_width, slide_height)

            blueprint_id = self._generate_id(f"{template_name}_{slide_num}_{layout_name}")

            blueprint = {
                'id': blueprint_id,
                'template': template_name,
                'slide_num': slide_num,
                'layout_name': layout_name,
                'element_count': len(elements),
                'elements': elements,
                'zones': zones,
                'margins': margins,
                'content_types': self._categorize_content_types(elements),
                'category': self._categorize_layout(elements, zones),
            }

            return blueprint

        except Exception as e:
            logger.warning(f"Failed to extract blueprint from slide {slide_num}: {e}")
            return None

    def _extract_element_info(self, shape, slide_width: float, slide_height: float) -> Optional[dict]:
        """Extract position and type info from a shape."""
        try:
            if not hasattr(shape, 'left') or shape.left is None:
                return None

            # Get shape type
            shape_type = 'unknown'
            try:
                if hasattr(shape, 'shape_type'):
                    shape_type = str(shape.shape_type).split('.')[-1]
            except:
                pass

            # Determine element category
            element_category = self._get_element_category(shape)

            # Calculate relative positions (0-1 scale)
            left = shape.left.inches / slide_width
            top = shape.top.inches / slide_height
            width = shape.width.inches / slide_width
            height = shape.height.inches / slide_height

            # Determine zone (top, middle, bottom, left, right, center)
            zone = self._get_zone_position(left, top, width, height)

            element = {
                'shape_type': shape_type,
                'category': element_category,
                'position': {
                    'left': round(left, 4),
                    'top': round(top, 4),
                    'width': round(width, 4),
                    'height': round(height, 4),
                    'right': round(left + width, 4),
                    'bottom': round(top + height, 4),
                },
                'position_inches': {
                    'left': round(shape.left.inches, 2),
                    'top': round(shape.top.inches, 2),
                    'width': round(shape.width.inches, 2),
                    'height': round(shape.height.inches, 2),
                },
                'zone': zone,
                'has_text': hasattr(shape, 'text') and bool(shape.text),
            }

            # Extract text sample if present
            if element['has_text']:
                text = shape.text[:50].replace('\n', ' ').strip()
                element['text_sample'] = text

            return element

        except Exception as e:
            logger.debug(f"Error extracting element info: {e}")
            return None

    def _get_element_category(self, shape) -> str:
        """Categorize the element type."""
        shape_type = 'unknown'
        try:
            if hasattr(shape, 'shape_type'):
                shape_type = str(shape.shape_type).split('.')[-1]
        except:
            pass

        if shape_type == 'PICTURE':
            return 'image'
        elif shape_type == 'CHART':
            return 'chart'
        elif shape_type == 'TABLE':
            return 'table'
        elif shape_type == 'GROUP':
            return 'diagram'
        elif shape_type in ('PLACEHOLDER', 'TEXT_BOX'):
            if hasattr(shape, 'text') and shape.text:
                text = shape.text.lower()
                if len(shape.text) < 100:
                    return 'title'
                return 'text_block'
            return 'placeholder'
        elif shape_type == 'AUTO_SHAPE':
            return 'shape'
        else:
            return 'other'

    def _get_zone_position(self, left: float, top: float, width: float, height: float) -> str:
        """Determine which zone an element is in."""
        center_x = left + width / 2
        center_y = top + height / 2

        # Vertical zone
        if center_y < 0.2:
            v_zone = 'top'
        elif center_y > 0.8:
            v_zone = 'bottom'
        else:
            v_zone = 'middle'

        # Horizontal zone
        if center_x < 0.33:
            h_zone = 'left'
        elif center_x > 0.67:
            h_zone = 'right'
        else:
            h_zone = 'center'

        return f"{v_zone}_{h_zone}"

    def _identify_zones(self, elements: List[dict], slide_width: float, slide_height: float) -> List[dict]:
        """Identify content zones from element positions."""
        zones = []

        # Group elements by zone
        zone_elements = defaultdict(list)
        for elem in elements:
            zone_elements[elem['zone']].append(elem)

        # Create zone definitions
        for zone_name, elems in zone_elements.items():
            if not elems:
                continue

            # Calculate zone bounds
            min_left = min(e['position']['left'] for e in elems)
            min_top = min(e['position']['top'] for e in elems)
            max_right = max(e['position']['right'] for e in elems)
            max_bottom = max(e['position']['bottom'] for e in elems)

            # Determine primary content type
            categories = [e['category'] for e in elems]
            primary_category = max(set(categories), key=categories.count)

            zone = {
                'name': zone_name,
                'bounds': {
                    'left': round(min_left, 4),
                    'top': round(min_top, 4),
                    'right': round(max_right, 4),
                    'bottom': round(max_bottom, 4),
                    'width': round(max_right - min_left, 4),
                    'height': round(max_bottom - min_top, 4),
                },
                'element_count': len(elems),
                'primary_content': primary_category,
            }
            zones.append(zone)

        return zones

    def _calculate_margins(self, elements: List[dict], slide_width: float, slide_height: float) -> dict:
        """Calculate margins from element positions."""
        if not elements:
            return {'left': 0, 'right': 0, 'top': 0, 'bottom': 0}

        # Find extreme positions
        min_left = min(e['position']['left'] for e in elements)
        max_right = max(e['position']['right'] for e in elements)
        min_top = min(e['position']['top'] for e in elements)
        max_bottom = max(e['position']['bottom'] for e in elements)

        return {
            'left': round(min_left, 4),
            'right': round(1 - max_right, 4),
            'top': round(min_top, 4),
            'bottom': round(1 - max_bottom, 4),
            'left_inches': round(min_left * slide_width, 2),
            'right_inches': round((1 - max_right) * slide_width, 2),
            'top_inches': round(min_top * slide_height, 2),
            'bottom_inches': round((1 - max_bottom) * slide_height, 2),
        }

    def _categorize_content_types(self, elements: List[dict]) -> dict:
        """Count content types in elements."""
        counts = defaultdict(int)
        for elem in elements:
            counts[elem['category']] += 1
        return dict(counts)

    def _categorize_layout(self, elements: List[dict], zones: List[dict]) -> str:
        """Categorize the overall layout type."""
        content_types = self._categorize_content_types(elements)
        zone_names = [z['name'] for z in zones]

        # Single content type layouts
        if content_types.get('chart', 0) > 0 and len(content_types) <= 2:
            if content_types.get('chart') == 1:
                return 'single_chart'
            elif content_types.get('chart') == 2:
                return 'two_charts'
            else:
                return 'multi_chart'

        if content_types.get('table', 0) > 0 and len(content_types) <= 2:
            return 'table_layout'

        if content_types.get('image', 0) > 0 and len(content_types) <= 2:
            return 'image_layout'

        # Layout patterns based on zones
        if 'left_middle' in zone_names and 'right_middle' in zone_names:
            return 'two_column'

        if len([z for z in zone_names if 'middle' in z]) >= 3:
            return 'multi_column'

        if content_types.get('title', 0) > 0 and content_types.get('text_block', 0) > 0:
            return 'title_content'

        if len(elements) <= 2:
            return 'minimal'

        return 'mixed_content'

    def _detect_grid_system(self, blueprint: dict, slide_width: float, slide_height: float) -> Optional[dict]:
        """Detect grid system from element positions."""
        elements = blueprint.get('elements', [])
        if len(elements) < 2:
            return None

        # Collect all unique x and y positions
        x_positions = set()
        y_positions = set()

        for elem in elements:
            pos = elem['position']
            x_positions.add(round(pos['left'], 2))
            x_positions.add(round(pos['right'], 2))
            y_positions.add(round(pos['top'], 2))
            y_positions.add(round(pos['bottom'], 2))

        x_positions = sorted(x_positions)
        y_positions = sorted(y_positions)

        # Detect column structure
        columns = self._detect_columns(x_positions)
        rows = self._detect_rows(y_positions)

        if columns <= 1 and rows <= 1:
            return None

        # Calculate gutters (spaces between columns/rows)
        column_gutters = self._calculate_gutters(x_positions, columns)
        row_gutters = self._calculate_gutters(y_positions, rows)

        grid_id = self._generate_id(f"grid_{blueprint['template']}_{columns}x{rows}")

        grid = {
            'id': grid_id,
            'columns': columns,
            'rows': rows,
            'column_positions': x_positions[:10],  # Limit for clarity
            'row_positions': y_positions[:10],
            'column_gutters': column_gutters,
            'row_gutters': row_gutters,
            'margins': blueprint.get('margins', {}),
            'usage_count': 1,
        }

        return grid

    def _detect_columns(self, x_positions: List[float]) -> int:
        """Detect number of columns from x positions."""
        if len(x_positions) < 2:
            return 1

        # Look for evenly spaced positions
        gaps = []
        for i in range(1, len(x_positions)):
            gaps.append(round(x_positions[i] - x_positions[i-1], 2))

        if not gaps:
            return 1

        # Find most common gap size (indicates column width)
        gap_counts = defaultdict(int)
        for gap in gaps:
            if gap > 0.05:  # Minimum 5% width
                gap_counts[gap] += 1

        if not gap_counts:
            return 1

        # Estimate columns based on positions
        # If positions are evenly distributed, we can infer column count
        content_width = max(x_positions) - min(x_positions)
        if content_width > 0.6:  # More than 60% of slide width used
            if len(x_positions) >= 6:
                return 3
            elif len(x_positions) >= 4:
                return 2

        return 1

    def _detect_rows(self, y_positions: List[float]) -> int:
        """Detect number of rows from y positions."""
        if len(y_positions) < 2:
            return 1

        # Similar logic to columns
        content_height = max(y_positions) - min(y_positions)
        if content_height > 0.5:  # More than 50% of slide height used
            if len(y_positions) >= 6:
                return 3
            elif len(y_positions) >= 4:
                return 2

        return 1

    def _calculate_gutters(self, positions: List[float], count: int) -> List[float]:
        """Calculate gutter sizes between columns/rows."""
        if count <= 1 or len(positions) < 3:
            return []

        # Look for small gaps that represent gutters
        gutters = []
        for i in range(1, len(positions)):
            gap = positions[i] - positions[i-1]
            if 0.01 < gap < 0.1:  # Between 1% and 10% - likely a gutter
                gutters.append(round(gap, 4))

        return gutters[:5]  # Return up to 5 gutters

    def _create_grid_key(self, grid: dict) -> str:
        """Create a key for grid deduplication."""
        return f"{grid['columns']}x{grid['rows']}_{round(grid['margins'].get('left', 0), 2)}"

    def _analyze_layout_patterns(self, blueprints: List[dict]) -> List[dict]:
        """Analyze patterns across multiple blueprints."""
        patterns = []

        # Group by category
        by_category = defaultdict(list)
        for bp in blueprints:
            by_category[bp.get('category', 'unknown')].append(bp)

        for category, bps in by_category.items():
            if len(bps) < 2:
                continue

            # Find common zone structures
            zone_structures = defaultdict(int)
            for bp in bps:
                zones = tuple(sorted(z['name'] for z in bp.get('zones', [])))
                zone_structures[zones] += 1

            # Most common structure
            if zone_structures:
                most_common = max(zone_structures.items(), key=lambda x: x[1])
                pattern = {
                    'category': category,
                    'zone_structure': list(most_common[0]),
                    'occurrence_count': most_common[1],
                    'total_slides': len(bps),
                }
                patterns.append(pattern)

        return patterns

    def _save_results(self, results: dict, template_name: str):
        """Save extraction results."""
        # Save individual blueprints
        for bp in results['blueprints']:
            bp_id = bp['id']
            filename = f"{bp_id}.json"
            filepath = self.layouts_dir / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(bp, f, indent=2)
            bp['filename'] = filename

            # Add to index
            existing = [x for x in self.index['blueprints'] if x.get('id') == bp_id]
            if not existing:
                self.index['blueprints'].append(bp)

        # Save grids
        for grid in results['grids']:
            grid_id = grid['id']
            filename = f"{grid_id}.json"
            filepath = self.grids_dir / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(grid, f, indent=2)
            grid['filename'] = filename

            # Add to index
            existing = [x for x in self.index['grids'] if x.get('id') == grid_id]
            if not existing:
                self.index['grids'].append(grid)

        # Save patterns to index
        for pattern in results['patterns']:
            existing = [p for p in self.index['patterns']
                       if p.get('category') == pattern.get('category')]
            if not existing:
                self.index['patterns'].append(pattern)

        self._save_index()

    def get_summary(self) -> dict:
        """Get summary of extracted layouts."""
        summary = {
            'total_blueprints': len(self.index['blueprints']),
            'total_grids': len(self.index['grids']),
            'total_patterns': len(self.index['patterns']),
            'by_category': defaultdict(int),
            'grid_types': defaultdict(int),
        }

        for bp in self.index['blueprints']:
            summary['by_category'][bp.get('category', 'unknown')] += 1

        for grid in self.index['grids']:
            grid_type = f"{grid['columns']}x{grid['rows']}"
            summary['grid_types'][grid_type] += 1

        summary['by_category'] = dict(summary['by_category'])
        summary['grid_types'] = dict(summary['grid_types'])

        return summary

    def search_blueprints(self, category: str = None, template: str = None,
                          min_elements: int = None, content_type: str = None) -> List[dict]:
        """Search layout blueprints."""
        results = []

        for bp in self.index['blueprints']:
            if category and bp.get('category') != category:
                continue
            if template and bp.get('template') != template:
                continue
            if min_elements and bp.get('element_count', 0) < min_elements:
                continue
            if content_type:
                content_types = bp.get('content_types', {})
                if content_type not in content_types or content_types[content_type] == 0:
                    continue
            results.append(bp)

        return results

    def search_grids(self, columns: int = None, rows: int = None) -> List[dict]:
        """Search grid systems."""
        results = []

        for grid in self.index['grids']:
            if columns and grid.get('columns') != columns:
                continue
            if rows and grid.get('rows') != rows:
                continue
            results.append(grid)

        return results

    def find_layout_for_content(self, content_description: dict) -> List[dict]:
        """
        Find suitable layouts for given content.

        Args:
            content_description: Dict with content counts, e.g.:
                {'chart': 2, 'table': 1, 'text_block': 1}

        Returns:
            List of matching blueprints
        """
        results = []

        for bp in self.index['blueprints']:
            content_types = bp.get('content_types', {})

            # Check if blueprint has capacity for all requested content
            matches = True
            for content_type, count in content_description.items():
                if content_types.get(content_type, 0) < count:
                    matches = False
                    break

            if matches:
                results.append(bp)

        # Sort by closest match (least excess capacity)
        def excess_capacity(bp):
            total_excess = 0
            content_types = bp.get('content_types', {})
            for content_type, count in content_description.items():
                total_excess += content_types.get(content_type, 0) - count
            return total_excess

        results.sort(key=excess_capacity)

        return results[:10]  # Return top 10 matches


def extract_layouts_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract layouts from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = LayoutBlueprintExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting layouts from: {pptx_path.name}")
        try:
            extractor.extract_template_layouts(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python layout_blueprint_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_layouts_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("LAYOUT EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
