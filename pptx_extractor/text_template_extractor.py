"""
Text Content Template Extractor - Extract text patterns and structures from PowerPoint templates.

This module extracts:
- Bullet point patterns (structure, levels, formatting)
- Title patterns (question-style, statement, etc.)
- Text block templates
- Callout/label patterns
- Paragraph formatting presets
"""

import json
import hashlib
import logging
import re
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any, Tuple
from collections import defaultdict

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class TextTemplateExtractor:
    """Extract text content templates from PowerPoint templates."""

    def __init__(self, output_dir: Path):
        """Initialize the text template extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories
        self.dirs = {
            'bullet_patterns': self.output_dir / 'text_templates' / 'bullet_patterns',
            'title_patterns': self.output_dir / 'text_templates' / 'title_patterns',
            'text_blocks': self.output_dir / 'text_templates' / 'text_blocks',
            'callouts': self.output_dir / 'text_templates' / 'callouts',
        }
        for d in self.dirs.values():
            d.mkdir(parents=True, exist_ok=True)

        # Text template index
        self.index_path = self.output_dir / 'text_templates' / 'text_template_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'bullet_patterns': [],
            'title_patterns': [],
            'text_blocks': [],
            'callouts': [],
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def extract_template_text_patterns(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all text patterns from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted text patterns
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting text patterns from: {template_name}")

        prs = Presentation(pptx_path)

        results = {
            'template': template_name,
            'bullet_patterns': [],
            'title_patterns': [],
            'text_blocks': [],
            'callouts': [],
        }

        # Pattern maps for deduplication
        bullet_map = {}
        title_map = {}
        text_block_map = {}
        callout_map = {}

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1

            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue

                text_frame = shape.text_frame

                if not text_frame.paragraphs:
                    continue

                # Analyze text content
                text_info = self._analyze_text_frame(text_frame, shape, slide_num, template_name)

                if text_info:
                    pattern_type = text_info.get('pattern_type')

                    if pattern_type == 'bullet_list':
                        key = self._create_bullet_key(text_info)
                        if key not in bullet_map:
                            bullet_map[key] = text_info
                        else:
                            bullet_map[key]['usage_count'] += 1

                    elif pattern_type == 'title':
                        key = self._create_title_key(text_info)
                        if key not in title_map:
                            title_map[key] = text_info
                        else:
                            title_map[key]['usage_count'] += 1

                    elif pattern_type == 'text_block':
                        key = self._create_text_block_key(text_info)
                        if key not in text_block_map:
                            text_block_map[key] = text_info
                        else:
                            text_block_map[key]['usage_count'] += 1

                    elif pattern_type == 'callout':
                        key = text_info.get('callout_type', 'generic')
                        if key not in callout_map:
                            callout_map[key] = text_info
                        else:
                            callout_map[key]['usage_count'] += 1

        # Process and save patterns
        self._save_patterns(bullet_map, 'bullet_patterns', results)
        self._save_patterns(title_map, 'title_patterns', results)
        self._save_patterns(text_block_map, 'text_blocks', results)
        self._save_patterns(callout_map, 'callouts', results)

        self._save_index()

        logger.info(f"Extracted: {len(results['bullet_patterns'])} bullet patterns, "
                   f"{len(results['title_patterns'])} title patterns, "
                   f"{len(results['text_blocks'])} text blocks, "
                   f"{len(results['callouts'])} callouts")

        return results

    def _analyze_text_frame(self, text_frame, shape, slide_num: int, template_name: str) -> Optional[dict]:
        """Analyze a text frame and categorize it."""
        try:
            paragraphs = list(text_frame.paragraphs)
            if not paragraphs:
                return None

            # Get full text
            full_text = '\n'.join(p.text for p in paragraphs if p.text)
            if not full_text.strip():
                return None

            # Extract paragraph structures
            para_structures = []
            max_level = 0
            has_bullets = False

            for para in paragraphs:
                if not para.text.strip():
                    continue

                para_info = self._extract_paragraph_info(para)
                para_structures.append(para_info)

                if para_info['level'] > max_level:
                    max_level = para_info['level']
                if para_info.get('has_bullet'):
                    has_bullets = True

            if not para_structures:
                return None

            # Determine pattern type
            pattern_type = self._determine_pattern_type(para_structures, shape, full_text)

            info = {
                'template': template_name,
                'slide_num': slide_num,
                'pattern_type': pattern_type,
                'paragraph_count': len(para_structures),
                'max_level': max_level,
                'has_bullets': has_bullets,
                'structure': para_structures,
                'sample_text': full_text[:200],
                'usage_count': 1,
            }

            # Add type-specific info
            if pattern_type == 'title':
                info['title_style'] = self._analyze_title_style(full_text, para_structures[0])

            elif pattern_type == 'bullet_list':
                info['bullet_style'] = self._analyze_bullet_style(para_structures)

            elif pattern_type == 'callout':
                info['callout_type'] = self._analyze_callout_type(shape, full_text)

            return info

        except Exception as e:
            logger.debug(f"Error analyzing text frame: {e}")
            return None

    def _extract_paragraph_info(self, para) -> dict:
        """Extract information from a paragraph."""
        info = {
            'level': 0,
            'has_bullet': False,
            'bullet_char': None,
            'alignment': 'left',
            'text_length': len(para.text) if para.text else 0,
            'word_count': len(para.text.split()) if para.text else 0,
        }

        try:
            # Level
            if hasattr(para, 'level') and para.level is not None:
                info['level'] = para.level

            # Alignment
            if para.alignment:
                info['alignment'] = str(para.alignment).split('.')[-1].lower()

            # Bullet (try to detect)
            if para.text:
                first_char = para.text.strip()[:1] if para.text.strip() else ''
                if first_char in ['•', '○', '●', '■', '□', '▪', '▫', '-', '*', '–', '►', '➤']:
                    info['has_bullet'] = True
                    info['bullet_char'] = first_char

            # Font info from first run
            if para.runs:
                run = para.runs[0]
                if run.font.size:
                    info['font_size'] = run.font.size.pt
                if run.font.name:
                    info['font_name'] = run.font.name
                info['bold'] = bool(run.font.bold)
                info['italic'] = bool(run.font.italic)

        except:
            pass

        return info

    def _determine_pattern_type(self, para_structures: List[dict], shape, full_text: str) -> str:
        """Determine the type of text pattern."""
        # Check for single short paragraph (likely title)
        if len(para_structures) == 1:
            para = para_structures[0]
            if para['word_count'] <= 10 and para.get('font_size', 0) >= 18:
                return 'title'
            if para['word_count'] <= 3:
                return 'callout'

        # Check for bullet list
        if any(p.get('has_bullet') or p.get('level', 0) > 0 for p in para_structures):
            return 'bullet_list'

        # Check if it's a multi-level list
        if len(set(p.get('level', 0) for p in para_structures)) > 1:
            return 'bullet_list'

        # Check for short callout
        total_words = sum(p['word_count'] for p in para_structures)
        if total_words <= 15 and len(para_structures) <= 2:
            return 'callout'

        # Check shape type for callout
        try:
            shape_type = str(shape.shape_type).split('.')[-1]
            if 'CALLOUT' in shape_type:
                return 'callout'
        except:
            pass

        return 'text_block'

    def _analyze_title_style(self, text: str, para_info: dict) -> dict:
        """Analyze title style."""
        style = {
            'format': 'statement',
            'case': 'mixed',
            'font_size': para_info.get('font_size'),
            'bold': para_info.get('bold', False),
            'alignment': para_info.get('alignment', 'left'),
        }

        # Determine format
        if text.strip().endswith('?'):
            style['format'] = 'question'
        elif text.strip().endswith(':'):
            style['format'] = 'header'
        elif ':' in text and text.index(':') < len(text) / 2:
            style['format'] = 'labeled'
        elif text.isupper():
            style['format'] = 'uppercase'
            style['case'] = 'upper'
        elif text.istitle():
            style['case'] = 'title'

        # Check for common patterns
        text_lower = text.lower()
        if text_lower.startswith('key '):
            style['format'] = 'key_point'
        elif any(text_lower.startswith(x) for x in ['summary', 'overview', 'agenda', 'contents']):
            style['format'] = 'section_header'

        return style

    def _analyze_bullet_style(self, para_structures: List[dict]) -> dict:
        """Analyze bullet list style."""
        style = {
            'total_items': len(para_structures),
            'max_depth': max(p.get('level', 0) for p in para_structures),
            'items_per_level': defaultdict(int),
            'bullet_chars': set(),
            'avg_words_per_item': 0,
        }

        total_words = 0
        for para in para_structures:
            level = para.get('level', 0)
            style['items_per_level'][level] += 1
            total_words += para.get('word_count', 0)

            if para.get('bullet_char'):
                style['bullet_chars'].add(para['bullet_char'])

        style['items_per_level'] = dict(style['items_per_level'])
        style['bullet_chars'] = list(style['bullet_chars'])
        style['avg_words_per_item'] = round(total_words / len(para_structures), 1) if para_structures else 0

        # Determine bullet pattern type
        if style['max_depth'] == 0:
            if style['total_items'] <= 3:
                style['pattern'] = 'short_list'
            elif style['total_items'] <= 5:
                style['pattern'] = 'key_points'
            else:
                style['pattern'] = 'long_list'
        else:
            if style['max_depth'] == 1:
                style['pattern'] = 'two_level'
            else:
                style['pattern'] = 'multi_level'

        return style

    def _analyze_callout_type(self, shape, text: str) -> str:
        """Analyze callout type."""
        try:
            shape_type = str(shape.shape_type).split('.')[-1]

            if 'CALLOUT' in shape_type:
                return 'speech_bubble'
        except:
            pass

        text_lower = text.lower()

        # Check for specific callout types
        if any(x in text_lower for x in ['note:', 'important:', 'tip:', 'hint:']):
            return 'note'
        elif any(x in text_lower for x in ['warning:', 'caution:', 'alert:']):
            return 'warning'
        elif any(x in text_lower for x in ['example:', 'e.g.:', 'for example:']):
            return 'example'
        elif text.startswith('"') and text.endswith('"'):
            return 'quote'
        elif len(text) <= 20:
            return 'label'
        else:
            return 'annotation'

    def _create_bullet_key(self, info: dict) -> str:
        """Create key for bullet pattern deduplication."""
        style = info.get('bullet_style', {})
        return f"{style.get('pattern', 'unknown')}_{style.get('max_depth', 0)}_{info.get('paragraph_count', 0)}"

    def _create_title_key(self, info: dict) -> str:
        """Create key for title pattern deduplication."""
        style = info.get('title_style', {})
        return f"{style.get('format', 'statement')}_{style.get('case', 'mixed')}"

    def _create_text_block_key(self, info: dict) -> str:
        """Create key for text block deduplication."""
        return f"{info.get('paragraph_count', 0)}_{info.get('max_level', 0)}"

    def _save_patterns(self, pattern_map: dict, pattern_type: str, results: dict):
        """Save patterns to files and index."""
        save_dir = self.dirs.get(pattern_type, self.dirs['text_blocks'])

        for key, pattern in pattern_map.items():
            pattern_id = self._generate_id(f"{pattern.get('template', '')}_{pattern_type}_{key}")
            pattern['id'] = pattern_id

            # Save pattern file
            filename = f"{pattern_id}.json"
            filepath = save_dir / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(pattern, f, indent=2)

            pattern['filename'] = filename
            results[pattern_type].append(pattern)

            # Add to index
            existing = [x for x in self.index[pattern_type] if x.get('id') == pattern_id]
            if not existing:
                self.index[pattern_type].append(pattern)

    def get_summary(self) -> dict:
        """Get summary of extracted text patterns."""
        summary = {
            'bullet_patterns': len(self.index['bullet_patterns']),
            'title_patterns': len(self.index['title_patterns']),
            'text_blocks': len(self.index['text_blocks']),
            'callouts': len(self.index['callouts']),
            'by_bullet_pattern': defaultdict(int),
            'by_title_format': defaultdict(int),
        }

        for bp in self.index['bullet_patterns']:
            pattern = bp.get('bullet_style', {}).get('pattern', 'unknown')
            summary['by_bullet_pattern'][pattern] += 1

        for tp in self.index['title_patterns']:
            format_type = tp.get('title_style', {}).get('format', 'statement')
            summary['by_title_format'][format_type] += 1

        summary['by_bullet_pattern'] = dict(summary['by_bullet_pattern'])
        summary['by_title_format'] = dict(summary['by_title_format'])

        return summary

    def search_bullet_patterns(self, pattern: str = None, max_depth: int = None,
                               template: str = None) -> List[dict]:
        """Search bullet patterns."""
        results = []

        for bp in self.index['bullet_patterns']:
            if template and bp.get('template') != template:
                continue
            if pattern and bp.get('bullet_style', {}).get('pattern') != pattern:
                continue
            if max_depth is not None and bp.get('bullet_style', {}).get('max_depth', 0) > max_depth:
                continue
            results.append(bp)

        return results

    def search_title_patterns(self, format_type: str = None, template: str = None) -> List[dict]:
        """Search title patterns."""
        results = []

        for tp in self.index['title_patterns']:
            if template and tp.get('template') != template:
                continue
            if format_type and tp.get('title_style', {}).get('format') != format_type:
                continue
            results.append(tp)

        return results

    def get_by_id(self, pattern_id: str) -> Optional[dict]:
        """Get a specific pattern by ID."""
        for category in ['bullet_patterns', 'title_patterns', 'text_blocks', 'callouts']:
            for pattern in self.index[category]:
                if pattern.get('id') == pattern_id:
                    return pattern
        return None

    def generate_bullet_template(self, items: int, depth: int = 0) -> dict:
        """Generate a bullet template with placeholder text."""
        template = {
            'type': 'bullet_list',
            'items': [],
        }

        for i in range(items):
            item = {
                'level': 0,
                'placeholder': f'Point {i + 1}',
            }
            template['items'].append(item)

            if depth > 0:
                for j in range(depth):
                    sub_item = {
                        'level': 1,
                        'placeholder': f'Sub-point {i + 1}.{j + 1}',
                    }
                    template['items'].append(sub_item)

        return template


def extract_text_patterns_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract text patterns from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = TextTemplateExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting text patterns from: {pptx_path.name}")
        try:
            extractor.extract_template_text_patterns(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python text_template_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_text_patterns_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("TEXT PATTERN EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
