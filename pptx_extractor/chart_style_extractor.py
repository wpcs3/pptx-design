"""
Chart Style Extractor - Extract chart formatting profiles from PowerPoint templates.

This module extracts chart styling separate from chart data:
- Series colors and fill patterns
- Axis formatting (fonts, number formats, gridlines)
- Legend settings (position, font)
- Data label settings
- Title formatting
- Plot area styling
"""

import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any
from collections import defaultdict
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


class ChartStyleExtractor:
    """Extract chart formatting profiles from PowerPoint templates."""

    # XML namespaces
    NAMESPACES = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    def __init__(self, output_dir: Path):
        """Initialize the chart style extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectory for chart styles
        self.styles_dir = self.output_dir / 'styles' / 'chart_styles'
        self.styles_dir.mkdir(parents=True, exist_ok=True)

        # Style index
        self.index_path = self.output_dir / 'styles' / 'chart_style_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'chart_styles': [],
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def extract_template_chart_styles(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all chart styles from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted chart styles
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting chart styles from: {template_name}")

        prs = Presentation(pptx_path)
        chart_styles = []
        style_map = {}  # For deduplication

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_chart:
                    continue

                try:
                    style = self._extract_chart_style(shape.chart, shape, template_name, slide_idx + 1)
                    if style:
                        # Create deduplication key based on core styling properties
                        style_key = self._create_style_key(style)
                        if style_key not in style_map:
                            style_map[style_key] = style
                        else:
                            # Increment usage count
                            style_map[style_key]['usage_count'] = style_map[style_key].get('usage_count', 1) + 1
                except Exception as e:
                    logger.warning(f"Failed to extract chart style from slide {slide_idx + 1}: {e}")

        # Convert to list and save
        for key, style in style_map.items():
            style_id = self._generate_id(f"{template_name}_{key}")
            style['id'] = style_id

            # Create style name
            chart_type = style.get('chart_type', 'unknown')
            style['style_name'] = f"{template_name}_{chart_type}_style"

            # Save style file
            filename = f"{style_id}.json"
            filepath = self.styles_dir / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(style, f, indent=2)

            style['filename'] = filename
            chart_styles.append(style)

            # Add to index
            existing = [x for x in self.index['chart_styles'] if x.get('id') == style_id]
            if not existing:
                self.index['chart_styles'].append(style)

        self._save_index()

        logger.info(f"Extracted {len(chart_styles)} unique chart styles")
        return {'chart_styles': chart_styles}

    def _create_style_key(self, style: dict) -> str:
        """Create a key for deduplication based on style properties."""
        key_parts = [
            style.get('chart_type', ''),
            str(style.get('series_colors', [])),
            str(style.get('category_axis', {}).get('font_size', '')),
            str(style.get('value_axis', {}).get('number_format', '')),
            str(style.get('legend', {}).get('position', '')),
        ]
        return '_'.join(key_parts)

    def _extract_chart_style(self, chart, shape, template_name: str, slide_num: int) -> Optional[dict]:
        """Extract comprehensive style information from a chart."""
        try:
            # Basic chart info
            chart_type = 'unknown'
            try:
                chart_type = str(chart.chart_type).split('.')[-1].lower() if chart.chart_type else 'unknown'
            except:
                pass

            style = {
                'chart_type': chart_type,
                'template': template_name,
                'slide_num': slide_num,
                'usage_count': 1,
                'series_colors': [],
                'series_fill_styles': [],
                'category_axis': {},
                'value_axis': {},
                'legend': {},
                'title': {},
                'data_labels': {},
                'plot_area': {},
                'gridlines': {},
            }

            # Extract series colors and fill styles
            style['series_colors'], style['series_fill_styles'] = self._extract_series_styles(chart)

            # Extract axis formatting
            style['category_axis'] = self._extract_axis_style(chart, 'category')
            style['value_axis'] = self._extract_axis_style(chart, 'value')

            # Extract legend settings
            style['legend'] = self._extract_legend_style(chart)

            # Extract title formatting
            style['title'] = self._extract_title_style(chart)

            # Extract data label settings
            style['data_labels'] = self._extract_data_label_style(chart)

            # Extract plot area styling
            style['plot_area'] = self._extract_plot_area_style(chart)

            # Extract gridline settings
            style['gridlines'] = self._extract_gridline_style(chart)

            return style

        except Exception as e:
            logger.debug(f"Error extracting chart style: {e}")
            return None

    def _extract_series_styles(self, chart) -> tuple:
        """Extract series colors and fill styles."""
        colors = []
        fill_styles = []

        try:
            for series in chart.series:
                color = None
                fill_style = 'solid'

                # Try to get series color
                try:
                    if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                        fill = series.format.fill
                        if fill.type is not None:
                            fill_type_str = str(fill.type)
                            if 'SOLID' in fill_type_str:
                                fill_style = 'solid'
                                if hasattr(fill, 'fore_color') and fill.fore_color:
                                    try:
                                        color = f"#{fill.fore_color.rgb}"
                                    except:
                                        pass
                            elif 'GRADIENT' in fill_type_str:
                                fill_style = 'gradient'
                            elif 'PATTERN' in fill_type_str:
                                fill_style = 'pattern'
                except:
                    pass

                colors.append(color)
                fill_styles.append(fill_style)

        except Exception as e:
            logger.debug(f"Error extracting series styles: {e}")

        return colors, fill_styles

    def _extract_axis_style(self, chart, axis_type: str) -> dict:
        """Extract axis formatting."""
        axis_style = {
            'visible': True,
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'font_color': None,
            'number_format': None,
            'major_gridlines': False,
            'minor_gridlines': False,
            'title': None,
            'title_font_size': None,
            'reverse_order': False,
            'logarithmic': False,
            'min_value': None,
            'max_value': None,
        }

        try:
            axis = None
            if axis_type == 'category':
                if hasattr(chart, 'category_axis'):
                    axis = chart.category_axis
            else:  # value axis
                if hasattr(chart, 'value_axis'):
                    axis = chart.value_axis

            if axis is None:
                return axis_style

            axis_style['visible'] = axis.visible if hasattr(axis, 'visible') else True

            # Font settings
            if hasattr(axis, 'tick_labels') and axis.tick_labels:
                tick_labels = axis.tick_labels
                if hasattr(tick_labels, 'font'):
                    font = tick_labels.font
                    if font.name:
                        axis_style['font_name'] = font.name
                    if font.size:
                        axis_style['font_size'] = font.size.pt
                    axis_style['font_bold'] = bool(font.bold)
                    if font.color and font.color.rgb:
                        try:
                            axis_style['font_color'] = f"#{font.color.rgb}"
                        except:
                            pass

                if hasattr(tick_labels, 'number_format'):
                    axis_style['number_format'] = tick_labels.number_format

            # Gridlines
            if hasattr(axis, 'has_major_gridlines'):
                axis_style['major_gridlines'] = axis.has_major_gridlines
            if hasattr(axis, 'has_minor_gridlines'):
                axis_style['minor_gridlines'] = axis.has_minor_gridlines

            # Title
            if hasattr(axis, 'has_title') and axis.has_title:
                if hasattr(axis, 'axis_title') and axis.axis_title:
                    try:
                        axis_style['title'] = axis.axis_title.text_frame.text
                    except:
                        pass

            # Scale settings (value axis only)
            if axis_type == 'value':
                if hasattr(axis, 'minimum_scale') and axis.minimum_scale is not None:
                    axis_style['min_value'] = axis.minimum_scale
                if hasattr(axis, 'maximum_scale') and axis.maximum_scale is not None:
                    axis_style['max_value'] = axis.maximum_scale

        except Exception as e:
            logger.debug(f"Error extracting {axis_type} axis style: {e}")

        return axis_style

    def _extract_legend_style(self, chart) -> dict:
        """Extract legend settings."""
        legend_style = {
            'visible': False,
            'position': None,
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'overlay': False,
        }

        try:
            if not hasattr(chart, 'has_legend') or not chart.has_legend:
                return legend_style

            legend_style['visible'] = True
            legend = chart.legend

            # Position
            if hasattr(legend, 'position'):
                legend_style['position'] = str(legend.position).split('.')[-1].lower()

            # Font
            if hasattr(legend, 'font'):
                font = legend.font
                if font.name:
                    legend_style['font_name'] = font.name
                if font.size:
                    legend_style['font_size'] = font.size.pt
                legend_style['font_bold'] = bool(font.bold)

            # Overlay
            if hasattr(legend, 'include_in_layout'):
                legend_style['overlay'] = not legend.include_in_layout

        except Exception as e:
            logger.debug(f"Error extracting legend style: {e}")

        return legend_style

    def _extract_title_style(self, chart) -> dict:
        """Extract chart title formatting."""
        title_style = {
            'visible': False,
            'text': None,
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'font_color': None,
            'overlay': False,
        }

        try:
            if not hasattr(chart, 'has_title') or not chart.has_title:
                return title_style

            title_style['visible'] = True

            if hasattr(chart, 'chart_title'):
                title = chart.chart_title
                if hasattr(title, 'text_frame') and title.text_frame:
                    try:
                        title_style['text'] = title.text_frame.text
                    except:
                        pass

                    # Extract font from first paragraph
                    if title.text_frame.paragraphs:
                        para = title.text_frame.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                            if run.font.name:
                                title_style['font_name'] = run.font.name
                            if run.font.size:
                                title_style['font_size'] = run.font.size.pt
                            title_style['font_bold'] = bool(run.font.bold)
                            if run.font.color and run.font.color.rgb:
                                try:
                                    title_style['font_color'] = f"#{run.font.color.rgb}"
                                except:
                                    pass

        except Exception as e:
            logger.debug(f"Error extracting title style: {e}")

        return title_style

    def _extract_data_label_style(self, chart) -> dict:
        """Extract data label settings."""
        label_style = {
            'visible': False,
            'position': None,
            'show_value': False,
            'show_category': False,
            'show_percentage': False,
            'show_series_name': False,
            'separator': None,
            'font_name': None,
            'font_size': None,
            'font_bold': False,
            'font_color': None,
            'number_format': None,
        }

        try:
            # Check plots for data labels
            if hasattr(chart, 'plots') and chart.plots:
                for plot in chart.plots:
                    if hasattr(plot, 'has_data_labels') and plot.has_data_labels:
                        label_style['visible'] = True
                        data_labels = plot.data_labels

                        # Position
                        if hasattr(data_labels, 'position'):
                            label_style['position'] = str(data_labels.position).split('.')[-1].lower()

                        # What to show
                        if hasattr(data_labels, 'show_value'):
                            label_style['show_value'] = data_labels.show_value
                        if hasattr(data_labels, 'show_category_name'):
                            label_style['show_category'] = data_labels.show_category_name
                        if hasattr(data_labels, 'show_percentage'):
                            label_style['show_percentage'] = data_labels.show_percentage
                        if hasattr(data_labels, 'show_series_name'):
                            label_style['show_series_name'] = data_labels.show_series_name

                        # Number format
                        if hasattr(data_labels, 'number_format'):
                            label_style['number_format'] = data_labels.number_format

                        # Font
                        if hasattr(data_labels, 'font'):
                            font = data_labels.font
                            if font.name:
                                label_style['font_name'] = font.name
                            if font.size:
                                label_style['font_size'] = font.size.pt
                            label_style['font_bold'] = bool(font.bold)

                        break  # Use first plot with data labels

        except Exception as e:
            logger.debug(f"Error extracting data label style: {e}")

        return label_style

    def _extract_plot_area_style(self, chart) -> dict:
        """Extract plot area styling."""
        plot_style = {
            'fill_type': None,
            'fill_color': None,
            'border_visible': False,
            'border_color': None,
            'border_width': None,
        }

        try:
            if hasattr(chart, 'plot_area') and chart.plot_area:
                plot_area = chart.plot_area

                # Fill
                if hasattr(plot_area, 'format') and hasattr(plot_area.format, 'fill'):
                    fill = plot_area.format.fill
                    if fill.type is not None:
                        plot_style['fill_type'] = str(fill.type).split('.')[-1].lower()
                        if hasattr(fill, 'fore_color') and fill.fore_color:
                            try:
                                plot_style['fill_color'] = f"#{fill.fore_color.rgb}"
                            except:
                                pass

                # Border
                if hasattr(plot_area, 'format') and hasattr(plot_area.format, 'line'):
                    line = plot_area.format.line
                    if line.fill.type is not None:
                        plot_style['border_visible'] = True
                        if hasattr(line, 'color') and line.color:
                            try:
                                plot_style['border_color'] = f"#{line.color.rgb}"
                            except:
                                pass
                        if hasattr(line, 'width') and line.width:
                            plot_style['border_width'] = line.width.pt

        except Exception as e:
            logger.debug(f"Error extracting plot area style: {e}")

        return plot_style

    def _extract_gridline_style(self, chart) -> dict:
        """Extract gridline settings."""
        gridline_style = {
            'major_category': False,
            'minor_category': False,
            'major_value': False,
            'minor_value': False,
            'line_color': None,
            'line_width': None,
            'line_style': None,
        }

        try:
            # Category axis gridlines
            if hasattr(chart, 'category_axis'):
                axis = chart.category_axis
                if hasattr(axis, 'has_major_gridlines'):
                    gridline_style['major_category'] = axis.has_major_gridlines
                if hasattr(axis, 'has_minor_gridlines'):
                    gridline_style['minor_category'] = axis.has_minor_gridlines

            # Value axis gridlines
            if hasattr(chart, 'value_axis'):
                axis = chart.value_axis
                if hasattr(axis, 'has_major_gridlines'):
                    gridline_style['major_value'] = axis.has_major_gridlines
                if hasattr(axis, 'has_minor_gridlines'):
                    gridline_style['minor_value'] = axis.has_minor_gridlines

                # Try to get gridline formatting
                if hasattr(axis, 'major_gridlines') and axis.has_major_gridlines:
                    try:
                        gridlines = axis.major_gridlines
                        if hasattr(gridlines, 'format') and hasattr(gridlines.format, 'line'):
                            line = gridlines.format.line
                            if hasattr(line, 'color') and line.color:
                                try:
                                    gridline_style['line_color'] = f"#{line.color.rgb}"
                                except:
                                    pass
                            if hasattr(line, 'width') and line.width:
                                gridline_style['line_width'] = line.width.pt
                    except:
                        pass

        except Exception as e:
            logger.debug(f"Error extracting gridline style: {e}")

        return gridline_style

    def get_summary(self) -> dict:
        """Get summary of extracted chart styles."""
        summary = {
            'total_styles': len(self.index['chart_styles']),
            'by_chart_type': defaultdict(int),
        }

        for style in self.index['chart_styles']:
            chart_type = style.get('chart_type', 'unknown')
            summary['by_chart_type'][chart_type] += 1

        summary['by_chart_type'] = dict(summary['by_chart_type'])
        return summary

    def search(self, chart_type: str = None, template: str = None,
               has_legend: bool = None, has_data_labels: bool = None) -> List[dict]:
        """Search chart styles."""
        results = []

        for style in self.index['chart_styles']:
            if chart_type and style.get('chart_type') != chart_type:
                continue
            if template and style.get('template') != template:
                continue
            if has_legend is not None:
                if style.get('legend', {}).get('visible') != has_legend:
                    continue
            if has_data_labels is not None:
                if style.get('data_labels', {}).get('visible') != has_data_labels:
                    continue
            results.append(style)

        return results

    def get_style_by_id(self, style_id: str) -> Optional[dict]:
        """Get a specific chart style by ID."""
        for style in self.index['chart_styles']:
            if style.get('id') == style_id:
                return style
        return None


def extract_chart_styles_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract chart styles from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = ChartStyleExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting chart styles from: {pptx_path.name}")
        try:
            extractor.extract_template_chart_styles(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python chart_style_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_chart_styles_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("CHART STYLE EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
