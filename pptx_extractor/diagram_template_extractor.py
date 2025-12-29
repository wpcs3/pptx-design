"""
Diagram Template Extractor - Extract shape combinations as reusable diagram templates.

This module extracts:
- Process flows (sequential shapes with arrows)
- Comparison matrices
- Org charts / hierarchies
- Timeline templates
- Cycle diagrams
- Custom grouped diagrams
"""

import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any, Tuple
from collections import defaultdict
import math

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class DiagramTemplateExtractor:
    """Extract diagram templates (shape combinations) from PowerPoint templates."""

    def __init__(self, output_dir: Path):
        """Initialize the diagram template extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories for diagram types
        self.dirs = {
            'process_flows': self.output_dir / 'diagrams' / 'process_flows',
            'matrices': self.output_dir / 'diagrams' / 'matrices',
            'hierarchies': self.output_dir / 'diagrams' / 'hierarchies',
            'timelines': self.output_dir / 'diagrams' / 'timelines',
            'cycles': self.output_dir / 'diagrams' / 'cycles',
            'custom': self.output_dir / 'diagrams' / 'custom',
        }
        for d in self.dirs.values():
            d.mkdir(parents=True, exist_ok=True)

        # Diagram index
        self.index_path = self.output_dir / 'diagrams' / 'diagram_template_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'templates': [],
            'categories': {
                'process_flows': [],
                'matrices': [],
                'hierarchies': [],
                'timelines': [],
                'cycles': [],
                'custom': [],
            },
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def extract_template_diagrams(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all diagram templates from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted diagram templates
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting diagram templates from: {template_name}")

        prs = Presentation(pptx_path)

        # Get slide dimensions
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches

        results = {
            'template': template_name,
            'diagrams': [],
            'by_category': defaultdict(list),
        }

        diagram_map = {}  # For deduplication

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1

            # Find groups (likely diagrams)
            for shape in slide.shapes:
                shape_type = self._get_shape_type(shape)

                if shape_type == 'GROUP':
                    diagram = self._extract_group_diagram(shape, slide_num, template_name,
                                                         slide_width, slide_height)
                    if diagram:
                        # Deduplicate by structure
                        diag_key = self._create_diagram_key(diagram)
                        if diag_key not in diagram_map:
                            diagram_map[diag_key] = diagram
                        else:
                            diagram_map[diag_key]['usage_count'] += 1

            # Also detect diagrams from ungrouped shape arrangements
            ungrouped_diagrams = self._detect_ungrouped_diagrams(slide, slide_num, template_name,
                                                                  slide_width, slide_height)
            for diagram in ungrouped_diagrams:
                diag_key = self._create_diagram_key(diagram)
                if diag_key not in diagram_map:
                    diagram_map[diag_key] = diagram

        # Process and save diagrams
        for key, diagram in diagram_map.items():
            diagram_id = self._generate_id(f"{template_name}_{key[:30]}")
            diagram['id'] = diagram_id

            category = diagram.get('category', 'custom')
            save_dir = self.dirs.get(category, self.dirs['custom'])

            # Save diagram file
            filename = f"{diagram_id}.json"
            filepath = save_dir / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(diagram, f, indent=2)

            diagram['filename'] = filename
            results['diagrams'].append(diagram)
            results['by_category'][category].append(diagram)

            # Add to index
            existing = [x for x in self.index['templates'] if x.get('id') == diagram_id]
            if not existing:
                self.index['templates'].append(diagram)
                self.index['categories'][category].append(diagram_id)

        self._save_index()

        logger.info(f"Extracted {len(results['diagrams'])} diagram templates")
        return results

    def _get_shape_type(self, shape) -> str:
        """Get shape type as string."""
        try:
            if hasattr(shape, 'shape_type'):
                return str(shape.shape_type).split('.')[-1]
        except:
            pass
        return 'UNKNOWN'

    def _extract_group_diagram(self, group_shape, slide_num: int, template_name: str,
                               slide_width: float, slide_height: float) -> Optional[dict]:
        """Extract diagram from a grouped shape."""
        try:
            shapes = []
            connectors = []
            text_elements = []

            def process_shapes(shape_collection, parent_left=0, parent_top=0):
                for shape in shape_collection:
                    shape_type = self._get_shape_type(shape)

                    # Calculate absolute position
                    try:
                        left = (shape.left.inches + parent_left) / slide_width
                        top = (shape.top.inches + parent_top) / slide_height
                        width = shape.width.inches / slide_width
                        height = shape.height.inches / slide_height
                    except:
                        continue

                    shape_info = {
                        'type': shape_type,
                        'auto_shape_type': self._get_auto_shape_type(shape),
                        'position': {
                            'left': round(left, 4),
                            'top': round(top, 4),
                            'width': round(width, 4),
                            'height': round(height, 4),
                        },
                        'has_text': False,
                        'text': '',
                    }

                    # Extract text
                    if hasattr(shape, 'text') and shape.text:
                        shape_info['has_text'] = True
                        shape_info['text'] = shape.text[:100]
                        text_elements.append(shape_info['text'])

                    # Extract colors
                    shape_info['colors'] = self._extract_shape_colors(shape)

                    # Categorize shape
                    if self._is_connector(shape, shape_type):
                        connectors.append(shape_info)
                    else:
                        shapes.append(shape_info)

                    # Recurse into nested groups
                    if shape_type == 'GROUP' and hasattr(shape, 'shapes'):
                        try:
                            process_shapes(shape.shapes, shape.left.inches, shape.top.inches)
                        except:
                            pass

            process_shapes(group_shape.shapes)

            if len(shapes) < 2:
                return None

            # Categorize the diagram
            category = self._categorize_diagram(shapes, connectors, text_elements)

            # Calculate overall bounds
            all_shapes = shapes + connectors
            bounds = self._calculate_bounds(all_shapes)

            diagram = {
                'template': template_name,
                'slide_num': slide_num,
                'category': category,
                'shape_count': len(shapes),
                'connector_count': len(connectors),
                'shapes': shapes,
                'connectors': connectors,
                'text_placeholders': text_elements[:10],  # First 10 text elements
                'bounds': bounds,
                'usage_count': 1,
            }

            return diagram

        except Exception as e:
            logger.warning(f"Failed to extract group diagram: {e}")
            return None

    def _detect_ungrouped_diagrams(self, slide, slide_num: int, template_name: str,
                                    slide_width: float, slide_height: float) -> List[dict]:
        """Detect diagrams from ungrouped shape arrangements."""
        diagrams = []

        # Collect all auto shapes
        auto_shapes = []
        for shape in slide.shapes:
            shape_type = self._get_shape_type(shape)
            if shape_type in ('AUTO_SHAPE', 'FREEFORM'):
                try:
                    left = shape.left.inches / slide_width
                    top = shape.top.inches / slide_height
                    width = shape.width.inches / slide_width
                    height = shape.height.inches / slide_height

                    shape_info = {
                        'type': shape_type,
                        'auto_shape_type': self._get_auto_shape_type(shape),
                        'position': {
                            'left': round(left, 4),
                            'top': round(top, 4),
                            'width': round(width, 4),
                            'height': round(height, 4),
                        },
                        'has_text': hasattr(shape, 'text') and bool(shape.text),
                        'text': shape.text[:100] if hasattr(shape, 'text') and shape.text else '',
                        'colors': self._extract_shape_colors(shape),
                    }
                    auto_shapes.append(shape_info)
                except:
                    continue

        # Detect patterns
        if len(auto_shapes) >= 3:
            # Check for horizontal alignment (process flow)
            horizontal_groups = self._find_aligned_shapes(auto_shapes, 'horizontal')
            for group in horizontal_groups:
                if len(group) >= 3:
                    diagram = {
                        'template': template_name,
                        'slide_num': slide_num,
                        'category': 'process_flows',
                        'shape_count': len(group),
                        'connector_count': 0,
                        'shapes': group,
                        'connectors': [],
                        'text_placeholders': [s['text'] for s in group if s['has_text']],
                        'bounds': self._calculate_bounds(group),
                        'usage_count': 1,
                    }
                    diagrams.append(diagram)

            # Check for grid arrangement (matrix)
            if len(auto_shapes) >= 4:
                grid = self._detect_grid_arrangement(auto_shapes)
                if grid:
                    diagram = {
                        'template': template_name,
                        'slide_num': slide_num,
                        'category': 'matrices',
                        'shape_count': len(grid['shapes']),
                        'connector_count': 0,
                        'shapes': grid['shapes'],
                        'connectors': [],
                        'grid_dimensions': grid['dimensions'],
                        'text_placeholders': [s['text'] for s in grid['shapes'] if s['has_text']],
                        'bounds': self._calculate_bounds(grid['shapes']),
                        'usage_count': 1,
                    }
                    diagrams.append(diagram)

        return diagrams

    def _get_auto_shape_type(self, shape) -> Optional[str]:
        """Get auto shape type name."""
        try:
            if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type:
                return str(shape.auto_shape_type).split('.')[-1]
        except:
            pass
        return None

    def _is_connector(self, shape, shape_type: str) -> bool:
        """Determine if shape is a connector/arrow."""
        if shape_type == 'CONNECTOR':
            return True

        auto_shape = self._get_auto_shape_type(shape)
        if auto_shape:
            auto_lower = auto_shape.lower()
            if any(x in auto_lower for x in ['arrow', 'connector', 'line']):
                return True

        return False

    def _extract_shape_colors(self, shape) -> List[str]:
        """Extract colors from a shape."""
        colors = []
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                    try:
                        colors.append(f"#{shape.fill.fore_color.rgb}")
                    except:
                        pass

            if hasattr(shape, 'line') and shape.line.fill.type is not None:
                try:
                    colors.append(f"#{shape.line.color.rgb}")
                except:
                    pass
        except:
            pass

        return colors

    def _categorize_diagram(self, shapes: List[dict], connectors: List[dict],
                           text_elements: List[str]) -> str:
        """Categorize the diagram type."""
        shape_count = len(shapes)
        connector_count = len(connectors)

        # Get auto shape types
        auto_types = [s.get('auto_shape_type', '') for s in shapes if s.get('auto_shape_type')]
        auto_types_str = ' '.join(auto_types).lower()

        # Check for cycle/circular arrangement
        if 'chevron' in auto_types_str or 'circular' in auto_types_str:
            return 'cycles'

        # Check for timeline
        if 'callout' in auto_types_str or connector_count >= shape_count - 1:
            positions = [s['position'] for s in shapes]
            if self._is_linear_arrangement(positions):
                return 'timelines'

        # Check for process flow (horizontal sequence with connectors)
        if connector_count >= shape_count / 2:
            return 'process_flows'

        # Check for hierarchy
        if self._is_hierarchical(shapes):
            return 'hierarchies'

        # Check for matrix
        if self._is_matrix(shapes):
            return 'matrices'

        # Check for process flow without connectors
        positions = [s['position'] for s in shapes]
        if self._is_linear_arrangement(positions):
            return 'process_flows'

        return 'custom'

    def _is_linear_arrangement(self, positions: List[dict]) -> bool:
        """Check if shapes are arranged linearly."""
        if len(positions) < 3:
            return False

        # Check horizontal alignment
        tops = [p['top'] for p in positions]
        top_variance = max(tops) - min(tops)
        if top_variance < 0.1:  # Within 10% variance
            return True

        # Check vertical alignment
        lefts = [p['left'] for p in positions]
        left_variance = max(lefts) - min(lefts)
        if left_variance < 0.1:
            return True

        return False

    def _is_hierarchical(self, shapes: List[dict]) -> bool:
        """Check if shapes form a hierarchy."""
        if len(shapes) < 3:
            return False

        # Group by vertical position (rows)
        rows = defaultdict(list)
        for s in shapes:
            row = round(s['position']['top'], 1)
            rows[row].append(s)

        # Hierarchy should have multiple rows with varying counts
        row_counts = [len(r) for r in rows.values()]
        if len(row_counts) >= 2:
            # Typical hierarchy: 1 at top, more below
            sorted_counts = sorted(row_counts)
            if sorted_counts[0] <= 2 and sorted_counts[-1] >= 2:
                return True

        return False

    def _is_matrix(self, shapes: List[dict]) -> bool:
        """Check if shapes form a matrix/grid."""
        if len(shapes) < 4:
            return False

        # Group by rows and columns
        rows = defaultdict(list)
        cols = defaultdict(list)

        for s in shapes:
            row = round(s['position']['top'], 1)
            col = round(s['position']['left'], 1)
            rows[row].append(s)
            cols[col].append(s)

        # Matrix should have multiple rows and columns
        if len(rows) >= 2 and len(cols) >= 2:
            # Check if relatively uniform
            row_counts = [len(r) for r in rows.values()]
            if max(row_counts) - min(row_counts) <= 1:
                return True

        return False

    def _find_aligned_shapes(self, shapes: List[dict], direction: str) -> List[List[dict]]:
        """Find groups of aligned shapes."""
        groups = []

        if direction == 'horizontal':
            # Group by similar top positions
            by_top = defaultdict(list)
            for s in shapes:
                top_key = round(s['position']['top'], 1)
                by_top[top_key].append(s)

            for top, group in by_top.items():
                if len(group) >= 3:
                    # Sort by left position
                    group.sort(key=lambda x: x['position']['left'])
                    groups.append(group)

        elif direction == 'vertical':
            # Group by similar left positions
            by_left = defaultdict(list)
            for s in shapes:
                left_key = round(s['position']['left'], 1)
                by_left[left_key].append(s)

            for left, group in by_left.items():
                if len(group) >= 3:
                    group.sort(key=lambda x: x['position']['top'])
                    groups.append(group)

        return groups

    def _detect_grid_arrangement(self, shapes: List[dict]) -> Optional[dict]:
        """Detect grid arrangement in shapes."""
        if len(shapes) < 4:
            return None

        # Group by rows
        rows = defaultdict(list)
        for s in shapes:
            row_key = round(s['position']['top'], 1)
            rows[row_key].append(s)

        if len(rows) < 2:
            return None

        # Check if rows have similar counts
        row_counts = [len(r) for r in rows.values()]
        if max(row_counts) - min(row_counts) > 1:
            return None

        # Collect grid shapes
        grid_shapes = []
        for row in rows.values():
            row.sort(key=lambda x: x['position']['left'])
            grid_shapes.extend(row)

        return {
            'shapes': grid_shapes,
            'dimensions': {
                'rows': len(rows),
                'cols': max(row_counts),
            }
        }

    def _calculate_bounds(self, shapes: List[dict]) -> dict:
        """Calculate bounding box for shapes."""
        if not shapes:
            return {'left': 0, 'top': 0, 'right': 0, 'bottom': 0, 'width': 0, 'height': 0}

        min_left = min(s['position']['left'] for s in shapes)
        min_top = min(s['position']['top'] for s in shapes)
        max_right = max(s['position']['left'] + s['position']['width'] for s in shapes)
        max_bottom = max(s['position']['top'] + s['position']['height'] for s in shapes)

        return {
            'left': round(min_left, 4),
            'top': round(min_top, 4),
            'right': round(max_right, 4),
            'bottom': round(max_bottom, 4),
            'width': round(max_right - min_left, 4),
            'height': round(max_bottom - min_top, 4),
        }

    def _create_diagram_key(self, diagram: dict) -> str:
        """Create a key for diagram deduplication."""
        shapes = diagram.get('shapes', [])
        shape_types = sorted([s.get('auto_shape_type', 'unknown') for s in shapes])
        return f"{diagram['category']}_{len(shapes)}_{','.join(shape_types[:5])}"

    def get_summary(self) -> dict:
        """Get summary of extracted diagrams."""
        summary = {
            'total_templates': len(self.index['templates']),
            'by_category': {},
        }

        for category, ids in self.index['categories'].items():
            summary['by_category'][category] = len(ids)

        return summary

    def search(self, category: str = None, template: str = None,
               min_shapes: int = None, max_shapes: int = None) -> List[dict]:
        """Search diagram templates."""
        results = []

        for diagram in self.index['templates']:
            if category and diagram.get('category') != category:
                continue
            if template and diagram.get('template') != template:
                continue
            if min_shapes and diagram.get('shape_count', 0) < min_shapes:
                continue
            if max_shapes and diagram.get('shape_count', 0) > max_shapes:
                continue
            results.append(diagram)

        return results

    def get_by_id(self, diagram_id: str) -> Optional[dict]:
        """Get a specific diagram template by ID."""
        for diagram in self.index['templates']:
            if diagram.get('id') == diagram_id:
                return diagram
        return None


def extract_diagrams_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract diagrams from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = DiagramTemplateExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting diagrams from: {pptx_path.name}")
        try:
            extractor.extract_template_diagrams(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python diagram_template_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_diagrams_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("DIAGRAM EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
