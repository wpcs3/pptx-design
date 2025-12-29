"""
Style Extractor - Extract color palettes, typography, and effect presets from PowerPoint templates.

This module extracts:
- Color palettes (theme colors, custom colors used)
- Typography presets (font combinations, heading/body styles)
- Effect presets (shadows, gradients, glows, 3D effects)
"""

import json
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any
from collections import defaultdict
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

logger = logging.getLogger(__name__)


class StyleExtractor:
    """Extract styles (colors, typography, effects) from PowerPoint templates."""

    # XML namespaces used in PowerPoint
    NAMESPACES = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    def __init__(self, output_dir: Path):
        """Initialize the style extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories for each style type
        self.dirs = {
            'color_palettes': self.output_dir / 'styles' / 'color_palettes',
            'typography': self.output_dir / 'styles' / 'typography',
            'effects': self.output_dir / 'styles' / 'effects',
            'gradients': self.output_dir / 'styles' / 'gradients',
            'shadows': self.output_dir / 'styles' / 'shadows',
        }
        for d in self.dirs.values():
            d.mkdir(parents=True, exist_ok=True)

        # Style index
        self.index_path = self.output_dir / 'styles' / 'style_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'color_palettes': [],
            'typography_presets': [],
            'effect_presets': [],
            'gradient_presets': [],
            'shadow_presets': [],
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def _rgb_to_hex(self, rgb) -> Optional[str]:
        """Convert RGB color to hex string."""
        try:
            if rgb is None:
                return None
            if isinstance(rgb, str):
                return rgb if rgb.startswith('#') else f'#{rgb}'
            if hasattr(rgb, 'rgb'):
                rgb = rgb.rgb
            if isinstance(rgb, RGBColor):
                return f'#{rgb}'
            return f'#{rgb:06X}' if isinstance(rgb, int) else str(rgb)
        except:
            return None

    def extract_template_styles(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract all styles from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted styles
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting styles from: {template_name}")

        prs = Presentation(pptx_path)

        results = {
            'template': template_name,
            'color_palettes': [],
            'typography_presets': [],
            'effect_presets': [],
            'gradient_presets': [],
            'shadow_presets': [],
        }

        # Extract theme colors
        color_palette = self._extract_color_palette(prs, template_name)
        if color_palette:
            results['color_palettes'].append(color_palette)

        # Extract typography from slides
        typography = self._extract_typography_presets(prs, template_name)
        results['typography_presets'].extend(typography)

        # Extract effects from slides
        effects = self._extract_effect_presets(prs, template_name)
        results['effect_presets'].extend(effects.get('effects', []))
        results['gradient_presets'].extend(effects.get('gradients', []))
        results['shadow_presets'].extend(effects.get('shadows', []))

        # Save to index
        self._merge_to_index(results)
        self._save_index()

        logger.info(f"Extracted: {len(results['color_palettes'])} palettes, "
                   f"{len(results['typography_presets'])} typography presets, "
                   f"{len(results['effect_presets'])} effect presets")

        return results

    def _extract_color_palette(self, prs: Presentation, template_name: str) -> Optional[dict]:
        """Extract color palette from theme."""
        try:
            # Get theme colors from the presentation
            theme_colors = {}
            custom_colors = set()

            # Try to access theme directly
            try:
                theme_part = prs.part.part_related_by('http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme')
                if theme_part:
                    theme_xml = theme_part.blob
                    root = ET.fromstring(theme_xml)

                    # Find color scheme
                    clr_scheme = root.find('.//a:clrScheme', self.NAMESPACES)
                    if clr_scheme:
                        for color_elem in clr_scheme:
                            color_name = color_elem.tag.split('}')[-1]

                            # Try different color specifications
                            srgb = color_elem.find('.//a:srgbClr', self.NAMESPACES)
                            sys_clr = color_elem.find('.//a:sysClr', self.NAMESPACES)

                            if srgb is not None:
                                hex_val = srgb.get('val', '')
                                theme_colors[color_name] = f'#{hex_val}'
                            elif sys_clr is not None:
                                last_clr = sys_clr.get('lastClr', '')
                                theme_colors[color_name] = f'#{last_clr}'
            except Exception as e:
                logger.debug(f"Could not extract theme XML: {e}")

            # Extract colors from shapes
            for slide in prs.slides:
                for shape in slide.shapes:
                    colors = self._extract_colors_from_shape(shape)
                    custom_colors.update(colors)

            if not theme_colors and not custom_colors:
                return None

            # Create palette
            palette_id = self._generate_id(f"{template_name}_colors")

            palette = {
                'id': palette_id,
                'name': f'{template_name}_palette',
                'template': template_name,
                'theme_colors': theme_colors,
                'custom_colors': list(custom_colors)[:20],  # Top 20 custom colors
                'primary': theme_colors.get('dk1', theme_colors.get('tx1', '#000000')),
                'secondary': theme_colors.get('dk2', theme_colors.get('tx2', '#444444')),
                'accent1': theme_colors.get('accent1', '#4472C4'),
                'accent2': theme_colors.get('accent2', '#ED7D31'),
                'accent3': theme_colors.get('accent3', '#A5A5A5'),
                'accent4': theme_colors.get('accent4', '#FFC000'),
                'accent5': theme_colors.get('accent5', '#5B9BD5'),
                'accent6': theme_colors.get('accent6', '#70AD47'),
                'background': theme_colors.get('lt1', '#FFFFFF'),
                'background_alt': theme_colors.get('lt2', '#E7E6E6'),
            }

            # Save palette file
            filename = f"{palette_id}.json"
            filepath = self.dirs['color_palettes'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(palette, f, indent=2)

            palette['filename'] = filename
            return palette

        except Exception as e:
            logger.warning(f"Failed to extract color palette: {e}")
            return None

    def _extract_colors_from_shape(self, shape) -> set:
        """Extract all colors used in a shape."""
        colors = set()
        try:
            # Fill color
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                    try:
                        hex_color = self._rgb_to_hex(shape.fill.fore_color.rgb)
                        if hex_color:
                            colors.add(hex_color)
                    except:
                        pass

            # Line color
            if hasattr(shape, 'line') and shape.line.fill.type is not None:
                try:
                    hex_color = self._rgb_to_hex(shape.line.color.rgb)
                    if hex_color:
                        colors.add(hex_color)
                except:
                    pass

            # Text color
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.color and run.font.color.rgb:
                            try:
                                hex_color = self._rgb_to_hex(run.font.color.rgb)
                                if hex_color:
                                    colors.add(hex_color)
                            except:
                                pass

        except Exception as e:
            logger.debug(f"Error extracting colors from shape: {e}")

        return colors

    def _extract_typography_presets(self, prs: Presentation, template_name: str) -> List[dict]:
        """Extract typography presets from slides."""
        typography_map = {}  # Use dict to deduplicate

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue

                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if not paragraph.runs:
                        continue

                    run = paragraph.runs[0]

                    # Extract font properties
                    font_name = None
                    font_size = None
                    bold = False
                    italic = False
                    color = None

                    try:
                        if run.font.name:
                            font_name = run.font.name
                        if run.font.size:
                            font_size = run.font.size.pt
                        bold = bool(run.font.bold)
                        italic = bool(run.font.italic)
                        if run.font.color and run.font.color.rgb:
                            color = self._rgb_to_hex(run.font.color.rgb)
                    except:
                        continue

                    if not font_name or not font_size:
                        continue

                    # Create typography key for deduplication
                    typo_key = f"{font_name}_{font_size}_{bold}_{italic}"

                    if typo_key not in typography_map:
                        # Determine preset type based on size
                        preset_type = 'body'
                        if font_size >= 28:
                            preset_type = 'title'
                        elif font_size >= 20:
                            preset_type = 'heading'
                        elif font_size >= 14:
                            preset_type = 'subheading'
                        elif font_size <= 10:
                            preset_type = 'caption'

                        # Extract paragraph formatting
                        line_spacing = 1.0
                        alignment = 'left'
                        try:
                            if paragraph.line_spacing:
                                line_spacing = paragraph.line_spacing
                            if paragraph.alignment:
                                alignment = str(paragraph.alignment).split('.')[-1].lower()
                        except:
                            pass

                        preset = {
                            'font_family': font_name,
                            'font_size_pt': font_size,
                            'bold': bold,
                            'italic': italic,
                            'color': color,
                            'line_spacing': line_spacing,
                            'alignment': alignment,
                            'preset_type': preset_type,
                            'template': template_name,
                            'usage_count': 1,
                        }
                        typography_map[typo_key] = preset
                    else:
                        typography_map[typo_key]['usage_count'] += 1

        # Convert to list and add IDs
        presets = []
        for key, preset in typography_map.items():
            preset_id = self._generate_id(f"{template_name}_{key}")
            preset['id'] = preset_id
            preset['preset_name'] = f"{preset['preset_type']}_{preset['font_family'].replace(' ', '_')}_{int(preset['font_size_pt'])}"

            # Save preset file
            filename = f"{preset_id}.json"
            filepath = self.dirs['typography'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(preset, f, indent=2)

            preset['filename'] = filename
            presets.append(preset)

        # Sort by usage count (most common first)
        presets.sort(key=lambda x: x['usage_count'], reverse=True)

        return presets[:20]  # Return top 20 presets

    def _extract_effect_presets(self, prs: Presentation, template_name: str) -> dict:
        """Extract effect presets (shadows, gradients, 3D effects)."""
        effects = []
        gradients = []
        shadows = []

        gradient_map = {}
        shadow_map = {}

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Extract gradient fills
                gradient = self._extract_gradient(shape, template_name)
                if gradient:
                    grad_key = json.dumps(gradient.get('stops', []), sort_keys=True)
                    if grad_key not in gradient_map:
                        gradient_map[grad_key] = gradient

                # Extract shadow effects
                shadow = self._extract_shadow(shape, template_name)
                if shadow:
                    shadow_key = f"{shadow.get('type')}_{shadow.get('blur_radius')}_{shadow.get('distance')}"
                    if shadow_key not in shadow_map:
                        shadow_map[shadow_key] = shadow

                # Extract 3D effects
                effect_3d = self._extract_3d_effect(shape, template_name)
                if effect_3d:
                    effects.append(effect_3d)

                # Extract glow effects
                glow = self._extract_glow_effect(shape, template_name)
                if glow:
                    effects.append(glow)

        # Convert gradient map to list with IDs
        for key, gradient in gradient_map.items():
            gradient_id = self._generate_id(f"grad_{template_name}_{key[:20]}")
            gradient['id'] = gradient_id

            filename = f"{gradient_id}.json"
            filepath = self.dirs['gradients'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(gradient, f, indent=2)

            gradient['filename'] = filename
            gradients.append(gradient)

        # Convert shadow map to list with IDs
        for key, shadow in shadow_map.items():
            shadow_id = self._generate_id(f"shadow_{template_name}_{key}")
            shadow['id'] = shadow_id

            filename = f"{shadow_id}.json"
            filepath = self.dirs['shadows'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(shadow, f, indent=2)

            shadow['filename'] = filename
            shadows.append(shadow)

        # Add IDs to general effects and save
        for i, effect in enumerate(effects):
            effect_id = self._generate_id(f"effect_{template_name}_{i}")
            effect['id'] = effect_id

            filename = f"{effect_id}.json"
            filepath = self.dirs['effects'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(effect, f, indent=2)

            effect['filename'] = filename

        return {
            'effects': effects[:20],
            'gradients': gradients[:20],
            'shadows': shadows[:20],
        }

    def _extract_gradient(self, shape, template_name: str) -> Optional[dict]:
        """Extract gradient fill from a shape."""
        try:
            if not hasattr(shape, 'fill'):
                return None

            fill = shape.fill
            if fill.type is None:
                return None

            # Check if gradient fill
            fill_type_str = str(fill.type)
            if 'GRADIENT' not in fill_type_str:
                return None

            gradient = {
                'type': 'linear',  # Default type
                'angle': 0,
                'stops': [],
                'template': template_name,
            }

            # Try to extract gradient stops from XML
            try:
                if hasattr(shape, '_element'):
                    spPr = shape._element.find('.//a:spPr', self.NAMESPACES)
                    if spPr is not None:
                        gradFill = spPr.find('.//a:gradFill', self.NAMESPACES)
                        if gradFill is not None:
                            # Get rotation angle
                            lin = gradFill.find('.//a:lin', self.NAMESPACES)
                            if lin is not None:
                                angle = lin.get('ang', '0')
                                gradient['angle'] = int(angle) / 60000  # Convert from EMUs

                            # Get color stops
                            gsLst = gradFill.find('.//a:gsLst', self.NAMESPACES)
                            if gsLst is not None:
                                for gs in gsLst.findall('.//a:gs', self.NAMESPACES):
                                    position = int(gs.get('pos', '0')) / 1000

                                    srgb = gs.find('.//a:srgbClr', self.NAMESPACES)
                                    if srgb is not None:
                                        color = f"#{srgb.get('val', '000000')}"
                                        gradient['stops'].append({
                                            'position': position,
                                            'color': color,
                                        })
            except:
                pass

            if gradient['stops']:
                return gradient
            return None

        except Exception as e:
            logger.debug(f"Error extracting gradient: {e}")
            return None

    def _extract_shadow(self, shape, template_name: str) -> Optional[dict]:
        """Extract shadow effect from a shape."""
        try:
            if not hasattr(shape, '_element'):
                return None

            # Look for effect list in shape properties
            spPr = shape._element.find('.//a:spPr', self.NAMESPACES)
            if spPr is None:
                return None

            effectLst = spPr.find('.//a:effectLst', self.NAMESPACES)
            if effectLst is None:
                return None

            # Check for outer shadow
            outerShdw = effectLst.find('.//a:outerShdw', self.NAMESPACES)
            if outerShdw is not None:
                shadow = {
                    'type': 'outer',
                    'blur_radius': int(outerShdw.get('blurRad', '0')) / 12700,  # EMU to pt
                    'distance': int(outerShdw.get('dist', '0')) / 12700,
                    'direction': int(outerShdw.get('dir', '0')) / 60000,  # EMU to degrees
                    'transparency': 1 - int(outerShdw.get('algn', '0')) / 100000,
                    'template': template_name,
                }

                # Get shadow color
                srgb = outerShdw.find('.//a:srgbClr', self.NAMESPACES)
                if srgb is not None:
                    shadow['color'] = f"#{srgb.get('val', '000000')}"

                    # Check for alpha
                    alpha = srgb.find('.//a:alpha', self.NAMESPACES)
                    if alpha is not None:
                        shadow['transparency'] = 1 - int(alpha.get('val', '100000')) / 100000

                return shadow

            # Check for inner shadow
            innerShdw = effectLst.find('.//a:innerShdw', self.NAMESPACES)
            if innerShdw is not None:
                return {
                    'type': 'inner',
                    'blur_radius': int(innerShdw.get('blurRad', '0')) / 12700,
                    'distance': int(innerShdw.get('dist', '0')) / 12700,
                    'direction': int(innerShdw.get('dir', '0')) / 60000,
                    'template': template_name,
                }

            return None

        except Exception as e:
            logger.debug(f"Error extracting shadow: {e}")
            return None

    def _extract_3d_effect(self, shape, template_name: str) -> Optional[dict]:
        """Extract 3D effect from a shape."""
        try:
            if not hasattr(shape, '_element'):
                return None

            spPr = shape._element.find('.//a:spPr', self.NAMESPACES)
            if spPr is None:
                return None

            # Check for scene 3D
            scene3d = spPr.find('.//a:scene3d', self.NAMESPACES)
            sp3d = spPr.find('.//a:sp3d', self.NAMESPACES)

            if scene3d is None and sp3d is None:
                return None

            effect = {
                'type': '3d',
                'template': template_name,
            }

            # Extract camera settings
            if scene3d is not None:
                camera = scene3d.find('.//a:camera', self.NAMESPACES)
                if camera is not None:
                    effect['camera_preset'] = camera.get('prst', 'orthographicFront')

                    rot = camera.find('.//a:rot', self.NAMESPACES)
                    if rot is not None:
                        effect['rotation'] = {
                            'lat': int(rot.get('lat', '0')) / 60000,
                            'lon': int(rot.get('lon', '0')) / 60000,
                            'rev': int(rot.get('rev', '0')) / 60000,
                        }

            # Extract 3D shape settings
            if sp3d is not None:
                effect['extrusion_height'] = int(sp3d.get('extrusionH', '0')) / 12700
                effect['contour_width'] = int(sp3d.get('contourW', '0')) / 12700

                bevelT = sp3d.find('.//a:bevelT', self.NAMESPACES)
                if bevelT is not None:
                    effect['bevel_top'] = {
                        'type': bevelT.get('prst', 'circle'),
                        'width': int(bevelT.get('w', '0')) / 12700,
                        'height': int(bevelT.get('h', '0')) / 12700,
                    }

            return effect if len(effect) > 2 else None

        except Exception as e:
            logger.debug(f"Error extracting 3D effect: {e}")
            return None

    def _extract_glow_effect(self, shape, template_name: str) -> Optional[dict]:
        """Extract glow effect from a shape."""
        try:
            if not hasattr(shape, '_element'):
                return None

            spPr = shape._element.find('.//a:spPr', self.NAMESPACES)
            if spPr is None:
                return None

            effectLst = spPr.find('.//a:effectLst', self.NAMESPACES)
            if effectLst is None:
                return None

            glow = effectLst.find('.//a:glow', self.NAMESPACES)
            if glow is None:
                return None

            effect = {
                'type': 'glow',
                'radius': int(glow.get('rad', '0')) / 12700,
                'template': template_name,
            }

            srgb = glow.find('.//a:srgbClr', self.NAMESPACES)
            if srgb is not None:
                effect['color'] = f"#{srgb.get('val', '000000')}"

                alpha = srgb.find('.//a:alpha', self.NAMESPACES)
                if alpha is not None:
                    effect['transparency'] = 1 - int(alpha.get('val', '100000')) / 100000

            return effect

        except Exception as e:
            logger.debug(f"Error extracting glow: {e}")
            return None

    def _merge_to_index(self, results: dict):
        """Merge extracted results to the main index."""
        for key in ['color_palettes', 'typography_presets', 'effect_presets',
                    'gradient_presets', 'shadow_presets']:
            for item in results.get(key, []):
                # Check if already exists
                existing = [x for x in self.index[key] if x.get('id') == item.get('id')]
                if not existing:
                    self.index[key].append(item)

    def get_summary(self) -> dict:
        """Get summary of extracted styles."""
        return {
            'color_palettes': len(self.index['color_palettes']),
            'typography_presets': len(self.index['typography_presets']),
            'effect_presets': len(self.index['effect_presets']),
            'gradient_presets': len(self.index['gradient_presets']),
            'shadow_presets': len(self.index['shadow_presets']),
        }

    def search_colors(self, query: str = None, template: str = None) -> List[dict]:
        """Search color palettes."""
        results = []
        for palette in self.index['color_palettes']:
            if template and palette.get('template') != template:
                continue
            if query:
                # Search in color values
                all_colors = list(palette.get('theme_colors', {}).values())
                all_colors.extend(palette.get('custom_colors', []))
                if not any(query.lower() in c.lower() for c in all_colors if c):
                    continue
            results.append(palette)
        return results

    def search_typography(self, font: str = None, preset_type: str = None,
                          template: str = None) -> List[dict]:
        """Search typography presets."""
        results = []
        for preset in self.index['typography_presets']:
            if template and preset.get('template') != template:
                continue
            if font and font.lower() not in preset.get('font_family', '').lower():
                continue
            if preset_type and preset.get('preset_type') != preset_type:
                continue
            results.append(preset)
        return results

    def search_effects(self, effect_type: str = None, template: str = None) -> List[dict]:
        """Search effect presets."""
        results = []

        # Search in all effect categories
        for category in ['effect_presets', 'gradient_presets', 'shadow_presets']:
            for effect in self.index[category]:
                if template and effect.get('template') != template:
                    continue
                if effect_type and effect.get('type') != effect_type:
                    continue
                results.append(effect)

        return results


def extract_styles_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract styles from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = StyleExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting styles from: {pptx_path.name}")
        try:
            extractor.extract_template_styles(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python style_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_styles_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("STYLE EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
