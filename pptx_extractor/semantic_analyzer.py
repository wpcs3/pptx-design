"""
LayoutLMv3 Semantic Slide Analyzer

Deep document understanding for PowerPoint slides using Microsoft's
LayoutLMv3 multimodal transformer model. Combines text, layout, and
visual features for semantic analysis.

Phase 4 Enhancement (2025-12-29):
- LayoutLMv3 integration for multimodal understanding
- Semantic element labeling (header, body, caption, etc.)
- Content purpose classification
- Template matching recommendations

Usage:
    from pptx_extractor.semantic_analyzer import SemanticSlideAnalyzer

    analyzer = SemanticSlideAnalyzer()
    result = analyzer.analyze("slide_image.png", text_boxes=[
        {"text": "Q4 Review", "bbox": [100, 50, 800, 150]},
        {"text": "Revenue increased 25%", "bbox": [100, 200, 700, 250]},
    ])
    print(result.semantic_labels)
    print(result.content_purpose)

Requirements:
    pip install transformers torch torchvision
    pip install pillow

Model Information:
    - Base: microsoft/layoutlmv3-base (125M params)
    - Large: microsoft/layoutlmv3-large (368M params)
    - Pre-trained on: IIT-CDIP, DocVQA, RVL-CDIP

Note:
    This module provides graceful fallbacks when dependencies are unavailable.
    Basic analysis works without deep learning models.
"""

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# Try to import deep learning dependencies
TRANSFORMERS_AVAILABLE = False
TORCH_AVAILABLE = False

try:
    import torch
    TORCH_AVAILABLE = True
    logger.info("PyTorch available")

    try:
        from transformers import (
            AutoProcessor,
            AutoModelForTokenClassification,
            AutoModelForSequenceClassification,
            LayoutLMv3Processor,
            LayoutLMv3ForTokenClassification,
        )
        TRANSFORMERS_AVAILABLE = True
        logger.info("Transformers library available")
    except ImportError:
        logger.info("Transformers library not installed")
except ImportError:
    logger.info("PyTorch not installed")


@dataclass
class TextBox:
    """A text box with position information."""
    text: str
    bbox: Tuple[int, int, int, int]  # x1, y1, x2, y2
    confidence: float = 1.0

    @property
    def width(self) -> int:
        return self.bbox[2] - self.bbox[0]

    @property
    def height(self) -> int:
        return self.bbox[3] - self.bbox[1]

    @property
    def area(self) -> int:
        return self.width * self.height

    @property
    def center(self) -> Tuple[float, float]:
        return (
            (self.bbox[0] + self.bbox[2]) / 2,
            (self.bbox[1] + self.bbox[3]) / 2
        )


@dataclass
class SemanticLabel:
    """Semantic label for a text element."""
    text: str
    label: str  # title, subtitle, header, body, bullet, caption, footer
    confidence: float
    bbox: Tuple[int, int, int, int]
    purpose: str = ""  # main_content, supporting, navigation, branding


@dataclass
class SemanticAnalysisResult:
    """Result of semantic slide analysis."""
    semantic_labels: List[SemanticLabel]
    content_purpose: str  # informative, persuasive, data_presentation, etc.
    suggested_template: str
    key_elements: Dict[str, str]  # title, main_point, etc.
    confidence: float
    analysis_method: str  # layoutlmv3, heuristic, basic

    @property
    def title(self) -> Optional[str]:
        """Extract the main title from labels."""
        titles = [l for l in self.semantic_labels if l.label == "title"]
        return titles[0].text if titles else None

    @property
    def body_text(self) -> List[str]:
        """Extract body text from labels."""
        return [l.text for l in self.semantic_labels if l.label in ("body", "bullet")]


# Semantic label definitions
SLIDE_ELEMENT_LABELS = {
    0: "title",
    1: "subtitle",
    2: "header",
    3: "body",
    4: "bullet",
    5: "caption",
    6: "footer",
    7: "other"
}

# Content purpose categories
CONTENT_PURPOSES = [
    "informative",      # General information sharing
    "persuasive",       # Sales, pitch, convincing
    "data_presentation", # Charts, metrics, analysis
    "instructional",    # How-to, process explanation
    "comparative",      # Comparison, pros/cons
    "summary",          # Executive summary, key takeaways
    "agenda",           # Meeting agenda, outline
    "introduction",     # Opening, title slides
    "conclusion"        # Closing, thank you slides
]


class SemanticSlideAnalyzer:
    """
    Deep semantic analyzer for slide content using LayoutLMv3.

    Provides multimodal understanding combining:
    - Text content and meaning
    - Visual layout and positioning
    - Document structure patterns

    Capabilities:
    - Label text elements (title, body, bullet, etc.)
    - Classify content purpose
    - Suggest matching templates
    - Extract key information
    """

    def __init__(
        self,
        model_name: str = "microsoft/layoutlmv3-base",
        use_gpu: bool = False,
        num_labels: int = 8
    ):
        """
        Initialize the semantic analyzer.

        Args:
            model_name: HuggingFace model identifier.
            use_gpu: Whether to use GPU for inference.
            num_labels: Number of semantic labels for classification.
        """
        self.model_name = model_name
        self.use_gpu = use_gpu
        self.num_labels = num_labels
        self.device = "cuda" if use_gpu and TORCH_AVAILABLE and torch.cuda.is_available() else "cpu"

        self.processor = None
        self.model = None

        if TRANSFORMERS_AVAILABLE and TORCH_AVAILABLE:
            self._load_model()

    def _load_model(self):
        """Load the LayoutLMv3 model and processor."""
        try:
            logger.info(f"Loading LayoutLMv3 model: {self.model_name}")

            self.processor = AutoProcessor.from_pretrained(
                self.model_name,
                apply_ocr=False  # We provide our own text boxes
            )

            # For token classification (labeling text elements)
            self.model = AutoModelForTokenClassification.from_pretrained(
                self.model_name,
                num_labels=self.num_labels
            )
            self.model.to(self.device)
            self.model.eval()

            logger.info(f"LayoutLMv3 model loaded on {self.device}")

        except Exception as e:
            logger.warning(f"Failed to load LayoutLMv3 model: {e}")
            self.processor = None
            self.model = None

    def analyze(
        self,
        image_path: str,
        text_boxes: List[Dict[str, Any]],
        normalize_coords: bool = True
    ) -> SemanticAnalysisResult:
        """
        Analyze slide semantics using LayoutLMv3.

        Args:
            image_path: Path to the slide image.
            text_boxes: List of text boxes with 'text' and 'bbox' keys.
                       bbox format: [x1, y1, x2, y2] in pixels.
            normalize_coords: Whether to normalize bbox to 0-1000 range.

        Returns:
            SemanticAnalysisResult with labels and analysis.
        """
        # Convert dict boxes to TextBox objects
        boxes = []
        for box in text_boxes:
            boxes.append(TextBox(
                text=box.get("text", ""),
                bbox=tuple(box.get("bbox", [0, 0, 0, 0])),
                confidence=box.get("confidence", 1.0)
            ))

        # Choose analysis method based on available resources
        if self.model and self.processor:
            return self._analyze_with_layoutlm(image_path, boxes, normalize_coords)
        else:
            return self._analyze_heuristic(image_path, boxes)

    def _analyze_with_layoutlm(
        self,
        image_path: str,
        text_boxes: List[TextBox],
        normalize_coords: bool
    ) -> SemanticAnalysisResult:
        """Analyze using LayoutLMv3 model."""
        try:
            from PIL import Image

            image = Image.open(image_path).convert("RGB")
            width, height = image.size

            # Prepare text and bounding boxes
            words = []
            boxes = []

            for box in text_boxes:
                # Split text into words for token-level processing
                box_words = box.text.split()
                for word in box_words:
                    words.append(word)

                    # Normalize bbox to 0-1000 range (LayoutLMv3 requirement)
                    if normalize_coords:
                        norm_bbox = [
                            int(box.bbox[0] * 1000 / width),
                            int(box.bbox[1] * 1000 / height),
                            int(box.bbox[2] * 1000 / width),
                            int(box.bbox[3] * 1000 / height),
                        ]
                    else:
                        norm_bbox = list(box.bbox)

                    boxes.append(norm_bbox)

            if not words:
                return self._analyze_heuristic(image_path, text_boxes)

            # Process with LayoutLMv3
            encoding = self.processor(
                image,
                words,
                boxes=boxes,
                return_tensors="pt",
                truncation=True,
                max_length=512
            )

            # Move to device
            encoding = {k: v.to(self.device) for k, v in encoding.items()}

            # Run inference
            with torch.no_grad():
                outputs = self.model(**encoding)
                predictions = torch.argmax(outputs.logits, dim=-1)
                confidences = torch.softmax(outputs.logits, dim=-1).max(dim=-1).values

            # Decode predictions to semantic labels
            semantic_labels = []
            pred_list = predictions[0].cpu().tolist()
            conf_list = confidences[0].cpu().tolist()

            word_idx = 0
            for box in text_boxes:
                box_words = box.text.split()
                if not box_words:
                    continue

                # Use first word's prediction for the box
                if word_idx < len(pred_list):
                    label_idx = pred_list[word_idx]
                    label = SLIDE_ELEMENT_LABELS.get(label_idx, "other")
                    confidence = conf_list[word_idx]

                    semantic_labels.append(SemanticLabel(
                        text=box.text,
                        label=label,
                        confidence=confidence,
                        bbox=box.bbox,
                        purpose=self._infer_purpose(label, box)
                    ))

                word_idx += len(box_words)

            # Determine overall content purpose and template
            content_purpose = self._classify_content_purpose(semantic_labels)
            suggested_template = self._suggest_template(semantic_labels, content_purpose)
            key_elements = self._extract_key_elements(semantic_labels)

            avg_confidence = (
                sum(l.confidence for l in semantic_labels) / len(semantic_labels)
                if semantic_labels else 0.0
            )

            return SemanticAnalysisResult(
                semantic_labels=semantic_labels,
                content_purpose=content_purpose,
                suggested_template=suggested_template,
                key_elements=key_elements,
                confidence=avg_confidence,
                analysis_method="layoutlmv3"
            )

        except Exception as e:
            logger.error(f"LayoutLMv3 analysis failed: {e}")
            return self._analyze_heuristic(image_path, text_boxes)

    def _analyze_heuristic(
        self,
        image_path: str,
        text_boxes: List[TextBox]
    ) -> SemanticAnalysisResult:
        """Fallback heuristic analysis without ML model."""
        try:
            from PIL import Image
            image = Image.open(image_path)
            width, height = image.size
        except Exception:
            width, height = 1920, 1080  # Assume standard size

        semantic_labels = []

        for box in text_boxes:
            label = self._heuristic_label(box, width, height)
            purpose = self._infer_purpose(label, box)

            semantic_labels.append(SemanticLabel(
                text=box.text,
                label=label,
                confidence=0.7,  # Lower confidence for heuristic
                bbox=box.bbox,
                purpose=purpose
            ))

        content_purpose = self._classify_content_purpose(semantic_labels)
        suggested_template = self._suggest_template(semantic_labels, content_purpose)
        key_elements = self._extract_key_elements(semantic_labels)

        return SemanticAnalysisResult(
            semantic_labels=semantic_labels,
            content_purpose=content_purpose,
            suggested_template=suggested_template,
            key_elements=key_elements,
            confidence=0.6,
            analysis_method="heuristic"
        )

    def _heuristic_label(
        self,
        box: TextBox,
        slide_width: int,
        slide_height: int
    ) -> str:
        """Assign label based on position and text characteristics."""
        rel_y = box.center[1] / slide_height
        rel_height = box.height / slide_height
        text_length = len(box.text)

        # Title: top of slide, larger text
        if rel_y < 0.2 and rel_height > 0.05:
            return "title"

        # Subtitle: below title, smaller
        if 0.15 < rel_y < 0.3 and text_length < 100:
            return "subtitle"

        # Footer: bottom of slide
        if rel_y > 0.9:
            return "footer"

        # Bullets: starts with bullet characters
        if box.text.strip().startswith(('-', '*', '•', '>')):
            return "bullet"

        # Caption: small text near figures (bottom half, short)
        if rel_y > 0.7 and text_length < 50:
            return "caption"

        # Default to body text
        return "body"

    def _infer_purpose(self, label: str, box: TextBox) -> str:
        """Infer the purpose of a text element."""
        if label in ("title", "subtitle"):
            return "main_content"
        elif label == "header":
            return "navigation"
        elif label in ("body", "bullet"):
            return "main_content"
        elif label == "caption":
            return "supporting"
        elif label == "footer":
            return "branding"
        return "other"

    def _classify_content_purpose(
        self,
        labels: List[SemanticLabel]
    ) -> str:
        """Classify the overall purpose of the slide content."""
        if not labels:
            return "informative"

        texts = " ".join(l.text.lower() for l in labels)

        # Check for purpose indicators
        if any(w in texts for w in ["agenda", "outline", "topics", "contents"]):
            return "agenda"

        if any(w in texts for w in ["thank", "questions", "contact", "end"]):
            return "conclusion"

        if any(w in texts for w in ["introduction", "welcome", "overview"]):
            return "introduction"

        if any(w in texts for w in ["vs", "versus", "comparison", "compare", "option"]):
            return "comparative"

        if any(w in texts for w in ["$", "%", "growth", "revenue", "metric"]):
            return "data_presentation"

        if any(w in texts for w in ["step", "process", "how to", "guide"]):
            return "instructional"

        if any(w in texts for w in ["key", "summary", "takeaway", "conclusion"]):
            return "summary"

        if any(w in texts for w in ["benefit", "why", "solution", "opportunity"]):
            return "persuasive"

        return "informative"

    def _suggest_template(
        self,
        labels: List[SemanticLabel],
        content_purpose: str
    ) -> str:
        """Suggest appropriate template based on content analysis."""
        # Map purposes to templates
        purpose_templates = {
            "agenda": "consulting_toolkit",
            "data_presentation": "market_analysis",
            "persuasive": "business_case",
            "comparative": "consulting_toolkit",
            "summary": "consulting_toolkit",
            "introduction": "default",
            "conclusion": "default",
        }

        return purpose_templates.get(content_purpose, "consulting_toolkit")

    def _extract_key_elements(
        self,
        labels: List[SemanticLabel]
    ) -> Dict[str, str]:
        """Extract key semantic elements from labels."""
        elements = {}

        for label in labels:
            if label.label == "title" and "title" not in elements:
                elements["title"] = label.text
            elif label.label == "subtitle" and "subtitle" not in elements:
                elements["subtitle"] = label.text
            elif label.label == "body" and "main_point" not in elements:
                elements["main_point"] = label.text[:100]

        return elements


def analyze_slide(
    image_path: str,
    text_boxes: List[Dict[str, Any]],
    use_gpu: bool = False
) -> SemanticAnalysisResult:
    """
    Convenience function to analyze a single slide.

    Args:
        image_path: Path to slide image.
        text_boxes: List of text boxes with 'text' and 'bbox' keys.
        use_gpu: Whether to use GPU.

    Returns:
        SemanticAnalysisResult.
    """
    analyzer = SemanticSlideAnalyzer(use_gpu=use_gpu)
    return analyzer.analyze(image_path, text_boxes)


def extract_text_boxes_from_pptx(pptx_path: str, slide_index: int) -> List[Dict[str, Any]]:
    """
    Extract text boxes from a PowerPoint slide.

    Args:
        pptx_path: Path to PPTX file.
        slide_index: Zero-based slide index.

    Returns:
        List of text box dictionaries with 'text' and 'bbox'.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation(pptx_path)
        slide = prs.slides[slide_index]

        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # Convert EMUs to pixels (assuming 96 DPI)
                    x = int(shape.left / Emu(914400) * 96)  # EMU to inches to pixels
                    y = int(shape.top / Emu(914400) * 96)
                    w = int(shape.width / Emu(914400) * 96)
                    h = int(shape.height / Emu(914400) * 96)

                    text_boxes.append({
                        "text": text,
                        "bbox": [x, y, x + w, y + h]
                    })

        return text_boxes

    except Exception as e:
        logger.error(f"Failed to extract text boxes: {e}")
        return []


# CLI interface
if __name__ == "__main__":
    import json
    import sys

    print("LayoutLMv3 Semantic Slide Analyzer")
    print(f"PyTorch available: {TORCH_AVAILABLE}")
    print(f"Transformers available: {TRANSFORMERS_AVAILABLE}")
    print()

    if len(sys.argv) > 1:
        image_path = sys.argv[1]

        # Demo text boxes (normally extracted from PPTX or OCR)
        demo_boxes = [
            {"text": "Q4 2025 Business Review", "bbox": [100, 50, 800, 120]},
            {"text": "Strategic Analysis and Recommendations", "bbox": [100, 130, 700, 170]},
            {"text": "Revenue grew 25% YoY", "bbox": [100, 250, 600, 290]},
            {"text": "Customer acquisition up 40%", "bbox": [100, 300, 600, 340]},
            {"text": "Market share expanded to 18%", "bbox": [100, 350, 600, 390]},
        ]

        print(f"Analyzing: {image_path}")
        print(f"Text boxes: {len(demo_boxes)}")
        print()

        result = analyze_slide(image_path, demo_boxes)

        print(f"Content Purpose: {result.content_purpose}")
        print(f"Suggested Template: {result.suggested_template}")
        print(f"Analysis Method: {result.analysis_method}")
        print(f"Confidence: {result.confidence:.2f}")
        print()

        print("Semantic Labels:")
        for label in result.semantic_labels:
            print(f"  [{label.label}] {label.text[:50]}...")
            print(f"      Purpose: {label.purpose}, Confidence: {label.confidence:.2f}")
        print()

        print("Key Elements:")
        for key, value in result.key_elements.items():
            print(f"  {key}: {value}")

    else:
        print("Usage: python -m pptx_extractor.semantic_analyzer <image_path>")
        print()
        print("Analyze slide semantics using LayoutLMv3.")
        print()
        print("Dependencies:")
        print("  pip install transformers torch torchvision")
        print("  pip install pillow")
