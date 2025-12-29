"""
Slide Sequence Extractor - Extract slide flow patterns and narrative structures.

This module extracts:
- Slide sequence patterns (executive summary, detail sections, appendix)
- Narrative flow structures
- Common slide progressions
- Section patterns
- Deck templates (complete presentation structures)
"""

import json
import hashlib
import logging
import re
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Any, Tuple
from collections import defaultdict

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class SequenceExtractor:
    """Extract slide sequence patterns and narrative structures."""

    # Common slide types based on content/layout
    SLIDE_TYPES = {
        'title': ['title', 'cover', 'front'],
        'agenda': ['agenda', 'contents', 'outline', 'index'],
        'section': ['section', 'divider', 'header'],
        'summary': ['summary', 'overview', 'key', 'takeaway', 'conclusion'],
        'detail': ['detail', 'analysis', 'deep dive'],
        'data': ['chart', 'graph', 'data', 'metrics'],
        'comparison': ['comparison', 'versus', 'vs', 'compare'],
        'timeline': ['timeline', 'roadmap', 'milestones'],
        'team': ['team', 'about', 'bio', 'leadership'],
        'appendix': ['appendix', 'backup', 'additional', 'reference'],
        'contact': ['contact', 'thank you', 'questions'],
    }

    def __init__(self, output_dir: Path):
        """Initialize the sequence extractor."""
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Create subdirectories
        self.dirs = {
            'sequences': self.output_dir / 'sequences' / 'patterns',
            'deck_templates': self.output_dir / 'sequences' / 'deck_templates',
        }
        for d in self.dirs.values():
            d.mkdir(parents=True, exist_ok=True)

        # Sequence index
        self.index_path = self.output_dir / 'sequences' / 'sequence_index.json'
        self.index = self._load_index()

    def _load_index(self) -> dict:
        """Load existing index or create new one."""
        if self.index_path.exists():
            with open(self.index_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            'metadata': {
                'created': datetime.now().isoformat(),
                'last_updated': datetime.now().isoformat(),
                'version': '1.0',
            },
            'sequences': [],
            'deck_templates': [],
            'common_patterns': [],
        }

    def _save_index(self):
        """Save index to file."""
        self.index['metadata']['last_updated'] = datetime.now().isoformat()
        self.index_path.parent.mkdir(parents=True, exist_ok=True)
        with open(self.index_path, 'w', encoding='utf-8') as f:
            json.dump(self.index, f, indent=2, ensure_ascii=False)

    def _generate_id(self, content: str) -> str:
        """Generate unique ID from content."""
        return hashlib.md5(content.encode()).hexdigest()[:12]

    def extract_template_sequences(self, pptx_path: Path, template_name: Optional[str] = None) -> dict:
        """
        Extract slide sequence patterns from a PowerPoint template.

        Args:
            pptx_path: Path to the PowerPoint file
            template_name: Optional name for the template

        Returns:
            Dictionary of extracted sequences
        """
        pptx_path = Path(pptx_path)
        template_name = template_name or pptx_path.stem

        logger.info(f"Extracting sequences from: {template_name}")

        prs = Presentation(pptx_path)

        # Analyze each slide
        slide_info = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1
            info = self._analyze_slide(slide, slide_num, template_name)
            slide_info.append(info)

        # Build deck template
        deck_template = self._build_deck_template(slide_info, template_name)

        # Detect sequences within the deck
        sequences = self._detect_sequences(slide_info, template_name)

        # Detect common patterns
        patterns = self._detect_common_patterns(sequences)

        results = {
            'template': template_name,
            'deck_template': deck_template,
            'sequences': sequences,
            'patterns': patterns,
            'slide_count': len(slide_info),
        }

        # Save results
        self._save_results(results, template_name)

        logger.info(f"Extracted {len(sequences)} sequences from {len(slide_info)} slides")

        return results

    def _analyze_slide(self, slide, slide_num: int, template_name: str) -> dict:
        """Analyze a slide to determine its type and purpose."""
        info = {
            'slide_num': slide_num,
            'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown',
            'slide_type': 'content',
            'title': '',
            'content_types': [],
            'word_count': 0,
            'element_count': 0,
        }

        # Count elements and extract content
        for shape in slide.shapes:
            info['element_count'] += 1
            shape_type = self._get_shape_type(shape)

            if shape_type == 'CHART':
                info['content_types'].append('chart')
            elif shape_type == 'TABLE':
                info['content_types'].append('table')
            elif shape_type == 'PICTURE':
                info['content_types'].append('image')
            elif shape_type == 'GROUP':
                info['content_types'].append('diagram')

            # Extract title
            if hasattr(shape, 'text') and shape.text:
                text = shape.text.strip()
                info['word_count'] += len(text.split())

                # Check if this might be the title
                if not info['title'] and len(text) < 100:
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                font_size = para.runs[0].font.size
                                if font_size and font_size.pt >= 18:
                                    info['title'] = text
                    except:
                        pass

        # If no title found, use first text
        if not info['title']:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    info['title'] = shape.text.strip()[:80]
                    break

        # Determine slide type
        info['slide_type'] = self._determine_slide_type(info)

        return info

    def _get_shape_type(self, shape) -> str:
        """Get shape type as string."""
        try:
            if hasattr(shape, 'shape_type'):
                return str(shape.shape_type).split('.')[-1]
        except:
            pass
        return 'UNKNOWN'

    def _determine_slide_type(self, info: dict) -> str:
        """Determine the type/purpose of a slide."""
        layout_name = info['layout_name'].lower()
        title = info['title'].lower()
        content_types = info['content_types']

        # Check layout name first
        for slide_type, keywords in self.SLIDE_TYPES.items():
            if any(kw in layout_name for kw in keywords):
                return slide_type

        # Check title
        for slide_type, keywords in self.SLIDE_TYPES.items():
            if any(kw in title for kw in keywords):
                return slide_type

        # Check content
        if 'chart' in content_types:
            if len(content_types) == 1 or content_types.count('chart') >= 2:
                return 'data'
        if 'table' in content_types:
            return 'comparison'
        if 'diagram' in content_types:
            return 'detail'

        # Check by position/characteristics
        if info['slide_num'] == 1:
            return 'title'
        if info['element_count'] <= 2:
            return 'section'

        return 'content'

    def _build_deck_template(self, slide_info: List[dict], template_name: str) -> dict:
        """Build a deck template from slide analysis."""
        # Extract the sequence of slide types
        type_sequence = [s['slide_type'] for s in slide_info]

        # Identify sections
        sections = self._identify_sections(slide_info)

        deck_template = {
            'id': self._generate_id(f"deck_{template_name}"),
            'template': template_name,
            'slide_count': len(slide_info),
            'type_sequence': type_sequence,
            'sections': sections,
            'structure': self._analyze_deck_structure(type_sequence),
        }

        return deck_template

    def _identify_sections(self, slide_info: List[dict]) -> List[dict]:
        """Identify sections within the deck."""
        sections = []
        current_section = None

        for i, slide in enumerate(slide_info):
            slide_type = slide['slide_type']

            # New section starts
            if slide_type in ('title', 'section', 'agenda'):
                if current_section:
                    current_section['end_slide'] = i
                    current_section['slide_count'] = i - current_section['start_slide']
                    sections.append(current_section)

                current_section = {
                    'name': slide['title'][:50] if slide['title'] else f'Section {len(sections) + 1}',
                    'start_slide': i + 1,
                    'type': slide_type,
                    'slides': [],
                }
            elif current_section:
                current_section['slides'].append({
                    'slide_num': slide['slide_num'],
                    'type': slide_type,
                    'title': slide['title'][:50] if slide['title'] else '',
                })

        # Close last section
        if current_section:
            current_section['end_slide'] = len(slide_info)
            current_section['slide_count'] = len(slide_info) - current_section['start_slide'] + 1
            sections.append(current_section)

        return sections

    def _analyze_deck_structure(self, type_sequence: List[str]) -> dict:
        """Analyze the overall deck structure."""
        structure = {
            'has_title': 'title' in type_sequence[:3],
            'has_agenda': 'agenda' in type_sequence[:5],
            'has_summary': 'summary' in type_sequence,
            'has_appendix': 'appendix' in type_sequence[-10:] if len(type_sequence) > 10 else False,
            'has_contact': 'contact' in type_sequence[-3:] if len(type_sequence) > 3 else False,
            'section_count': type_sequence.count('section'),
            'data_slide_count': type_sequence.count('data'),
            'comparison_count': type_sequence.count('comparison'),
            'structure_type': 'unknown',
        }

        # Determine structure type
        if structure['has_title'] and structure['has_agenda'] and structure['has_summary']:
            structure['structure_type'] = 'executive_presentation'
        elif structure['data_slide_count'] >= len(type_sequence) * 0.5:
            structure['structure_type'] = 'data_heavy'
        elif structure['section_count'] >= 3:
            structure['structure_type'] = 'multi_section'
        elif len(type_sequence) <= 10:
            structure['structure_type'] = 'brief'
        else:
            structure['structure_type'] = 'detailed'

        return structure

    def _detect_sequences(self, slide_info: List[dict], template_name: str) -> List[dict]:
        """Detect reusable sequences within the slides."""
        sequences = []

        # Look for common sequence patterns
        i = 0
        while i < len(slide_info):
            # Check for section + content sequence
            if slide_info[i]['slide_type'] == 'section':
                seq = self._extract_section_sequence(slide_info, i)
                if seq:
                    seq['template'] = template_name
                    sequences.append(seq)
                    i += seq['length']
                    continue

            # Check for data comparison sequence (multiple data slides)
            if slide_info[i]['slide_type'] == 'data':
                seq = self._extract_data_sequence(slide_info, i)
                if seq:
                    seq['template'] = template_name
                    sequences.append(seq)
                    i += seq['length']
                    continue

            # Check for summary sequence
            if slide_info[i]['slide_type'] == 'summary':
                seq = self._extract_summary_sequence(slide_info, i)
                if seq:
                    seq['template'] = template_name
                    sequences.append(seq)
                    i += seq['length']
                    continue

            i += 1

        return sequences

    def _extract_section_sequence(self, slide_info: List[dict], start_idx: int) -> Optional[dict]:
        """Extract a section + content sequence."""
        if start_idx >= len(slide_info):
            return None

        sequence = {
            'type': 'section_content',
            'slides': [],
            'length': 0,
        }

        # Add section header
        sequence['slides'].append({
            'type': slide_info[start_idx]['slide_type'],
            'purpose': 'section_header',
        })

        # Add following content slides until next section or end
        for i in range(start_idx + 1, min(start_idx + 20, len(slide_info))):
            slide_type = slide_info[i]['slide_type']

            if slide_type in ('section', 'appendix', 'contact'):
                break

            sequence['slides'].append({
                'type': slide_type,
                'purpose': 'content',
            })

        sequence['length'] = len(sequence['slides'])

        if sequence['length'] < 2:
            return None

        return sequence

    def _extract_data_sequence(self, slide_info: List[dict], start_idx: int) -> Optional[dict]:
        """Extract a data visualization sequence."""
        if start_idx >= len(slide_info):
            return None

        sequence = {
            'type': 'data_analysis',
            'slides': [],
            'length': 0,
        }

        for i in range(start_idx, min(start_idx + 10, len(slide_info))):
            slide_type = slide_info[i]['slide_type']

            if slide_type not in ('data', 'comparison', 'content'):
                break

            sequence['slides'].append({
                'type': slide_type,
                'purpose': 'data_point' if slide_type == 'data' else 'analysis',
            })

        sequence['length'] = len(sequence['slides'])

        if sequence['length'] < 2:
            return None

        return sequence

    def _extract_summary_sequence(self, slide_info: List[dict], start_idx: int) -> Optional[dict]:
        """Extract a summary/conclusion sequence."""
        sequence = {
            'type': 'conclusion',
            'slides': [],
            'length': 0,
        }

        for i in range(start_idx, min(start_idx + 5, len(slide_info))):
            slide_type = slide_info[i]['slide_type']
            sequence['slides'].append({
                'type': slide_type,
                'purpose': 'wrap_up',
            })

            if slide_type in ('contact', 'appendix'):
                break

        sequence['length'] = len(sequence['slides'])

        if sequence['length'] < 1:
            return None

        return sequence

    def _detect_common_patterns(self, sequences: List[dict]) -> List[dict]:
        """Detect common patterns across sequences."""
        pattern_counts = defaultdict(int)

        for seq in sequences:
            # Create pattern signature
            types = [s['type'] for s in seq.get('slides', [])]
            pattern_key = '_'.join(types[:5])  # First 5 types
            pattern_counts[pattern_key] += 1

        # Convert to list of patterns
        patterns = []
        for pattern_key, count in pattern_counts.items():
            if count >= 2:  # Only include patterns that appear multiple times
                patterns.append({
                    'pattern': pattern_key.split('_'),
                    'occurrence_count': count,
                })

        patterns.sort(key=lambda x: x['occurrence_count'], reverse=True)
        return patterns[:10]  # Top 10 patterns

    def _save_results(self, results: dict, template_name: str):
        """Save extraction results."""
        # Save deck template
        deck = results['deck_template']
        filename = f"{deck['id']}.json"
        filepath = self.dirs['deck_templates'] / filename
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(deck, f, indent=2)

        deck['filename'] = filename

        # Add to index
        existing = [x for x in self.index['deck_templates'] if x.get('id') == deck['id']]
        if not existing:
            self.index['deck_templates'].append(deck)

        # Save sequences
        for seq in results['sequences']:
            seq_id = self._generate_id(f"{template_name}_{seq['type']}_{len(seq['slides'])}")
            seq['id'] = seq_id

            filename = f"{seq_id}.json"
            filepath = self.dirs['sequences'] / filename
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(seq, f, indent=2)

            seq['filename'] = filename

            existing = [x for x in self.index['sequences'] if x.get('id') == seq_id]
            if not existing:
                self.index['sequences'].append(seq)

        # Save common patterns
        self.index['common_patterns'] = results['patterns']

        self._save_index()

    def get_summary(self) -> dict:
        """Get summary of extracted sequences."""
        summary = {
            'deck_templates': len(self.index['deck_templates']),
            'sequences': len(self.index['sequences']),
            'common_patterns': len(self.index['common_patterns']),
            'by_sequence_type': defaultdict(int),
        }

        for seq in self.index['sequences']:
            summary['by_sequence_type'][seq.get('type', 'unknown')] += 1

        summary['by_sequence_type'] = dict(summary['by_sequence_type'])

        return summary

    def search_sequences(self, sequence_type: str = None, template: str = None,
                         min_length: int = None) -> List[dict]:
        """Search sequences."""
        results = []

        for seq in self.index['sequences']:
            if sequence_type and seq.get('type') != sequence_type:
                continue
            if template and seq.get('template') != template:
                continue
            if min_length and seq.get('length', 0) < min_length:
                continue
            results.append(seq)

        return results

    def get_deck_template(self, template_name: str) -> Optional[dict]:
        """Get deck template by template name."""
        for deck in self.index['deck_templates']:
            if deck.get('template') == template_name:
                return deck
        return None

    def generate_deck_outline(self, structure_type: str, topic: str = 'Presentation') -> dict:
        """
        Generate a deck outline based on structure type.

        Args:
            structure_type: Type of structure (executive_presentation, data_heavy, etc.)
            topic: Topic/title for the presentation

        Returns:
            Deck outline with slide placeholders
        """
        outlines = {
            'executive_presentation': [
                {'type': 'title', 'title': topic},
                {'type': 'agenda', 'title': 'Agenda'},
                {'type': 'summary', 'title': 'Executive Summary'},
                {'type': 'section', 'title': 'Background'},
                {'type': 'content', 'title': 'Current Situation'},
                {'type': 'section', 'title': 'Analysis'},
                {'type': 'data', 'title': 'Key Metrics'},
                {'type': 'comparison', 'title': 'Comparison'},
                {'type': 'section', 'title': 'Recommendations'},
                {'type': 'content', 'title': 'Proposed Solution'},
                {'type': 'timeline', 'title': 'Implementation Roadmap'},
                {'type': 'summary', 'title': 'Key Takeaways'},
                {'type': 'contact', 'title': 'Questions & Next Steps'},
            ],
            'data_heavy': [
                {'type': 'title', 'title': topic},
                {'type': 'summary', 'title': 'Key Findings'},
                {'type': 'data', 'title': 'Overview Metrics'},
                {'type': 'data', 'title': 'Trend Analysis'},
                {'type': 'data', 'title': 'Segment Analysis'},
                {'type': 'comparison', 'title': 'Benchmark Comparison'},
                {'type': 'data', 'title': 'Detailed Breakdown'},
                {'type': 'summary', 'title': 'Conclusions'},
                {'type': 'appendix', 'title': 'Appendix: Data Sources'},
            ],
            'brief': [
                {'type': 'title', 'title': topic},
                {'type': 'content', 'title': 'Overview'},
                {'type': 'content', 'title': 'Key Points'},
                {'type': 'data', 'title': 'Supporting Data'},
                {'type': 'summary', 'title': 'Summary'},
            ],
            'multi_section': [
                {'type': 'title', 'title': topic},
                {'type': 'agenda', 'title': 'Contents'},
                {'type': 'section', 'title': 'Section 1'},
                {'type': 'content', 'title': 'Content 1'},
                {'type': 'section', 'title': 'Section 2'},
                {'type': 'content', 'title': 'Content 2'},
                {'type': 'section', 'title': 'Section 3'},
                {'type': 'content', 'title': 'Content 3'},
                {'type': 'summary', 'title': 'Summary'},
            ],
        }

        outline = outlines.get(structure_type, outlines['brief'])

        return {
            'structure_type': structure_type,
            'topic': topic,
            'slides': outline,
            'slide_count': len(outline),
        }


def extract_sequences_from_templates(templates_dir: Path, output_dir: Path) -> dict:
    """
    Extract sequences from all templates in a directory.

    Args:
        templates_dir: Directory containing PPTX templates
        output_dir: Output directory for the library

    Returns:
        Summary of extraction
    """
    extractor = SequenceExtractor(output_dir)

    # Find all PPTX files
    pptx_files = list(templates_dir.rglob('*.pptx'))

    print(f"Found {len(pptx_files)} PowerPoint files")

    for pptx_path in pptx_files:
        print(f"\nExtracting sequences from: {pptx_path.name}")
        try:
            extractor.extract_template_sequences(pptx_path)
        except Exception as e:
            print(f"  Error: {e}")

    extractor._save_index()
    return extractor.get_summary()


if __name__ == '__main__':
    import sys

    if len(sys.argv) < 3:
        print("Usage: python sequence_extractor.py <templates_dir> <output_dir>")
        sys.exit(1)

    templates_dir = Path(sys.argv[1])
    output_dir = Path(sys.argv[2])

    summary = extract_sequences_from_templates(templates_dir, output_dir)

    print("\n" + "="*50)
    print("SEQUENCE EXTRACTION COMPLETE")
    print("="*50)
    print(json.dumps(summary, indent=2))
